"""
Top-level autonomous pipeline router.

Accepts a markdown file, an existing .pptx, a plan dict, or a plan JSON file
and produces a Merck-branded .pptx.
"""

import colorsys
import json
import pathlib

from merck_pptx.build_from_plan import build_from_plan
from merck_pptx import llm


def _shape_color_name(shape) -> str | None:
    """Return a named Merck-palette color category for a shape's fill, or None."""
    try:
        fill = shape.fill
        if fill.type is None:
            return None
        rgb = fill.fore_color.rgb
        r, g, b = rgb.r / 255.0, rgb.g / 255.0, rgb.b / 255.0
        h, s, v = colorsys.rgb_to_hsv(r, g, b)
        # Near-white, very light, or very dark — skip
        if v > 0.92 or v < 0.08:
            return None
        # Low saturation → gray
        if s < 0.18:
            return "gray"
        # Classify by hue
        if h < 0.04 or h >= 0.92:
            return "red"
        if h < 0.08:
            return "orange"
        if h < 0.18:
            return "yellow"
        if h < 0.42:
            return "green"
        if h < 0.56:
            return "teal"
        if h < 0.68:
            return "blue"
        if h < 0.82:
            return "purple"
        return "pink"
    except Exception:
        return None


def _extract_pptx_structured(source_path: pathlib.Path) -> list:
    """Extract per-slide structured content from a PPTX, preserving text verbatim.

    Returns a list of dicts, one per slide:
      {
        "slide":          int,      # 1-based slide index
        "title":          str,      # title placeholder text (verbatim)
        "text_blocks": [            # all non-title text shapes, in top-to-bottom order
          {
            "text":       str,      # full raw text
            "bullets":    [str],    # paragraphs when >1
            "is_figure":  bool,     # True when text starts with [FIGURE]
            "color":      str|null, # Merck palette name: gray/teal/green/yellow/orange/
                                    #   red/blue/purple/pink — or null if no distinctive fill
          }
        ],
        "has_figures":    bool,     # True if any non-text shape or [FIGURE] block
        "color_sequence": [str],    # ordered list of non-null color names from content shapes
      }

    Retries up to 3 times with a short delay to handle transient OneDrive file locks.
    On PermissionError (file open in PowerPoint / OneDrive sync lock) copies to a
    temp file first so python-pptx can open a read-only snapshot.
    """
    import shutil
    import tempfile
    import time
    from pptx import Presentation
    from pptx.exc import PackageNotFoundError

    def _open_pptx(path: pathlib.Path):
        try:
            return Presentation(str(path))
        except PermissionError:
            # File is locked (open in PowerPoint or OneDrive sync); copy to temp
            tmp = pathlib.Path(tempfile.mktemp(suffix=".pptx"))
            shutil.copy2(str(path), str(tmp))
            try:
                return Presentation(str(tmp))
            finally:
                try:
                    tmp.unlink(missing_ok=True)
                except OSError:
                    pass

    last_exc: Exception = RuntimeError("no attempts made")
    for attempt in range(3):
        try:
            prs = _open_pptx(source_path)
            break
        except (MemoryError, KeyboardInterrupt):
            raise
        except (PackageNotFoundError, OSError, IOError) as exc:
            last_exc = exc
            if attempt < 2:
                time.sleep(1.5)
        except Exception as exc:
            last_exc = exc
            if attempt < 2:
                time.sleep(1.5)
    else:
        raise last_exc

    result = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        entry = {
            "slide":          slide_idx,
            "title":          "",
            "text_blocks":    [],
            "has_figures":    False,
            "color_sequence": [],
        }

        # Sort shapes top-to-bottom then left-to-right so reading order is
        # preserved for two-column and other multi-column source layouts.
        # Shapes with top=None are pushed to the end (unpositioned shapes).
        shapes_sorted = sorted(
            slide.shapes,
            key=lambda s: (
                s.top  if s.top  is not None else float("inf"),
                s.left if s.left is not None else 0,
            ),
        )

        for shape in shapes_sorted:
            if not shape.has_text_frame:
                entry["has_figures"] = True
                continue

            text = shape.text_frame.text.strip()
            if not text:
                continue

            is_fig = text.upper().startswith("[FIGURE]")
            if is_fig:
                entry["has_figures"] = True

            color = _shape_color_name(shape)

            is_title = (
                shape.is_placeholder
                and shape.placeholder_format.idx == 0
                and not entry["title"]
            )
            if is_title:
                entry["title"] = text
            else:
                bullets = [
                    p.text.strip()
                    for p in shape.text_frame.paragraphs
                    if p.text.strip()
                ]
                block = {
                    "text":      text,
                    "bullets":   bullets if len(bullets) > 1 else [],
                    "is_figure": is_fig,
                    "color":     color,
                }
                entry["text_blocks"].append(block)
                if color:
                    entry["color_sequence"].append(color)

        if not entry["title"] and entry["text_blocks"]:
            entry["title"] = entry["text_blocks"].pop(0)["text"]

        result.append(entry)

    if not result:
        raise ValueError(
            f"No slides were extracted from '{source_path}'. "
            f"The file may be empty or corrupted."
        )

    return result


def generate_deck(source, output_path, meta=None, save_plan=None) -> str:
    """
    Autonomous end-to-end pipeline. Any input type → Merck-branded .pptx.

    Parameters
    ----------
    source : dict | str | pathlib.Path
        One of:
          - A plan dict  →  built directly (no LLM)
          - A .json file →  loaded and built directly (no LLM)
          - A .md file   →  LLM converts content to plan, then builds
          - A .pptx file →  slides extracted, LLM converts to plan, then builds
    output_path : str | pathlib.Path
        Destination .pptx path.
    meta : dict | None
        Required when source is a .md or .pptx file. Contains the 6 gate answers:
        region, deck_label, classification, month_year, audience, deck_style.
        Ignored when source is a dict or .json (meta is already embedded in the plan).
    save_plan : str | pathlib.Path | None
        If set, the LLM-generated plan dict is written to this path as JSON before
        building. Ignored for dict/.json inputs (no LLM is called). Useful for
        inspecting or re-using the plan with the `build` command.

    Returns
    -------
    str
        Absolute path of the saved .pptx.

    Raises
    ------
    ValueError
        If meta is None when source is .md or .pptx.
    EnvironmentError
        If AIP_BASE_URL or AIP_TOKEN are not set when LLM is required.
    """
    import warnings as _warnings

    # Plan dict — direct build, no LLM
    if isinstance(source, dict):
        if save_plan is not None:
            _warnings.warn(
                "save_plan is ignored when source is a dict (no LLM call is made).",
                stacklevel=2,
            )
        return build_from_plan(source, output_path)

    # Resolve to a Path for all file-based inputs (str or pathlib.Path)
    if isinstance(source, (str, pathlib.Path)):
        source_path = pathlib.Path(source)
    else:
        raise ValueError(
            f"Unsupported source type: {type(source).__name__}. "
            f"Expected a dict, or a path (str or pathlib.Path) to a .json, .md, or .pptx file."
        )

    suffix = source_path.suffix.lower()

    # .json plan file — direct build, no LLM
    if suffix == ".json":
        if save_plan is not None:
            _warnings.warn(
                "save_plan is ignored when source is a .json file (no LLM call is made).",
                stacklevel=2,
            )
        return build_from_plan(source_path, output_path)

    # .md or .pptx — LLM path
    if meta is None:
        raise ValueError(
            "meta is required when source is a .md or .pptx file. "
            "Provide a dict with region, deck_label, classification, "
            "month_year, audience, and deck_style."
        )

    if suffix == ".md":
        raw_content = source_path.read_text(encoding="utf-8")
        plan = llm.generate_plan(raw_content, meta)
    elif suffix == ".pptx":
        slides = _extract_pptx_structured(source_path)
        plan = llm.generate_plan_from_pptx(slides, meta)
    elif suffix == ".ppt":
        raise ValueError(
            "Legacy .ppt (PowerPoint 97-2003 binary) format is not supported. "
            "Open the file in PowerPoint, save as .pptx, then retry."
        )
    else:
        raise ValueError(
            f"Unsupported file type '{suffix}'. "
            f"Expected .json (plan), .md (markdown), or .pptx."
        )

    if save_plan is not None:
        save_path = pathlib.Path(save_plan)
        save_path.parent.mkdir(parents=True, exist_ok=True)
        with open(save_path, "w", encoding="utf-8") as f:
            json.dump(plan, f, ensure_ascii=False, indent=2)

    return build_from_plan(plan, output_path)
