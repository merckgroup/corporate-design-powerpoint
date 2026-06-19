from __future__ import annotations
import math
import sys
from typing import Optional
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt
from ._ml_constants import (
    FONT_HEAD, FONT_BODY,
    MERCK_PURPLE, MERCK_GOLD, MERCK_YELLOW, MERCK_BLUE, MERCK_AQUA,
    WHITE, INK_DARK, INK_GRAY, LIGHT_GRAY, PANEL_LIGHT,
    PURPLE_DEEP, PURPLE_MUTED, BAD_RED, GOOD_GREEN,
    ACT_PURPLE, LY_CYAN, OP_LIME, FC_PINK, DEV_POS_BLUE, DEV_NEG_RED,
    CHART_PALETTE, PHASE_1_COLOR, PHASE_2_COLOR, PHASE_3_COLOR,
    CONTENT_X, CONTENT_Y, CONTENT_Y_SUBTITLE, CONTENT_W, CONTENT_H,
    SOURCE_Y, SOURCE_H, TAKEAWAY_Y, TAKEAWAY_H, PHASE_Y, PHASE_H,
    FOOTER_Y, FOOTER_H, FOOTER_TEXT_Y,
    TITLE_X, TITLE_Y_NUMBERED, TITLE_Y_UNNUMBERED, TITLE_W, TITLE_H,
    SUB_X, SUB_W, SUB_H, SUB_GAP,
    AUTO_PROMOTE_EXECUTIVE,
    _palette_for, _rgb_tuple, _is_dark,
)
from ._ml_primitives import (
    rect, rounded, oval, circle, line, hairline, txt, _add_run,
    _freeform_poly, _emu, _apply_fill, _apply_border,
)
from ._ml_icons import draw_icon, ICON_REGISTRY
from ._ml_deck import (
    _new_slide, _intro_layout, _divider_layout,
    _cover_picture_layout, _populate_placeholder,
)
from ._ml_chrome import (
    apply_chrome, _phase_progress, _draw_card, _bulleted_list,
    stub_and_flag, footnotes_block, _section_marker,
    _top_chrome, _bottom_chrome, _gold_square_bullet,
    _tracked, _track_letters, _format_section_number, _pad_int,
    _render_action_title, _source_line,
    statement_card, in_slide_section,
    _takeaway_band, _superscript,
)
from ._ml_charts import (
    add_slope_chart, add_dot_plot, add_marimekko, add_waterfall,
    add_small_multiples, add_simple_bar, _render_chart,
)
from ._ml_helpers import _style_or_promote, _tone_color, _rag_color, _norm_key

# ===========================================================================
# Layout: CHART SLIDE
# ===========================================================================

def _auto_callout_for_chart(chart):
    """Generate a fallback callout for a chart when none was provided.

    Returns a single callout dict {x, y, label, direction} or None. Aimed at
    the first highlighted/relevant data point. The agent normally overrides
    this, but it ensures every data slide has at least one annotation.
    """
    if not chart:
        return None
    ctype = chart.get("type", "")
    data = chart.get("data", {})
    if ctype == "slope":
        items = data.get("items", [])
        highlights = data.get("highlight_indices") or []
        if not items or not highlights:
            return None
        idx = highlights[0]
        if idx >= len(items):
            return None
        item = items[idx]
        label = item[0] if isinstance(item, (list, tuple)) else item.get("label", "")
        try:
            before, after = item[1], item[2]
        except Exception:
            before, after = 0, 0
        delta = float(after) - float(before)
        direction = "decline" if delta < 0 else "gain"
        return {
            "x": Inches(10.4), "y": Inches(4.7),
            "label": f"Largest {direction} at: {label}",
            "direction": "up_left",
        }
    if ctype == "waterfall":
        bars = data.get("bars", [])
        if not bars:
            return None
        start_v, end_v = None, None
        for b in bars:
            t = b.get("type", "")
            if t == "start" and start_v is None:
                start_v = float(b.get("value", 0))
            if t == "end":
                end_v = float(b.get("value", 0))
        if start_v is not None and end_v is not None:
            delta = end_v - start_v
            sign = "+" if delta >= 0 else ""
            return {
                "x": Inches(6.5), "y": Inches(4.0),
                "label": f"Net change: {sign}{delta:g}",
                "direction": "up_right",
            }
        return None
    if ctype == "dot":
        items = data.get("items", [])
        if not items:
            return None
        # Top performer: max value.
        def _val(it):
            try:
                return float(str(it[1]).replace("%", "").strip())
            except Exception:
                return 0.0
        top = max(items, key=_val)
        label = top[0] if isinstance(top, (list, tuple)) else top.get("label", "")
        return {
            "x": Inches(10.4), "y": Inches(4.0),
            "label": f"Top performer: {label}",
            "direction": "up_left",
        }
    if ctype == "marimekko":
        cols = data.get("columns", [])
        if not cols:
            return None
        # Largest column by weight.
        top = max(cols, key=lambda c: c.get("weight", 0))
        return {
            "x": Inches(2.4), "y": Inches(6.0),
            "label": f"Largest column: {top.get('label', '')}",
            "direction": "down_left",
        }
    return None


def build_chart_slide(prs, meta, action_title=None, chart=None, takeaway=None,
                      source=None, callouts=None, category=None, subtitle=None,
                      style="merck_executive", page=None, total=None,
                      section_number=None, footnotes=None,
                      methodology_note=None, content=None):
    """Chart slide. chart={'type','data'}, optional callouts and footnotes.

    If callouts is None and the chart type benefits from a fallback callout,
    a single auto-callout is generated from the chart data.
    """
    if content:
        if "chart"    in content: chart    = content["chart"]
        if "callouts" in content: callouts = content["callouts"]
        if "takeaway" in content: takeaway = content["takeaway"]
    style = _style_or_promote(category, style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 subtitle=subtitle, takeaway=takeaway, source=source,
                 page=page, total=total, palette=style,
                 section_number=section_number,
                 methodology_note=methodology_note)

    chart_y = Inches(2.85) if subtitle else Inches(2.50)
    chart_h = Inches(6.50) - chart_y
    chart_rendered = _render_chart(slide, chart, Inches(0.65), chart_y,
                                   Inches(12.0), chart_h, style)
    if not chart_rendered:
        # Show a placeholder instead of a silently blank content area.
        txt(slide, Inches(0.65), chart_y + Inches(0.50), Inches(12.0),
            Inches(1.20),
            "[Chart data not available — populate chart.type and chart.data in the plan]",
            sz=13, color=INK_GRAY, italic=True, font=FONT_BODY)

    # Auto-fallback callout when none provided.
    effective_callouts = callouts
    if effective_callouts is None:
        auto = _auto_callout_for_chart(chart)
        if auto:
            effective_callouts = [auto]
    if effective_callouts:
        for c in effective_callouts:
            if c.get("x") is not None or c.get("y") is not None:
                # Slide-absolute EMU coordinates (from _auto_callout_for_chart or
                # hand-crafted plans that use the internal x/y keys directly).
                cx = c.get("x") or Inches(6.0)
                cy = c.get("y") or chart_y + Inches(1.5)
            else:
                # LLM plan schema: x_in/y_in are chart-relative float inches.
                # Add chart_y so the annotation lands inside the chart area, not
                # on top of the action title (which sits above chart_y).
                cx = Inches(float(c.get("x_in", 6.0)))
                cy = chart_y + Inches(float(c.get("y_in", 1.5)))
            stub_and_flag(slide, cx, cy,
                          c.get("label") or c.get("text", ""), style,
                          direction=c.get("direction", "up_right"))

    # Footnotes block above the source line (chrome already drew source/takeaway).
    if footnotes:
        footnotes_block(slide, footnotes, style)
    return slide



# ===========================================================================
# Layout: WATERFALL
# ===========================================================================

def build_waterfall_slide(prs, meta, action_title=None, bars=None, takeaway=None,
                          source=None, category=None, style="merck_executive",
                          page=None, total=None, section_number=None,
                          footnotes=None, methodology_note=None,
                          subtitle=None, content=None):
    if content:
        if "bars"      in content: bars      = content["bars"]
        if "takeaway"  in content: takeaway  = content["takeaway"]
        if "footnotes" in content: footnotes = content["footnotes"]
    style = _style_or_promote(category, style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 subtitle=subtitle,
                 takeaway=takeaway, source=source, page=page, total=total,
                 palette=style, section_number=section_number, methodology_note=methodology_note)
    chart_y = Inches(2.95) if subtitle else Inches(2.55)
    chart_h = Inches(6.50) - chart_y
    add_waterfall(slide, Inches(0.65), chart_y, Inches(12.0), chart_h,
                  bars or [], style)
    if footnotes:
        footnotes_block(slide, footnotes, style)
    return slide



# ===========================================================================
# Layout: HERO STAT
# ===========================================================================

def build_hero_stat(prs, meta, stat=None, context=None, source=None, category=None,
                    style="merck_storytelling", page=None, total=None,
                    section_number=None, methodology_note=None,
                    subtitle=None, action_title=None, content=None):
    """One huge stat with a one-line context. Defaults to storytelling.
    section_number, methodology_note, subtitle, and action_title are accepted
    for caller consistency but are ignored; hero_stat slides do not show a
    section marker, subtitle, or action title.
    """
    if content:
        if "stat"    in content: stat    = content["stat"]
        if "context" in content: context = content["context"]
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    _top_chrome(slide, meta, category, style, page=page, total=total)
    _bottom_chrome(slide, meta, category, page, total, style)

    stat = stat or {}
    value = stat.get("value", "")
    label = stat.get("label", "")

    dark = _is_dark(style)
    number_color = pal["hot"] if dark else MERCK_PURPLE
    label_color = WHITE if dark else INK_DARK
    context_color = PANEL_LIGHT if dark else INK_GRAY

    # Huge number centered.
    txt(slide, Inches(0.65), Inches(2.40), Inches(12.0), Inches(2.20),
        value, sz=110, color=number_color, bold=True, font=FONT_HEAD,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # One-line context underneath.
    txt(slide, Inches(0.65), Inches(4.75), Inches(12.0), Inches(0.50),
        label, sz=18, color=label_color, font=FONT_BODY,
        align=PP_ALIGN.CENTER)
    if context:
        txt(slide, Inches(0.65), Inches(5.25), Inches(12.0), Inches(1.30),
            context, sz=13, color=context_color, italic=True,
            font=FONT_BODY, align=PP_ALIGN.CENTER)
    # Theme-accent rule under the number.
    rule_color = pal["highlight"]
    hairline(slide,
             Inches(6.665) - Inches(1.2), Inches(4.55),
             Inches(2.4), Emu(int(Pt(2.5))), rule_color)
    return slide



# ===========================================================================
# Layout: STAT STRIP (Recall 3-up / 4-up stat cards)
# ===========================================================================

def build_stat_strip(prs, meta, action_title=None, stats=None, takeaway=None,
                     source=None, category=None, subtitle=None,
                     style="merck_executive", page=None, total=None,
                     section_number=None, methodology_note=None, content=None):
    """Recall-style 3-up or 4-up stat strip with cream cards.

    stats: list of 3 or 4 dicts: {value, label, body}.
      - value: hero number string in Merck Web 56pt MERCK_PURPLE bold.
      - label: SMALL UPPERCASE TRACKED label in MERCK_GOLD Verdana 10pt.
      - body: body copy in INK_GRAY Verdana 11pt.

    All cards in the row share the same height (equal-height enforcement).
    """
    if content:
        if "stats"    in content: stats    = content["stats"]
        if "takeaway" in content: takeaway = content["takeaway"]
    style = _style_or_promote(category, style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 subtitle=subtitle, takeaway=takeaway, source=source,
                 page=page, total=total, palette=style,
                 section_number=section_number, methodology_note=methodology_note)

    cards = [s for s in (stats or []) if s]
    if not cards:
        return slide
    # Clamp to a maximum of 4 cards.
    if len(cards) > 4:
        cards = cards[:4]
    n = len(cards)

    zone_top = Inches(2.95) if subtitle else Inches(2.50)
    card_h = Inches(2.40)
    gutter = Inches(0.18)
    total_w = CONTENT_W
    card_w = (total_w - gutter * (n - 1)) / n

    for i, s in enumerate(cards):
        cx = CONTENT_X + i * (card_w + gutter)
        cy = zone_top
        # Cream card with theme-accent top stripe + hairline border.
        body = rounded(slide, cx, cy, card_w, card_h, fill=PANEL_LIGHT)
        _apply_border(body, LIGHT_GRAY, Pt(0.5))
        rounded(slide, cx, cy, card_w, Inches(0.06), fill=pal["highlight"], adj=50000)

        pad = Inches(0.22)
        # Big hero number — font size scales with string length to prevent overlap.
        value = s.get("value", "")
        val_str = str(value)
        if len(val_str) <= 5:
            val_sz, num_h = 52, Inches(1.05)
        elif len(val_str) <= 8:
            val_sz, num_h = 40, Inches(0.85)
        else:
            val_sz, num_h = 32, Inches(0.72)
        num_y = cy + Inches(0.30)
        txt(slide, cx + pad, num_y, card_w - pad * 2, num_h,
            val_str, sz=val_sz, color=pal["accent"], bold=True,
            font=FONT_HEAD, anchor=MSO_ANCHOR.TOP)

        # Tracked uppercase label in theme accent.
        # Tracking is skipped for labels > 12 chars to prevent mid-word wrapping.
        label = s.get("label", "")
        label_y = num_y + num_h + Inches(0.06)
        label_h = Inches(0.40)
        label_text = _track_letters(label) if len(label) <= 12 else label.upper()
        txt(slide, cx + pad, label_y, card_w - pad * 2, label_h,
            label_text, sz=10, color=pal["highlight"], bold=True,
            font=FONT_BODY, anchor=MSO_ANCHOR.TOP)

        # Body in INK_GRAY Verdana 11pt.
        body_text = s.get("body", "")
        body_y = label_y + Inches(0.46)
        body_h = (cy + card_h) - body_y - Inches(0.18)
        if body_h < Inches(0.30):
            body_h = Inches(0.30)
        txt(slide, cx + pad, body_y, card_w - pad * 2, body_h,
            str(body_text), sz=11, color=INK_GRAY, font=FONT_BODY)

    return slide



# ===========================================================================
# Layout: DONUT / RING CHART
# ===========================================================================

def build_donut_chart(prs, meta, action_title=None, segments=None,
                      center_label=None, center_value=None,
                      legend_title=None, takeaway="", source=None,
                      subtitle=None, methodology_note=None,
                      style="merck_executive", page=None, total=None,
                      section_number=None, category=None):
    """Donut ring chart (native pptx DOUGHNUT) with legend.

    segments:     list of dicts — label (str), value (float), color (tuple opt)
    center_value: str — big number shown in the donut hole
    center_label: str — small descriptor below center_value
    legend_title: str — optional legend heading
    """
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    segments = list(segments or [])[:8]
    if not segments:
        return slide

    dark = _is_dark(style)
    bg = pal["bg"]
    DEFAULT_COLORS = [MERCK_PURPLE, MERCK_BLUE, MERCK_GOLD, MERCK_AQUA,
                      GOOD_GREEN, BAD_RED, PURPLE_MUTED, MERCK_YELLOW]
    total_val = sum(float(s.get("value", 0)) for s in segments) or 1.0

    from pptx.chart.data import CategoryChartData as _CCD
    chart_data = _CCD()
    chart_data.categories = [str(s.get("label", "")) for s in segments]
    chart_data.add_series("", [float(s.get("value", 0)) for s in segments])

    outer_r = Inches(1.90)
    cx, cy  = Inches(4.20), Inches(3.85)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        _emu(cx - outer_r), _emu(cy - outer_r),
        _emu(outer_r * 2),  _emu(outer_r * 2),
        chart_data,
    )
    ch = chart_shape.chart
    ch.has_legend = False
    ch.has_title  = False

    series = ch.series[0]
    for i, s in enumerate(segments):
        col = s.get("color") or DEFAULT_COLORS[i % len(DEFAULT_COLORS)]
        pt  = series.points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = _rgb_tuple(col)

    try:
        from lxml import etree as _et
        plot_el = ch.plots[0]._element
        for el in plot_el.iter():
            if el.tag.endswith("}holeSize"):
                el.set("val", "55")
                break
        else:
            ns = "http://schemas.openxmlformats.org/drawingml/2006/chart"
            _et.SubElement(plot_el, f"{{{ns}}}holeSize").set("val", "55")
    except Exception:
        pass

    inner_r = Inches(1.00)
    circle(slide, cx - inner_r, cy - inner_r, inner_r * 2, fill=bg)

    val_color = MERCK_YELLOW if dark else MERCK_PURPLE
    lab_color = PANEL_LIGHT  if dark else INK_GRAY
    if center_value:
        txt(slide, cx - inner_r, cy - Inches(0.46), inner_r * 2, Inches(0.72),
            str(center_value), sz=32, color=val_color, bold=True,
            font=FONT_HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if center_label:
        txt(slide, cx - inner_r, cy + Inches(0.20), inner_r * 2, Inches(0.30),
            str(center_label), sz=10, color=lab_color,
            font=FONT_BODY, align=PP_ALIGN.CENTER)

    leg_x, leg_y = Inches(8.60), Inches(2.10)
    leg_item_h   = Inches(0.44)
    dot_sz       = Inches(0.18)
    leg_color    = PANEL_LIGHT if dark else INK_DARK

    if legend_title:
        txt(slide, leg_x, leg_y - Inches(0.36), Inches(4.5), Inches(0.28),
            str(legend_title).upper(), sz=9, color=MERCK_GOLD, bold=True,
            font=FONT_BODY)

    for i, s in enumerate(segments):
        col = s.get("color") or DEFAULT_COLORS[i % len(DEFAULT_COLORS)]
        iy  = leg_y + i * leg_item_h
        circle(slide, leg_x, iy + (leg_item_h - dot_sz) / 2, dot_sz, fill=col)
        pct = f"{float(s.get('value', 0)) / total_val * 100:.0f}%"
        txt(slide, leg_x + dot_sz + Inches(0.12), iy,
            Inches(4.0), leg_item_h,
            f"{pct}  {s.get('label', '')}",
            sz=11, color=leg_color, font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE)

    return slide



# ===========================================================================
# Layout: RISK HEAT MAP
# ===========================================================================

def build_risk_heatmap(prs, meta, action_title=None, risks=None,
                       x_label="LIKELIHOOD", y_label="IMPACT",
                       takeaway="", source=None, subtitle=None,
                       methodology_note=None, style="merck_executive",
                       page=None, total=None, section_number=None,
                       category=None):
    """Probability x Impact heat map with numbered risk dots and legend.

    risks: list of dicts with keys:
        label      (str) — short risk name
        likelihood (int 1-5) — x-axis position
        impact     (int 1-5) — y-axis position
        color      (tuple, opt) — dot color override
    """
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    risks = list(risks or [])[:15]
    dark  = _is_dark(style)

    # grid_y must clear the action title chrome (single-line ~1.50, two-line
    # can reach ~1.90).  Use the same subtitle-aware offset other layouts use.
    GRID_N  = 5
    grid_x  = Inches(1.60)
    grid_y  = Inches(2.55) if not subtitle else Inches(2.95)
    grid_w  = Inches(6.20)
    # Bottom boundary: leave room for takeaway bar (~0.55) + source (~0.35).
    grid_h  = Inches(6.50) - grid_y - Inches(0.90)
    cell_w  = grid_w / GRID_N
    cell_h  = grid_h / GRID_N

    def _heat(row, col):
        s = row + col
        if s <= 3: return (0xC8, 0xE6, 0xC9)
        if s <= 5: return (0xFF, 0xF1, 0x7F)
        if s <= 7: return (0xFF, 0xCC, 0x66)
        return           (0xFF, 0x99, 0x99)

    for row in range(GRID_N):
        for col in range(GRID_N):
            cx = grid_x + col * cell_w
            cy = grid_y + (GRID_N - 1 - row) * cell_h
            rect(slide, cx, cy, cell_w, cell_h,
                 fill=_heat(row, col), border=WHITE, border_w=Pt(0.5))

    ax_col   = PANEL_LIGHT if dark else INK_GRAY
    ax_lbls  = ["Very Low", "Low", "Medium", "High", "Very High"]
    for j in range(GRID_N):
        txt(slide, grid_x + j * cell_w, grid_y + grid_h + Inches(0.04),
            cell_w, Inches(0.24), ax_lbls[j],
            sz=8, color=ax_col, font=FONT_BODY, align=PP_ALIGN.CENTER)
        txt(slide, Inches(0.10), grid_y + (GRID_N - 1 - j) * cell_h,
            Inches(1.40), cell_h, ax_lbls[j],
            sz=8, color=ax_col, font=FONT_BODY,
            align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    txt(slide, grid_x, grid_y + grid_h + Inches(0.30), grid_w, Inches(0.24),
        x_label.upper(), sz=9, color=ax_col, bold=True,
        font=FONT_BODY, align=PP_ALIGN.CENTER)

    RISK_COLORS = [MERCK_PURPLE, MERCK_BLUE, BAD_RED, GOOD_GREEN, MERCK_GOLD,
                   PURPLE_MUTED, MERCK_AQUA, INK_DARK, MERCK_YELLOW,
                   MERCK_BLUE, PURPLE_MUTED, GOOD_GREEN, BAD_RED, MERCK_GOLD, MERCK_PURPLE]
    dot_r = Inches(0.18)
    leg_x = Inches(8.30)
    leg_y = Inches(1.78)
    leg_h = Inches(0.34)
    txt(slide, leg_x, leg_y - Inches(0.28), Inches(4.8), Inches(0.24),
        "RISK REGISTER", sz=9,
        color=MERCK_GOLD if dark else MERCK_PURPLE, bold=True, font=FONT_BODY)

    for i, risk in enumerate(risks):
        col = risk.get("color") or RISK_COLORS[i % len(RISK_COLORS)]
        try:
            lk = max(1, min(5, int(float(risk.get("likelihood", 3)))))
        except (TypeError, ValueError):
            lk = 3
        try:
            im = max(1, min(5, int(float(risk.get("impact", 3)))))
        except (TypeError, ValueError):
            im = 3
        dx  = grid_x + (lk - 1) * cell_w + cell_w / 2
        dy  = grid_y + (GRID_N - im) * cell_h + cell_h / 2
        circle(slide, dx - dot_r, dy - dot_r, dot_r * 2, fill=col)
        txt(slide, dx - dot_r, dy - dot_r, dot_r * 2, dot_r * 2,
            str(i + 1), sz=8, color=WHITE, bold=True, font=FONT_BODY,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        ry = leg_y + i * (leg_h + Inches(0.04))
        circle(slide, leg_x, ry + (leg_h - dot_r * 2) / 2, dot_r * 2, fill=col)
        txt(slide, leg_x + dot_r * 2 + Inches(0.10), ry,
            Inches(4.0), leg_h,
            f"{i + 1}.  {risk.get('label', '')}",
            sz=9, color=PANEL_LIGHT if dark else INK_DARK,
            font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE)

    return slide


# ===========================================================================
# Layout: RADAR / SPIDER CHART (hand-drawn)
# ===========================================================================

def build_radar_chart(prs, meta, action_title=None, axes=None,
                      series=None, takeaway="", source=None, subtitle=None,
                      methodology_note=None, style="merck_executive",
                      page=None, total=None, section_number=None,
                      category=None):
    """Multi-axis spider chart drawn with polygon lines.

    axes:   list of str — axis labels (4-8 axes)
    series: list of dicts with keys:
        label  (str) — series name
        values (list of float 0-100) — one per axis
        color  (tuple, opt)
    """
    import math as _math
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    axes   = list(axes   or [])[:8]
    series = list(series or [])[:4]
    n_axes = len(axes)
    if n_axes < 3:
        return slide

    dark = _is_dark(style)
    cx, cy   = Inches(5.00), Inches(3.85)
    R        = Inches(2.20)
    N_RINGS  = 4
    S_COLORS = [MERCK_PURPLE, MERCK_GOLD, MERCK_BLUE, GOOD_GREEN]

    def _pt(axis_i, frac):
        angle = _math.radians(-90 + 360 / n_axes * axis_i)
        return (int(cx + frac * R * _math.cos(angle)),
                int(cy + frac * R * _math.sin(angle)))

    for ring in range(1, N_RINGS + 1):
        pts = [_pt(j, ring / N_RINGS) for j in range(n_axes)]
        for k in range(n_axes):
            line(slide, pts[k][0], pts[k][1],
                 pts[(k + 1) % n_axes][0], pts[(k + 1) % n_axes][1],
                 color=LIGHT_GRAY, weight=Pt(0.5))

    for j in range(n_axes):
        x2, y2 = _pt(j, 1.0)
        line(slide, int(cx), int(cy), x2, y2, color=LIGHT_GRAY, weight=Pt(0.75))

    for j, lbl in enumerate(axes):
        ax, ay = _pt(j, 1.20)
        txt(slide, ax - Inches(0.70), ay - Inches(0.20),
            Inches(1.40), Inches(0.42), lbl,
            sz=9, color=PANEL_LIGHT if dark else INK_DARK,
            font=FONT_BODY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    for si, ser in enumerate(series):
        vals = [float(v) for v in ser.get("values", [])]
        if len(vals) < n_axes:
            vals += [0.0] * (n_axes - len(vals))
        col = ser.get("color") or S_COLORS[si % len(S_COLORS)]
        pts = [_pt(j, max(0.0, min(1.0, vals[j] / 100.0)))
               for j in range(n_axes)]
        _freeform_poly(slide, pts, fill=col, border=col, border_w=Pt(1.5))

    leg_x, leg_y, leg_h = Inches(10.60), Inches(2.30), Inches(0.36)
    for si, ser in enumerate(series):
        col = ser.get("color") or S_COLORS[si % len(S_COLORS)]
        ry  = leg_y + si * (leg_h + Inches(0.10))
        rect(slide, leg_x, ry + Inches(0.10),
             Inches(0.20), Inches(0.16), fill=col)
        txt(slide, leg_x + Inches(0.28), ry, Inches(2.30), leg_h,
            str(ser.get("label", "")),
            sz=9, color=PANEL_LIGHT if dark else INK_DARK,
            font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE)

    return slide


