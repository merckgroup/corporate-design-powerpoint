"""Merck Layouts v3: a premium python-pptx slide layout library.

This module is the rendering engine called by an LLM agent that authors decks.
It produces consistent, premium-feel 16:9 slides across three locked styles:
merck_executive, merck_corporate, merck_storytelling.

v3 ports the visual rhythm and chrome of the Recall Vision reference deck onto
the Merck LE rich purple brand: rich purple primary, gold and yellow accents,
cream PANEL_LIGHT card panels, Merck Web for hero/serif type, Verdana for
body and chrome. Bug fixes for waterfall (floating bars) and slope chart
(anti-collision labels) are baked in.

Public API (preserved):

  Deck lifecycle:
    open_deck(base_path=None)
    save_deck(prs, output_path)

  Palette helpers:
    PALETTES
    rgb(palette_name, key)

  Primitives:
    rect, rounded, oval, circle, line, hairline, txt

  Vector icon system:
    icon_chart_bar, icon_chart_line, icon_chart_pie,
    icon_arrow_up, icon_arrow_down, icon_arrow_right,
    icon_check, icon_x, icon_alert, icon_info, icon_target,
    icon_gear, icon_users, icon_calendar, icon_clock,
    icon_lightbulb, icon_lock, icon_globe, icon_search,
    icon_money, icon_trending_up, icon_trending_down,
    icon_shield, icon_flag, icon_doc
    draw_icon(slide, name, x, y, size, color)

  Universal chrome:
    merck_stub, category_tag, action_title, subtitle_line,
    takeaway_band, source_line, page_number, apply_chrome,
    footnotes_block

  Craft helpers:
    stub_and_flag, decimal_align

  Premium chart helpers:
    add_slope_chart, add_dot_plot, add_marimekko,
    add_waterfall, add_small_multiples, add_simple_bar

  Layout functions:
    build_cover, build_exec_summary, build_agenda,
    build_section_divider, build_chart_slide, build_two_column,
    build_three_column, build_2x2_matrix, build_phase_process,
    build_vertical_numbered, build_waterfall_slide,
    build_decision_rows, build_gantt, build_hero_stat,
    build_close, build_stat_strip, build_before_after
"""

from __future__ import annotations

import os
from typing import Optional

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt


# ===========================================================================
# Fonts and palette
# ===========================================================================

FONT_HEAD = "Merck Web"
FONT_BODY = "Verdana"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Merck Corporate Design official palette (exact hex from brand guidelines).
# Full set:  violet, blue, green, red, pink, lightblue, lightgreen, yellow,
#            palepink, paleblue, palegreen, paleyellow, grey
# ---------------------------------------------------------------------------
MERCK_PURPLE  = (0x50, 0x32, 0x91)   # #503291 — violet     PRIMARY
MERCK_BLUE    = (0x0F, 0x69, 0xAF)   # #0F69AF — blue       SECONDARY
MERCK_GOLD    = (0xEB, 0x3C, 0x96)   # #EB3C96 — pink       TERTIARY (was MERCK_GOLD)
PURPLE_DEEP   = (0x3A, 0x24, 0x68)   # #3A2468 — footer / dark backgrounds
PURPLE_MUTED  = (0x7D, 0x74, 0xA0)   # #7D74A0 — muted / separator
MERCK_YELLOW  = (0xFF, 0xC8, 0x32)   # #FFC832 — yellow
MERCK_AQUA    = (0x96, 0xD7, 0xD2)   # #96D7D2 — paleblue   (sensitive-blue)
LIGHT_GRAY    = (0xE0, 0xE0, 0xE0)   # #E0E0E0 — rules / borders
PANEL_LIGHT   = (0xF4, 0xF2, 0xF8)   # #F4F2F8 — card panel background
WHITE         = (0xFF, 0xFF, 0xFF)
INK_DARK      = (0x1A, 0x16, 0x26)   # #1A1626 — primary body text
INK_GRAY      = (0x55, 0x5D, 0x6E)   # #555D6E — secondary text / sources
BAD_RED       = (0xE6, 0x1E, 0x50)   # #E61E50 — red
GOOD_GREEN    = (0x14, 0x9B, 0x5F)   # #149B5F — green

# Full official Merck Corporate Design colour names (for theme palettes).
_MC_VIOLET     = MERCK_PURPLE          # #503291
_MC_BLUE       = MERCK_BLUE            # #0F69AF
_MC_GREEN      = GOOD_GREEN            # #149B5F
_MC_RED        = BAD_RED               # #E61E50
_MC_PINK       = MERCK_GOLD            # #EB3C96
_MC_LIGHTBLUE  = (0x2D, 0xBE, 0xCD)   # #2DBECD
_MC_LIGHTGREEN = (0xA5, 0xCD, 0x50)   # #A5CD50  (= OP_LIME)
_MC_YELLOW     = MERCK_YELLOW          # #FFC832
_MC_PALEPINK   = (0xE1, 0xC3, 0xCD)   # #E1C3CD
_MC_PALEBLUE   = MERCK_AQUA            # #96D7D2
_MC_PALEGREEN  = (0xB4, 0xDC, 0x96)   # #B4DC96
_MC_PALEYELLOW = (0xFF, 0xDC, 0xB9)   # #FFDCB9
_MC_GREY       = (0x99, 0x99, 0x99)   # #999999

# Liquid Carbon LS dashboard data-series palette (from LS dashboard legend).
# Use these for chart series, RAG indicators, and data comparisons.
ACT_PURPLE   = MERCK_PURPLE          # #503291 — Actual (same as primary)
LY_CYAN      = (0x2D, 0xBE, 0xCD)   # #2DBECD — Last Year
OP_LIME      = (0xA5, 0xCD, 0x50)   # #A5CD50 — Operating Plan
FC_PINK      = (0xEB, 0x3C, 0x96)   # #EB3C96 — Forecast
DEV_POS_BLUE = (0x0F, 0x69, 0xAF)   # #0F69AF — Deviation positive
DEV_NEG_RED  = BAD_RED               # #E61E50 — Deviation negative (same as BAD_RED)

# Liquid Carbon 12-color chart palette (resolved from tailwind-colors.json).
# chart-1 through chart-12 match the LC setup.json chart token order.
# chart-8 (sensitive-blue) and chart-10 (sensitive-yellow) are custom tokens
# approximated from the brand palette.
CHART_1  = (0xDE, 0x7A, 0x21)   # #DE7A21 — orange-600
CHART_2  = (0x12, 0x87, 0x9D)   # #12879D — teal-600
CHART_3  = (0x18, 0x42, 0x51)   # #184251 — cyan-900
CHART_4  = (0xFE, 0xC7, 0x31)   # #FEC731 — amber-400
CHART_5  = (0xF7, 0xA2, 0x16)   # #F7A216 — amber-500
CHART_6  = (0x0E, 0x69, 0xAF)   # #0E69AF — blue-700
CHART_7  = (0xE6, 0x1E, 0x50)   # #E61E50 — red-600 (same as BAD_RED)
CHART_8  = (0x96, 0xD7, 0xD2)   # #96D7D2 — sensitive-blue (exact from vibrant-m)
CHART_9  = (0x14, 0x9B, 0x5F)   # #149B5F — green-500
CHART_10 = (0xFF, 0xDC, 0xB9)   # #FFDCB9 — sensitive-yellow (exact from vibrant-m)
CHART_11 = (0x9D, 0x80, 0xE6)   # #9D80E6 — violet-400
CHART_12 = (0x8C, 0x22, 0x35)   # #8C2235 — red-800

# Ordered list for cycling through chart series automatically.
CHART_PALETTE = [
    CHART_1, CHART_2, CHART_3, CHART_4,
    CHART_5, CHART_6, CHART_7, CHART_8,
    CHART_9, CHART_10, CHART_11, CHART_12,
]

PHASE_1_COLOR = MERCK_PURPLE
PHASE_2_COLOR = MERCK_AQUA
PHASE_3_COLOR = MERCK_GOLD


PALETTES = {
    "merck_executive": {
        "bg":        WHITE,
        "ink":       INK_DARK,
        "ink_2":     INK_GRAY,
        "ink_3":     LIGHT_GRAY,
        "accent":    MERCK_PURPLE,
        "accent_2":  PURPLE_DEEP,
        "accent_3":  LY_CYAN,
        "highlight": MERCK_GOLD,
        "hot":       MERCK_YELLOW,
        "rule":      LIGHT_GRAY,
        "panel":     PANEL_LIGHT,
        "muted":     PURPLE_MUTED,
        "good":      GOOD_GREEN,
        "warn":      MERCK_YELLOW,
        "bad":       BAD_RED,
        "lime":      OP_LIME,
    },
    "merck_corporate": {
        "bg":        WHITE,
        "ink":       INK_DARK,
        "ink_2":     INK_GRAY,
        "ink_3":     LIGHT_GRAY,
        "accent":    MERCK_PURPLE,
        "accent_2":  PURPLE_DEEP,
        "accent_3":  LY_CYAN,
        "highlight": MERCK_GOLD,
        "hot":       MERCK_YELLOW,
        "rule":      LIGHT_GRAY,
        "panel":     PANEL_LIGHT,
        "muted":     PURPLE_MUTED,
        "good":      GOOD_GREEN,
        "warn":      MERCK_YELLOW,
        "bad":       BAD_RED,
        "lime":      OP_LIME,
    },
    "merck_storytelling": {
        "bg":        MERCK_PURPLE,
        "ink":       WHITE,
        "ink_2":     PANEL_LIGHT,
        "ink_3":     PURPLE_MUTED,
        "accent":    MERCK_GOLD,
        "accent_2":  MERCK_YELLOW,
        "accent_3":  LY_CYAN,
        "highlight": MERCK_GOLD,
        "hot":       MERCK_YELLOW,
        "rule":      PURPLE_MUTED,
        "panel":     PANEL_LIGHT,
        "muted":     PURPLE_MUTED,
        "good":      GOOD_GREEN,
        "warn":      MERCK_YELLOW,
        "bad":       BAD_RED,
        "lime":      OP_LIME,
    },
    # ------------------------------------------------------------------
    # Official Merck Corporate Design themes (6 variants).
    # Theme names match the empower library folder names exactly.
    # Light themes use WHITE content bg; dark themes use MERCK_PURPLE.
    # ------------------------------------------------------------------
    "functional": {       # lightgreen bg cover, teal accent, organic cells
        "bg":        WHITE,
        "ink":       INK_DARK,
        "ink_2":     INK_GRAY,
        "ink_3":     LIGHT_GRAY,
        "accent":    _MC_LIGHTBLUE,
        "accent_2":  MERCK_PURPLE,
        "accent_3":  _MC_LIGHTGREEN,
        "highlight": _MC_LIGHTBLUE,
        "hot":       _MC_LIGHTGREEN,
        "rule":      LIGHT_GRAY,
        "panel":     (0xF0, 0xF9, 0xFA),    # very pale teal
        "muted":     _MC_PALEBLUE,
        "good":      GOOD_GREEN,
        "warn":      MERCK_YELLOW,
        "bad":       BAD_RED,
        "lime":      _MC_LIGHTGREEN,
    },
    "organic": {          # paleyellow bg cover, red accent, warm life-science
        "bg":        WHITE,
        "ink":       INK_DARK,
        "ink_2":     INK_GRAY,
        "ink_3":     LIGHT_GRAY,
        "accent":    _MC_RED,
        "accent_2":  MERCK_PURPLE,
        "accent_3":  _MC_PALEYELLOW,
        "highlight": _MC_RED,
        "hot":       _MC_PALEYELLOW,
        "rule":      LIGHT_GRAY,
        "panel":     (0xFD, 0xF5, 0xED),    # very pale warm
        "muted":     _MC_PALEPINK,
        "good":      GOOD_GREEN,
        "warn":      MERCK_YELLOW,
        "bad":       _MC_RED,
        "lime":      _MC_LIGHTGREEN,
    },
    "plastic": {          # lightgreen bg cover, pink accent — current default
        "bg":        WHITE,
        "ink":       INK_DARK,
        "ink_2":     INK_GRAY,
        "ink_3":     LIGHT_GRAY,
        "accent":    MERCK_PURPLE,
        "accent_2":  PURPLE_DEEP,
        "accent_3":  LY_CYAN,
        "highlight": _MC_PINK,
        "hot":       _MC_LIGHTGREEN,
        "rule":      LIGHT_GRAY,
        "panel":     PANEL_LIGHT,
        "muted":     PURPLE_MUTED,
        "good":      GOOD_GREEN,
        "warn":      MERCK_YELLOW,
        "bad":       BAD_RED,
        "lime":      OP_LIME,
    },
    "synthetic": {        # dark violet bg, yellow accent, industrial
        "bg":        MERCK_PURPLE,
        "ink":       WHITE,
        "ink_2":     PANEL_LIGHT,
        "ink_3":     PURPLE_MUTED,
        "accent":    _MC_YELLOW,
        "accent_2":  _MC_LIGHTBLUE,
        "accent_3":  _MC_LIGHTGREEN,
        "highlight": _MC_YELLOW,
        "hot":       _MC_YELLOW,      # yellow: section number + takeaway band
        "rule":      PURPLE_MUTED,
        "panel":     (0x3F, 0x28, 0x70),    # dark panel (lighter than bg)
        "muted":     PURPLE_MUTED,
        "good":      GOOD_GREEN,
        "warn":      _MC_YELLOW,
        "bad":       BAD_RED,
        "lime":      _MC_LIGHTGREEN,
    },
    "technical": {        # paleyellow bg, teal/lightblue, angular shapes
        "bg":        WHITE,
        "ink":       INK_DARK,
        "ink_2":     INK_GRAY,
        "ink_3":     LIGHT_GRAY,
        "accent":    _MC_LIGHTBLUE,
        "accent_2":  MERCK_PURPLE,
        "accent_3":  _MC_PALEYELLOW,
        "highlight": _MC_LIGHTBLUE,
        "hot":       _MC_PALEYELLOW,
        "rule":      LIGHT_GRAY,
        "panel":     (0xF0, 0xF9, 0xFA),    # very pale teal (same as functional)
        "muted":     _MC_PALEBLUE,
        "good":      GOOD_GREEN,
        "warn":      MERCK_YELLOW,
        "bad":       BAD_RED,
        "lime":      _MC_LIGHTGREEN,
    },
    "electronics": {      # dark violet bg, yellow accent, photo placeholder on cover
        "bg":        MERCK_PURPLE,
        "ink":       WHITE,
        "ink_2":     PANEL_LIGHT,
        "ink_3":     PURPLE_MUTED,
        "accent":    _MC_YELLOW,
        "accent_2":  _MC_LIGHTBLUE,
        "accent_3":  _MC_LIGHTGREEN,
        "highlight": _MC_YELLOW,
        "hot":       _MC_YELLOW,      # yellow: section number + takeaway band
        "rule":      PURPLE_MUTED,
        "panel":     (0x3F, 0x28, 0x70),
        "muted":     PURPLE_MUTED,
        "good":      GOOD_GREEN,
        "warn":      _MC_YELLOW,
        "bad":       BAD_RED,
        "lime":      _MC_LIGHTGREEN,
    },
}

AUTO_PROMOTE_EXECUTIVE = {
    "Executive Summary",
    "Recommendation",
    "Decision Request",
    "Risk",
    "Tradeoff",
}


# ===========================================================================
# Layout constants
# ===========================================================================

# Top chrome (universal)
GOLD_RULE_X = Inches(0.65)
GOLD_RULE_Y = Inches(0.30)
GOLD_RULE_W = Inches(3.5)
GOLD_RULE_H = Inches(0.04)

CLASS_BADGE_X = Inches(11.5)
CLASS_BADGE_Y = Inches(0.20)
CLASS_BADGE_W = Inches(1.7)

BREADCRUMB_X = Inches(0.65)
BREADCRUMB_Y = Inches(0.80)
BREADCRUMB_W = Inches(9.0)
BREADCRUMB_H = Inches(0.30)

# Section marker zone (numbered content slides). Top breadcrumb removed,
# so positions shifted up by 0.90" to reclaim that vertical space.
SECTION_CIRCLE_X = Inches(0.65)
SECTION_CIRCLE_Y = Inches(0.40)
SECTION_CIRCLE_D = Inches(0.55)

SECTION_TAG_X = Inches(1.35)
SECTION_TAG_Y = Inches(0.52)
SECTION_TAG_W = Inches(8.0)
SECTION_TAG_H = Inches(0.30)

# Action title zone (32pt bold serif takes ~0.55in/line; allow 2 lines + slack)
TITLE_X = Inches(0.65)
TITLE_Y_NUMBERED = Inches(1.15)
TITLE_Y_UNNUMBERED = Inches(0.45)
TITLE_W = Inches(12.0)
TITLE_H = Inches(1.30)

# Subtitle zone (placed beneath the title block)
SUB_X = Inches(0.65)
SUB_W = Inches(12.0)
SUB_H = Inches(0.40)
SUB_GAP = Inches(1.30)   # distance from title_y to subtitle_y

# Content zone (between title block and source/takeaway/footer chrome)
CONTENT_X = Inches(0.65)
CONTENT_Y = Inches(2.55)
CONTENT_Y_SUBTITLE = Inches(2.95)
CONTENT_W = Inches(12.0)
CONTENT_H = Inches(4.00)

# Below-content chrome positions (above the footer band)
SOURCE_Y    = Inches(6.55)
SOURCE_H    = Inches(0.22)
TAKEAWAY_Y  = Inches(6.83)
TAKEAWAY_H  = Inches(0.22)

# Phase progress (above footer band)
PHASE_Y = Inches(6.25)
PHASE_H = Inches(0.55)

# Footer band. Reduced from 0.40" to 0.30" height (top dropped from 7.10
# to 7.20) so the band no longer overlaps the master slide's Merck "M"
# logo, which spans y=6.73 to y=7.15. Band now sits cleanly below the
# logo at 7.20 → 7.50. Text inside re-centered to 7.24.
FOOTER_Y = Inches(7.20)
FOOTER_H = Inches(0.30)
FOOTER_TEXT_Y = Inches(7.24)


# ===========================================================================
# Color helpers
# ===========================================================================

def _rgb_tuple(t) -> RGBColor:
    if isinstance(t, RGBColor):
        return t
    return RGBColor(*t)


def _palette_for(style: str) -> dict:
    return PALETTES.get(style, PALETTES["merck_executive"])


_DARK_STYLES = frozenset({"merck_storytelling", "synthetic", "electronics"})


def _is_dark(style: str) -> bool:
    return style in _DARK_STYLES


def rgb(palette_name: str, key: str) -> RGBColor:
    """Return an RGBColor from a named palette and key. Falls back to ink."""
    pal = PALETTES.get(palette_name, PALETTES["merck_executive"])
    tup = pal.get(key) or pal.get("ink") or (0, 0, 0)
    return RGBColor(*tup)


# ===========================================================================
# Deck lifecycle
# ===========================================================================

def open_deck(base_path: Optional[str] = None) -> Presentation:
    """Open a Presentation. Reuses the Merck themed base pptx if it can be
    found at the given path or any common location; else falls back to a
    default python-pptx deck and warns to stderr.

    Common search paths tried when base_path is missing or not found:
      - the path passed in
      - Merck_Themed_Base_v1.pptx in the current working directory
      - /mnt/data/Merck_Themed_Base_v1.pptx (MyGPT sandbox upload location)
      - the directory containing merck_layouts.py
    """
    candidates = []
    if base_path:
        candidates.append(base_path)
    # Common fallback locations so a wrong-path call still finds the base.
    fname = "Merck_Themed_Base_v1.pptx"
    candidates.extend([
        fname,
        os.path.join("/mnt/data", fname),
        os.path.join(os.path.dirname(os.path.abspath(__file__)), fname),
    ])

    prs = None
    used = None
    for cand in candidates:
        if not cand:
            continue
        if os.path.isfile(cand):
            try:
                prs = Presentation(cand)
                used = cand
                # Strip pre-existing slides so the base provides only theme
                # and layouts; the caller adds slides from scratch.
                sldIdLst = prs.slides._sldIdLst  # noqa
                for sldId in list(sldIdLst):
                    rId = sldId.get(qn("r:id"))
                    try:
                        prs.part.drop_rel(rId)
                    except Exception:
                        pass
                    sldIdLst.remove(sldId)
                break
            except Exception:
                prs = None
                continue

    if prs is None:
        import sys
        print("WARNING: Merck_Themed_Base_v1.pptx not found in any search "
              "path. Falling back to default python-pptx layouts. Slides may "
              "look generic and PowerPoint may show a 'repairs' warning on "
              "open. Place the base pptx next to merck_layouts.py or pass "
              "the correct path to open_deck().",
              file=sys.stderr)
        prs = Presentation()

    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def save_deck(prs: Presentation, output_path: str) -> str:
    """Persist the deck and return the path."""
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    return output_path


def _blank_layout(prs: Presentation):
    """Pick the most placeholder-free layout we can find.

    Order of preference:
      1. ANY layout whose name is exactly "Blank" (case-insensitive).
         Merck_Themed_Base_v1.pptx has this at index 9; default python-pptx
         has it at index 6. Searching by name handles both cleanly.
      2. The layout with the FEWEST placeholders (avoids "Vertical Title and
         Text" and other heavy default layouts that trigger PowerPoint
         repair warnings when chrome is drawn on top).
      3. Last resort: the first layout. Never the last layout (which is the
         default python-pptx "Vertical Title and Text" and looked wrong on
         every slide of the prior run).
    """
    layouts = list(prs.slide_layouts)
    if not layouts:
        return None
    # 1. Search by name.
    for layout in layouts:
        if str(layout.name).strip().lower() == "blank":
            return layout
    # 2. Fewest placeholders.
    best = None
    best_count = None
    for layout in layouts:
        try:
            count = len(list(layout.placeholders))
        except Exception:
            continue
        if best_count is None or count < best_count:
            best = layout
            best_count = count
            if count == 0:
                return layout
    if best is not None:
        return best
    # 3. First layout (NOT last).
    return layouts[0]


def _new_slide(prs: Presentation, bg_color=None):
    slide = prs.slides.add_slide(_blank_layout(prs))
    if bg_color is not None:
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _rgb_tuple(bg_color)
    return slide


def _intro_layout(prs: Presentation):
    """Pick the themed cover layout from Merck_Themed_Base_v1.pptx.

    Returns the 'Title' layout (index 1 in the Merck base template) which
    provides the diagonal green/yellow cover design, the auto-disclaimer
    text, EMD logos, and three placeholders: title, subtitle, name/date.

    Falls back to None when running against a generic deck (no themed base)
    so build_cover can route to its chrome-from-scratch path.
    """
    if prs is None:
        return None
    # 1. Match by exact name 'Title' (Merck template).
    for layout in prs.slide_layouts:
        if str(layout.name).strip().lower() == "title":
            # Sanity: must have all 3 cover placeholders.
            try:
                placeholders = list(layout.placeholders)
            except Exception:
                continue
            has_title = any(
                ph.placeholder_format.idx == 0 for ph in placeholders
            )
            has_subtitle = any(
                ph.placeholder_format.idx == 1 for ph in placeholders
            )
            if has_title and has_subtitle:
                return layout
    # 2. No themed cover available — caller will fall back.
    return None


def _divider_layout(prs: Presentation):
    """Return the template's native 'Divider' layout, or None.

    The Merck base template includes a branded section-divider layout (index 10,
    named 'Divider') that carries the organic blob shapes and colour scheme.
    Using it makes dividers visually consistent with covers from the same template.

    Layout must have placeholder idx=0 (TITLE — section number) and idx=13
    (BODY — chapter title).  'Divider plain' (no number) is intentionally skipped.
    """
    if prs is None:
        return None
    for layout in prs.slide_layouts:
        name = str(layout.name).strip().lower()
        if name != "divider":          # skip "divider plain" and others
            continue
        try:
            placeholders = list(layout.placeholders)
        except Exception:
            continue
        has_title = any(ph.placeholder_format.idx == 0  for ph in placeholders)
        has_body  = any(ph.placeholder_format.idx == 13 for ph in placeholders)
        if has_title and has_body:
            return layout
    return None


def _cover_picture_layout(prs: Presentation):
    """Return the 'Title with picture' cover layout for EMD Electronics, or None.

    This layout (index 2 in the Merck base template, named 'Title with picture')
    includes a PICTURE placeholder (idx=20) that the user fills in PowerPoint
    after generation.  It is only used when color_theme == 'electronics'.
    """
    if prs is None:
        return None
    for layout in prs.slide_layouts:
        name = str(layout.name).strip().lower()
        if "picture" not in name:
            continue
        try:
            placeholders = list(layout.placeholders)
        except Exception:
            continue
        has_title    = any(ph.placeholder_format.idx == 0  for ph in placeholders)
        has_subtitle = any(ph.placeholder_format.idx == 1  for ph in placeholders)
        has_picture  = any(ph.placeholder_format.idx == 20 for ph in placeholders)
        if has_title and has_subtitle and has_picture:
            return layout
    return None


def add_image(slide, image_path, x, y, w=None, h=None):
    """Insert a user-provided image. Returns the picture shape, or None
    if the file is missing or unreadable. Aspect ratio preserved when only
    one of w/h is supplied.

    Path is validated to be a local file (not a URL) for safety. The
    canonical runner passes user-uploaded paths from `content.image.path`.
    """
    if not image_path:
        return None
    if not os.path.isfile(image_path):
        return None
    try:
        if w is not None and h is not None:
            return slide.shapes.add_picture(image_path,
                                            _emu(x), _emu(y),
                                            width=_emu(w), height=_emu(h))
        if w is not None:
            return slide.shapes.add_picture(image_path,
                                            _emu(x), _emu(y),
                                            width=_emu(w))
        if h is not None:
            return slide.shapes.add_picture(image_path,
                                            _emu(x), _emu(y),
                                            height=_emu(h))
        return slide.shapes.add_picture(image_path, _emu(x), _emu(y))
    except Exception:
        return None


def add_slide_jump_hyperlink(shape, target_slide):
    """Make a shape's text a clickable jump to target_slide. Works by
    creating a slide-jump relationship on the source slide part and
    wrapping each <a:r> text run with an <a:hlinkClick> element pointing
    at that relationship via r:id, plus action='ppaction://hlinksldjump'.

    Used by the canonical runner to make agenda chapter rows clickable.
    """
    from lxml import etree
    if not shape.has_text_frame:
        return
    src_part = shape.part
    try:
        rId = src_part.relate_to(
            target_slide.part,
            "http://schemas.openxmlformats.org/officeDocument/2006/"
            "relationships/slide"
        )
    except Exception:
        return
    tx_body = shape.text_frame._txBody
    for r in tx_body.iter(qn("a:r")):
        rPr = r.find(qn("a:rPr"))
        if rPr is None:
            rPr = etree.SubElement(r, qn("a:rPr"))
            # rPr must be the FIRST child of <a:r> per OOXML schema.
            r.insert(0, rPr)
        # Remove any existing hlinkClick before adding ours.
        for existing in rPr.findall(qn("a:hlinkClick")):
            rPr.remove(existing)
        hlink = etree.SubElement(rPr, qn("a:hlinkClick"))
        hlink.set(qn("r:id"), rId)
        hlink.set("action", "ppaction://hlinksldjump")


def add_speaker_notes(slide, notes):
    """Add speaker notes to a slide. Multi-line text becomes multiple
    paragraphs. Safe to call with None / empty string (no-op).

    Used by the canonical runner — agents emit `notes` per slide and the
    runner lifts them through to PowerPoint's Notes pane for presenters.
    """
    if not notes:
        return
    text = str(notes).strip()
    if not text:
        return
    try:
        notes_tf = slide.notes_slide.notes_text_frame
    except Exception:
        return
    notes_tf.text = ""
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = notes_tf.paragraphs[0] if i == 0 else notes_tf.add_paragraph()
        p.text = line


def _populate_placeholder(layout_ph_idx, slide, text, *, font=None, sz=None,
                          color=None, bold=None, italic=None):
    """Find a placeholder by idx on the slide and set its text + formatting.

    Returns the placeholder shape (or None if not found). When text contains
    newlines, each line becomes a separate paragraph.
    """
    if text is None:
        return None
    target = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == layout_ph_idx:
            target = ph
            break
    if target is None or not target.has_text_frame:
        return None
    tf = target.text_frame
    tf.clear()
    lines = str(text).split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        if font is not None:
            run.font.name = font
        if sz is not None:
            run.font.size = Pt(sz)
        if bold is not None:
            run.font.bold = bool(bold)
        if italic is not None:
            run.font.italic = bool(italic)
        if color is not None:
            run.font.color.rgb = _rgb_tuple(color)
    return target


# ===========================================================================
# Primitives
# ===========================================================================

def _apply_fill(shape, fill):
    if fill is None:
        shape.fill.background()
        return
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb_tuple(fill)


def _apply_border(shape, color, weight=None):
    line_obj = shape.line
    if color is None:
        line_obj.fill.background()
        return
    line_obj.color.rgb = _rgb_tuple(color)
    line_obj.width = weight if weight is not None else Pt(0.5)


def _emu(v):
    """Coerce a coordinate or extent value to an integer EMU.

    OOXML's xsd:long requires plain integers; lxml will serialize a Python
    float as '2138340.0', which PowerPoint rejects and shows as a 'needs
    repair' warning on open. python-pptx's add_shape / add_textbox happen
    to auto-cast, but add_connector does NOT — and any arithmetic involving
    `/` (Python 3 true division) silently produces a float that propagates
    all the way to the XML. This helper is the single coercion point used
    by every primitive below.
    """
    if v is None:
        return 0
    try:
        return int(v)
    except (TypeError, ValueError):
        return 0


def rect(slide, x, y, w, h, fill=None, border=None, border_w=None):
    """Sharp rectangle. Default has no border, no fill."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 _emu(x), _emu(y), _emu(w), _emu(h))
    shp.shadow.inherit = False
    _apply_fill(shp, fill)
    _apply_border(shp, border, border_w if border_w is not None else Pt(0.5))
    if shp.has_text_frame:
        shp.text_frame.text = ""
        shp.text_frame.margin_left = Inches(0.05)
        shp.text_frame.margin_right = Inches(0.05)
        shp.text_frame.margin_top = Inches(0.02)
        shp.text_frame.margin_bottom = Inches(0.02)
    return shp


def rounded(slide, x, y, w, h, fill=None, adj=6000):
    """Rounded rectangle. adj controls corner radius (0 to 50000)."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 _emu(x), _emu(y), _emu(w), _emu(h))
    shp.shadow.inherit = False
    _apply_fill(shp, fill)
    _apply_border(shp, None)
    try:
        shp.adjustments[0] = max(0, min(adj, 50000)) / 100000.0
    except Exception:
        pass
    if shp.has_text_frame:
        shp.text_frame.margin_left = Inches(0.1)
        shp.text_frame.margin_right = Inches(0.1)
        shp.text_frame.margin_top = Inches(0.05)
        shp.text_frame.margin_bottom = Inches(0.05)
    return shp


def oval(slide, x, y, w, h, fill=None):
    """Oval; no border by default."""
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 _emu(x), _emu(y), _emu(w), _emu(h))
    shp.shadow.inherit = False
    _apply_fill(shp, fill)
    _apply_border(shp, None)
    return shp


def circle(slide, x, y, size, fill=None):
    """Circle of given side length."""
    return oval(slide, x, y, size, size, fill=fill)


def line(slide, x1, y1, x2, y2, color, weight=Pt(0.5)):
    """Straight line connector. Coords are forced to int EMU — add_connector
    does not auto-cast, and float coords produce 'needs repair' warnings."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      _emu(x1), _emu(y1),
                                      _emu(x2), _emu(y2))
    conn.shadow.inherit = False
    conn.line.color.rgb = _rgb_tuple(color)
    conn.line.width = weight
    return conn


def hairline(slide, x, y, w, h, color):
    """Thin separator drawn as a no-border filled rectangle."""
    shp = rect(slide, x, y, w, h, fill=color)
    _apply_border(shp, None)
    return shp


def txt(slide, x, y, w, h, text, sz=14, color=None, bold=False, italic=False,
        align=PP_ALIGN.LEFT, font=FONT_BODY, anchor=MSO_ANCHOR.TOP):
    """Transparent textbox with one formatted run."""
    box = slide.shapes.add_textbox(_emu(x), _emu(y), _emu(w), _emu(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = "" if text is None else str(text)
    run.font.name = font
    run.font.size = Pt(sz)
    run.font.bold = bool(bold)
    run.font.italic = bool(italic)
    if color is not None:
        run.font.color.rgb = _rgb_tuple(color)
    return box


def _add_run(paragraph, text, sz=14, color=None, bold=False, italic=False,
             font=FONT_BODY):
    run = paragraph.add_run()
    run.text = "" if text is None else str(text)
    run.font.name = font
    run.font.size = Pt(sz)
    run.font.bold = bool(bold)
    run.font.italic = bool(italic)
    if color is not None:
        run.font.color.rgb = _rgb_tuple(color)
    return run


# ===========================================================================
# Freeform helper
# ===========================================================================

def _freeform_poly(slide, points, fill=None, border=None, border_w=None):
    """Closed freeform polygon from a list of (x, y) EMU points."""
    if not points:
        return None
    builder = slide.shapes.build_freeform(points[0][0], points[0][1])
    for px, py in points[1:]:
        builder.add_line_segments([(px, py)], close=False)
    builder.add_line_segments([(points[0][0], points[0][1])], close=True)
    shp = builder.convert_to_shape()
    shp.shadow.inherit = False
    _apply_fill(shp, fill)
    if border is not None:
        _apply_border(shp, border, border_w if border_w is not None else Pt(0.5))
    else:
        _apply_border(shp, None)
    return shp


# ===========================================================================
# Vector icons (unchanged structure; pen-quality silhouettes)
# ===========================================================================

def icon_chart_bar(slide, x, y, size, color):
    n = 3
    gap = size / 9
    bar_w = (size - (n + 1) * gap) / n
    base_y = y + size * 0.95
    heights = [size * 0.40, size * 0.60, size * 0.85]
    for i in range(n):
        bx = x + gap + i * (bar_w + gap)
        by = base_y - heights[i]
        rect(slide, bx, by, bar_w, heights[i], fill=color)
    hairline(slide, x + size * 0.05, base_y, size * 0.90, Emu(int(Pt(0.75))), color)


def icon_chart_line(slide, x, y, size, color):
    pts = [
        (x + size * 0.08, y + size * 0.80),
        (x + size * 0.32, y + size * 0.55),
        (x + size * 0.55, y + size * 0.65),
        (x + size * 0.92, y + size * 0.20),
    ]
    for a, b in zip(pts, pts[1:]):
        line(slide, a[0], a[1], b[0], b[1], color, Pt(1.5))
    dot = size * 0.10
    circle(slide, pts[-1][0] - dot / 2, pts[-1][1] - dot / 2, dot, fill=color)


def icon_chart_pie(slide, x, y, size, color):
    oval(slide, x, y, size, size, fill=color)
    cx = x + size / 2
    cy = y + size / 2
    pts = [(cx, cy), (cx + size / 2, cy), (cx + size / 2, cy + size / 2),
           (cx, cy + size / 2)]
    _freeform_poly(slide, pts, fill=WHITE)


def icon_arrow_up(slide, x, y, size, color):
    pts = [
        (x + size / 2, y + size * 0.08),
        (x + size * 0.90, y + size * 0.65),
        (x + size * 0.65, y + size * 0.65),
        (x + size * 0.65, y + size * 0.92),
        (x + size * 0.35, y + size * 0.92),
        (x + size * 0.35, y + size * 0.65),
        (x + size * 0.10, y + size * 0.65),
    ]
    _freeform_poly(slide, pts, fill=color)


def icon_arrow_down(slide, x, y, size, color):
    pts = [
        (x + size / 2, y + size * 0.92),
        (x + size * 0.90, y + size * 0.35),
        (x + size * 0.65, y + size * 0.35),
        (x + size * 0.65, y + size * 0.08),
        (x + size * 0.35, y + size * 0.08),
        (x + size * 0.35, y + size * 0.35),
        (x + size * 0.10, y + size * 0.35),
    ]
    _freeform_poly(slide, pts, fill=color)


def icon_arrow_right(slide, x, y, size, color):
    pts = [
        (x + size * 0.92, y + size / 2),
        (x + size * 0.55, y + size * 0.10),
        (x + size * 0.55, y + size * 0.35),
        (x + size * 0.08, y + size * 0.35),
        (x + size * 0.08, y + size * 0.65),
        (x + size * 0.55, y + size * 0.65),
        (x + size * 0.55, y + size * 0.90),
    ]
    _freeform_poly(slide, pts, fill=color)


def icon_check(slide, x, y, size, color):
    p1 = (x + size * 0.10, y + size * 0.55)
    p2 = (x + size * 0.40, y + size * 0.85)
    p3 = (x + size * 0.90, y + size * 0.18)
    line(slide, p1[0], p1[1], p2[0], p2[1], color, Pt(2.0))
    line(slide, p2[0], p2[1], p3[0], p3[1], color, Pt(2.0))


def icon_x(slide, x, y, size, color):
    line(slide, x + size * 0.12, y + size * 0.12, x + size * 0.88, y + size * 0.88, color, Pt(2.0))
    line(slide, x + size * 0.88, y + size * 0.12, x + size * 0.12, y + size * 0.88, color, Pt(2.0))


def icon_alert(slide, x, y, size, color):
    pts = [(x + size / 2, y + size * 0.08), (x + size * 0.95, y + size * 0.90),
           (x + size * 0.05, y + size * 0.90)]
    _freeform_poly(slide, pts, fill=color)
    bar_w = size * 0.10
    bar_h = size * 0.32
    bx = x + size / 2 - bar_w / 2
    by = y + size * 0.35
    rect(slide, bx, by, bar_w, bar_h, fill=WHITE)
    dot_size = size * 0.10
    circle(slide, x + size / 2 - dot_size / 2, y + size * 0.74, dot_size, fill=WHITE)


def icon_info(slide, x, y, size, color):
    circle(slide, x, y, size, fill=color)
    dot = size * 0.10
    circle(slide, x + size / 2 - dot / 2, y + size * 0.22, dot, fill=WHITE)
    bar_w = size * 0.10
    bar_h = size * 0.36
    rect(slide, x + size / 2 - bar_w / 2, y + size * 0.42, bar_w, bar_h, fill=WHITE)


def icon_target(slide, x, y, size, color):
    outer = size
    mid = size * 0.66
    inner = size * 0.30
    circle(slide, x, y, outer, fill=color)
    circle(slide, x + (outer - mid) / 2, y + (outer - mid) / 2, mid, fill=WHITE)
    circle(slide, x + (outer - inner) / 2, y + (outer - inner) / 2, inner, fill=color)


def icon_gear(slide, x, y, size, color):
    cx = x + size / 2
    cy = y + size / 2
    circle(slide, x, y, size, fill=color)
    hub = size * 0.32
    circle(slide, cx - hub / 2, cy - hub / 2, hub, fill=WHITE)
    notch_w = size * 0.14
    notch_h = size * 0.16
    positions = [
        (cx - notch_w / 2, y - notch_h * 0.3),
        (cx - notch_w / 2, y + size - notch_h * 0.7),
        (x - notch_w * 0.3, cy - notch_h / 2),
        (x + size - notch_w * 0.7, cy - notch_h / 2),
    ]
    for px, py in positions:
        rect(slide, px, py, notch_w, notch_h, fill=color)


def icon_users(slide, x, y, size, color):
    head_d = size * 0.32
    body_w = size * 0.50
    body_h = size * 0.32
    circle(slide, x + size * 0.30, y + size * 0.08, head_d, fill=color)
    rect(slide, x + size * 0.18, y + size * 0.46, body_w, body_h, fill=color)
    circle(slide, x + size * 0.55, y + size * 0.18, head_d, fill=color)
    rect(slide, x + size * 0.43, y + size * 0.56, body_w, body_h, fill=color)


def icon_calendar(slide, x, y, size, color):
    body_y = y + size * 0.18
    body_h = size * 0.80
    rect(slide, x + size * 0.05, body_y, size * 0.90, body_h, fill=color)
    rect(slide, x + size * 0.08, body_y + size * 0.18, size * 0.84, body_h - size * 0.22, fill=WHITE)
    rect(slide, x + size * 0.22, y + size * 0.04, size * 0.08, size * 0.22, fill=color)
    rect(slide, x + size * 0.70, y + size * 0.04, size * 0.08, size * 0.22, fill=color)
    hairline(slide, x + size * 0.15, body_y + body_h * 0.55, size * 0.70, Emu(int(Pt(0.75))), color)


def icon_clock(slide, x, y, size, color):
    circle(slide, x, y, size, fill=color)
    inner = size * 0.86
    off = (size - inner) / 2
    circle(slide, x + off, y + off, inner, fill=WHITE)
    cx = x + size / 2
    cy = y + size / 2
    line(slide, cx, cy, cx, cy - size * 0.28, color, Pt(1.5))
    line(slide, cx, cy, cx + size * 0.32, cy, color, Pt(1.0))
    dot = size * 0.08
    circle(slide, cx - dot / 2, cy - dot / 2, dot, fill=color)


def icon_lightbulb(slide, x, y, size, color):
    circle(slide, x + size * 0.18, y + size * 0.06, size * 0.64, fill=color)
    rect(slide, x + size * 0.34, y + size * 0.70, size * 0.32, size * 0.12, fill=color)
    rect(slide, x + size * 0.38, y + size * 0.82, size * 0.24, size * 0.10, fill=color)


def icon_lock(slide, x, y, size, color):
    rect(slide, x + size * 0.15, y + size * 0.42, size * 0.70, size * 0.50, fill=color)
    rect(slide, x + size * 0.25, y + size * 0.18, size * 0.08, size * 0.30, fill=color)
    rect(slide, x + size * 0.67, y + size * 0.18, size * 0.08, size * 0.30, fill=color)
    rect(slide, x + size * 0.25, y + size * 0.15, size * 0.50, size * 0.08, fill=color)
    kh = size * 0.10
    circle(slide, x + size / 2 - kh / 2, y + size * 0.55, kh, fill=WHITE)


def icon_globe(slide, x, y, size, color):
    circle(slide, x, y, size, fill=color)
    cx = x + size / 2
    cy = y + size / 2
    line(slide, x + size * 0.05, cy, x + size * 0.95, cy, WHITE, Pt(0.75))
    line(slide, cx, y + size * 0.05, cx, y + size * 0.95, WHITE, Pt(0.75))


def icon_search(slide, x, y, size, color):
    ring_d = size * 0.70
    rx = x
    ry = y
    circle(slide, rx, ry, ring_d, fill=color)
    hole = ring_d * 0.66
    off = (ring_d - hole) / 2
    circle(slide, rx + off, ry + off, hole, fill=WHITE)
    line(slide, rx + ring_d * 0.78, ry + ring_d * 0.78,
         x + size * 0.98, y + size * 0.98, color, Pt(2.0))


def icon_money(slide, x, y, size, color):
    circle(slide, x, y, size, fill=color)
    bar_w = size * 0.08
    rect(slide, x + size / 2 - bar_w / 2, y + size * 0.18, bar_w, size * 0.64, fill=WHITE)
    rect(slide, x + size * 0.30, y + size * 0.30, size * 0.40, size * 0.08, fill=WHITE)
    rect(slide, x + size * 0.30, y + size * 0.46, size * 0.40, size * 0.08, fill=WHITE)
    rect(slide, x + size * 0.30, y + size * 0.62, size * 0.40, size * 0.08, fill=WHITE)


def icon_trending_up(slide, x, y, size, color):
    pts = [
        (x + size * 0.05, y + size * 0.85),
        (x + size * 0.35, y + size * 0.55),
        (x + size * 0.55, y + size * 0.70),
        (x + size * 0.90, y + size * 0.20),
    ]
    for a, b in zip(pts, pts[1:]):
        line(slide, a[0], a[1], b[0], b[1], color, Pt(1.75))
    tip = pts[-1]
    head = [(tip[0], tip[1]), (tip[0] - size * 0.20, tip[1] + size * 0.02),
            (tip[0] - size * 0.02, tip[1] + size * 0.20)]
    _freeform_poly(slide, head, fill=color)


def icon_trending_down(slide, x, y, size, color):
    pts = [
        (x + size * 0.05, y + size * 0.20),
        (x + size * 0.35, y + size * 0.50),
        (x + size * 0.55, y + size * 0.40),
        (x + size * 0.90, y + size * 0.85),
    ]
    for a, b in zip(pts, pts[1:]):
        line(slide, a[0], a[1], b[0], b[1], color, Pt(1.75))
    tip = pts[-1]
    head = [(tip[0], tip[1]), (tip[0] - size * 0.20, tip[1] - size * 0.02),
            (tip[0] - size * 0.02, tip[1] - size * 0.20)]
    _freeform_poly(slide, head, fill=color)


def icon_shield(slide, x, y, size, color):
    pts = [
        (x + size * 0.50, y + size * 0.05),
        (x + size * 0.92, y + size * 0.20),
        (x + size * 0.92, y + size * 0.55),
        (x + size * 0.50, y + size * 0.95),
        (x + size * 0.08, y + size * 0.55),
        (x + size * 0.08, y + size * 0.20),
    ]
    _freeform_poly(slide, pts, fill=color)


def icon_flag(slide, x, y, size, color):
    rect(slide, x + size * 0.18, y + size * 0.08, size * 0.06, size * 0.84, fill=color)
    pts = [(x + size * 0.24, y + size * 0.10), (x + size * 0.92, y + size * 0.28),
           (x + size * 0.24, y + size * 0.46)]
    _freeform_poly(slide, pts, fill=color)


def icon_doc(slide, x, y, size, color):
    pts = [(x + size * 0.18, y + size * 0.08), (x + size * 0.66, y + size * 0.08),
           (x + size * 0.86, y + size * 0.30), (x + size * 0.86, y + size * 0.92),
           (x + size * 0.18, y + size * 0.92)]
    _freeform_poly(slide, pts, fill=color)
    fold = [(x + size * 0.66, y + size * 0.08), (x + size * 0.66, y + size * 0.30),
            (x + size * 0.86, y + size * 0.30)]
    _freeform_poly(slide, fold, fill=WHITE)
    for k in range(3):
        hairline(slide, x + size * 0.28, y + size * (0.46 + k * 0.12),
                 size * 0.50, Emu(int(Pt(0.75))), WHITE)


ICON_REGISTRY = {
    "chart_bar": icon_chart_bar,
    "chart_line": icon_chart_line,
    "chart_pie": icon_chart_pie,
    "arrow_up": icon_arrow_up,
    "arrow_down": icon_arrow_down,
    "arrow_right": icon_arrow_right,
    "check": icon_check,
    "x": icon_x,
    "alert": icon_alert,
    "info": icon_info,
    "target": icon_target,
    "gear": icon_gear,
    "users": icon_users,
    "calendar": icon_calendar,
    "clock": icon_clock,
    "lightbulb": icon_lightbulb,
    "lock": icon_lock,
    "globe": icon_globe,
    "search": icon_search,
    "money": icon_money,
    "trending_up": icon_trending_up,
    "trending_down": icon_trending_down,
    "shield": icon_shield,
    "flag": icon_flag,
    "doc": icon_doc,
}


def draw_icon(slide, name, x, y, size, color):
    fn = ICON_REGISTRY.get(name, icon_target)
    fn(slide, x, y, size, color)


# ===========================================================================
# Tracked-letter helper (simulates letter-spacing)
# ===========================================================================

def _tracked(text: str) -> str:
    """Insert spaces between letters to simulate uppercase tracking."""
    if not text:
        return ""
    words = str(text).upper().split()
    spaced = [" ".join(list(w)) for w in words]
    return "   ".join(spaced)


def _track_letters(text: str) -> str:
    """Alias for _tracked, exposed for layouts that prefer this name."""
    return _tracked(text)


def _estimate_card_content_h(col: dict, default_h):
    """Estimate vertical height needed for a card based on label + title + items.

    col supports keys: label, title (or body for two/three col), body (for stat
    strip), items (list), milestone, timestamp. Returns an Emu value.
    """
    h = Inches(0.40)  # padding top + label
    title_text = col.get("title") or col.get("body_title") or col.get("body") or ""
    if title_text:
        # ~1 line per 32 chars at 16pt.
        title_lines = max(1, (len(str(title_text)) // 32) + (1 if len(str(title_text)) % 32 else 0))
        h += Inches(0.42) * title_lines
    body_text = col.get("desc") or ""
    if body_text:
        body_lines = max(1, (len(str(body_text)) // 50) + (1 if len(str(body_text)) % 50 else 0))
        h += Inches(0.24) * body_lines
    items = col.get("items") or []
    if items:
        h += Inches(0.05)
        h += Inches(0.32) * len(items)
    if col.get("milestone") or col.get("timestamp"):
        h += Inches(0.50)
    h += Inches(0.30)  # bottom padding
    if h < default_h:
        return default_h
    return h


def _compute_row_card_height(cols, default_h):
    """Equal-height enforcement: return the max needed across all cards.

    Cards in a row must visually look like a row of equals; cards with less
    content show whitespace at the bottom rather than auto-shrinking.
    """
    if not cols:
        return default_h
    heights = [_estimate_card_content_h(c or {}, default_h) for c in cols]
    return max(heights)


def _pad_int(n) -> str:
    try:
        return f"{int(n):02d}"
    except Exception:
        return str(n)


def _format_section_number(num):
    """Render the section-circle text. Pure ints zero-pad to 2 digits to match
    the agenda chips ('01', '02', ..., '10'). Strings with letters pass through
    verbatim to support sub-labels ('3A', '3B', '12B')."""
    if num is None:
        return ""
    s = str(num).strip()
    if not s:
        return ""
    try:
        return f"{int(s):02d}"
    except (ValueError, TypeError):
        return s


# ===========================================================================
# Universal chrome
# ===========================================================================

def _top_chrome(slide, meta, category, palette, top_bar=False,
                page=None, total=None):
    """Top chrome: classification badge + thin progress bar.

    The previous gold rule + breadcrumb ('DECK_LABEL • CATEGORY • MONTH YEAR')
    duplicated information that's already in the footer band (deck_label +
    category) and in the section circle (category tag) — three copies of
    the same words at the top of every content slide. Replaced with:

      - Classification badge (top-right, compliance requirement)
      - Thin progress bar across the top (page / total), filled in
        MERCK_GOLD on light palettes / MERCK_YELLOW on dark. Subtle
        signal of "how far through the deck am I" — the only navigation
        question audiences actually have. Skipped on the cover (when the
        full-width gold top_bar is used instead).

    When meta['cover_top_bar'] is True OR top_bar=True, a full-width gold
    bar renders at y=0. Used for cover slides only; progress bar is
    suppressed in that case.
    """
    pal = _palette_for(palette)
    meta = meta or {}
    dark = _is_dark(palette)

    # NOTE: meta['cover_top_bar'] is intentionally NOT read here — only the
    # explicit top_bar=True argument enables the cover gold-bar treatment.
    # build_cover passes top_bar=True when meta['cover_top_bar'] is set, so
    # content slides never accidentally inherit the cover bar even when the
    # meta key is present in the deck-wide meta dict.
    if bool(top_bar):
        # Full-width gold bar at the very top of the slide (cover only).
        hairline(slide, Inches(0.0), Inches(0.0), SLIDE_W, Inches(0.20),
                 MERCK_GOLD)
    elif page is not None and total:
        # Thin progress bar at the very top: light track + gold fill.
        try:
            p = int(page)
            t = int(total)
        except (TypeError, ValueError):
            p, t = 0, 0
        if t > 0 and p > 0:
            bar_x = Inches(0.65)
            bar_y = Inches(0.10)
            bar_w = Inches(12.0)
            bar_h = Inches(0.05)
            track_color = PURPLE_MUTED if dark else LIGHT_GRAY
            fill_color = pal["hot"] if dark else pal["highlight"]
            hairline(slide, bar_x, bar_y, bar_w, bar_h, track_color)
            # Filled portion proportional to page / total. Min 1px visible.
            frac = max(0.0, min(1.0, p / t))
            fill_w_emu = int(bar_w * frac)
            if fill_w_emu > 0:
                hairline(slide, bar_x, bar_y, fill_w_emu, bar_h, fill_color)

    # Classification badge (top-right, two stacked lines). Kept for
    # compliance — internal vs confidential should be visible at a glance.
    classification = meta.get("classification") or "Internal"
    badge = slide.shapes.add_textbox(CLASS_BADGE_X, CLASS_BADGE_Y,
                                     CLASS_BADGE_W, Inches(0.50))
    tf = badge.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    _add_run(p1, "Classification:", sz=9, color=MERCK_GOLD, font=FONT_BODY)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    _add_run(p2, str(classification).upper(), sz=9, color=MERCK_GOLD,
             bold=True, font=FONT_BODY)


def _bottom_chrome(slide, meta, category, page, total, palette):
    """Footer text boxes: deck label + category (left) and page number (right).

    No band is drawn. Both elements use Rich Purple at 8 pt per CD guidelines.
    Custom text boxes are used (not native placeholder injection) because the
    content is per-slide dynamic.
    """
    meta = meta or {}

    deck_label = meta.get("deck_label", "")
    cat = category or ""
    left_parts = [p for p in [deck_label, cat] if p]
    # On dark-background slides (merck_storytelling) Rich Purple is invisible —
    # use a light off-white so the footer remains readable.
    _ftr_dark = _is_dark(palette)
    footer_color = MERCK_PURPLE if not _ftr_dark else PANEL_LIGHT

    if left_parts:
        txt(slide, Inches(0.65), FOOTER_TEXT_Y, Inches(10.0), Inches(0.20),
            "   •   ".join(left_parts), sz=8,
            color=footer_color, bold=False,
            font=FONT_BODY, anchor=MSO_ANCHOR.TOP)

    if page is not None:
        page_text = _pad_int(page)
        if total:
            page_text = f"{_pad_int(page)} / {_pad_int(total)}"
        txt(slide, Inches(12.0), FOOTER_TEXT_Y, Inches(1.0), Inches(0.20),
            page_text, sz=8, color=footer_color,
            bold=False, font=FONT_BODY, align=PP_ALIGN.RIGHT)


_ICON_DISPATCH = {
    "lightbulb": "icon_lightbulb", "alert": "icon_alert",
    "info": "icon_info", "target": "icon_target",
    "check": "icon_check", "flag": "icon_flag",
    "shield": "icon_shield", "lock": "icon_lock",
    "gear": "icon_gear", "users": "icon_users",
    "chart_bar": "icon_chart_bar", "trending_up": "icon_trending_up",
}

def _section_marker(slide, number, category, palette, icon=None):
    """Numbered circle + uppercase category tag.

    The circle is only drawn when it has visible content (a number or icon) —
    an empty filled circle with nothing inside is omitted. The category tag
    always renders independently of the circle.
    """
    dark = _is_dark(palette)
    circle_fill = MERCK_PURPLE if not dark else MERCK_YELLOW
    number_color = WHITE if not dark else INK_DARK
    tag_color = PURPLE_MUTED if not dark else MERCK_GOLD

    has_number = number is not None and str(number) != ""
    if has_number or icon:
        circle(slide, SECTION_CIRCLE_X, SECTION_CIRCLE_Y,
               SECTION_CIRCLE_D, fill=circle_fill)
        if has_number:
            num_sz = 22 if len(str(number)) <= 2 else 14
            txt(slide, SECTION_CIRCLE_X, SECTION_CIRCLE_Y,
                SECTION_CIRCLE_D, SECTION_CIRCLE_D,
                str(number), sz=num_sz, color=number_color, bold=True,
                font=FONT_BODY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if category:
        txt(slide, SECTION_TAG_X, SECTION_TAG_Y,
            SECTION_TAG_W, SECTION_TAG_H,
            _tracked(category), sz=11, color=tag_color, bold=True,
            font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE)


def _render_action_title(slide, x, y, w, h, content, palette,
                         size=22, italic_color=None, base_color=None):
    """Render action title that may be a string OR list of (text, italic_bool) tuples."""
    pal = _palette_for(palette)
    dark = _is_dark(palette)
    if base_color is None:
        base_color = WHITE if dark else MERCK_PURPLE
    if italic_color is None:
        italic_color = pal["hot"]

    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    # Auto-shrink font if the text is too long to fit — prevents visible clipping.
    try:
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    if isinstance(content, (list, tuple)) and content and \
       all(isinstance(seg, (list, tuple)) and len(seg) == 2 for seg in content):
        for seg_text, seg_italic in content:
            color = italic_color if seg_italic else base_color
            _add_run(p, seg_text, sz=size, color=color, bold=False,
                     italic=bool(seg_italic), font=(FONT_HEAD if seg_italic else FONT_BODY))
    else:
        text = "" if content is None else str(content)
        _add_run(p, text, sz=size, color=base_color, bold=False,
                 italic=False, font=FONT_BODY)
    return box


def _subtitle(slide, x, y, w, h, text, palette):
    if not text:
        return None
    pal = _palette_for(palette)
    color = pal["ink_2"] if not _is_dark(palette) else PANEL_LIGHT
    return txt(slide, x, y, w, h, text, sz=13, color=color, italic=True,
               font=FONT_BODY, anchor=MSO_ANCHOR.TOP)


def _source_line(slide, x, y, w, h, text, palette):
    if not text:
        return None
    pal = _palette_for(palette)
    color = pal["ink_2"] if not _is_dark(palette) else PANEL_LIGHT
    body = str(text)
    if not body.lower().startswith("source"):
        body = f"Source: {body}"
    return txt(slide, x, y, w, h, body, sz=8, color=color, italic=True,
               font=FONT_BODY)


def _methodology_note(slide, x, y, w, h, text, palette):
    """Render a small italic methodology note in INK_GRAY Verdana 9pt."""
    if not text:
        return None
    pal = _palette_for(palette)
    color = pal["ink_2"] if not _is_dark(palette) else PANEL_LIGHT
    body = str(text)
    return txt(slide, x, y, w, h, body, sz=9, color=color, italic=True,
               font=FONT_BODY)


def _takeaway_band(slide, text, palette):
    """Bottom takeaway band placed just above the footer band.

    Single consistent render across all slides: 9pt Verdana bold, mixed case,
    no letter-tracking, no uppercase. The earlier letter-tracked-uppercase
    treatment (the "McKinsey caps look") inflated visual width ~2.2x and
    silently clipped takeaways over ~50 chars at the right edge. Bands are
    decorated by their colored fill (yellow / purple / gold) — typography
    doesn't need to also shout.

    Band geometry: starts at x=0.65, width=11.1 → ends at x=11.75. This is
    0.07" short of the Merck "M" logo placed at x=11.82 in the master slide
    chrome; the previous 12.0" width covered the logo on every content
    slide.

    Author contract: takeaway should be 40-90 characters. Up to ~120 chars
    fit at sz=9 mixed case across the 10.7" usable band width (11.1 minus
    2x 0.20" margins); beyond that, word-wrap pushes to a second line that
    the 0.22" band height clips. The runner validates this.
    """
    if not text:
        return None
    pal  = _palette_for(palette)
    dark = _is_dark(palette)
    # Band fill: use the theme's primary accent colour so each theme's
    # takeaway band is visually distinct.  Dark themes use "hot" (their
    # bright accent on dark bg); light themes use "accent".
    fill = pal["hot"] if dark else pal["accent"]
    # Text colour: pick white or near-black based on perceived luminance of
    # the fill so readability is maintained for all 6 themes.
    r, g, b = fill
    luminance = 0.299 * r + 0.587 * g + 0.114 * b
    text_color = PURPLE_DEEP if luminance > 140 else WHITE
    band = rounded(slide, Inches(0.65), TAKEAWAY_Y, Inches(11.1), TAKEAWAY_H,
                   fill=fill, adj=50000)
    tf = band.text_frame
    tf.margin_left = Inches(0.20)
    tf.margin_right = Inches(0.20)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _add_run(p, str(text).strip(), sz=9, color=text_color, bold=True,
             font=FONT_BODY)
    return band


# ===========================================================================
# Backward-compatible chrome surface
# ===========================================================================

def merck_stub(slide, palette):
    """Legacy: no-op (kept for API parity)."""
    return None


def category_tag(slide, text, palette):
    """Legacy: small uppercase tag at the breadcrumb position."""
    if not text:
        return None
    color = MERCK_GOLD
    return txt(slide, BREADCRUMB_X, BREADCRUMB_Y, BREADCRUMB_W, BREADCRUMB_H,
               _tracked(text), sz=10, color=color, bold=True, font=FONT_BODY,
               anchor=MSO_ANCHOR.MIDDLE)


def action_title(slide, text, palette):
    """Legacy entry point: render action title at TITLE_Y_NUMBERED."""
    return _render_action_title(slide, TITLE_X, TITLE_Y_NUMBERED,
                                TITLE_W, TITLE_H, text, palette)


def subtitle_line(slide, text, palette):
    """Legacy entry point: render subtitle below the action title."""
    return _subtitle(slide, SUB_X, Inches(3.05), SUB_W, SUB_H, text, palette)


def takeaway_band(slide, text, palette):
    return _takeaway_band(slide, text, palette)


def source_line(slide, text, palette):
    return _source_line(slide, Inches(0.65), Inches(7.22),
                        Inches(11.0), Inches(0.20), text, palette)


def footnotes_block(slide, items, palette, y=None):
    """Render numbered footnotes just above the source line.

    items: list of (number, text) tuples. Renders as a single textbox with
    each footnote on its own paragraph. Numbers are rendered as small
    superscript-style runs. Verdana italic 9pt INK_GRAY.
    """
    if not items:
        return None
    pal = _palette_for(palette)
    color = pal["ink_2"] if not _is_dark(palette) else PANEL_LIGHT
    if y is None:
        # Place above SOURCE_Y, but stack upward if many footnotes.
        n = len(items)
        block_h = Inches(0.18) * n
        y = SOURCE_Y - block_h - Inches(0.02)
    box = slide.shapes.add_textbox(Inches(0.65), y, Inches(11.5),
                                   Inches(0.20) * max(len(items), 1) + Inches(0.05))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for i, item in enumerate(items):
        try:
            number, text = item[0], item[1]
        except Exception:
            continue
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(1)
        # Small superscript-like number then italic body text.
        _add_run(p, f"{number}  ", sz=8, color=MERCK_GOLD, bold=True,
                 font=FONT_BODY)
        _add_run(p, str(text), sz=9, color=color, italic=True, font=FONT_BODY)
    return box


def _superscript(slide, x, y, number, color=None):
    """Tiny superscript number used next to chart data labels."""
    txt(slide, x, y, Inches(0.30), Inches(0.18),
        str(number), sz=8, color=color or MERCK_GOLD, bold=True,
        font=FONT_BODY)


def page_number(slide, n, total, palette):
    if n is None:
        return None
    page_text = _pad_int(n)
    if total:
        page_text = f"{_pad_int(n)} / {_pad_int(total)}"
    return txt(slide, Inches(12.0), FOOTER_TEXT_Y, Inches(1.0), Inches(0.20),
               page_text, sz=8, color=MERCK_PURPLE, bold=False, font=FONT_BODY,
               align=PP_ALIGN.RIGHT)


def apply_chrome(slide, meta, action_title, category=None, subtitle=None,
                 takeaway=None, source=None, page=None, total=None,
                 palette="merck_executive", section_number=None,
                 show_section_marker=True, methodology_note=None,
                 section_icon=None):
    """Lay down the universal chrome.

    action_title may be either a string or a list of (text, italic_bool)
    tuples. When show_section_marker is True (default) and category is set,
    a section number marker is rendered.

    section_number: number to show in the circle. Each content slide should
        have a UNIQUE sequential number (1, 2, 3, ...) that matches its entry
        in the agenda. Not the page number.

    section_icon: if provided (e.g. "lightbulb"), the circle renders an icon
        instead of a number. Used for exec_summary and other slides where the
        circle anchors on a concept rather than a sequence position.

    methodology_note: optional italic note rendered just above the source
    line (or in its place if source is empty). 9pt INK_GRAY italic.
    """
    _top_chrome(slide, meta, category, palette, page=page, total=total)
    _bottom_chrome(slide, meta, category, page, total, palette)

    title_y = TITLE_Y_NUMBERED
    if show_section_marker and category:
        if section_icon:
            _section_marker(slide, None, category, palette,
                            icon=section_icon)
            title_y = TITLE_Y_NUMBERED
        else:
            num = section_number
            if num is not None:
                _section_marker(slide,
                                _format_section_number(num),
                                category, palette)
                title_y = TITLE_Y_NUMBERED
            else:
                title_y = TITLE_Y_UNNUMBERED
    else:
        title_y = TITLE_Y_UNNUMBERED

    if action_title:
        _render_action_title(slide, TITLE_X, title_y, TITLE_W, TITLE_H,
                             action_title, palette)
    if subtitle:
        sub_y = title_y + SUB_GAP
        _subtitle(slide, SUB_X, sub_y, SUB_W, SUB_H, subtitle, palette)

    if takeaway:
        _takeaway_band(slide, takeaway, palette)
    # Source + methodology note placement: when both are supplied, the
    # methodology note sits ABOVE the source line (so it doesn't collide with
    # the takeaway band at TAKEAWAY_Y). When only methodology_note is
    # supplied, it occupies the source line slot.
    if source and methodology_note:
        _methodology_note(slide, Inches(0.65),
                          SOURCE_Y - Inches(0.22),
                          Inches(12.0), Inches(0.20),
                          methodology_note, palette)
        _source_line(slide, Inches(0.65), SOURCE_Y,
                     Inches(12.0), SOURCE_H, source, palette)
    elif source:
        _source_line(slide, Inches(0.65), SOURCE_Y,
                     Inches(12.0), SOURCE_H, source, palette)
    elif methodology_note:
        _methodology_note(slide, Inches(0.65), SOURCE_Y,
                          Inches(12.0), SOURCE_H, methodology_note, palette)
    return slide


# ===========================================================================
# Craft helpers
# ===========================================================================

def stub_and_flag(slide, anchor_x, anchor_y, label, palette,
                  direction="up_right",
                  label_offset_x=Inches(0.4),
                  label_offset_y=Inches(-0.4)):
    """Draw a dot at the anchor, an elbow connector, and a small label."""
    pal = _palette_for(palette)
    color = pal["highlight"] if not _is_dark(palette) else pal["hot"]
    dot_size = Inches(0.08)
    circle(slide, anchor_x - dot_size / 2, anchor_y - dot_size / 2,
           dot_size, fill=color)

    if direction == "up_left":
        label_x = anchor_x - Inches(2.0)
        label_y = anchor_y - Inches(0.40)
    elif direction == "down_right":
        label_x = anchor_x + Inches(0.40)
        label_y = anchor_y + Inches(0.30)
    elif direction == "down_left":
        label_x = anchor_x - Inches(2.0)
        label_y = anchor_y + Inches(0.30)
    else:
        label_x = anchor_x + Inches(0.40)
        label_y = anchor_y - Inches(0.40)

    conn = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW,
                                      _emu(anchor_x), _emu(anchor_y),
                                      _emu(label_x), _emu(label_y))
    conn.shadow.inherit = False
    conn.line.color.rgb = _rgb_tuple(color)
    conn.line.width = Pt(0.75)
    box_w = Inches(2.0)
    box_h = Inches(0.30)
    txt(slide, label_x, label_y - box_h / 2,
        box_w, box_h, label, sz=9, color=color, bold=True, font=FONT_BODY)
    return conn


def decimal_align(text_value: str):
    if text_value is None:
        return ("", "")
    s = str(text_value).strip()
    if not s:
        return ("", "")
    if "." in s:
        whole, dec = s.split(".", 1)
        return (whole, "." + dec)
    return (s, "")


# ===========================================================================
# Card pattern (cream PANEL_LIGHT + gold top stripe)
# ===========================================================================



def _draw_card(slide, x, y, w, h, palette, highlighted=False,
               stripe_color=None):
    """Draw a Recall-style card. Returns (top_inside_y, body_color, label_color)."""
    pal = _palette_for(palette)
    if highlighted:
        card_fill = PURPLE_DEEP if _is_dark(palette) else MERCK_PURPLE
        body = rounded(slide, x, y, w, h, fill=card_fill)
        stripe = stripe_color or pal["hot"]
        rounded(slide, x, y, w, Inches(0.06), fill=stripe, adj=50000)
        return WHITE, PANEL_LIGHT, stripe
    # Default card: panel fill + theme-accent stripe + hairline border.
    body = rounded(slide, x, y, w, h, fill=PANEL_LIGHT)
    _apply_border(body, LIGHT_GRAY, Pt(0.5))
    stripe = stripe_color or pal["highlight"]
    rounded(slide, x, y, w, Inches(0.06), fill=stripe, adj=50000)
    return MERCK_PURPLE, INK_GRAY, stripe


def _gold_square_bullet(slide, x, y, size=Inches(0.10), color=None):
    rect(slide, x, y, size, size, fill=color or MERCK_GOLD)


def _bulleted_list(slide, x, y, w, h, items, palette,
                   text_color=None, sz=11, bullet_color=None):
    """Render a list with bullet hierarchy and CD-compliant paragraph spacing.

    Indent detection: items starting with 2+ spaces are rendered as level-3
    sub-bullets (en-dash, slightly smaller, indented) per CD bullet hierarchy.
    All other items are level-2 bullets (▪).  Level-1 (plain text, no bullet)
    is not used here — these are always body list items.

    Paragraph spacing follows CD spec: 3 pt before each paragraph (except the
    first), 3 pt after every paragraph, line spacing 1.05×.
    """
    if not items:
        return
    pal = _palette_for(palette)
    if text_color is None:
        text_color = pal["ink"]
    bcol = bullet_color or pal["highlight"]

    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)

    for i, it in enumerate(items):
        raw = str(it)
        is_sub = raw.startswith("  ")          # 2+ leading spaces → level-3
        text   = raw.lstrip() if is_sub else raw

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(3)             # CD: 3 pt before non-first paragraphs
        p.space_after = Pt(3)                  # CD: 3 pt after every paragraph

        # CD line spacing: 1.05× multiple via XML (no python-pptx shortcut for multiples)
        from lxml import etree as _et
        pPr = p._p.get_or_add_pPr()
        for existing in pPr.findall(qn("a:lnSpc")):
            pPr.remove(existing)
        lnSpc = _et.SubElement(pPr, qn("a:lnSpc"))
        spcPct = _et.SubElement(lnSpc, qn("a:spcPct"))
        spcPct.set("val", "105000")            # 105 % = 1.05×

        if is_sub:
            # Level 3: indented en-dash bullet, slightly smaller
            sub_sz = max(sz - 1, 9)
            _add_run(p, "    –  ", sz=sub_sz, color=bcol, bold=False, font=FONT_BODY)
            _add_run(p, text,     sz=sub_sz, color=text_color, font=FONT_BODY)
        else:
            # Level 2: filled square bullet
            _add_run(p, "▪  ", sz=sz, color=bcol, bold=True, font=FONT_BODY)
            _add_run(p, text, sz=sz, color=text_color, font=FONT_BODY)


# ===========================================================================
# Statement card primitive (left or top stripe)
# ===========================================================================

def statement_card(slide, x, y, w, h, label, body, palette,
                   accent="left", stripe_color=None):
    """Pale-lavender rounded statement card with a theme-accent stripe.

    accent="left" draws a vertical stripe down the left edge; accent="top"
    draws a horizontal stripe across the top. Label sits at the top-left
    inside the card in theme accent colour uppercase tracked Verdana 10pt; body in
    INK_DARK Verdana 14pt bold (auto-shrinks to 13pt for longer bodies).
    """
    pal = _palette_for(palette)
    fill = PANEL_LIGHT
    stripe = stripe_color or pal["highlight"]
    card = rounded(slide, x, y, w, h, fill=fill)
    _apply_border(card, LIGHT_GRAY, Pt(0.5))

    if accent == "top":
        rounded(slide, x, y, w, Inches(0.06), fill=stripe, adj=50000)
        text_inset_x = x + Inches(0.30)
        text_w = w - Inches(0.60)
        label_y = y + Inches(0.18)
        body_y = y + Inches(0.50)
    else:
        # Default: left stripe
        rect(slide, x, y, Inches(0.06), h, fill=stripe)
        text_inset_x = x + Inches(0.24)
        text_w = w - Inches(0.44)
        label_y = y + Inches(0.16)
        body_y = y + Inches(0.48)

    if label:
        txt(slide, text_inset_x, label_y, text_w, Inches(0.24),
            _track_letters(label), sz=10, color=pal["highlight"], bold=True,
            font=FONT_BODY)
    if body:
        body_sz = 14 if len(str(body)) <= 180 else 13
        body_h = (y + h) - body_y - Inches(0.16)
        if body_h < Inches(0.30):
            body_h = Inches(0.30)
        txt(slide, text_inset_x, body_y, text_w, body_h,
            str(body), sz=body_sz, color=INK_DARK, bold=True,
            font=FONT_BODY)
    return card


def in_slide_section(slide, y, label, palette, x=None, w=None, align="center"):
    """Render an in-slide sub-section heading.

    Small theme-accent uppercase tracked Verdana 11pt label used to break a
    content zone into sub-sections (e.g. four cards, then THREE PILLARS
    heading, then three more cards).
    """
    if not label:
        return None
    pal = _palette_for(palette)
    x_use = x if x is not None else Inches(0.65)
    w_use = w if w is not None else Inches(12.0)
    align_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }
    a = align_map.get(align, PP_ALIGN.CENTER)
    return txt(slide, x_use, y, w_use, Inches(0.28),
               _track_letters(label), sz=11, color=pal["highlight"], bold=True,
               font=FONT_BODY, align=a, anchor=MSO_ANCHOR.MIDDLE)


# ===========================================================================
# Phase progress strip (used on cover)
# ===========================================================================

def _phase_progress(slide, phases, palette):
    """Phase progress indicator above the footer band.

    phases: list of dicts {label, desc, current(bool)}.
    """
    if not phases:
        return
    pal = _palette_for(palette)
    dark = _is_dark(palette)
    n = len(phases)
    margin = Inches(0.65)
    avail = SLIDE_W - 2 * margin
    seg = avail / n

    # Draw connecting hairline first.
    cy = PHASE_Y + Inches(0.05)
    line(slide, margin + seg / 2, cy, margin + seg * (n - 0.5), cy,
         PURPLE_MUTED, Pt(0.5))

    for i, ph in enumerate(phases):
        current = bool(ph.get("current"))
        cx = margin + seg * i + seg / 2
        dot_d = Inches(0.20)
        dot_color = pal["highlight"] if current else PURPLE_MUTED
        circle(slide, cx - dot_d / 2, cy - dot_d / 2, dot_d, fill=dot_color)

        label = ph.get("label", "")
        desc = ph.get("desc", "")
        label_color = (WHITE if dark else INK_DARK) if current else PURPLE_MUTED
        desc_color = pal["highlight"] if current else PURPLE_MUTED

        # Label below the dot.
        label_y = cy + Inches(0.18)
        txt(slide, cx - seg / 2 + Inches(0.10), label_y,
            seg - Inches(0.20), Inches(0.22),
            _tracked(label), sz=9, color=label_color, bold=True,
            font=FONT_BODY, align=PP_ALIGN.CENTER)
        # Description.
        if desc:
            txt(slide, cx - seg / 2 + Inches(0.10), label_y + Inches(0.22),
                seg - Inches(0.20), Inches(0.20),
                desc, sz=8, color=desc_color, font=FONT_BODY,
                align=PP_ALIGN.CENTER)


# ===========================================================================
# Chart helpers
# ===========================================================================

def _scale(value, vmin, vmax, lo, hi):
    if vmax == vmin:
        return lo + (hi - lo) / 2
    return lo + (value - vmin) * (hi - lo) / (vmax - vmin)


def add_slope_chart(slide, x, y, w, h, before_label, after_label, items,
                    palette, highlight_indices=None):
    """Slope chart with anti-collision label stacking on both sides."""
    pal = _palette_for(palette)
    if not items:
        return
    highlight_set = set(highlight_indices or [])

    label_gutter = Inches(1.7)
    left_axis = x + label_gutter
    right_axis = x + w - label_gutter
    plot_top = y + Inches(0.65)
    plot_bot = y + h - Inches(0.20)

    vals = []
    for it in items:
        vals.append(it[1])
        vals.append(it[2])
    vmin = min(vals)
    vmax = max(vals)
    if vmin == vmax:
        vmin -= 1
        vmax += 1

    # Gold "Before" / "After" headers at top of axes.
    txt(slide, left_axis - Inches(1.1), plot_top - Inches(0.55),
        Inches(2.2), Inches(0.30),
        _tracked(before_label), sz=11, color=pal["highlight"], bold=True,
        font=FONT_BODY, align=PP_ALIGN.CENTER)
    txt(slide, right_axis - Inches(1.1), plot_top - Inches(0.55),
        Inches(2.2), Inches(0.30),
        _tracked(after_label), sz=11, color=pal["highlight"], bold=True,
        font=FONT_BODY, align=PP_ALIGN.CENTER)

    # Axis hairlines.
    hairline(slide, left_axis, plot_top, Emu(int(Pt(0.75))),
             plot_bot - plot_top, LIGHT_GRAY)
    hairline(slide, right_axis, plot_top, Emu(int(Pt(0.75))),
             plot_bot - plot_top, LIGHT_GRAY)

    # Natural positions on each side, with index.
    left_entries = []
    right_entries = []
    for idx, item in enumerate(items):
        label, bv, av = item[0], item[1], item[2]
        y1 = _scale(bv, vmin, vmax, plot_bot, plot_top)
        y2 = _scale(av, vmin, vmax, plot_bot, plot_top)
        left_entries.append({"idx": idx, "label": label, "value": bv, "y": y1})
        right_entries.append({"idx": idx, "label": label, "value": av, "y": y2})

    def _enforce_gap(entries, min_gap_emu, lo, hi):
        if not entries:
            return entries
        entries = sorted(entries, key=lambda e: e["y"])
        # First pass: push down if too close to prior.
        for i in range(1, len(entries)):
            if entries[i]["y"] - entries[i - 1]["y"] < min_gap_emu:
                entries[i]["y"] = entries[i - 1]["y"] + min_gap_emu
        # Clamp to hi.
        if entries[-1]["y"] > hi:
            entries[-1]["y"] = hi
            for i in range(len(entries) - 2, -1, -1):
                if entries[i + 1]["y"] - entries[i]["y"] < min_gap_emu:
                    entries[i]["y"] = entries[i + 1]["y"] - min_gap_emu
        # Clamp to lo.
        if entries[0]["y"] < lo:
            entries[0]["y"] = lo
            for i in range(1, len(entries)):
                if entries[i]["y"] - entries[i - 1]["y"] < min_gap_emu:
                    entries[i]["y"] = entries[i - 1]["y"] + min_gap_emu
        return entries

    min_gap = Inches(0.24)
    label_lo = plot_top - Inches(0.10)
    label_hi = plot_bot + Inches(0.10)
    left_entries = _enforce_gap(left_entries, min_gap, label_lo, label_hi)
    right_entries = _enforce_gap(right_entries, min_gap, label_lo, label_hi)

    # Map back by idx for label positions.
    left_label_y = {e["idx"]: e["y"] for e in left_entries}
    right_label_y = {e["idx"]: e["y"] for e in right_entries}

    # Draw slope lines and endpoints using NATURAL positions (not adjusted).
    for idx, item in enumerate(items):
        label, bv, av = item[0], item[1], item[2]
        y1 = _scale(bv, vmin, vmax, plot_bot, plot_top)
        y2 = _scale(av, vmin, vmax, plot_bot, plot_top)
        highlighted = idx in highlight_set
        if highlighted:
            line_color = MERCK_YELLOW
            weight = Pt(2.5)
            dot_d = Inches(0.14)
        else:
            line_color = PURPLE_MUTED
            weight = Pt(1.25)
            dot_d = Inches(0.10)
        line(slide, left_axis, y1, right_axis, y2, line_color, weight)
        circle(slide, left_axis - dot_d / 2, y1 - dot_d / 2, dot_d, fill=line_color)
        circle(slide, right_axis - dot_d / 2, y2 - dot_d / 2, dot_d, fill=line_color)

    # Render labels at adjusted positions.
    for idx, item in enumerate(items):
        label, bv, av = item[0], item[1], item[2]
        highlighted = idx in highlight_set
        text_color = MERCK_PURPLE if highlighted else PURPLE_MUTED
        bold = highlighted
        # Left label: "label  value", right-aligned to left axis.
        ly = left_label_y[idx]
        txt(slide, x, ly - Inches(0.11), left_axis - x - Inches(0.20),
            Inches(0.22),
            f"{label}   {bv}", sz=10, color=text_color, bold=bold,
            font=FONT_BODY, align=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE)
        # Right label: "value  label", left-aligned to right axis.
        ry = right_label_y[idx]
        txt(slide, right_axis + Inches(0.20), ry - Inches(0.11),
            x + w - right_axis - Inches(0.20), Inches(0.22),
            f"{av}   {label}", sz=10, color=text_color, bold=bold,
            font=FONT_BODY, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.MIDDLE)


def add_dot_plot(slide, x, y, w, h, categories, items, palette):
    """Horizontal dot plot."""
    pal = _palette_for(palette)
    if not items:
        return
    left = x + Inches(2.2)
    right = x + w - Inches(0.4)
    axis_y = y + h - Inches(0.40)
    top = y + Inches(0.20)
    rows = len(items)
    row_h = (axis_y - top) / max(rows, 1)

    hairline(slide, left, axis_y, right - left, Emu(int(Pt(0.75))), LIGHT_GRAY)
    try:
        cat_nums = [float(str(c).replace("%", "").strip()) for c in categories]
        vmin = min(cat_nums)
        vmax = max(cat_nums)
    except Exception:
        cat_nums = list(range(len(categories)))
        vmin = 0
        vmax = len(categories) - 1
    for i, c in enumerate(categories):
        tx = _scale(cat_nums[i], vmin, vmax, left, right)
        line(slide, tx, axis_y, tx, axis_y + Inches(0.08), LIGHT_GRAY, Pt(0.5))
        txt(slide, tx - Inches(0.4), axis_y + Inches(0.10),
            Inches(0.8), Inches(0.24),
            str(c), sz=9, color=INK_GRAY, font=FONT_BODY, align=PP_ALIGN.CENTER)

    for i, item in enumerate(items):
        label, value = item[0], item[1]
        ry = top + row_h * i + row_h / 2
        line(slide, left, ry, right, ry, LIGHT_GRAY, Pt(0.5))
        txt(slide, x, ry - Inches(0.13), left - x - Inches(0.10),
            Inches(0.26),
            str(label), sz=10, color=INK_DARK, font=FONT_BODY,
            align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        try:
            v = float(str(value).replace("%", "").strip())
        except Exception:
            v = vmin
        dx = _scale(v, vmin, vmax, left, right)
        dot = Inches(0.18)
        circle(slide, dx - dot / 2, ry - dot / 2, dot, fill=MERCK_PURPLE)
        txt(slide, dx + Inches(0.12), ry - Inches(0.12),
            Inches(1.0), Inches(0.24),
            str(value), sz=9, color=PURPLE_DEEP, bold=True, font=FONT_BODY,
            anchor=MSO_ANCHOR.MIDDLE)


def add_marimekko(slide, x, y, w, h, columns, palette):
    """Mosaic chart with proportional column widths and stacked segments.

    Each column has a top label (uppercase tracked) and a bottom share label.
    Segments inside each column are filled with the brand palette in order;
    a label and percentage are rendered inside if the segment is tall enough,
    otherwise an external callout-style label is drawn beside the column.
    A consistent legend appears below the chart.
    """
    pal = _palette_for(palette)
    if not columns:
        return
    palette_colors = CHART_PALETTE
    total_weight = sum(max(c.get("weight", 0), 0) for c in columns) or 1.0
    label_h = Inches(0.30)
    share_h = Inches(0.26)
    legend_h = Inches(0.30)
    top = y + label_h
    bot = y + h - share_h - legend_h - Inches(0.10)
    col_h = bot - top
    gap = Inches(0.03)
    n_cols = len(columns)

    # Build a unified legend from segment label order across columns.
    legend_keys = []
    seen = set()
    for col in columns:
        for seg in col.get("segments", []):
            key = seg[0]
            if key not in seen:
                seen.add(key)
                legend_keys.append(key)
    color_by_key = {k: palette_colors[i % len(palette_colors)]
                    for i, k in enumerate(legend_keys)}

    cx = x
    for ci, col in enumerate(columns):
        weight = max(col.get("weight", 0), 0) / total_weight
        cw = w * weight - (gap if ci < n_cols - 1 else 0)
        # Column header in MERCK_PURPLE tracked bold.
        txt(slide, cx, y, cw + (gap if ci < n_cols - 1 else 0), label_h,
            _tracked(col.get("label", "")), sz=10,
            color=MERCK_PURPLE, bold=True, font=FONT_BODY,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        segs = col.get("segments", [])
        seg_total = sum(max(s[1], 0) for s in segs) or 1.0

        sy = top
        for si, seg in enumerate(segs):
            seg_label, pct = seg[0], max(seg[1], 0)
            sh = col_h * (pct / seg_total)
            color = color_by_key.get(seg_label,
                                     palette_colors[si % len(palette_colors)])
            rect(slide, cx, sy, cw, sh, fill=color)
            # Inside label only when the segment is tall AND wide enough.
            if sh >= Inches(0.40) and cw >= Inches(0.80):
                txt(slide, cx + Inches(0.06), sy, cw - Inches(0.12), sh,
                    f"{seg_label}  {int(round(pct))}%", sz=10, color=WHITE,
                    bold=True, font=FONT_BODY, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
            elif sh >= Inches(0.18) and cw >= Inches(0.60):
                txt(slide, cx + Inches(0.04), sy, cw - Inches(0.08), sh,
                    f"{int(round(pct))}%", sz=9, color=WHITE,
                    bold=True, font=FONT_BODY, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
            sy += sh

        # Share label below the column (% of total weight).
        share_pct = int(round(weight * 100))
        txt(slide, cx, bot + Inches(0.06),
            cw + (gap if ci < n_cols - 1 else 0), share_h,
            f"{share_pct}% of total", sz=9,
            color=INK_GRAY, italic=True, font=FONT_BODY, align=PP_ALIGN.CENTER)
        cx += cw + gap

    # Unified legend below the share row.
    legend_y = bot + share_h + Inches(0.12)
    if legend_keys:
        swatch = Inches(0.16)
        # Estimate text widths so chips don't overlap.
        chip_pad = Inches(0.20)
        chips = []
        for k in legend_keys:
            chip_w = swatch + Inches(0.08) + max(Inches(0.6),
                                                 Inches(0.10) * (len(str(k)) + 2))
            chips.append((k, chip_w))
        total_chip_w = sum(c[1] for c in chips) + chip_pad * (len(chips) - 1)
        lx = x + (w - total_chip_w) / 2
        for k, chip_w in chips:
            rect(slide, lx, legend_y + (legend_h - swatch) / 2,
                 swatch, swatch, fill=color_by_key.get(k, MERCK_PURPLE))
            txt(slide, lx + swatch + Inches(0.06), legend_y,
                chip_w - swatch - Inches(0.06), legend_h,
                str(k), sz=9, color=INK_DARK, font=FONT_BODY,
                anchor=MSO_ANCHOR.MIDDLE)
            lx += chip_w + chip_pad


def add_waterfall(slide, x, y, w, h, bars, palette):
    """Waterfall chart with floating bars at running totals.

    bars: list of {label, value, type} where type is start, up, down, end.
    """
    pal = _palette_for(palette)
    if not bars:
        return
    n = len(bars)
    plot_top = y + Inches(0.55)
    plot_bot = y + h - Inches(0.60)
    gap = Inches(0.10)
    bar_w = (w - gap * (n + 1)) / n

    # Compute running totals.
    running = 0.0
    levels = []  # (low, high) for each bar
    running_totals = []
    for b in bars:
        t = b.get("type", "up")
        v = float(b.get("value", 0))
        if t == "start":
            low, high = 0.0, v
            running = v
        elif t == "end":
            low, high = 0.0, v
            running = v
        elif t == "down":
            high = running
            running -= abs(v)
            low = running
        else:  # up
            low = running
            running += v
            high = running
        levels.append((low, high))
        running_totals.append(running)

    all_y = []
    for lo, hi in levels:
        all_y.append(lo)
        all_y.append(hi)
    vmin = min(all_y + [0.0])
    vmax = max(all_y + [0.0]) * 1.10
    if vmax == vmin:
        vmax = vmin + 1.0

    baseline_y = _scale(0, vmin, vmax, plot_bot, plot_top)
    hairline(slide, x, baseline_y, w, Emu(int(Pt(0.75))), LIGHT_GRAY)

    # Connector dotted lines between adjacent bars at their running totals.
    bx = x + gap
    for i, b in enumerate(bars):
        if i < n - 1:
            t = b.get("type", "up")
            next_t = bars[i + 1].get("type", "up")
            connect_y = _scale(running_totals[i], vmin, vmax, plot_bot, plot_top)
            cx1 = bx + bar_w
            cx2 = bx + bar_w + gap
            # Use a thin dashed-look hairline by drawing tiny segments.
            seg_w = Inches(0.05)
            ax = cx1
            while ax < cx2:
                ax2 = min(ax + seg_w, cx2)
                line(slide, ax, connect_y, ax2, connect_y, PURPLE_MUTED, Pt(0.5))
                ax += seg_w * 2
        bx += bar_w + gap

    # Draw bars.
    bx = x + gap
    for i, b in enumerate(bars):
        t = b.get("type", "up")
        label = b.get("label", "")
        v = float(b.get("value", 0))
        low, high = levels[i]
        y_high = _scale(high, vmin, vmax, plot_bot, plot_top)
        y_low = _scale(low, vmin, vmax, plot_bot, plot_top)
        top_px = min(y_high, y_low)
        bar_h = abs(y_high - y_low)
        bar_h = max(bar_h, Inches(0.04))

        if t in ("start", "end"):
            color = MERCK_PURPLE
            shown = f"{v:g}"
            label_above = True
            label_color = MERCK_PURPLE
        elif t == "down":
            color = BAD_RED
            shown = f"-{abs(v):g}"
            label_above = False
            label_color = BAD_RED
        else:  # up — GOOD_GREEN for positive/growth bars (CD traffic-light convention)
            color = GOOD_GREEN
            shown = f"+{v:g}"
            label_above = True
            label_color = MERCK_PURPLE

        rect(slide, bx, top_px, bar_w, bar_h, fill=color)

        # Value label (always one line; above for up/start/end, below for down).
        if label_above:
            ly = top_px - Inches(0.32)
            anchor = MSO_ANCHOR.BOTTOM
        else:
            ly = top_px + bar_h + Inches(0.04)
            anchor = MSO_ANCHOR.TOP
        txt(slide, bx - Inches(0.30), ly, bar_w + Inches(0.60), Inches(0.28),
            shown, sz=10, color=label_color, bold=True, font=FONT_BODY,
            align=PP_ALIGN.CENTER, anchor=anchor)

        # X-axis label (with optional footnote superscript).
        fn = b.get("footnote") if isinstance(b, dict) else None
        label_text = str(label)
        if fn is not None:
            label_text = f"{label_text} {fn}"
        txt(slide, bx - Inches(0.20), plot_bot + Inches(0.08),
            bar_w + Inches(0.40), Inches(0.40),
            label_text, sz=9, color=INK_DARK, font=FONT_BODY,
            align=PP_ALIGN.CENTER)
        bx += bar_w + gap


def add_small_multiples(slide, x, y, w, h, datasets, palette, grid=(2, 3)):
    """Grid of mini line charts with titles and shared visual rhythm.

    Each cell is a small panel with:
      - a card chrome background (PANEL_LIGHT)
      - a title in MERCK_PURPLE bold
      - a baseline hairline at the bottom of the plot
      - left-edge value labels (min, max) for context
      - a MERCK_PURPLE line trace with a MERCK_GOLD endpoint dot
      - the latest value as a small label at the endpoint
    """
    pal = _palette_for(palette)
    if not datasets:
        return
    rows, cols = grid
    cell_pad = Inches(0.10)
    cell_w = (w - cell_pad * (cols - 1)) / cols
    cell_h = (h - cell_pad * (rows - 1)) / rows
    pad_x = Inches(0.16)
    title_h = Inches(0.30)
    bottom_h = Inches(0.10)

    for i, ds in enumerate(datasets):
        r = i // cols
        c = i % cols
        if r >= rows:
            break
        cx = x + c * (cell_w + cell_pad)
        cy = y + r * (cell_h + cell_pad)

        # Cell chrome.
        card = rounded(slide, cx, cy, cell_w, cell_h, fill=PANEL_LIGHT)
        _apply_border(card, LIGHT_GRAY, Pt(0.5))
        rect(slide, cx, cy, cell_w, Inches(0.04), fill=pal["highlight"])

        title = ds[0]
        values = list(ds[1] or [])

        # Cell title.
        txt(slide, cx + pad_x, cy + Inches(0.10), cell_w - pad_x * 2,
            title_h, title, sz=10, color=MERCK_PURPLE, bold=True,
            font=FONT_BODY)

        # Plot area.
        px = cx + pad_x + Inches(0.32)  # leave room for left value labels
        py = cy + title_h + Inches(0.12)
        pw = cell_w - pad_x * 2 - Inches(0.34)
        ph = cell_h - title_h - bottom_h - Inches(0.22)
        # Baseline hairline.
        hairline(slide, px, py + ph, pw, Emu(int(Pt(0.5))), LIGHT_GRAY)

        if len(values) >= 2:
            vmin = min(values)
            vmax = max(values)
            if vmin == vmax:
                vmin -= 1
                vmax += 1
            # Y-axis labels (max top, min bottom) on the left.
            txt(slide, cx + Inches(0.05), py - Inches(0.06),
                pad_x + Inches(0.22), Inches(0.18),
                f"{vmax:g}", sz=7, color=INK_GRAY, font=FONT_BODY,
                align=PP_ALIGN.RIGHT)
            txt(slide, cx + Inches(0.05), py + ph - Inches(0.10),
                pad_x + Inches(0.22), Inches(0.18),
                f"{vmin:g}", sz=7, color=INK_GRAY, font=FONT_BODY,
                align=PP_ALIGN.RIGHT)
            pts = []
            for j, v in enumerate(values):
                gx = px + (pw * j / (len(values) - 1))
                gy = _scale(v, vmin, vmax, py + ph - Inches(0.04),
                            py + Inches(0.04))
                pts.append((gx, gy))
            for a, b in zip(pts, pts[1:]):
                line(slide, a[0], a[1], b[0], b[1], MERCK_PURPLE, Pt(1.5))
            # Start dot.
            sdot = Inches(0.07)
            circle(slide, pts[0][0] - sdot / 2, pts[0][1] - sdot / 2,
                   sdot, fill=PURPLE_MUTED)
            # End dot (highlight).
            dot = Inches(0.10)
            circle(slide, pts[-1][0] - dot / 2, pts[-1][1] - dot / 2,
                   dot, fill=pal["highlight"])
            # Endpoint value label, placed inside the cell.
            label_w = Inches(0.70)
            lx = pts[-1][0] - label_w + Inches(0.05)
            ly = pts[-1][1] - Inches(0.26)
            if ly < py:
                ly = py + Inches(0.02)
            txt(slide, lx, ly, label_w, Inches(0.22),
                f"{values[-1]:g}", sz=9, color=PURPLE_DEEP, bold=True,
                font=FONT_BODY, align=PP_ALIGN.RIGHT)


def add_simple_bar(slide, x, y, w, h, items=None, palette=None,
                   horizontal=False, highlight_index=None,
                   series=None, categories=None, chart_type_name="column"):
    """Simple bar / column / line chart as a NATIVE PowerPoint chart.

    Accepts two data formats:
      Legacy format  — items: list of [label, value] pairs (single series).
      Standard format — categories: list of str  +  series: list of
                        {name: str, values: list of float}.

    chart_type_name: "column" | "bar" | "line"  (ignored for legacy callers
    that set horizontal=True, which stays as BAR_CLUSTERED).
    """
    _CHART_TYPES = {
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "bar":    XL_CHART_TYPE.BAR_CLUSTERED,
        "line":   XL_CHART_TYPE.LINE,
    }

    # ------------------------------------------------------------------
    # Normalise to (cats, series_list) regardless of which format arrived.
    # ------------------------------------------------------------------
    if series is not None and categories is not None:
        # Standard {categories, series} format.
        cats        = [str(c) for c in categories]
        series_list = list(series)
    elif items:
        # Legacy [label, value] pairs → single series named "Value".
        cats        = [str(it[0]) for it in items]
        series_list = [{"name": "Value", "values": [float(it[1]) for it in items]}]
    else:
        return None

    if not cats or not series_list:
        return None

    chart_data = CategoryChartData()
    chart_data.categories = cats
    for s in series_list:
        vals = [float(v) for v in s.get("values", [])]
        chart_data.add_series(str(s.get("name", "")), vals)

    # ------------------------------------------------------------------
    # Determine PowerPoint chart type.
    # ------------------------------------------------------------------
    if horizontal:
        xl_type = XL_CHART_TYPE.BAR_CLUSTERED
    else:
        xl_type = _CHART_TYPES.get(str(chart_type_name).lower(),
                                   XL_CHART_TYPE.COLUMN_CLUSTERED)

    chart_shape = slide.shapes.add_chart(
        xl_type, _emu(x), _emu(y), _emu(w), _emu(h), chart_data
    )
    chart = chart_shape.chart

    # Clean styling — no legend when single series; show for multi-series.
    chart.has_title = False
    chart.has_legend = len(series_list) > 1
    if chart.has_legend:
        try:
            chart.legend.position    = 2  # BOTTOM
            chart.legend.include_in_layout = False
        except Exception:
            pass

    try:
        for axis in (chart.category_axis, chart.value_axis):
            try:
                axis.tick_labels.font.size = Pt(10)
                axis.tick_labels.font.name = FONT_BODY
            except Exception:
                pass
    except Exception:
        pass

    # ------------------------------------------------------------------
    # Color each series.  For a single series use the highlight logic;
    # for multiple series rotate through the Merck palette.
    # ------------------------------------------------------------------
    # Use the official Merck Corporate Design 12-color chart palette so
    # multi-series charts follow CD specifications (orange / teal / dark-cyan …).
    _MULTI_SERIES_COLORS = CHART_PALETTE
    for si, s_obj in enumerate(chart.series):
        series_color = _MULTI_SERIES_COLORS[si % len(_MULTI_SERIES_COLORS)]
        try:
            s_obj.format.line.color.rgb = _rgb_tuple(series_color)
            s_obj.format.fill.solid()
            s_obj.format.fill.fore_color.rgb = _rgb_tuple(series_color)
        except Exception:
            pass
        if si == 0 and highlight_index is not None:
            for i, pt in enumerate(s_obj.points):
                fill = pt.format.fill
                fill.solid()
                fill.fore_color.rgb = _rgb_tuple(
                    MERCK_GOLD if i == int(highlight_index) else PURPLE_MUTED
                )
        elif si == 0 and len(series_list) == 1:
            for pt in s_obj.points:
                fill = pt.format.fill
                fill.solid()
                fill.fore_color.rgb = _rgb_tuple(MERCK_PURPLE)

    # Data labels — only on bar/column charts (not line).
    if xl_type not in (XL_CHART_TYPE.LINE,):
        try:
            from pptx.enum.chart import XL_LABEL_POSITION
            plot = chart.plots[0]
            plot.has_data_labels = True
            dl = plot.data_labels
            dl.show_value = True
            dl.font.size = Pt(9)
            dl.font.bold = True
            try:
                dl.font.color.rgb = _rgb_tuple(INK_DARK)
            except Exception:
                pass
            try:
                dl.position = XL_LABEL_POSITION.OUTSIDE_END
            except Exception:
                pass
        except Exception:
            pass

    return chart_shape


# Chart dispatcher used by build_chart_slide

def _render_chart(slide, chart, x, y, w, h, palette):
    """Render chart content onto slide. Returns True if a chart was drawn."""
    if not chart:
        return False
    ctype = chart.get("type", "column")
    data = chart.get("data", {})

    if ctype == "slope":
        add_slope_chart(slide, x, y, w, h,
                        data.get("before_label", "Before"),
                        data.get("after_label", "After"),
                        data.get("items", []),
                        palette,
                        highlight_indices=data.get("highlight_indices"))
        return True
    elif ctype == "dot":
        add_dot_plot(slide, x, y, w, h,
                     data.get("categories", ["0", "25", "50", "75", "100"]),
                     data.get("items", []), palette)
        return True
    elif ctype == "marimekko":
        add_marimekko(slide, x, y, w, h, data.get("columns", []), palette)
        return True
    elif ctype == "waterfall":
        add_waterfall(slide, x, y, w, h, data.get("bars", []), palette)
        return True
    elif ctype == "small_multiples":
        add_small_multiples(slide, x, y, w, h,
                            data.get("datasets", []), palette,
                            grid=tuple(data.get("grid", (2, 3))))
        return True
    elif ctype in ("bar", "column", "line"):
        # Standard {categories, series} format — preferred for LLM-generated plans.
        # Falls back to legacy [label, value] items format when series is absent.
        series_data  = data.get("series")
        categories   = data.get("categories")
        legacy_items = data.get("items")
        if series_data and categories:
            add_simple_bar(slide, x, y, w, h,
                           categories=categories,
                           series=series_data,
                           palette=palette,
                           horizontal=(ctype == "bar"),
                           highlight_index=data.get("highlight_index"),
                           chart_type_name=ctype)
            return True
        elif legacy_items:
            add_simple_bar(slide, x, y, w, h,
                           items=legacy_items,
                           palette=palette,
                           horizontal=(ctype == "bar"),
                           highlight_index=data.get("highlight_index"),
                           chart_type_name=ctype)
            return True
        else:
            import sys
            print(f"WARNING: chart_slide received type='{ctype}' but no "
                  f"'series'+'categories' or 'items' data found — slide "
                  f"will render without a chart.", file=sys.stderr)
            return False
    else:
        # Unknown type: attempt legacy items format and warn.
        import sys
        items = data.get("items", [])
        if items:
            add_simple_bar(slide, x, y, w, h,
                           items=items, palette=palette,
                           horizontal=bool(data.get("horizontal", False)),
                           highlight_index=data.get("highlight_index"))
            return True
        else:
            print(f"WARNING: chart_slide received unknown chart type "
                  f"'{ctype}' with no renderable data.", file=sys.stderr)
            return False


# ===========================================================================
# Style helpers
# ===========================================================================

def _style_or_promote(category, style):
    # Never override an explicit storytelling style — visual consistency
    # across the deck matters more than per-category promotion.
    if style == "merck_storytelling":
        return style
    if category and str(category) in AUTO_PROMOTE_EXECUTIVE:
        return "merck_executive"
    return style


def _tone_color(tone, palette):
    if tone == "positive":
        return GOOD_GREEN
    if tone == "negative":
        return BAD_RED
    v = str(tone).strip().lower()
    # Decision-row semantic tones (plan schema: approve / discuss / note).
    if v in ("approve", "approved", "yes", "accept", "green", "g", "good"):
        return GOOD_GREEN
    if v in ("discuss", "review", "consider", "pending", "open",
             "amber", "yellow", "a"):
        return MERCK_YELLOW
    if v in ("note", "fyi", "info", "neutral", "observation"):
        return PURPLE_MUTED
    if v in ("red", "r"):
        from_rag = _rag_color(tone)
        return from_rag
    return MERCK_PURPLE


# ===========================================================================
# Layout: COVER
# ===========================================================================

def _draw_cover_keymessages_grid(slide, cards, grid_top, grid_h, card_style):
    """N-aware key_messages grid for cover slides.

    - 1 card: full-width single row
    - 2 cards: two columns, one row
    - 3 cards: three columns, one row (centered, not 2x2 with orphan)
    - 4 cards: 2x2 grid

    Avoids the asymmetric look of 3 cards in a 2x2 grid (top-left, top-right,
    bottom-left, with the bottom-right empty).
    """
    if not cards:
        return
    cards = list(cards)[:4]
    n = len(cards)
    gap = Inches(0.30)
    pad = Inches(0.20)
    left = Inches(0.65)

    if n == 4:
        cols, rows = 2, 2
    elif n == 3:
        cols, rows = 3, 1
    elif n == 2:
        cols, rows = 2, 1
    else:
        cols, rows = 1, 1

    col_w = (Inches(12.0) - gap * (cols - 1)) / cols
    row_h = (grid_h - gap * (rows - 1)) / rows

    for i, m in enumerate(cards):
        r = i // cols
        c = i % cols
        cx = left + c * (col_w + gap)
        cy = grid_top + r * (row_h + gap)
        _draw_card(slide, cx, cy, col_w, row_h, card_style, highlighted=False)
        _gold_square_bullet(slide, cx + pad, cy + Inches(0.22))
        title_box = slide.shapes.add_textbox(
            cx + pad + Inches(0.20), cy + Inches(0.15),
            col_w - pad * 2 - Inches(0.25), Inches(0.36))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.0)
        tf.margin_bottom = Inches(0.0)
        p = tf.paragraphs[0]
        _add_run(p, m.get("label", ""), sz=14, color=MERCK_PURPLE,
                 bold=True, font=FONT_BODY)
        txt(slide, cx + pad, cy + Inches(0.55), col_w - pad * 2,
            row_h - Inches(0.65),
            m.get("body", ""), sz=11, color=INK_GRAY, font=FONT_BODY)


def _draw_cover_authors(slide, authors, dark=False):
    """Stack a list of {name, title} byline entries below the subtitle area.

    Authors render one per line from y ~5.20 to y ~6.10. Name in bold (white
    on dark covers, INK_DARK on light), title in lighter color on the same
    line, separated by a small gap.
    """
    if not authors:
        return
    rows = list(authors)[:5]
    line_h = Inches(0.32)
    start_y = Inches(5.20)
    if dark:
        name_color = WHITE
        title_color = PANEL_LIGHT
    else:
        name_color = INK_DARK
        title_color = INK_GRAY
    for i, a in enumerate(rows):
        ry = start_y + line_h * i
        name = a.get("name", "") if isinstance(a, dict) else str(a)
        atitle = a.get("title", "") if isinstance(a, dict) else ""
        box = slide.shapes.add_textbox(Inches(0.65), ry, Inches(12.0), line_h)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.02)
        tf.margin_right = Inches(0.02)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        if name:
            _add_run(p, name, sz=12, color=name_color, bold=True,
                     font=FONT_BODY)
        if atitle:
            _add_run(p, "   ", sz=12, color=name_color, font=FONT_BODY)
            _add_run(p, atitle, sz=11, color=title_color, font=FONT_BODY)


def build_cover(prs, meta, title=None, subtitle="", style="merck_executive",
                key_messages=None, phases=None, action_title=None,
                page=None, total=None, section_number=None, category=None,
                authors=None, top_bar=False, methodology_note=None,
                content=None, color_theme=None):
    """Cover slide.

    Preferred path: when the Merck themed base is loaded, this uses the
    'Title' layout (index 1) directly. The layout provides the diagonal
    green/yellow design, EMD/MilliporeSigma/EMD-Electronics logos, the auto
    disclaimer ('The businesses of Merck KGaA...'), and the classification
    badge. The function populates three placeholders:
      - Title 1 (idx 0): action_title (keep short, ≤60 chars)
      - Subtitle 2 (idx 1): subtitle (main subtitle, multi-line OK)
      - Name/Date (idx 10): name + month_year on two lines

    Electronics theme path: uses the 'Title with picture' layout (index 2).
    The PICTURE placeholder (idx 20) is left empty so the user can add their
    own image in PowerPoint after generation.

    Fallback path: when no themed layout is found (generic python-pptx
    deck), draws the cover chrome from scratch — the legacy behavior.
    key_messages / phases / authors-as-grid render only in fallback mode.

    title (also accepted as action_title) may be a string OR a list of
    (text, italic_bool) tuples — italic runs render in MERCK_YELLOW (only
    honored on the fallback path; the themed layout uses its own type).
    """
    if content:
        if "authors"      in content: authors      = content["authors"]
        if "key_messages" in content: key_messages = content["key_messages"]
        if "phases"       in content: phases       = content["phases"]
        if "title"        in content: title        = content["title"]
        if "subtitle"     in content: subtitle     = content["subtitle"]
    if title is None:
        title = action_title if action_title is not None else ""

    # For Electronics theme, try the 'Title with picture' layout first.
    theme_lower = str(color_theme or "").lower()
    if theme_lower == "electronics":
        pic_layout = _cover_picture_layout(prs)
        if pic_layout is not None:
            import sys
            print(
                "INFO: Electronics cover — image placeholder (idx 20) left "
                "empty for manual editing in PowerPoint.",
                file=sys.stderr,
            )
            intro_layout = pic_layout
        else:
            intro_layout = _intro_layout(prs)
    else:
        intro_layout = _intro_layout(prs)

    # Themed-layout path: drop the visual chrome onto the template.
    if intro_layout is not None:
        slide = prs.slides.add_slide(intro_layout)

        # Dark color themes (synthetic, electronics) need a dark slide background.
        # The template's "light panel" freeform is made transparent by
        # _apply_color_theme(), so the slide's own background fills the gap.
        # This must be set before placeholders are populated.
        if theme_lower in ("synthetic", "electronics"):
            try:
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = _rgb_tuple(MERCK_PURPLE)
            except Exception:
                pass

        # Flatten tuple titles to a single string; the template's Title 1
        # uses its own typography so we don't try to honor italic runs.
        if isinstance(title, (list, tuple)) and title and \
           all(isinstance(seg, (list, tuple)) and len(seg) == 2
               for seg in title):
            title_text = "".join(seg[0] for seg in title)
        else:
            title_text = str(title) if title else ""

        # Split title at the schema-convention ';' separator: everything before
        # the first ';' is the deck title; everything after is the subtitle.
        # This prevents long combined strings from overflowing the title box.
        title_part, _, subtitle_part = title_text.partition(";")
        title_part    = title_part.strip()
        subtitle_part = subtitle_part.strip()

        ph_title = _populate_placeholder(0, slide, title_part)
        if ph_title is not None and ph_title.has_text_frame:
            try:
                from pptx.enum.text import MSO_AUTO_SIZE
                ph_title.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except Exception:
                pass

        # Theme-dependent body text color for subtitle / name-date placeholders.
        # The USA template can inherit teal from its theme master; explicitly
        # override so these placeholders always use a CD-compliant neutral.
        _cover_dark = theme_lower in ("synthetic", "electronics")
        _body_color = PANEL_LIGHT if _cover_dark else INK_GRAY

        # Populate subtitle placeholder (idx 1) with subtitle_part if the
        # slide had no explicit subtitle and the title contained a ';' split.
        if subtitle_part and not subtitle:
            subtitle = subtitle_part
        if subtitle:
            ph_sub = _populate_placeholder(1, slide, subtitle,
                                           font=FONT_BODY, color=_body_color)
            if ph_sub is not None and ph_sub.has_text_frame:
                ph_sub.text_frame.word_wrap = True
                try:
                    from pptx.enum.text import MSO_AUTO_SIZE
                    ph_sub.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                except Exception:
                    pass
        # Name/Date placeholder (idx 10): two lines (author name; month/year).
        name_line = ""
        if authors:
            first = authors[0] if isinstance(authors, list) else authors
            if isinstance(first, dict):
                name_line = first.get("name", "")
        if not name_line:
            name_line = (meta or {}).get("author") or ""
        date_line = (meta or {}).get("month_year", "")
        nd_text = "\n".join(s for s in [name_line, date_line] if s)
        if nd_text:
            _populate_placeholder(10, slide, nd_text,
                                  font=FONT_BODY, color=_body_color)
        return slide

    # Fallback: no themed base. Build the cover from scratch (legacy path).
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])

    if style == "merck_storytelling":
        # Dark cover: full-bleed purple + top chrome + huge italic hero.
        _top_chrome(slide, meta, category, style,
                    top_bar=top_bar or bool((meta or {}).get("cover_top_bar")))
        # Hero title.
        _render_action_title(slide, Inches(0.65), Inches(1.50),
                             Inches(12.0), Inches(1.60), title, style,
                             size=44, italic_color=MERCK_YELLOW,
                             base_color=WHITE)
        # Subtitle.
        if subtitle:
            txt(slide, Inches(0.65), Inches(3.10), Inches(12.0), Inches(0.50),
                subtitle, sz=15, color=PANEL_LIGHT, italic=True,
                font=FONT_BODY)

        # N-aware key_messages grid (1/2/3/4 cards).
        if key_messages:
            _draw_cover_keymessages_grid(slide, key_messages,
                                         Inches(4.10), Inches(2.10),
                                         "merck_executive")

        # Authors byline (stacked, name bold + title in lighter color).
        if authors:
            _draw_cover_authors(slide, authors, dark=True)

        # Phase progress strip.
        if phases:
            _phase_progress(slide, phases, style)

        # Bottom footer band.
        _bottom_chrome(slide, meta, None, None, None, style)

    else:
        # Light cover: white bg.
        _top_chrome(slide, meta, category, style,
                    top_bar=top_bar or bool((meta or {}).get("cover_top_bar")))
        # Hero title in purple bold (with optional italic emphasis in yellow).
        _render_action_title(slide, Inches(0.65), Inches(1.80),
                             Inches(12.0), Inches(1.70), title, style,
                             size=44, italic_color=MERCK_YELLOW,
                             base_color=MERCK_PURPLE)
        # Theme-accent rule under title.
        hairline(slide, Inches(0.65), Inches(3.55), Inches(2.6),
                 Emu(int(Pt(2.5))), pal["highlight"])
        # Subtitle gray italic.
        if subtitle:
            txt(slide, Inches(0.65), Inches(3.75), Inches(12.0), Inches(0.50),
                subtitle, sz=16, color=INK_GRAY, italic=True, font=FONT_BODY)

        # N-aware key_messages grid (1/2/3/4 cards).
        if key_messages:
            _draw_cover_keymessages_grid(slide, key_messages,
                                         Inches(4.20), Inches(1.80), style)

        # Authors byline (stacked, name bold + title in lighter color).
        if authors:
            _draw_cover_authors(slide, authors, dark=False)

        # Deck label and month/year strip near bottom.
        deck_label = (meta or {}).get("deck_label", "")
        month_year = (meta or {}).get("month_year", "")
        if deck_label or month_year:
            parts = []
            if deck_label:
                parts.append(_tracked(deck_label))
            if month_year:
                parts.append(_tracked(month_year))
            sep = "   •   "
            txt(slide, Inches(0.65), Inches(5.95), Inches(12.0), Inches(0.30),
                sep.join(parts), sz=10, color=PURPLE_MUTED, bold=True,
                font=FONT_BODY)

        # Phase strip.
        if phases:
            _phase_progress(slide, phases, style)

        # Footer band.
        _bottom_chrome(slide, meta, None, None, None, style)

    return slide


# ===========================================================================
# Layout: EXEC SUMMARY
# ===========================================================================

def build_exec_summary(prs, meta, action_title=None, key_messages=None, takeaway=None,
                       source=None, page=None, total=None, style=None,
                       category=None, section_number=None,
                       methodology_note=None, subtitle=None, content=None):
    """Executive Summary: numbered rows with label + body.
    Auto-promoted to executive style (unless overridden by an exec-eligible
    category). Category defaults to 'Executive Summary'.
    """
    if content:
        if "key_messages" in content: key_messages = content["key_messages"]
        if "takeaway"     in content: takeaway     = content["takeaway"]
    style = _style_or_promote(category or "Executive Summary", style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title,
                 category="Executive Summary",
                 subtitle=subtitle,
                 takeaway=takeaway,
                 source=source, page=page, total=total, palette=style,
                 section_number=section_number,
                 methodology_note=methodology_note)

    items = (key_messages or [])[:5]
    n = max(len(items), 1)
    # Content area below the title block.
    zone_top = Inches(2.95) if subtitle else Inches(2.55)
    zone_bot = Inches(6.50)
    zone_h = zone_bot - zone_top
    row_gap = Inches(0.10)
    row_h = (zone_h - row_gap * (n - 1)) / n
    for i, msg in enumerate(items):
        ry = zone_top + i * (row_h + row_gap)
        num_w = Inches(0.65)
        txt(slide, Inches(0.65), ry + Inches(0.02), num_w, row_h,
            _pad_int(i + 1), sz=24, color=MERCK_PURPLE, bold=True,
            font=FONT_BODY, anchor=MSO_ANCHOR.TOP)
        hairline(slide, Inches(0.65) + num_w + Inches(0.10),
                 ry + row_h - Emu(int(Pt(0.5))),
                 Inches(12.0) - num_w - Inches(0.10),
                 Emu(int(Pt(0.5))), LIGHT_GRAY)
        label = msg.get("label", "")
        body = msg.get("body", "")
        body_x = Inches(0.65) + num_w + Inches(0.20)
        body_w = Inches(12.0) - num_w - Inches(0.30)
        if label:
            txt(slide, body_x, ry + Inches(0.02), body_w, Inches(0.30),
                label, sz=13, color=MERCK_PURPLE, bold=True, font=FONT_BODY)
        txt(slide, body_x, ry + Inches(0.34), body_w, row_h - Inches(0.40),
            body, sz=10, color=INK_GRAY, font=FONT_BODY)
    return slide


# ===========================================================================
# Layout: AGENDA
# ===========================================================================

def build_agenda(prs, meta, chapters=None, style="merck_executive",
                 page=None, total=None, action_title=None,
                 section_number=None, category=None,
                 methodology_note=None, content=None):
    # action_title is honored when provided; otherwise the title defaults to
    # "INDEX". methodology_note is accepted for caller consistency and ignored.
    """Agenda / chapters slide (no section marker; INDEX-style title)."""
    if content:
        if "chapters" in content: chapters = content["chapters"]
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    _top_chrome(slide, meta, None, style, page=page, total=total)
    _bottom_chrome(slide, meta, "Index", page, total, style)

    # Big "INDEX" title (or supplied action_title).
    dark = _is_dark(style)
    title_color = WHITE if dark else MERCK_PURPLE
    title_text = action_title if action_title else "INDEX"
    heading_box = txt(slide, Inches(0.65), Inches(1.30), Inches(12.0), Inches(1.20),
                      str(title_text).upper(), sz=44, color=title_color, bold=True,
                      font=FONT_HEAD)
    try:
        from pptx.enum.text import MSO_AUTO_SIZE
        heading_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass

    items = (chapters or [])[:12]
    if not items:
        return slide
    # Two-column layout (left half + right half).
    half = len(items) // 2 + len(items) % 2
    col_top = Inches(2.80)  # heading box ends at 1.30+1.20=2.50; 0.30 gap
    col_h = Inches(3.95)
    col_w = (Inches(12.0) - Inches(0.50)) / 2
    row_gap = Inches(0.12)
    rows_left = items[:half]
    rows_right = items[half:]

    def _draw_col(rows, x0):
        if not rows:
            return
        n = len(rows)
        compact = (n >= 5)
        title_sz = 12 if compact else 14
        sub_sz = 9 if compact else 10
        row_h = (col_h - row_gap * (n - 1)) / n
        for i, ch in enumerate(rows):
            ry = col_top + i * (row_h + row_gap)
            number = ch.get("number") or _pad_int(i + 1)
            title = ch.get("title", "")
            sub = ch.get("subtitle", "")
            # Number in purple bold. Named with the chapter number so the
            # canonical runner can attach a slide-jump hyperlink in a
            # post-pass after all target slides exist.
            num_color   = MERCK_PURPLE if not dark else pal["hot"]
            num_shape = txt(slide, x0, ry, Inches(0.60), Inches(0.40),
                str(number), sz=title_sz, color=num_color, bold=True,
                font=FONT_BODY)
            try:
                num_shape.name = f"AgendaChapter_{number}"
            except Exception:
                pass
            # Title in INK_DARK bold (also named for hyperlink targeting).
            # Height is dynamic so long titles can wrap to 2 lines without clipping.
            title_h = min(row_h * 0.60, Inches(0.52))
            title_shape = txt(slide, x0 + Inches(0.55), ry,
                col_w - Inches(0.55),
                title_h,
                title, sz=title_sz,
                color=INK_DARK if not dark else WHITE,
                bold=True, font=FONT_BODY)
            try:
                title_shape.name = f"AgendaChapter_{number}"
            except Exception:
                pass
            if sub:
                sub_y = ry + title_h + Inches(0.02)
                sub_h = max(row_h - title_h - Inches(0.06), Inches(0.18))
                txt(slide, x0 + Inches(0.55), sub_y,
                    col_w - Inches(0.55),
                    sub_h,
                    sub, sz=sub_sz,
                    color=INK_GRAY if not dark else PANEL_LIGHT,
                    italic=True, font=FONT_BODY)
    _draw_col(rows_left, Inches(0.65))
    _draw_col(rows_right, Inches(0.65) + col_w + Inches(0.50))
    return slide


# ===========================================================================
# Layout: SECTION DIVIDER
# ===========================================================================

def build_section_divider(prs, meta, number=None, title=None, style="merck_executive",
                          page=None, total=None, action_title=None,
                          section_number=None, category=None,
                          takeaway=None, source=None,
                          methodology_note=None, content=None):
    """Section divider with a section number and chapter title.

    Preferred path: uses the template's native 'Divider' layout (branded
    organic blob shapes, colour scheme from the theme template file).  The
    section number goes into placeholder idx=0 and the chapter title into
    idx=13.

    Fallback path (no 'Divider' layout found): draws the divider from scratch
    using the programmatic chrome — a large serif number at left and an italic
    title at right with a gold rule below.

    title (also accepted as action_title) is the section heading.
    methodology_note is accepted for caller consistency and ignored.
    """
    if content:
        if "number" in content: number = content["number"]
        if "title"  in content: title  = content["title"]
    if title is None:
        title = action_title if action_title is not None else ""

    # Format the section number as a zero-padded two-digit string.
    num_str = ""
    if number is not None:
        try:
            num_str = f"{int(number):02d}"
        except (TypeError, ValueError):
            num_str = str(number)

    pal = _palette_for(style)

    # ------------------------------------------------------------------
    # Preferred path: native template Divider layout.
    # ------------------------------------------------------------------
    # For dark themes (synthetic, electronics) the native Divider layout's
    # full-slide background rectangle uses scheme:accent1, which is yellow
    # for those themes — this would make the divider entirely yellow.
    # Use the programmatic fallback instead so the palette bg (violet) is
    # applied correctly.
    divider_layout = _divider_layout(prs) if not _is_dark(style) else None
    if divider_layout is not None:
        slide = prs.slides.add_slide(divider_layout)
        # Populate number placeholder (idx 0).
        if num_str:
            ph0 = _populate_placeholder(0, slide, num_str)
            if ph0 is not None and ph0.has_text_frame:
                ph0.text_frame.word_wrap = True
                try:
                    from pptx.enum.text import MSO_AUTO_SIZE
                    ph0.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                except Exception:
                    pass
        # Populate chapter title placeholder (idx 13).
        ph13 = _populate_placeholder(13, slide, str(title))
        if ph13 is not None and ph13.has_text_frame:
            ph13.text_frame.word_wrap = True
            try:
                from pptx.enum.text import MSO_AUTO_SIZE
                ph13.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except Exception:
                pass
        # Add bottom chrome (classification + page number).
        _bottom_chrome(slide, meta, "Section", page, total, style)
        return slide

    # ------------------------------------------------------------------
    # Fallback: programmatic renderer (no themed template available).
    # ------------------------------------------------------------------
    slide = _new_slide(prs, bg_color=pal["bg"])
    _top_chrome(slide, meta, None, style, page=page, total=total)
    _bottom_chrome(slide, meta, "Section", page, total, style)

    dark = _is_dark(style)
    number_color  = pal["hot"]       if dark else pal["accent"]
    title_color   = WHITE            if dark else INK_DARK
    accent_color  = pal["highlight"]            # was hardcoded MERCK_GOLD

    # Huge serif number at left.
    txt(slide, Inches(0.65), Inches(1.80), Inches(5.0), Inches(3.20),
        str(num_str or number or ""), sz=140, color=number_color, bold=True,
        font=FONT_HEAD, anchor=MSO_ANCHOR.MIDDLE)
    # Title to the right (uppercase tracking-wide).
    txt(slide, Inches(5.50), Inches(3.10), Inches(7.5), Inches(0.40),
        _tracked(title), sz=13, color=accent_color, bold=True, font=FONT_BODY,
        anchor=MSO_ANCHOR.TOP)
    # Italic phrase below.
    txt(slide, Inches(5.50), Inches(3.55), Inches(7.5), Inches(1.20),
        str(title), sz=32, color=title_color, italic=True, bold=True,
        font=FONT_HEAD)
    # Accent rule — colour from palette (was hardcoded MERCK_GOLD).
    hairline(slide, Inches(5.50), Inches(5.05), Inches(2.4),
             Emu(int(Pt(2.5))), accent_color)

    # Emit takeaway band and/or source line if supplied.
    if takeaway:
        _takeaway_band(slide, takeaway, style)
    if source:
        _source_line(slide, Inches(0.65), SOURCE_Y, Inches(12.0), SOURCE_H,
                     source, style)
    return slide


# ===========================================================================
# Layout: CHART SLIDE
# ===========================================================================

def _auto_callout_for_chart(chart):
    """Generate a fallback callout for a chart when none was provided.

    Returns a single callout dict {x, y, label, direction} or None. Aimed at
    the first highlighted/relevant data point. The agent normally overrides
    this, but it ensures every data slide has at least one annotation.
    """
    if not chart:
        return None
    ctype = chart.get("type", "")
    data = chart.get("data", {})
    if ctype == "slope":
        items = data.get("items", [])
        highlights = data.get("highlight_indices") or []
        if not items or not highlights:
            return None
        idx = highlights[0]
        if idx >= len(items):
            return None
        item = items[idx]
        label = item[0] if isinstance(item, (list, tuple)) else item.get("label", "")
        try:
            before, after = item[1], item[2]
        except Exception:
            before, after = 0, 0
        delta = float(after) - float(before)
        direction = "decline" if delta < 0 else "gain"
        return {
            "x": Inches(10.4), "y": Inches(4.7),
            "label": f"Largest {direction} at: {label}",
            "direction": "up_left",
        }
    if ctype == "waterfall":
        bars = data.get("bars", [])
        if not bars:
            return None
        start_v, end_v = None, None
        for b in bars:
            t = b.get("type", "")
            if t == "start" and start_v is None:
                start_v = float(b.get("value", 0))
            if t == "end":
                end_v = float(b.get("value", 0))
        if start_v is not None and end_v is not None:
            delta = end_v - start_v
            sign = "+" if delta >= 0 else ""
            return {
                "x": Inches(6.5), "y": Inches(4.0),
                "label": f"Net change: {sign}{delta:g}",
                "direction": "up_right",
            }
        return None
    if ctype == "dot":
        items = data.get("items", [])
        if not items:
            return None
        # Top performer: max value.
        def _val(it):
            try:
                return float(str(it[1]).replace("%", "").strip())
            except Exception:
                return 0.0
        top = max(items, key=_val)
        label = top[0] if isinstance(top, (list, tuple)) else top.get("label", "")
        return {
            "x": Inches(10.4), "y": Inches(4.0),
            "label": f"Top performer: {label}",
            "direction": "up_left",
        }
    if ctype == "marimekko":
        cols = data.get("columns", [])
        if not cols:
            return None
        # Largest column by weight.
        top = max(cols, key=lambda c: c.get("weight", 0))
        return {
            "x": Inches(2.4), "y": Inches(6.0),
            "label": f"Largest column: {top.get('label', '')}",
            "direction": "down_left",
        }
    return None


def build_chart_slide(prs, meta, action_title=None, chart=None, takeaway=None,
                      source=None, callouts=None, category=None, subtitle=None,
                      style="merck_executive", page=None, total=None,
                      section_number=None, footnotes=None,
                      methodology_note=None, content=None):
    """Chart slide. chart={'type','data'}, optional callouts and footnotes.

    If callouts is None and the chart type benefits from a fallback callout,
    a single auto-callout is generated from the chart data.
    """
    if content:
        if "chart"    in content: chart    = content["chart"]
        if "callouts" in content: callouts = content["callouts"]
        if "takeaway" in content: takeaway = content["takeaway"]
    style = _style_or_promote(category, style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 subtitle=subtitle, takeaway=takeaway, source=source,
                 page=page, total=total, palette=style,
                 section_number=section_number,
                 methodology_note=methodology_note)

    chart_y = Inches(2.85) if subtitle else Inches(2.50)
    chart_h = Inches(6.50) - chart_y
    chart_rendered = _render_chart(slide, chart, Inches(0.65), chart_y,
                                   Inches(12.0), chart_h, style)
    if not chart_rendered:
        # Show a placeholder instead of a silently blank content area.
        txt(slide, Inches(0.65), chart_y + Inches(0.50), Inches(12.0),
            Inches(1.20),
            "[Chart data not available — populate chart.type and chart.data in the plan]",
            sz=13, color=INK_GRAY, italic=True, font=FONT_BODY)

    # Auto-fallback callout when none provided.
    effective_callouts = callouts
    if effective_callouts is None:
        auto = _auto_callout_for_chart(chart)
        if auto:
            effective_callouts = [auto]
    if effective_callouts:
        for c in effective_callouts:
            if c.get("x") is not None or c.get("y") is not None:
                # Slide-absolute EMU coordinates (from _auto_callout_for_chart or
                # hand-crafted plans that use the internal x/y keys directly).
                cx = c.get("x") or Inches(6.0)
                cy = c.get("y") or chart_y + Inches(1.5)
            else:
                # LLM plan schema: x_in/y_in are chart-relative float inches.
                # Add chart_y so the annotation lands inside the chart area, not
                # on top of the action title (which sits above chart_y).
                cx = Inches(float(c.get("x_in", 6.0)))
                cy = chart_y + Inches(float(c.get("y_in", 1.5)))
            stub_and_flag(slide, cx, cy,
                          c.get("label") or c.get("text", ""), style,
                          direction=c.get("direction", "up_right"))

    # Footnotes block above the source line (chrome already drew source/takeaway).
    if footnotes:
        footnotes_block(slide, footnotes, style)
    return slide


# ===========================================================================
# Framework presets for non-matrix layouts.
# Set "framework" in content to auto-populate columns/items/phases/spokes.
# ===========================================================================

_COLUMN_PRESETS = {
    "4ps": {
        "columns": [
            {"label": "PRODUCT",   "body": "What you sell"},
            {"label": "PRICE",     "body": "What you charge"},
            {"label": "PLACE",     "body": "How you distribute"},
            {"label": "PROMOTION", "body": "How you communicate"},
        ],
    },
    "value_disciplines": {
        "columns": [
            {"label": "PRODUCT LEADERSHIP",    "body": "New products, new markets, experimentation"},
            {"label": "OPERATIONAL EXCELLENCE", "body": "Efficiency, cost, volume, narrow product lines"},
            {"label": "CUSTOMER INTIMACY",      "body": "Deep customer focus, high expertise in chosen areas"},
        ],
    },
    "balanced_scorecard": {
        "columns": [
            {"label": "FINANCIAL",  "body": "How do we look to shareholders?"},
            {"label": "CUSTOMERS",  "body": "How do customers see us?"},
            {"label": "PROCESSES",  "body": "What must we excel at?"},
            {"label": "LEARNING",   "body": "How can we improve and create value?"},
        ],
    },
}

_NUMBERED_PRESETS = {
    "kotter_8": {
        "items": [
            {"number": "1", "title": "Create Urgency",         "body": "Help others see the need for change"},
            {"number": "2", "title": "Build the Coalition",    "body": "Assemble a group with power to lead change"},
            {"number": "3", "title": "Form Strategic Vision",  "body": "Shape a vision and initiatives"},
            {"number": "4", "title": "Enlist a Volunteer Army","body": "Raise a large force of people"},
            {"number": "5", "title": "Enable Action",          "body": "Remove barriers to change"},
            {"number": "6", "title": "Generate Short-Term Wins","body": "Recognise and reward progress"},
            {"number": "7", "title": "Sustain Acceleration",   "body": "Press harder and faster"},
            {"number": "8", "title": "Institute Change",       "body": "Reinforce new behaviours"},
        ],
    },
    "adkar": {
        "items": [
            {"number": "1", "title": "Awareness",   "body": "Of the need for change"},
            {"number": "2", "title": "Desire",      "body": "To participate and support the change"},
            {"number": "3", "title": "Knowledge",   "body": "On how to change"},
            {"number": "4", "title": "Ability",     "body": "To implement required skills and behaviours"},
            {"number": "5", "title": "Reinforcement","body": "To sustain the change"},
        ],
    },
}

_CIRCULAR_PRESETS = {
    "pdca": {
        "phases": [
            {"label": "PLAN",  "body": "Identify the problem and plan the change"},
            {"label": "DO",    "body": "Implement the change on a small scale"},
            {"label": "CHECK", "body": "Analyse results and identify learnings"},
            {"label": "ACT",   "body": "Scale up or adjust based on findings"},
        ],
    },
    "dmaic": {
        "phases": [
            {"label": "DEFINE",   "body": "Define the problem and goals"},
            {"label": "MEASURE",  "body": "Measure current process performance"},
            {"label": "ANALYSE",  "body": "Identify root causes"},
            {"label": "IMPROVE",  "body": "Implement and verify solutions"},
            {"label": "CONTROL",  "body": "Sustain the improvements"},
        ],
    },
}

_SPOKE_PRESETS = {
    "porters_5": {
        "hub": {"label": "Competitive Rivalry", "body": "Intensity among existing competitors"},
        "spokes": [
            {"label": "Supplier Power",     "body": "Bargaining power of suppliers"},
            {"label": "Buyer Power",        "body": "Bargaining power of customers"},
            {"label": "Threat of New Entry","body": "Ease of entry for new competitors"},
            {"label": "Substitution Risk",  "body": "Threat from substitute products"},
        ],
    },
}


# ===========================================================================
# Layout: TWO COLUMN
# ===========================================================================

def _two_or_three_column_card(slide, x, y, w, h, col, palette):
    """Render a column card with dark-bar header + tone dot treatment.

    Both normal and highlighted cards use PANEL_LIGHT fill. The PURPLE_DEEP
    header bar blends into the card chrome. A small dot signals tone:
    PURPLE_MUTED (neutral/normal) or theme hot colour (highlighted).

    Highlight derives from tone="positive" automatically.
    An explicit "highlighted" key overrides tone when present.
    """
    pal = _palette_for(palette)
    if "highlighted" in col:
        highlighted = bool(col["highlighted"])
    else:
        highlighted = col.get("tone") == "positive"

    # Card background — always PANEL_LIGHT.
    card = rounded(slide, x, y, w, h, fill=PANEL_LIGHT)
    _apply_border(card, LIGHT_GRAY, Pt(0.5))

    # Header bar: MERCK_PURPLE (#503291) normal, LY_CYAN (#2DBECD) highlighted
    HDR_H    = Inches(0.40)
    DOT_SZ   = Inches(0.16)
    bar_fill = LY_CYAN if highlighted else MERCK_PURPLE
    lbl_col  = PURPLE_DEEP if highlighted else WHITE   # dark text on cyan, white on purple
    rect(slide, x, y, w, HDR_H, fill=bar_fill)

    # Tone dot
    dot_col = pal["hot"] if highlighted else PURPLE_MUTED
    dot_x   = x + Inches(0.16)
    dot_y   = y + (HDR_H - DOT_SZ) / 2
    circle(slide, dot_x, dot_y, DOT_SZ, fill=dot_col)

    # Label text on bar — accept "label" (internal name) or "header" (schema name).
    label = col.get("label") or col.get("header", "")
    if label:
        hdr_box = txt(slide, dot_x + DOT_SZ + Inches(0.10), y,
                      w - DOT_SZ - Inches(0.40), HDR_H,
                      _tracked(label), sz=9, color=lbl_col, bold=True,
                      font=FONT_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        try:
            from pptx.enum.text import MSO_AUTO_SIZE
            hdr_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except Exception:
            pass

    # Body content below the bar
    pad   = Inches(0.20)
    cur_y = y + HDR_H + Inches(0.18)

    body_text = col.get("body", "")
    items     = col.get("items") or col.get("bullets") or []

    if body_text:
        hero_h = Inches(0.95)
        if len(str(body_text)) > 60:
            hero_h = Inches(1.30)
        if len(str(body_text)) > 120:
            hero_h = Inches(1.55)
        hero_sz = 16 if len(str(body_text)) <= 60 else 14
        txt(slide, x + pad, cur_y, w - pad * 2, hero_h,
            body_text, sz=hero_sz, color=MERCK_PURPLE, bold=True,
            font=FONT_BODY)
        cur_y += hero_h + Inches(0.05)

    if items:
        avail_h = (y + h) - cur_y - Inches(0.50)
        _bulleted_list(slide, x + pad, cur_y, w - pad * 2,
                       max(avail_h, Inches(0.30)),
                       items, palette, text_color=INK_DARK, sz=11,
                       bullet_color=pal["highlight"])

    # Timestamp pill at bottom
    ts = col.get("timestamp")
    if ts:
        pill_w = Inches(1.6)
        pill_h = Inches(0.32)
        pill_x = x + pad
        pill_y = y + h - pill_h - Inches(0.18)
        rect(slide, pill_x, pill_y, pill_w, pill_h,
             fill=pal["hot"] if highlighted else PURPLE_MUTED)
        txt(slide, pill_x, pill_y, pill_w, pill_h,
            str(ts), sz=10,
            color=PURPLE_DEEP if highlighted else WHITE,
            bold=True, font=FONT_BODY,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def build_two_column(prs, meta, action_title=None, left=None, right=None, takeaway=None,
                     source=None, category=None, style="merck_executive",
                     page=None, total=None, section_number=None,
                     methodology_note=None, subtitle=None, content=None):
    if content:
        if "left"     in content: left     = content["left"]
        if "right"    in content: right    = content["right"]
        if "takeaway" in content: takeaway = content["takeaway"]
    style = _style_or_promote(category, style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 subtitle=subtitle,
                 takeaway=takeaway, source=source, page=page, total=total,
                 palette=style, section_number=section_number,
                 methodology_note=methodology_note)
    zone_top = Inches(2.95) if subtitle else Inches(2.55)
    zone_h = Inches(6.50) - zone_top
    gap = Inches(0.30)
    col_w = (Inches(12.0) - gap) / 2
    _two_or_three_column_card(slide, Inches(0.65), zone_top, col_w, zone_h,
                              left or {}, style)
    _two_or_three_column_card(slide, Inches(0.65) + col_w + gap, zone_top,
                              col_w, zone_h, right or {}, style)
    return slide


def build_three_column(prs, meta, action_title=None, columns=None, takeaway=None,
                       source=None, category=None, style="merck_executive",
                       page=None, total=None, section_number=None,
                       methodology_note=None, subtitle=None, content=None):
    if content:
        fw = content.get("framework")
        if fw and fw in _COLUMN_PRESETS and not columns:
            columns = _COLUMN_PRESETS[fw]["columns"]
        if "columns"  in content: columns  = content["columns"]
        if "takeaway" in content: takeaway = content["takeaway"]
    style = _style_or_promote(category, style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 subtitle=subtitle,
                 takeaway=takeaway, source=source, page=page, total=total,
                 palette=style, section_number=section_number,
                 methodology_note=methodology_note)
    cols = (columns or [])[:3]
    while len(cols) < 3:
        cols.append({})
    zone_top = Inches(2.95) if subtitle else Inches(2.55)
    zone_h = Inches(6.50) - zone_top
    gap = Inches(0.25)
    col_w = (Inches(12.0) - gap * 2) / 3
    for i, col in enumerate(cols):
        cx = Inches(0.65) + i * (col_w + gap)
        _two_or_three_column_card(slide, cx, zone_top, col_w, zone_h, col, style)
    return slide


# ===========================================================================
# Layout: 2x2 MATRIX
# ===========================================================================

# Framework presets: set "framework" in content to auto-configure axes and quadrants.
# Explicit x_axis/y_axis/quadrants in content always override the preset.
_MATRIX_PRESETS = {
    "bcg": {
        "x_axis": "Relative Market Share",
        "y_axis": "Market Growth Rate",
        "quadrants": {
            "top_left":     {"label": "Stars",           "highlighted": True},
            "top_right":    {"label": "Question Marks"},
            "bottom_left":  {"label": "Cash Cows"},
            "bottom_right": {"label": "Dogs"},
        },
    },
    "swot": {
        "x_axis": "Internal / External",
        "y_axis": "Positive / Negative",
        "quadrants": {
            "top_left":     {"label": "Strengths",      "highlighted": True},
            "top_right":    {"label": "Opportunities"},
            "bottom_left":  {"label": "Weaknesses"},
            "bottom_right": {"label": "Threats"},
        },
    },
    "ansoff": {
        "x_axis": "Products (Existing → New)",
        "y_axis": "Markets (Existing → New)",
        "quadrants": {
            "top_left":     {"label": "Market Development"},
            "top_right":    {"label": "Diversification",   "highlighted": True},
            "bottom_left":  {"label": "Market Penetration"},
            "bottom_right": {"label": "Product Development"},
        },
    },
    "risk": {
        "x_axis": "Probability (Low → High)",
        "y_axis": "Impact (Low → High)",
        "quadrants": {
            "top_left":     {"label": "Monitor"},
            "top_right":    {"label": "Critical",          "highlighted": True},
            "bottom_left":  {"label": "Accept"},
            "bottom_right": {"label": "Mitigate"},
        },
    },
}


def build_2x2_matrix(prs, meta, action_title=None, x_axis=None, y_axis=None, quadrants=None,
                     takeaway=None, source=None, category=None,
                     style="merck_executive", page=None, total=None,
                     section_number=None, methodology_note=None,
                     subtitle=None, content=None):
    if content:
        framework = content.get("framework")
        if framework and framework in _MATRIX_PRESETS:
            preset = _MATRIX_PRESETS[framework]
            if not x_axis:    x_axis    = preset.get("x_axis")
            if not y_axis:    y_axis    = preset.get("y_axis")
            if not quadrants: quadrants = preset.get("quadrants")
        if "x_axis"    in content: x_axis    = content["x_axis"]
        if "y_axis"    in content: y_axis    = content["y_axis"]
        if "quadrants" in content: quadrants = content["quadrants"]
        if "takeaway"  in content: takeaway  = content["takeaway"]
    style = _style_or_promote(category, style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 subtitle=subtitle,
                 takeaway=takeaway, source=source, page=page, total=total,
                 palette=style, section_number=section_number,
                 methodology_note=methodology_note)

    margin_left = Inches(0.55)
    margin_bottom = Inches(0.55)
    zone_top = Inches(2.95) if subtitle else Inches(2.55)
    zone_h = Inches(6.50) - zone_top
    plot_x = Inches(0.65) + margin_left
    plot_y = zone_top + Inches(0.05)
    plot_w = Inches(12.0) - margin_left
    plot_h = zone_h - margin_bottom - Inches(0.10)
    cell_w = plot_w / 2
    cell_h = plot_h / 2

    cells = ["top_left", "top_right", "bottom_left", "bottom_right"]
    positions = [
        (plot_x, plot_y),
        (plot_x + cell_w, plot_y),
        (plot_x, plot_y + cell_h),
        (plot_x + cell_w, plot_y + cell_h),
    ]
    for k, key in enumerate(cells):
        cx, cy = positions[k]
        q = (quadrants or {}).get(key, {})
        highlighted = bool(q.get("highlighted"))
        if highlighted:
            quad_fill = PURPLE_DEEP if _is_dark(style) else MERCK_PURPLE
            rounded(slide, cx, cy, cell_w, cell_h, fill=quad_fill)
            rounded(slide, cx, cy, cell_w, Inches(0.06), fill=pal["hot"], adj=50000)
            tcolor = WHITE
            btext = PANEL_LIGHT
            label_color = pal["hot"]
        else:
            cell = rounded(slide, cx, cy, cell_w, cell_h, fill=PANEL_LIGHT)
            _apply_border(cell, LIGHT_GRAY, Pt(0.5))
            rounded(slide, cx, cy, cell_w, Inches(0.06), fill=pal["highlight"], adj=50000)
            tcolor = MERCK_PURPLE
            btext = INK_DARK
            label_color = pal["highlight"]
        label = q.get("label", "")
        items = q.get("items") or q.get("bullets") or []
        body  = q.get("body", "")
        pad = Inches(0.18)
        txt(slide, cx + pad, cy + Inches(0.18), cell_w - pad * 2, Inches(0.26),
            _tracked(label), sz=10, color=label_color, bold=True,
            font=FONT_BODY)
        if items:
            _bulleted_list(slide, cx + pad, cy + Inches(0.50),
                           cell_w - pad * 2, cell_h - Inches(0.65),
                           items, palette=style, text_color=btext, sz=10,
                           bullet_color=label_color)
        elif body:
            txt(slide, cx + pad, cy + Inches(0.50),
                cell_w - pad * 2, cell_h - Inches(0.65),
                str(body), sz=11, color=btext, font=FONT_BODY,
                anchor=MSO_ANCHOR.TOP)

    # Axes lines (light).
    line(slide, plot_x, plot_y + plot_h, plot_x + plot_w, plot_y + plot_h,
         PURPLE_MUTED, Pt(0.75))
    line(slide, plot_x, plot_y, plot_x, plot_y + plot_h, PURPLE_MUTED, Pt(0.75))
    ah = Inches(0.12)
    _freeform_poly(slide,
                   [(plot_x + plot_w + ah, plot_y + plot_h),
                    (plot_x + plot_w, plot_y + plot_h - ah / 2),
                    (plot_x + plot_w, plot_y + plot_h + ah / 2)],
                   fill=PURPLE_MUTED)
    _freeform_poly(slide,
                   [(plot_x, plot_y - ah),
                    (plot_x - ah / 2, plot_y),
                    (plot_x + ah / 2, plot_y)],
                   fill=PURPLE_MUTED)

    x_ax = x_axis or {}
    y_ax = y_axis or {}
    txt(slide, plot_x, plot_y + plot_h + Inches(0.08), plot_w, Inches(0.26),
        _tracked(x_ax.get("label", "")), sz=10, color=pal["highlight"], bold=True,
        font=FONT_BODY, align=PP_ALIGN.CENTER)
    txt(slide, plot_x, plot_y + plot_h + Inches(0.30), Inches(1.5),
        Inches(0.22),
        x_ax.get("low", ""), sz=9, color=PURPLE_MUTED, font=FONT_BODY)
    txt(slide, plot_x + plot_w - Inches(1.5), plot_y + plot_h + Inches(0.30),
        Inches(1.5), Inches(0.22),
        x_ax.get("high", ""), sz=9, color=PURPLE_MUTED, font=FONT_BODY,
        align=PP_ALIGN.RIGHT)
    txt(slide, Inches(0.65), plot_y, margin_left, plot_h,
        _tracked(y_ax.get("label", "")), sz=10, color=pal["highlight"], bold=True,
        font=FONT_BODY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return slide


# ===========================================================================
# Layout: PHASE PROCESS
# ===========================================================================

def build_phase_process(prs, meta, action_title=None, phases=None, takeaway=None,
                        source=None, category=None, style="merck_executive",
                        page=None, total=None, section_number=None,
                        highlight_index=None, show_arrows=True,
                        methodology_note=None, subtitle=None, content=None):
    """Horizontal phased flow as a series of cards with arrow connectors.

    Each phase dict: {label, title, body, milestone, highlighted}.
    highlight_index (0-based int): when set, only this index is highlighted
        and overrides per-phase 'highlighted' flags.
    show_arrows (bool): draw a MERCK_GOLD right-arrow in the gutter between
        adjacent cards (default True).

    All cards are rendered at the same height (equal-height enforcement) so
    they line up as a row of equals.
    """
    if content:
        if "phases"          in content: phases          = content["phases"]
        if "highlight_index" in content: highlight_index = content["highlight_index"]
        if "show_arrows"     in content: show_arrows     = content["show_arrows"]
        if "takeaway"        in content: takeaway        = content["takeaway"]
    style = _style_or_promote(category, style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 subtitle=subtitle,
                 takeaway=takeaway, source=source, page=page, total=total,
                 palette=style, section_number=section_number,
                 methodology_note=methodology_note)

    phs = (phases or [])[:5]
    if not phs:
        return slide
    n = len(phs)
    zone_top = Inches(2.95) if subtitle else Inches(2.55)
    zone_h = Inches(6.50) - zone_top
    # Reserve a top strip for the "STEP a -> STEP b -> ..." breadcrumb.
    breadcrumb_h = Inches(0.32)
    cards_top = zone_top + breadcrumb_h + Inches(0.10)
    base_cards_h = zone_h - breadcrumb_h - Inches(0.10)
    arrow_w = Inches(0.36) if show_arrows else Inches(0.18)
    total_w = Inches(12.0)
    card_w = (total_w - arrow_w * (n - 1)) / n

    # Equal-height enforcement based on content estimates.
    estimated = []
    for ph in phs:
        # Per-card content height estimate based on title length, body, milestone.
        est_h = Inches(0.50)  # top label
        title_len = len(str(ph.get("title", "")))
        if title_len:
            title_lines = max(1, (title_len // 22) + (1 if title_len % 22 else 0))
            est_h += Inches(0.32) * title_lines
        body_len = len(str(ph.get("body", "")))
        if body_len:
            body_lines = max(1, (body_len // 32) + (1 if body_len % 32 else 0))
            est_h += Inches(0.22) * body_lines
        if ph.get("milestone"):
            est_h += Inches(0.50)
        est_h += Inches(0.40)  # bottom padding
        estimated.append(est_h)
    cards_h = max(base_cards_h, max(estimated))
    # Clamp so cards don't run past the source line.
    max_h_avail = Inches(6.40) - cards_top
    if cards_h > max_h_avail:
        cards_h = max_h_avail

    # Breadcrumb across the top: "STEP 1 -> STEP 2 -> ...".
    crumb_parts = []
    for i, ph in enumerate(phs):
        lbl = ph.get("label", f"STEP {i + 1}")
        crumb_parts.append(_tracked(lbl))
    crumb = "   →   ".join(crumb_parts)
    txt(slide, Inches(0.65), zone_top, total_w, breadcrumb_h,
        crumb, sz=10, color=pal["highlight"], bold=True, font=FONT_BODY,
        anchor=MSO_ANCHOR.MIDDLE)

    for i, ph in enumerate(phs):
        cx = Inches(0.65) + i * (card_w + arrow_w)
        cy = cards_top
        if "highlighted" in ph:
            highlighted = bool(ph["highlighted"])
        elif ph.get("status") == "current":
            highlighted = True
        elif highlight_index is not None:
            highlighted = (i == int(highlight_index))
        else:
            highlighted = False
        _draw_card(slide, cx, cy, card_w, cards_h, style,
                   highlighted=highlighted)

        pad = Inches(0.20)
        # Tiny step label inside card top.
        label_color = pal["hot"] if highlighted else pal["highlight"]
        sub_color = WHITE if highlighted else INK_GRAY
        title_color = WHITE if highlighted else MERCK_PURPLE
        txt(slide, cx + pad, cy + Inches(0.20), card_w - pad * 2,
            Inches(0.22),
            _tracked(ph.get("label", f"STEP {i + 1}")), sz=9,
            color=label_color, bold=True, font=FONT_BODY)

        title = ph.get("title", "")
        if title:
            txt(slide, cx + pad, cy + Inches(0.50), card_w - pad * 2,
                Inches(0.95),
                title, sz=15, color=title_color, bold=True, font=FONT_BODY)

        # Decide milestone treatment before drawing body so we can shrink the
        # body region when a long milestone caption needs vertical room.
        ms = ph.get("milestone")
        ms_text = str(ms) if ms else ""
        ms_is_long = bool(ms_text) and len(ms_text) > 32
        # Default body bottom margin: 0.60" reserves room for the 0.32" pill +
        # 0.18" gap + 0.10" slack. For long milestones we reserve 1.00" so the
        # 0.80"-tall caption sits cleanly below the body.
        body_bottom_margin = Inches(1.00) if ms_is_long else Inches(0.60)

        body = ph.get("body", "")
        if body:
            body_h = cards_h - Inches(1.50) - body_bottom_margin
            if body_h < Inches(0.30):
                body_h = Inches(0.30)
            txt(slide, cx + pad, cy + Inches(1.50),
                card_w - pad * 2, body_h,
                body, sz=10, color=sub_color, font=FONT_BODY)

        # Milestone at the bottom. Short text (<=32 chars) renders as a
        # colored pill. Long text renders as an italic gold caption (no fill)
        # so it can wrap to 2-3 lines without overlapping body text — the
        # previous fixed 0.32" pill clipped long milestones into the body,
        # producing strikethrough-style render artifacts.
        if ms:
            block_w = card_w - pad * 2
            block_x = cx + pad
            if ms_is_long:
                cap_h = Inches(0.80)
                cap_y = cy + cards_h - cap_h - Inches(0.10)
                cap_color = pal["hot"] if highlighted else pal["highlight"]
                txt(slide, block_x, cap_y, block_w, cap_h,
                    ms_text, sz=9, color=cap_color, italic=True, bold=True,
                    font=FONT_BODY, align=PP_ALIGN.LEFT,
                    anchor=MSO_ANCHOR.TOP)
            else:
                pill_h = Inches(0.32)
                pill_y = cy + cards_h - pill_h - Inches(0.18)
                if highlighted:
                    if _is_dark(style):
                        pill_fill = PURPLE_DEEP
                        pill_text = WHITE
                    else:
                        pill_fill = pal["hot"]
                        pill_text = PURPLE_DEEP
                else:
                    if _is_dark(style):
                        pill_fill = PURPLE_DEEP
                    else:
                        pill_fill = MERCK_PURPLE
                    pill_text = WHITE
                rect(slide, block_x, pill_y, block_w, pill_h, fill=pill_fill)
                txt(slide, block_x, pill_y, block_w, pill_h,
                    ms_text, sz=10, color=pill_text, bold=True,
                    font=FONT_BODY, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)

        # Arrow connector in the gutter to the next card.
        if i < n - 1 and show_arrows:
            ax = cx + card_w
            ay = cy + cards_h / 2
            ar_w = Inches(0.30)
            ar_h = Inches(0.40)
            ar_x = ax + (arrow_w - ar_w) / 2
            ar_y = ay - ar_h / 2
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                           ar_x, ar_y, ar_w, ar_h)
            arrow.shadow.inherit = False
            _apply_fill(arrow, pal["highlight"])
            _apply_border(arrow, None)
    return slide


# ===========================================================================
# Layout: VERTICAL NUMBERED
# ===========================================================================

def build_vertical_numbered(prs, meta, action_title=None, items=None, takeaway=None,
                            source=None, category=None, style="merck_executive",
                            page=None, total=None, section_number=None,
                            methodology_note=None, subtitle=None, content=None):
    if content:
        fw = content.get("framework")
        if fw and fw in _NUMBERED_PRESETS and not items:
            items = _NUMBERED_PRESETS[fw]["items"]
        if "items"    in content: items    = content["items"]
        if "takeaway" in content: takeaway = content["takeaway"]
    style = _style_or_promote(category, style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 subtitle=subtitle,
                 takeaway=takeaway, source=source, page=page, total=total,
                 palette=style, section_number=section_number, methodology_note=methodology_note)

    rows = (items or [])[:5]
    n = max(len(rows), 1)
    # Scale number and title font sizes down for dense layouts so badges
    # don't overflow into adjacent rows.
    num_sz   = 34 if n >= 5 else (38 if n == 4 else 44)
    title_sz = 13 if n >= 5 else (14 if n == 4 else 15)
    zone_top = Inches(2.95) if subtitle else Inches(2.55)
    zone_h = Inches(6.50) - zone_top
    gap = Inches(0.12)
    row_h = (zone_h - gap * (n - 1)) / n
    for i, row in enumerate(rows):
        ry = zone_top + i * (row_h + gap)
        num_w = Inches(0.90)
        # Number badge in MERCK_PURPLE.
        txt(slide, Inches(0.65), ry, num_w, row_h, str(i + 1),
            sz=num_sz, color=MERCK_PURPLE, bold=True, font=FONT_BODY,
            anchor=MSO_ANCHOR.TOP)
        bx = Inches(0.65) + num_w + Inches(0.25)
        bw = Inches(12.0) - num_w - Inches(0.35)
        txt(slide, bx, ry + Inches(0.04), bw, Inches(0.36),
            row.get("title", ""), sz=title_sz, color=INK_DARK, bold=True,
            font=FONT_BODY)
        txt(slide, bx, ry + Inches(0.42), bw, row_h - Inches(0.46),
            row.get("body", ""), sz=11, color=INK_GRAY, font=FONT_BODY)
        if i < n - 1:
            hairline(slide, Inches(0.65),
                     ry + row_h + gap / 2 - Emu(int(Pt(0.5))),
                     Inches(12.0), Emu(int(Pt(0.5))), LIGHT_GRAY)
    return slide


# ===========================================================================
# Layout: WATERFALL
# ===========================================================================

def build_waterfall_slide(prs, meta, action_title=None, bars=None, takeaway=None,
                          source=None, category=None, style="merck_executive",
                          page=None, total=None, section_number=None,
                          footnotes=None, methodology_note=None,
                          subtitle=None, content=None):
    if content:
        if "bars"      in content: bars      = content["bars"]
        if "takeaway"  in content: takeaway  = content["takeaway"]
        if "footnotes" in content: footnotes = content["footnotes"]
    style = _style_or_promote(category, style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 subtitle=subtitle,
                 takeaway=takeaway, source=source, page=page, total=total,
                 palette=style, section_number=section_number, methodology_note=methodology_note)
    chart_y = Inches(2.95) if subtitle else Inches(2.55)
    chart_h = Inches(6.50) - chart_y
    add_waterfall(slide, Inches(0.65), chart_y, Inches(12.0), chart_h,
                  bars or [], style)
    if footnotes:
        footnotes_block(slide, footnotes, style)
    return slide


# ===========================================================================
# Layout: DECISION ROWS
# ===========================================================================

def build_decision_rows(prs, meta, action_title=None, decisions=None, takeaway=None,
                        source=None, category=None, page=None, total=None,
                        style=None, section_number=None, methodology_note=None,
                        subtitle=None, content=None):
    """Auto-promoted to executive style (unless overridden by an exec-eligible
    category). The style kwarg is honored when the category does not trigger
    auto-promotion.
    """
    if content:
        if "decisions" in content: decisions = content["decisions"]
        if "takeaway"  in content: takeaway  = content["takeaway"]
    style = _style_or_promote(category or "Decision Request", style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title,
                 category=category or "Decision Request",
                 subtitle=subtitle,
                 takeaway=takeaway, source=source, page=page, total=total,
                 palette=style, section_number=section_number, methodology_note=methodology_note)

    rows = (decisions or [])[:5]
    n = max(len(rows), 1)
    zone_top = Inches(2.95) if subtitle else Inches(2.55)
    zone_h = Inches(6.50) - zone_top
    gap = Inches(0.12)
    row_h = (zone_h - gap * (n - 1)) / n
    for i, d in enumerate(rows):
        ry = zone_top + i * (row_h + gap)
        # Card chrome (panel + hairline + theme-accent stripe on top).
        card = rect(slide, Inches(0.65), ry, Inches(12.0), row_h,
                    fill=PANEL_LIGHT)
        _apply_border(card, LIGHT_GRAY, Pt(0.5))
        rect(slide, Inches(0.65), ry, Inches(12.0), Inches(0.06),
             fill=pal["highlight"])
        # Number circle.
        tone = d.get("tone", "neutral")
        tone_color = _tone_color(tone, style)
        circ_d = Inches(0.55)
        cx = Inches(0.95)
        cy = ry + (row_h - circ_d) / 2
        circle(slide, cx, cy, circ_d, fill=tone_color)
        txt(slide, cx, cy, circ_d, circ_d, str(d.get("number") or (i + 1)),
            sz=18, color=WHITE, bold=True, font=FONT_BODY,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Title + description.
        bx = cx + circ_d + Inches(0.25)
        owner = d.get("owner", "")
        owner_w = Inches(2.4) if owner else Inches(0.0)
        bw = Inches(12.0) - (bx - Inches(0.65)) - owner_w - Inches(0.20)
        title_y = ry + Inches(0.20)
        txt(slide, bx, title_y, bw, Inches(0.34),
            d.get("title") or d.get("decision", ""), sz=14, color=MERCK_PURPLE, bold=True,
            font=FONT_BODY)
        txt(slide, bx, title_y + Inches(0.36), bw,
            row_h - Inches(0.60),
            d.get("desc") or d.get("body") or d.get("description", ""), sz=11, color=INK_GRAY, font=FONT_BODY)
        if owner:
            txt(slide, Inches(0.65) + Inches(12.0) - owner_w - Inches(0.15),
                ry, owner_w, row_h,
                f"Owner: {owner}", sz=10, color=PURPLE_MUTED, italic=True,
                font=FONT_BODY, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    return slide


# ===========================================================================
# Layout: GANTT
# ===========================================================================

def build_gantt(prs, meta, action_title=None, rows=None, quarters=None, takeaway=None,
                source=None, category=None, style="merck_executive",
                page=None, total=None, section_number=None,
                methodology_note=None, subtitle=None, content=None):
    if content:
        if "rows"     in content: rows     = content["rows"]
        if "quarters" in content: quarters = content["quarters"]
        if "takeaway" in content: takeaway = content["takeaway"]
    style = _style_or_promote(category, style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 subtitle=subtitle,
                 takeaway=takeaway, source=source, page=page, total=total,
                 palette=style, section_number=section_number, methodology_note=methodology_note)

    zone_top = Inches(2.95) if subtitle else Inches(2.55)
    zone_h = Inches(6.50) - zone_top
    qs = list(quarters or ["Q1", "Q2", "Q3", "Q4"])
    qn_count = max(len(qs), 1)
    label_w = Inches(2.6)
    grid_x = Inches(0.65) + label_w
    grid_w = Inches(12.0) - label_w
    header_h = Inches(0.40)
    header_y = zone_top

    col_w = grid_w / qn_count
    for i, q in enumerate(qs):
        txt(slide, grid_x + i * col_w, header_y, col_w, header_h,
            _tracked(str(q)), sz=10, color=pal["highlight"], bold=True,
            font=FONT_BODY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        line(slide, grid_x + i * col_w, header_y + header_h,
             grid_x + i * col_w, zone_top + zone_h,
             LIGHT_GRAY, Pt(0.5))
    line(slide, grid_x + qn_count * col_w, header_y + header_h,
         grid_x + qn_count * col_w, zone_top + zone_h,
         LIGHT_GRAY, Pt(0.5))

    body_y = header_y + header_h
    body_h = zone_h - header_h
    rs = rows or []
    n = max(len(rs), 1)
    row_h = body_h / n
    for i, r in enumerate(rs):
        ry = body_y + i * row_h
        txt(slide, Inches(0.65), ry, label_w - Inches(0.10), row_h,
            r.get("label", ""), sz=11, color=INK_DARK, font=FONT_BODY,
            anchor=MSO_ANCHOR.MIDDLE)
        hairline(slide, grid_x, ry + row_h - Emu(int(Pt(0.5))),
                 grid_w, Emu(int(Pt(0.5))), LIGHT_GRAY)
        try:
            sq = float(r.get("start_q", 1))
        except (TypeError, ValueError):
            sq = 1.0
        try:
            du = float(r.get("duration_q", 1))
        except (TypeError, ValueError):
            du = 1.0
        bar_x = grid_x + (sq - 1) * col_w
        bar_w = du * col_w
        if bar_x < grid_x:
            bar_w = bar_w - (grid_x - bar_x)
            bar_x = grid_x
        if bar_x + bar_w > grid_x + grid_w:
            bar_w = grid_x + grid_w - bar_x
        bar_w = max(bar_w, Inches(0.10))
        tone_color = _tone_color(r.get("tone", "neutral"), style)
        pad = Inches(0.10)
        rect(slide, bar_x, ry + pad, bar_w, row_h - pad * 2, fill=tone_color)
    return slide


# ===========================================================================
# Layout: HERO STAT
# ===========================================================================

def build_hero_stat(prs, meta, stat=None, context=None, source=None, category=None,
                    style="merck_storytelling", page=None, total=None,
                    section_number=None, methodology_note=None,
                    subtitle=None, action_title=None, content=None):
    """One huge stat with a one-line context. Defaults to storytelling.
    section_number, methodology_note, subtitle, and action_title are accepted
    for caller consistency but are ignored; hero_stat slides do not show a
    section marker, subtitle, or action title.
    """
    if content:
        if "stat"    in content: stat    = content["stat"]
        if "context" in content: context = content["context"]
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    _top_chrome(slide, meta, category, style, page=page, total=total)
    _bottom_chrome(slide, meta, category, page, total, style)

    stat = stat or {}
    value = stat.get("value", "")
    label = stat.get("label", "")

    dark = _is_dark(style)
    number_color = pal["hot"] if dark else MERCK_PURPLE
    label_color = WHITE if dark else INK_DARK
    context_color = PANEL_LIGHT if dark else INK_GRAY

    # Huge number centered.
    txt(slide, Inches(0.65), Inches(2.40), Inches(12.0), Inches(2.20),
        value, sz=110, color=number_color, bold=True, font=FONT_HEAD,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # One-line context underneath.
    txt(slide, Inches(0.65), Inches(4.75), Inches(12.0), Inches(0.50),
        label, sz=18, color=label_color, font=FONT_BODY,
        align=PP_ALIGN.CENTER)
    if context:
        txt(slide, Inches(0.65), Inches(5.25), Inches(12.0), Inches(1.30),
            context, sz=13, color=context_color, italic=True,
            font=FONT_BODY, align=PP_ALIGN.CENTER)
    # Theme-accent rule under the number.
    rule_color = pal["highlight"]
    hairline(slide,
             Inches(6.665) - Inches(1.2), Inches(4.55),
             Inches(2.4), Emu(int(Pt(2.5))), rule_color)
    return slide


# ===========================================================================
# Layout: CLOSE
# ===========================================================================

def build_close(prs, meta, action_statement=None, style="merck_executive",
                page=None, total=None, section_number=None,
                category=None, takeaway=None, source=None,
                action_title=None, methodology_note=None, content=None):
    """Close slide. Action statement in big italic Merck Web, gold rule.
    page and total are accepted for caller consistency; close slides typically
    omit the page number visually but the kwargs prevent a TypeError on callers
    that pass page=N, total=N as part of their consistent slide-call pattern.
    action_title is accepted as an alias for action_statement.
    methodology_note is accepted for caller consistency and ignored.
    """
    if content:
        if "action_statement" in content: action_statement = content["action_statement"]
        if "takeaway"         in content: takeaway         = content["takeaway"]
    if action_statement is None:
        action_statement = action_title
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    # Close slides: progress bar shows full (deck ends here) when total known.
    close_p = total if total is not None else page
    _top_chrome(slide, meta, None, style, page=close_p, total=total)
    _bottom_chrome(slide, meta, "Close", None, None, style)

    dark = _is_dark(style)
    rule_color = pal["highlight"]
    text_color = WHITE if dark else MERCK_PURPLE

    # Gold accent rule above.
    hairline(slide, Inches(0.65), Inches(2.80), Inches(2.6),
             Emu(int(Pt(2.5))), rule_color)

    # Action statement text box — use the full safe zone from rule to source line.
    stmt_top = Inches(3.10)
    stmt_h   = SOURCE_Y - Inches(0.15) - stmt_top  # ≈ 3.30 in (3.10 → 6.40)

    # Adaptive font: scale down for long statements so text fits without clipping.
    _stmt = action_statement or ""
    _n = (len(_stmt) if isinstance(_stmt, str)
          else sum(len(s) for s, _ in _stmt if isinstance(s, str)))
    sz_stmt = 38 if _n <= 100 else (28 if _n <= 200 else 20)

    # Action statement in italic display font. Supports either a string OR a list
    # of (text, italic_bool) tuples for mixed-emphasis runs.
    if isinstance(action_statement, (list, tuple)) and action_statement and \
       all(isinstance(seg, (list, tuple)) and len(seg) == 2
           for seg in action_statement):
        box = slide.shapes.add_textbox(Inches(0.65), stmt_top,
                                       Inches(12.0), stmt_h)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.02)
        tf.margin_right = Inches(0.02)
        p = tf.paragraphs[0]
        for seg_text, seg_italic in action_statement:
            seg_color = pal["hot"] if seg_italic else text_color
            _add_run(p, seg_text, sz=sz_stmt, color=seg_color, bold=True,
                     italic=True, font=FONT_HEAD)
    else:
        txt(slide, Inches(0.65), stmt_top, Inches(12.0), stmt_h,
            action_statement or "", sz=sz_stmt, color=text_color, italic=True,
            bold=True, font=FONT_HEAD)
    # Small tagline beneath: deck label + month/year.
    deck_label = (meta or {}).get("deck_label", "")
    month_year = (meta or {}).get("month_year", "")
    parts = []
    if deck_label:
        parts.append(_tracked(deck_label))
    if month_year:
        parts.append(_tracked(month_year))
    if parts:
        sep = "   •   "
        tag_color = PANEL_LIGHT if dark else PURPLE_MUTED
        # When a takeaway band is present, the takeaway sits at 6.83; move the
        # tagline above the gold rule (which is at 2.80) so it doesn't collide.
        tag_y = Inches(2.50) if takeaway else Inches(6.40)
        txt(slide, Inches(0.65), tag_y, Inches(12.0), Inches(0.30),
            sep.join(parts), sz=10, color=tag_color, bold=True,
            font=FONT_BODY)

    # Emit takeaway band and/or source line if supplied.
    if takeaway:
        _takeaway_band(slide, takeaway, style)
    if source:
        _source_line(slide, Inches(0.65), SOURCE_Y, Inches(12.0), SOURCE_H,
                     source, style)
    return slide


# ===========================================================================
# Layout: STAT STRIP (Recall 3-up / 4-up stat cards)
# ===========================================================================

def build_stat_strip(prs, meta, action_title=None, stats=None, takeaway=None,
                     source=None, category=None, subtitle=None,
                     style="merck_executive", page=None, total=None,
                     section_number=None, methodology_note=None, content=None):
    """Recall-style 3-up or 4-up stat strip with cream cards.

    stats: list of 3 or 4 dicts: {value, label, body}.
      - value: hero number string in Merck Web 56pt MERCK_PURPLE bold.
      - label: SMALL UPPERCASE TRACKED label in MERCK_GOLD Verdana 10pt.
      - body: body copy in INK_GRAY Verdana 11pt.

    All cards in the row share the same height (equal-height enforcement).
    """
    if content:
        if "stats"    in content: stats    = content["stats"]
        if "takeaway" in content: takeaway = content["takeaway"]
    style = _style_or_promote(category, style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 subtitle=subtitle, takeaway=takeaway, source=source,
                 page=page, total=total, palette=style,
                 section_number=section_number, methodology_note=methodology_note)

    cards = [s for s in (stats or []) if s]
    if not cards:
        return slide
    # Clamp to a maximum of 4 cards.
    if len(cards) > 4:
        cards = cards[:4]
    n = len(cards)

    zone_top = Inches(2.95) if subtitle else Inches(2.50)
    card_h = Inches(2.40)
    gutter = Inches(0.18)
    total_w = CONTENT_W
    card_w = (total_w - gutter * (n - 1)) / n

    for i, s in enumerate(cards):
        cx = CONTENT_X + i * (card_w + gutter)
        cy = zone_top
        # Cream card with theme-accent top stripe + hairline border.
        body = rounded(slide, cx, cy, card_w, card_h, fill=PANEL_LIGHT)
        _apply_border(body, LIGHT_GRAY, Pt(0.5))
        rounded(slide, cx, cy, card_w, Inches(0.06), fill=pal["highlight"], adj=50000)

        pad = Inches(0.22)
        # Big hero number — font size scales with string length to prevent overlap.
        value = s.get("value", "")
        val_str = str(value)
        if len(val_str) <= 5:
            val_sz, num_h = 52, Inches(1.05)
        elif len(val_str) <= 8:
            val_sz, num_h = 40, Inches(0.85)
        else:
            val_sz, num_h = 32, Inches(0.72)
        num_y = cy + Inches(0.30)
        txt(slide, cx + pad, num_y, card_w - pad * 2, num_h,
            val_str, sz=val_sz, color=MERCK_PURPLE, bold=True,
            font=FONT_HEAD, anchor=MSO_ANCHOR.TOP)

        # Tracked uppercase label in theme accent.
        # Tracking is skipped for labels > 12 chars to prevent mid-word wrapping.
        label = s.get("label", "")
        label_y = num_y + num_h + Inches(0.06)
        label_h = Inches(0.40)
        label_text = _track_letters(label) if len(label) <= 12 else label.upper()
        txt(slide, cx + pad, label_y, card_w - pad * 2, label_h,
            label_text, sz=10, color=pal["highlight"], bold=True,
            font=FONT_BODY, anchor=MSO_ANCHOR.TOP)

        # Body in INK_GRAY Verdana 11pt.
        body_text = s.get("body", "")
        body_y = label_y + Inches(0.46)
        body_h = (cy + card_h) - body_y - Inches(0.18)
        if body_h < Inches(0.30):
            body_h = Inches(0.30)
        txt(slide, cx + pad, body_y, card_w - pad * 2, body_h,
            str(body_text), sz=11, color=INK_GRAY, font=FONT_BODY)

    return slide


# ===========================================================================
# Layout: BEFORE / AFTER (Today vs Tomorrow with gold arrow)
# ===========================================================================

def build_before_after(prs, meta, action_title=None, before=None, after=None, takeaway=None,
                       source=None, category=None, subtitle=None,
                       style="merck_executive", page=None, total=None,
                       section_number=None,
                       before_label="TODAY", after_label="TOMORROW", methodology_note=None,
                       content=None):
    """Two cards side by side with a gold right-arrow between them.

    before, after: dicts {title, items} where title is the headline (Verdana
    18pt bold MERCK_PURPLE) and items is a small bullet list.

    The BEFORE card is muted (PANEL_LIGHT bg, hairline border, PURPLE_MUTED
    text). The AFTER card is highlighted (PANEL_LIGHT bg, MERCK_GOLD top
    accent stripe, MERCK_GOLD label, MERCK_PURPLE bold title, INK_DARK body).

    A MERCK_GOLD right-arrow shape sits in the gutter between the cards.
    """
    if content:
        if "before"       in content: before       = content["before"]
        if "after"        in content: after        = content["after"]
        if "takeaway"     in content: takeaway     = content["takeaway"]
        if "before_label" in content: before_label = content["before_label"]
        if "after_label"  in content: after_label  = content["after_label"]
    style = _style_or_promote(category, style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 subtitle=subtitle, takeaway=takeaway, source=source,
                 page=page, total=total, palette=style,
                 section_number=section_number, methodology_note=methodology_note)

    zone_top = Inches(2.95) if subtitle else Inches(2.50)
    # Equal-height enforcement.
    base_h = Inches(2.55)
    estimated = []
    for col in (before or {}, after or {}):
        est = Inches(0.60)  # label + title
        title_len = len(str(col.get("title", "")))
        if title_len:
            title_lines = max(1, (title_len // 30) + (1 if title_len % 30 else 0))
            est += Inches(0.40) * title_lines
        items = col.get("items") or []
        est += Inches(0.10) + Inches(0.32) * len(items)
        est += Inches(0.30)
        estimated.append(est)
    card_h = max(base_h, max(estimated))
    # Clamp.
    max_h_avail = Inches(6.40) - zone_top
    if card_h > max_h_avail:
        card_h = max_h_avail

    # Card geometry: x=0.65 w=5.7 for BEFORE; x=7.10 w=5.7 for AFTER.
    before_x = Inches(0.65)
    before_w = Inches(5.70)
    after_x = Inches(7.10)
    after_w = Inches(5.70)

    # --- BEFORE card (muted) ---
    bcard = rounded(slide, before_x, zone_top, before_w, card_h,
                 fill=PANEL_LIGHT)
    _apply_border(bcard, LIGHT_GRAY, Pt(0.5))
    # No gold top stripe on the muted card.
    pad = Inches(0.30)
    txt(slide, before_x + pad, zone_top + Inches(0.22),
        before_w - pad * 2, Inches(0.26),
        _track_letters(before_label), sz=10, color=PURPLE_MUTED, bold=True,
        font=FONT_BODY)
    btitle = (before or {}).get("title", "")
    txt(slide, before_x + pad, zone_top + Inches(0.60),
        before_w - pad * 2, Inches(0.95),
        str(btitle), sz=18, color=PURPLE_MUTED, bold=True, font=FONT_BODY)
    bitems = (before or {}).get("items") or (before or {}).get("bullets") or []
    if bitems:
        _bulleted_list(slide, before_x + pad, zone_top + Inches(1.60),
                       before_w - pad * 2, card_h - Inches(1.85),
                       bitems, palette=style,
                       text_color=PURPLE_MUTED, sz=12,
                       bullet_color=PURPLE_MUTED)

    # --- AFTER card (highlighted) ---
    acard = rounded(slide, after_x, zone_top, after_w, card_h, fill=PANEL_LIGHT)
    _apply_border(acard, LIGHT_GRAY, Pt(0.5))
    rounded(slide, after_x, zone_top, after_w, Inches(0.06), fill=pal["highlight"], adj=50000)
    txt(slide, after_x + pad, zone_top + Inches(0.22),
        after_w - pad * 2, Inches(0.26),
        _track_letters(after_label), sz=10, color=pal["highlight"], bold=True,
        font=FONT_BODY)
    atitle = (after or {}).get("title", "")
    txt(slide, after_x + pad, zone_top + Inches(0.60),
        after_w - pad * 2, Inches(0.95),
        str(atitle), sz=18, color=MERCK_PURPLE, bold=True, font=FONT_BODY)
    aitems = (after or {}).get("items") or (after or {}).get("bullets") or []
    if aitems:
        _bulleted_list(slide, after_x + pad, zone_top + Inches(1.60),
                       after_w - pad * 2, card_h - Inches(1.85),
                       aitems, palette=style,
                       text_color=INK_DARK, sz=12,
                       bullet_color=pal["highlight"])

    # --- Theme-accent right-arrow shape between the cards ---
    ar_w = Inches(0.65)
    ar_h = Inches(0.85)
    ar_x = (before_x + before_w + after_x) / 2 - ar_w / 2
    ar_y = zone_top + (card_h - ar_h) / 2
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   ar_x, ar_y, ar_w, ar_h)
    arrow.shadow.inherit = False
    _apply_fill(arrow, pal["highlight"])
    _apply_border(arrow, None)

    return slide


# End of module


# ===========================================================================
# Tier A: New layout functions
# ===========================================================================

def _draw_check_mark(slide, cx, cy, size, color):
    """Draw a check-mark inside a circle by composing two line segments."""
    half = size / 2
    arm = size * 0.30
    # Down-stroke of check
    s1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    _emu(cx + half - arm),
                                    _emu(cy + half),
                                    _emu(cx + half - arm * 0.15),
                                    _emu(cy + half + arm * 0.6))
    s1.line.color.rgb = _rgb_tuple(color)
    s1.line.width = Pt(2.25)
    # Up-stroke of check
    s2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    _emu(cx + half - arm * 0.15),
                                    _emu(cy + half + arm * 0.6),
                                    _emu(cx + half + arm * 0.95),
                                    _emu(cy + half - arm * 0.7))
    s2.line.color.rgb = _rgb_tuple(color)
    s2.line.width = Pt(2.25)


def build_milestone_timeline(prs, meta, action_title=None, milestones=None, takeaway=None,
                             source=None, methodology_note=None,
                             category=None, subtitle=None,
                             style="merck_executive", page=None, total=None,
                             section_number=None, content=None):
    """Horizontal milestone timeline.

    milestones = list of {"date","title","body","status"} where status is one
    of "done", "current", "future".
    """
    if content:
        if "milestones" in content: milestones = content["milestones"]
        if "takeaway"   in content: takeaway   = content["takeaway"]
    style = _style_or_promote(category, style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 subtitle=subtitle, takeaway=takeaway, source=source,
                 page=page, total=total, palette=style,
                 section_number=section_number,
                 methodology_note=methodology_note)

    items = list(milestones or [])[:6]
    n = max(len(items), 1)
    zone_x = CONTENT_X
    zone_w = CONTENT_W
    col_w = zone_w / n
    line_y = Inches(4.10)
    circle_d = Inches(0.42)
    date_y = Inches(3.25)
    title_y = Inches(3.55)
    body_y = Inches(4.75)

    # Pre-pass: figure out last done index for line coloring.
    statuses = [str(it.get("status", "future")).lower() for it in items]
    last_done = -1
    for i, s in enumerate(statuses):
        if s == "done":
            last_done = i

    # Draw connecting lines between adjacent circles.
    for i in range(n - 1):
        cx1 = zone_x + col_w * i + col_w / 2
        cx2 = zone_x + col_w * (i + 1) + col_w / 2
        s_left = statuses[i]
        s_right = statuses[i + 1]
        if s_left == "done" and (s_right in ("done", "current")):
            line_color = MERCK_PURPLE
        else:
            line_color = LIGHT_GRAY
        rect(slide, cx1 + circle_d / 2, line_y + circle_d / 2 - Emu(int(Pt(1.0))),
             cx2 - cx1 - circle_d, Emu(int(Pt(2.0))),
             fill=line_color)

    # Draw each milestone column.
    for i, m in enumerate(items):
        col_cx = zone_x + col_w * i + col_w / 2
        st = statuses[i]

        # Date label.
        txt(slide, zone_x + col_w * i, date_y, col_w, Inches(0.26),
            _track_letters(str(m.get("date", "")).upper()),
            sz=10, color=pal["highlight"], bold=True, font=FONT_BODY,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Title — use white on dark-background styles so text is readable.
        dark = _is_dark(style)
        title_color = pal["highlight"] if st == "current" else (WHITE if dark else MERCK_PURPLE)
        txt(slide, zone_x + col_w * i, title_y, col_w, Inches(0.40),
            str(m.get("title", "")), sz=13, color=title_color, bold=True,
            font=FONT_BODY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Circle.
        circle_x = col_cx - circle_d / 2
        circle_y = line_y
        if st == "done":
            circle(slide, circle_x, circle_y, circle_d, fill=MERCK_PURPLE)
            _draw_check_mark(slide, circle_x, circle_y, circle_d, WHITE)
        elif st == "current":
            circle(slide, circle_x, circle_y, circle_d, fill=pal["highlight"])
        else:
            outer = circle(slide, circle_x, circle_y, circle_d, fill=pal["bg"])
            _apply_border(outer, LIGHT_GRAY, Pt(1.25))

        # Body text.
        txt(slide, zone_x + col_w * i + Inches(0.10), body_y,
            col_w - Inches(0.20), Inches(1.50),
            str(m.get("body", "")), sz=10, color=INK_GRAY, font=FONT_BODY,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    return slide


def _rag_color(value):
    v = str(value or "").strip().upper()
    if v in ("RED", "R", "HIGH"):
        return BAD_RED
    if v in ("AMBER", "A", "YELLOW", "MEDIUM", "MED"):
        return MERCK_GOLD
    if v in ("GREEN", "G", "LOW", "OK"):
        return GOOD_GREEN
    return PURPLE_MUTED


def _norm_key(s):
    return "".join(ch for ch in str(s).lower() if ch.isalnum())


def build_status_table(prs, meta, action_title=None, columns=None, rows=None, takeaway=None,
                       source=None, methodology_note=None,
                       category=None, subtitle=None,
                       style="merck_executive", page=None, total=None,
                       section_number=None, content=None):
    """RAG status / dependency table with header row + RAG pills.
    Caps at 7 rows; extra rows are dropped.
    """
    if content:
        if "columns"  in content: columns  = content["columns"]
        if "rows"     in content: rows     = content["rows"]
        if "takeaway" in content: takeaway = content["takeaway"]
    style = _style_or_promote(category, style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 subtitle=subtitle, takeaway=takeaway, source=source,
                 page=page, total=total, palette=style,
                 section_number=section_number,
                 methodology_note=methodology_note)

    cols = list(columns or [])
    if not cols:
        return slide
    rs = list(rows or [])[:7]

    # Detect RAG columns by header text.
    rag_set = {"status", "rag", "health", "severity",
               "priority", "risk level", "rag status"}
    rag_cols = [i for i, c in enumerate(cols)
                if str(c).strip().lower() in rag_set]

    # Compute column widths. RAG columns are narrower.
    zone_x = CONTENT_X
    zone_w = CONTENT_W
    zone_y = Inches(2.55) if not subtitle else Inches(2.95)
    n = len(cols)
    rag_w = Inches(1.20)
    non_rag_count = n - len(rag_cols)
    remaining = zone_w - rag_w * len(rag_cols)
    first_col_share = 0.40 if non_rag_count >= 2 else 1.0
    other_col_share = (1.0 - first_col_share) / max(non_rag_count - 1, 1)
    col_widths = []
    seen_first = False
    for i, c in enumerate(cols):
        if i in rag_cols:
            col_widths.append(rag_w)
        else:
            if not seen_first:
                col_widths.append(remaining * first_col_share)
                seen_first = True
            else:
                col_widths.append(remaining * other_col_share)

    # Header row.
    header_h = Inches(0.42)
    avail_h = SOURCE_Y - zone_y - header_h - Inches(0.15)
    n_rows = max(len(rs), 1)
    row_h = min(Inches(0.78), avail_h / n_rows)
    row_h = max(row_h, Inches(0.34))  # floor for legibility
    rect(slide, zone_x, zone_y, zone_w, header_h, fill=MERCK_PURPLE)
    cx = zone_x
    for i, c in enumerate(cols):
        txt(slide, cx + Inches(0.18), zone_y, col_widths[i] - Inches(0.20),
            header_h, _track_letters(str(c).upper()),
            sz=10, color=WHITE, bold=True, font=FONT_BODY,
            anchor=MSO_ANCHOR.MIDDLE)
        cx = cx + col_widths[i]

    # Body rows.
    PANEL_LAVENDER = RGBColor(0xF7, 0xF4, 0xFA)
    for ri, row in enumerate(rs):
        ry = zone_y + header_h + row_h * ri
        rect(slide, zone_x, ry, zone_w, row_h, fill=tuple(PANEL_LAVENDER))
        # Bottom hairline.
        hairline(slide, zone_x, ry + row_h - Emu(int(Pt(0.5))),
                 zone_w, Emu(int(Pt(0.5))), LIGHT_GRAY)
        cx = zone_x
        first_text_col = True
        for i, c in enumerate(cols):
            col_key = _norm_key(c)
            val = ""
            for k, v in row.items():
                if _norm_key(k) == col_key:
                    val = v
                    break
            if not val:
                # Fallback to first-word match or exact-lowercase match.
                _parts = str(c).strip().lower().split() if c else []
                first_word = _parts[0] if _parts else ""
                val = row.get(first_word) or row.get(str(c).strip().lower()) or ""
            cell_w = col_widths[i]
            if i in rag_cols:
                # RAG pill.
                pill_w = Inches(0.85)
                pill_h = Inches(0.35)
                pill_x = cx + (cell_w - pill_w) / 2
                pill_y = ry + (row_h - pill_h) / 2
                pill = rounded(slide, pill_x, pill_y, pill_w, pill_h,
                               fill=_rag_color(val), adj=5000)
                _apply_border(pill, None)
                txt(slide, pill_x, pill_y, pill_w, pill_h,
                    str(val).upper(), sz=10, color=WHITE, bold=True,
                    font=FONT_BODY, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
            else:
                # Text cell. First column anchors TOP so long labels don't clip.
                color = MERCK_PURPLE if i == 0 else INK_DARK
                bold = (i == 0)
                first_text_col = False
                anchor = MSO_ANCHOR.TOP if i == 0 else MSO_ANCHOR.MIDDLE
                cell_sz = 9 if (i == 0 and len(str(val)) > 30) else 10
                txt(slide, cx + Inches(0.18), ry + Inches(0.06),
                    cell_w - Inches(0.20), row_h - Inches(0.08),
                    str(val), sz=cell_sz, color=color, bold=bold,
                    font=FONT_BODY, anchor=anchor)
            cx = cx + cell_w
    return slide


def _dashed_line(slide, x1, y1, x2, y2, color, weight=Pt(1.0)):
    """Add a dashed connector between two points."""
    from lxml import etree
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      _emu(x1), _emu(y1),
                                      _emu(x2), _emu(y2))
    conn.line.color.rgb = _rgb_tuple(color)
    conn.line.width = weight
    # Set dash style via XML.
    ln = conn.line._get_or_add_ln()
    pr_dash = ln.find(qn("a:prstDash"))
    if pr_dash is None:
        pr_dash = etree.SubElement(ln, qn("a:prstDash"))
    pr_dash.set("val", "dash")
    return conn


def build_hub_spoke(prs, meta, action_title=None, hub=None, spokes=None, takeaway=None,
                    source=None, methodology_note=None,
                    category=None, subtitle=None,
                    style="merck_executive", page=None, total=None,
                    section_number=None, content=None):
    """Center hub oval + 4 corner cards with dashed connectors."""
    if content:
        fw = content.get("framework")
        if fw and fw in _SPOKE_PRESETS and not hub and not spokes:
            hub    = _SPOKE_PRESETS[fw]["hub"]
            spokes = _SPOKE_PRESETS[fw]["spokes"]
        if "hub"      in content: hub      = content["hub"]
        if "spokes"   in content: spokes   = content["spokes"]
        if "takeaway" in content: takeaway = content["takeaway"]
    style = _style_or_promote(category, style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 subtitle=subtitle, takeaway=takeaway, source=source,
                 page=page, total=total, palette=style,
                 section_number=section_number,
                 methodology_note=methodology_note)

    sp = list(spokes or [])[:4]

    # Center hub oval.
    hub_w = Inches(2.80)
    hub_h = Inches(2.00)
    hub_x = Inches(13.333) / 2 - hub_w / 2
    hub_y = Inches(4.55) - hub_h / 2
    hub_shape = oval(slide, hub_x, hub_y, hub_w, hub_h, fill=MERCK_PURPLE)
    _apply_border(hub_shape, pal["highlight"], Pt(2.25))

    # Hub contents.
    hub_label = (hub or {}).get("label", "")
    hub_title = (hub or {}).get("title", "")
    hub_sub = (hub or {}).get("subtitle", "")
    if hub_label:
        txt(slide, hub_x, hub_y + Inches(0.35), hub_w, Inches(0.30),
            _track_letters(str(hub_label).upper()),
            sz=10, color=pal["highlight"], bold=True, font=FONT_BODY,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if hub_title:
        txt(slide, hub_x, hub_y + Inches(0.70), hub_w, Inches(0.55),
            str(hub_title), sz=18, color=WHITE, bold=True,
            font=FONT_BODY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if hub_sub:
        txt(slide, hub_x, hub_y + Inches(1.30), hub_w, Inches(0.45),
            str(hub_sub), sz=11, color=WHITE, italic=True,
            font=FONT_BODY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Four corner cards.
    card_w = Inches(3.40)
    card_h = Inches(1.55)
    positions = [
        (Inches(0.55), Inches(3.20)),                   # top-left
        (Inches(13.333) - Inches(0.55) - card_w,        # top-right
         Inches(3.20)),
        (Inches(0.55), Inches(5.40)),                   # bottom-left
        (Inches(13.333) - Inches(0.55) - card_w,        # bottom-right
         Inches(5.40)),
    ]
    hub_cx = hub_x + hub_w / 2
    hub_cy = hub_y + hub_h / 2
    for i, item in enumerate(sp):
        if i >= 4:
            break
        cx, cy = positions[i]
        card = rounded(slide, cx, cy, card_w, card_h, fill=pal["panel"])
        _apply_border(card, LIGHT_GRAY, Pt(0.5))
        # Theme-accent top stripe.
        rounded(slide, cx, cy, card_w, Inches(0.06), fill=pal["highlight"], adj=50000)
        # Title + body.
        pad = Inches(0.22)
        txt(slide, cx + pad, cy + Inches(0.20),
            card_w - pad * 2, Inches(0.45),
            str(item.get("title", "")), sz=13, color=MERCK_PURPLE, bold=True,
            font=FONT_BODY)
        txt(slide, cx + pad, cy + Inches(0.72),
            card_w - pad * 2, card_h - Inches(0.80),
            str(item.get("body", "")), sz=10, color=INK_GRAY, font=FONT_BODY)
        # Dashed connector from card edge nearest hub to hub edge.
        if cx < hub_x:
            edge_x = cx + card_w
            edge_y = cy + card_h / 2
            hub_edge_x = hub_x
            hub_edge_y = hub_cy
        else:
            edge_x = cx
            edge_y = cy + card_h / 2
            hub_edge_x = hub_x + hub_w
            hub_edge_y = hub_cy
        _dashed_line(slide, edge_x, edge_y, hub_edge_x, hub_edge_y,
                     PURPLE_MUTED, weight=Pt(0.75))
    return slide


def build_pillar_detail(prs, meta, action_title=None, pillar_number=None, pillar_label=None,
                        owner=None, sections=None, takeaway=None,
                        source=None, methodology_note=None,
                        category=None, subtitle=None,
                        style="merck_executive", page=None, total=None,
                        section_number=None, content=None):
    """Dedicated page per pillar: left purple panel with huge number, right
    column with sectioned content.

    pillar_number: e.g. "02"
    pillar_label:  e.g. "PILLAR"
    owner: {"label": "OWNER", "name": "Natalia Mocan, Future Platform Architecture"}
    sections: list of {"label", "body"} dicts
    """
    if content:
        if "pillar_number" in content: pillar_number = content["pillar_number"]
        if "pillar_label"  in content: pillar_label  = content["pillar_label"]
        if "owner"         in content: owner         = content["owner"]
        if "sections"      in content: sections      = content["sections"]
        if "takeaway"      in content: takeaway      = content["takeaway"]
    style = _style_or_promote(category, style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 subtitle=subtitle, takeaway=takeaway, source=source,
                 page=page, total=total, palette=style,
                 section_number=section_number,
                 methodology_note=methodology_note)

    # Left purple panel.
    panel_x = Inches(0.65)
    panel_y = Inches(2.95) if subtitle else Inches(2.55)
    panel_w = Inches(3.40)
    # Panel bottom snaps just above source / takeaway band.
    panel_h = SOURCE_Y - panel_y - Inches(0.30)
    rect(slide, panel_x, panel_y, panel_w, panel_h, fill=MERCK_PURPLE)

    # Pillar label at top.
    txt(slide, panel_x + Inches(0.30), panel_y + Inches(0.25),
        panel_w - Inches(0.40), Inches(0.30),
        _track_letters(str(pillar_label or "PILLAR").upper()),
        sz=11, color=pal["highlight"], bold=True, font=FONT_BODY)

    # Huge number (Merck Web — hero moment).
    txt(slide, panel_x, panel_y + Inches(0.65),
        panel_w, panel_h - Inches(1.80),
        str(pillar_number or ""), sz=120, color=WHITE, bold=True,
        font=FONT_HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Owner block at panel bottom.
    if owner:
        owner_y = panel_y + panel_h - Inches(0.95)
        owner_label = (owner or {}).get("label", "OWNER")
        owner_name = (owner or {}).get("name", "")
        txt(slide, panel_x + Inches(0.30), owner_y,
            panel_w - Inches(0.40), Inches(0.30),
            _track_letters(str(owner_label).upper()),
            sz=10, color=pal["highlight"], bold=True, font=FONT_BODY)
        txt(slide, panel_x + Inches(0.30), owner_y + Inches(0.30),
            panel_w - Inches(0.40), Inches(0.55),
            str(owner_name), sz=13, color=WHITE, bold=True,
            font=FONT_BODY)

    # Right column: sectioned content blocks.
    sec = list(sections or [])[:5]
    right_x = panel_x + panel_w + Inches(0.40)
    right_w = Inches(13.333) - right_x - Inches(0.65)
    right_y = panel_y
    right_h = panel_h
    if sec:
        block_gap = Inches(0.18)
        block_h = (right_h - block_gap * (len(sec) - 1)) / len(sec)
        for i, s in enumerate(sec):
            by = right_y + (block_h + block_gap) * i
            stripe_w = Inches(0.05)
            rect(slide, right_x, by, stripe_w, block_h, fill=pal["highlight"])
            txt(slide, right_x + stripe_w + Inches(0.18), by,
                right_w - stripe_w - Inches(0.20), Inches(0.28),
                _track_letters(str(s.get("label", "")).upper()),
                sz=10, color=pal["highlight"], bold=True, font=FONT_BODY)
            txt(slide, right_x + stripe_w + Inches(0.18),
                by + Inches(0.32),
                right_w - stripe_w - Inches(0.20),
                block_h - Inches(0.36),
                str(s.get("body", "")), sz=11, color=INK_DARK,
                font=FONT_BODY)
    return slide



# ===========================================================================
# build_four_column: 4-up cards in a row
# ===========================================================================

def build_four_column(prs, meta, action_title=None, columns=None, takeaway=None,
                      source=None, category=None, style="merck_executive",
                      page=None, total=None, section_number=None,
                      methodology_note=None, subtitle=None, content=None):
    """4-up columns. Same card chrome as three_column.
    columns = list of dicts: {label, title, body, tone, items}.
    Set "framework" in content for auto-configured presets (e.g. "4ps", "balanced_scorecard").
    """
    if content:
        fw = content.get("framework")
        if fw and fw in _COLUMN_PRESETS and not columns:
            columns = _COLUMN_PRESETS[fw]["columns"]
        if "columns"  in content: columns  = content["columns"]
        if "takeaway" in content: takeaway = content["takeaway"]
    style = _style_or_promote(category, style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category, subtitle=subtitle,
                 takeaway=takeaway, source=source, page=page, total=total,
                 palette=style, section_number=section_number,
                 methodology_note=methodology_note)
    cols = (columns or [])[:4]
    while len(cols) < 4:
        cols.append({})
    zone_top = Inches(2.55) if not subtitle else Inches(2.95)
    zone_h = Inches(6.30) - zone_top
    gap = Inches(0.18)
    col_w = (Inches(12.0) - gap * 3) / 4
    for i, col in enumerate(cols):
        cx = Inches(0.65) + i * (col_w + gap)
        _two_or_three_column_card(slide, cx, zone_top, col_w, zone_h, col, style)
    return slide


# ===========================================================================
# build_label_rows: colored label card on left + body on right per row
# ===========================================================================

def build_label_rows(prs, meta, action_title=None, rows=None, takeaway=None,
                     source=None, category=None, style="merck_executive",
                     page=None, total=None, section_number=None,
                     methodology_note=None, subtitle=None,
                     label_color=None, content=None):
    """Row-per-item layout with a colored LABEL CARD on the left and BODY text on
    the right. Use for diagnosis / root cause / risk taxonomy slides where the
    name of each item is the anchor.

    rows = list of {"label": str, "body": str} (3-6 rows).
    label_color: color tuple for the label card fill. Defaults to MERCK_PURPLE.
        Set explicitly for a one-off accent (e.g. BAD_RED for negative tone).
    """
    if content:
        if "rows"        in content: rows        = content["rows"]
        if "takeaway"    in content: takeaway    = content["takeaway"]
        if "label_color" in content: label_color = content["label_color"]
    style = _style_or_promote(category, style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category, subtitle=subtitle,
                 takeaway=takeaway, source=source, page=page, total=total,
                 palette=style, section_number=section_number,
                 methodology_note=methodology_note)

    rs = list(rows or [])[:6]
    n = max(len(rs), 1)
    zone_top = Inches(2.55) if not subtitle else Inches(2.95)
    zone_h = SOURCE_Y - zone_top - Inches(0.20)
    gap = Inches(0.14)
    row_h = (zone_h - gap * (n - 1)) / n

    label_w = Inches(3.10)
    body_x = Inches(0.65) + label_w + Inches(0.30)
    body_w = Inches(12.0) - label_w - Inches(0.30)
    fill_color = label_color if label_color is not None else MERCK_PURPLE
    # On dark / storytelling bg, body text needs to be light.
    body_color = WHITE if _is_dark(style) else INK_DARK

    for i, row in enumerate(rs):
        ry = zone_top + i * (row_h + gap)
        # Label card: colored fill, rounded corners, white bold text centered.
        card = rounded(slide, Inches(0.65), ry, label_w, row_h,
                       fill=fill_color, adj=2000)
        _apply_border(card, None)
        txt(slide, Inches(0.65) + Inches(0.12), ry,
            label_w - Inches(0.24), row_h,
            str(row.get("label", "")), sz=14, color=WHITE, bold=True,
            font=FONT_BODY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Body text on right.
        txt(slide, body_x, ry, body_w, row_h,
            str(row.get("body", "")), sz=12, color=body_color,
            font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE)
        # Hairline below row (except last).
        if i < n - 1:
            rule_color = PURPLE_MUTED if _is_dark(style) else LIGHT_GRAY
            hairline(slide, Inches(0.65), ry + row_h + gap / 2 - Emu(int(Pt(0.5))),
                     Inches(12.0), Emu(int(Pt(0.5))), rule_color)
    return slide



def build_circular_flow(prs, meta, action_title, phases, takeaway="", source=None,
                        subtitle=None, methodology_note=None,
                        style="merck_corporate", page=None, total=None,
                        section_number=None, category=None):
    """Circular process flow with 2-8 phases arranged around a centre point.

    phases: list of dicts with keys: label (str), body (str), icon (optional str icon name).
    Returns the slide.
    """
    import math

    style = _style_or_promote(category, style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 subtitle=subtitle, takeaway=takeaway, source=source,
                 page=page, total=total, palette=style,
                 section_number=section_number,
                 methodology_note=methodology_note)

    n = max(2, min(8, len(phases)))
    cx = SLIDE_W / 2
    cy = Inches(4.0)
    radius = Inches(2.0)
    circle_size = Inches(0.9)
    half = circle_size / 2

    # Compute centre positions for each phase
    positions = []
    for i in range(n):
        angle = (2 * math.pi * i / n) - math.pi / 2
        px = cx + radius * math.cos(angle)
        py = cy + radius * math.sin(angle)
        positions.append((px, py))

    # Draw connector lines between adjacent phase circles
    for i in range(n):
        x1, y1 = positions[i]
        x2, y2 = positions[(i + 1) % n]
        # Shorten line endpoints to circle edge
        dx = x2 - x1
        dy = y2 - y1
        dist = math.hypot(dx, dy)
        if dist > 0:
            ux, uy = dx / dist, dy / dist
            edge = half + Inches(0.05)
            lx1 = x1 + ux * edge
            ly1 = y1 + uy * edge
            lx2 = x2 - ux * edge
            ly2 = y2 - uy * edge
            line(slide, lx1, ly1, lx2, ly2,
                 color=pal.get("accent", MERCK_PURPLE), weight=Pt(1.5))

    # Draw phase circles and labels
    body_box_h = Inches(0.55)
    for i, phase in enumerate(phases[:n]):
        px, py = positions[i]
        cx_shape = px - half
        cy_shape = py - half

        # Filled circle
        circle(slide, cx_shape, cy_shape, circle_size,
               fill=pal.get("accent", MERCK_PURPLE))

        # Label inside circle (number + short label)
        label_text = str(phase.get("label", str(i + 1)))
        txt(slide, cx_shape, cy_shape, circle_size, circle_size,
            label_text, sz=10, color=WHITE, bold=True,
            align=PP_ALIGN.CENTER, font=FONT_BODY,
            anchor=MSO_ANCHOR.MIDDLE)

        # Body text below circle
        body_text = str(phase.get("body", ""))
        if body_text:
            txt(slide, cx_shape - Inches(0.3), py + half + Inches(0.06),
                circle_size + Inches(0.6), body_box_h,
                body_text, sz=11, color=pal.get("ink", INK_DARK),
                align=PP_ALIGN.CENTER, font=FONT_BODY,
                anchor=MSO_ANCHOR.TOP)

    return slide


def build_org_chart(prs, meta, action_title, root, children,
                    takeaway="", source=None, subtitle=None,
                    methodology_note=None, style="merck_corporate",
                    page=None, total=None, section_number=None, category=None):
    """Two-level org chart.

    root: dict with keys: name (str), title (str, optional).
    children: list of dicts with keys: name (str), title (str, optional),
              reports (list of dicts with name/title, optional).
    Returns the slide.
    """
    pal = PALETTES.get(style, PALETTES["merck_corporate"])
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    children = children[:6]
    n = len(children)

    # Root box geometry
    root_w = Inches(2.5)
    root_h = Inches(0.7)
    root_x = Inches(13.333 / 2) - root_w / 2
    root_y = Inches(2.0)
    root_cx = root_x + root_w / 2   # horizontal centre of root

    root_box = rounded(slide, root_x, root_y, root_w, root_h,
                       fill=MERCK_PURPLE, adj=6000)
    txt(slide, root_x, root_y, root_w, root_h,
        str(root.get("name", "")),
        sz=13, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER, font=FONT_BODY,
        anchor=MSO_ANCHOR.MIDDLE)

    # Optional title text beneath root box
    root_title = str(root.get("title", ""))
    if root_title:
        txt(slide, root_x, root_y + root_h + Inches(0.04),
            root_w, Inches(0.25),
            root_title, sz=9, color=pal.get("ink", INK_DARK),
            align=PP_ALIGN.CENTER, font=FONT_BODY,
            anchor=MSO_ANCHOR.TOP)

    if n == 0:
        return slide

    # Child row geometry
    child_w = Inches(1.7)
    child_h = Inches(0.65)
    child_y = Inches(3.4)
    available_w = Inches(12.0)
    col_w = available_w / n
    start_x = (Inches(13.333) - available_w) / 2

    # Vertical trunk: root bottom to midpoint row
    trunk_top_y = root_y + root_h
    trunk_bot_y = child_y                          # top of child boxes
    mid_y = trunk_top_y + (trunk_bot_y - trunk_top_y) / 2

    # Draw trunk line from root bottom down to mid_y
    line(slide, root_cx, trunk_top_y, root_cx, mid_y,
         color=PURPLE_MUTED, weight=Pt(1.0))

    for i, child in enumerate(children):
        child_cx = start_x + col_w * i + col_w / 2
        child_x = child_cx - child_w / 2

        # Horizontal branch at mid_y
        line(slide, root_cx, mid_y, child_cx, mid_y,
             color=PURPLE_MUTED, weight=Pt(1.0))
        # Vertical drop from mid_y to child top
        line(slide, child_cx, mid_y, child_cx, child_y,
             color=PURPLE_MUTED, weight=Pt(1.0))

        # Child box
        rounded(slide, child_x, child_y, child_w, child_h,
                fill=PANEL_LIGHT, adj=6000)
        txt(slide, child_x, child_y, child_w, child_h,
            str(child.get("name", "")),
            sz=11, color=MERCK_PURPLE, bold=True,
            align=PP_ALIGN.CENTER, font=FONT_BODY,
            anchor=MSO_ANCHOR.MIDDLE)

        child_title = str(child.get("title", ""))
        if child_title:
            txt(slide, child_x, child_y + child_h + Inches(0.03),
                child_w, Inches(0.22),
                child_title, sz=8, color=pal.get("ink", INK_DARK),
                align=PP_ALIGN.CENTER, font=FONT_BODY,
                anchor=MSO_ANCHOR.TOP)

        # Grandchild boxes
        reports = list(child.get("reports", []))[:2]
        if reports:
            gc_w = Inches(1.4)
            gc_h = Inches(0.5)
            gc_y = Inches(4.5)
            gc_spacing = Inches(0.12)
            total_gc_w = len(reports) * gc_w + (len(reports) - 1) * gc_spacing
            gc_start_x = child_cx - total_gc_w / 2

            # Vertical connector from child bottom to grandchild row
            line(slide, child_cx, child_y + child_h, child_cx, gc_y,
                 color=LIGHT_GRAY, weight=Pt(0.75))

            for j, rep in enumerate(reports):
                gx = gc_start_x + j * (gc_w + gc_spacing)
                gc_cx = gx + gc_w / 2

                # Horizontal branch at gc_y
                if len(reports) > 1:
                    line(slide, child_cx, gc_y, gc_cx, gc_y,
                         color=LIGHT_GRAY, weight=Pt(0.75))

                rounded(slide, gx, gc_y, gc_w, gc_h,
                        fill=LIGHT_GRAY, adj=6000)
                txt(slide, gx, gc_y, gc_w, gc_h,
                    str(rep.get("name", "")),
                    sz=9, color=INK_DARK, bold=False,
                    align=PP_ALIGN.CENTER, font=FONT_BODY,
                    anchor=MSO_ANCHOR.MIDDLE)

    return slide


# ---------------------------------------------------------------------------
# build_topic_set — 3–6 topic cards with numbered circles
# ---------------------------------------------------------------------------

def build_topic_set(prs, meta, action_title, topics, takeaway="", source=None,
                    subtitle=None, methodology_note=None, style="merck_corporate",
                    page=None, total=None, section_number=None, category=None):
    """3-6 topic cards with numbered circles, title and body.

    topics: list of dicts with keys:
        label (str) — short uppercase label shown above title
        title (str) — bold topic title
        body  (str) — 1-2 sentence description
        icon  (str, optional) — icon name for draw_icon; if None, shows sequential number
    Returns the slide.
    """
    pal = PALETTES.get(style, PALETTES["merck_corporate"])
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    topics = list(topics)[:6]
    n = max(len(topics), 1)

    # Decide row layout
    if n <= 3:
        rows = [topics]
    else:
        rows = [topics[:3], topics[3:]]

    row_count = len(rows)
    zone_x = Inches(0.65)
    zone_y = Inches(2.4)
    zone_w = Inches(12.0)
    zone_h = Inches(3.8)

    gap = Inches(0.25)
    row_h = zone_h / row_count

    circle_size = Inches(0.75)
    circle_r = circle_size / 2

    global_idx = 0

    for row_idx, row_topics in enumerate(rows):
        rn = len(row_topics)
        card_w = (zone_w - gap * (rn - 1)) / rn
        row_y = zone_y + row_idx * row_h

        for col_idx, topic in enumerate(row_topics):
            global_idx += 1
            card_x = zone_x + col_idx * (card_w + gap)

            # Circle centred horizontally on the card
            circ_x = card_x + card_w / 2 - circle_r
            circ_y = row_y
            circle(slide, circ_x, circ_y, circle_size, fill=MERCK_PURPLE)

            # Number or icon inside circle
            icon_name = topic.get("icon")
            if icon_name:
                icon_margin = Inches(0.12)
                draw_icon(slide, icon_name,
                          circ_x + icon_margin,
                          circ_y + icon_margin,
                          circle_size - icon_margin * 2,
                          color=WHITE)
            else:
                txt(slide, circ_x, circ_y, circle_size, circle_size,
                    str(global_idx),
                    sz=18, color=WHITE, bold=True,
                    align=PP_ALIGN.CENTER, font=FONT_BODY,
                    anchor=MSO_ANCHOR.MIDDLE)

            cursor_y = circ_y + circle_size + Inches(0.1)

            # Label
            label_text = str(topic.get("label", "")).upper()
            if label_text:
                label_h = Inches(0.22)
                txt(slide, card_x, cursor_y, card_w, label_h,
                    label_text,
                    sz=9, color=pal["highlight"], bold=False,
                    align=PP_ALIGN.CENTER, font=FONT_BODY,
                    anchor=MSO_ANCHOR.TOP)
                cursor_y += label_h + Inches(0.04)

            # Title
            title_h = Inches(0.36)
            txt(slide, card_x, cursor_y, card_w, title_h,
                str(topic.get("title", "")),
                sz=13, color=MERCK_PURPLE, bold=True,
                align=PP_ALIGN.CENTER, font=FONT_BODY,
                anchor=MSO_ANCHOR.TOP)
            cursor_y += title_h + Inches(0.08)

            # Body
            remaining_h = (row_y + row_h) - cursor_y - Inches(0.1)
            txt(slide, card_x, cursor_y, card_w, max(remaining_h, Inches(0.4)),
                str(topic.get("body", "")),
                sz=10, color=INK_DARK, bold=False,
                align=PP_ALIGN.CENTER, font=FONT_BODY,
                anchor=MSO_ANCHOR.TOP)

    return slide


# ---------------------------------------------------------------------------
# build_arrow_chain — 3–5 horizontal boxes connected by triangle arrows
# ---------------------------------------------------------------------------

def build_arrow_chain(prs, meta, action_title, steps, consequence=None,
                      takeaway="", source=None, subtitle=None,
                      methodology_note=None, style="merck_corporate",
                      page=None, total=None, section_number=None, category=None):
    """Horizontal arrow chain: 3-5 steps + optional consequence box.

    steps: list of dicts with keys:
        label (str) — uppercase short label
        body (str) — 1-2 sentence description
        highlighted (bool, optional) — if True, uses MERCK_PURPLE fill + white text
    consequence: dict with label and body (optional right-end conclusion box,
                 MERCK_YELLOW fill)
    Returns the slide.
    """
    pal = PALETTES.get(style, PALETTES["merck_corporate"])
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    steps = list(steps)[:5]
    n_steps = max(len(steps), 1)
    n_arrows = n_steps - 1
    has_consequence = consequence is not None

    # Layout geometry
    zone_x = Inches(0.55)
    zone_y = Inches(2.3)
    zone_w = Inches(12.2)
    box_h = Inches(2.4)

    arrow_w = Inches(0.35)
    # consequence box is 1.0 * box_w wide; preceded by a slightly larger arrow
    consequence_w_ratio = 1.1
    consequence_arrow_w = Inches(0.42)

    # Total "units" of width consumed:
    #   n_steps boxes  +  n_arrows regular arrows
    #   + (consequence_arrow_w + consequence_w_ratio * box_w) if consequence
    # Solve for box_w:
    #   zone_w = n_steps*box_w + n_arrows*arrow_w
    #            [+ consequence_arrow_w + consequence_w_ratio*box_w]
    if has_consequence:
        # zone_w = (n_steps + consequence_w_ratio)*box_w
        #        + n_arrows*arrow_w + consequence_arrow_w
        box_w = (zone_w - n_arrows * arrow_w - consequence_arrow_w) / (
            n_steps + consequence_w_ratio)
        consequence_box_w = box_w * consequence_w_ratio
    else:
        box_w = (zone_w - n_arrows * arrow_w) / n_steps
        consequence_box_w = 0.0

    # Helper: draw a right-pointing filled triangle (arrow connector)
    def _draw_triangle(ax, ay, aw, ah, fill_color):
        """Draw filled right-pointing triangle at (ax, ay) with size (aw x ah)."""
        tip_x = ax + aw
        mid_y = ay + ah / 2
        pts = [
            (int(ax), int(ay)),
            (int(ax), int(ay + ah)),
            (int(tip_x), int(mid_y)),
        ]
        _freeform_poly(slide, pts, fill=fill_color)

    cursor_x = zone_x
    arrow_h = Inches(0.55)  # triangle height (vertical span)
    arrow_y = zone_y + (box_h - arrow_h) / 2  # vertically centred

    for i, step in enumerate(steps):
        highlighted = bool(step.get("highlighted", False))

        # Box fill and text colors
        box_fill = MERCK_PURPLE if highlighted else PANEL_LIGHT
        label_color = WHITE if highlighted else pal["highlight"]
        body_color = WHITE if highlighted else INK_DARK

        # Draw box
        rounded(slide, cursor_x, zone_y, box_w, box_h, fill=box_fill)

        # Label — uppercase, tracked at sz=9
        label_text = str(step.get("label", "")).upper()
        label_pad_x = Inches(0.12)
        label_pad_y = Inches(0.14)
        label_h = Inches(0.25)
        if label_text:
            txt(slide,
                cursor_x + label_pad_x,
                zone_y + label_pad_y,
                box_w - label_pad_x * 2,
                label_h,
                label_text,
                sz=9, color=label_color, bold=False,
                align=PP_ALIGN.LEFT, font=FONT_BODY,
                anchor=MSO_ANCHOR.TOP)

        # Body text
        body_top = zone_y + label_pad_y + label_h + Inches(0.08)
        body_h = box_h - label_pad_y - label_h - Inches(0.08) - Inches(0.12)
        txt(slide,
            cursor_x + label_pad_x,
            body_top,
            box_w - label_pad_x * 2,
            body_h,
            str(step.get("body", "")),
            sz=11, color=body_color, bold=False,
            align=PP_ALIGN.LEFT, font=FONT_BODY,
            anchor=MSO_ANCHOR.TOP)

        cursor_x += box_w

        # Arrow between boxes (not after last step unless consequence follows)
        if i < n_steps - 1:
            _draw_triangle(cursor_x, arrow_y, arrow_w, arrow_h, MERCK_PURPLE)
            cursor_x += arrow_w
        elif has_consequence:
            # Larger arrow before consequence box
            _draw_triangle(cursor_x, arrow_y - Inches(0.04),
                           consequence_arrow_w, arrow_h + Inches(0.08),
                           MERCK_PURPLE)
            cursor_x += consequence_arrow_w

    # Consequence box
    if has_consequence:
        con_label = str(consequence.get("label", "")).upper()
        con_body = str(consequence.get("body", ""))

        rounded(slide, cursor_x, zone_y, consequence_box_w, box_h,
             fill=pal["hot"])

        label_pad_x = Inches(0.12)
        label_pad_y = Inches(0.14)
        label_h = Inches(0.25)
        if con_label:
            txt(slide,
                cursor_x + label_pad_x,
                zone_y + label_pad_y,
                consequence_box_w - label_pad_x * 2,
                label_h,
                con_label,
                sz=9, color=PURPLE_DEEP, bold=True,
                align=PP_ALIGN.LEFT, font=FONT_BODY,
                anchor=MSO_ANCHOR.TOP)

        body_top = zone_y + label_pad_y + label_h + Inches(0.08)
        body_h = box_h - label_pad_y - label_h - Inches(0.08) - Inches(0.12)
        txt(slide,
            cursor_x + label_pad_x,
            body_top,
            consequence_box_w - label_pad_x * 2,
            body_h,
            con_body,
            sz=11, color=PURPLE_DEEP, bold=True,
            align=PP_ALIGN.LEFT, font=FONT_BODY,
            anchor=MSO_ANCHOR.TOP)

    return slide


# ===========================================================================
# Layout: PULL QUOTE / STATEMENT
# ===========================================================================

def build_pull_quote(prs, meta, quote=None, attribution=None,
                     context=None, takeaway="", source=None, subtitle=None,
                     methodology_note=None, style="merck_storytelling",
                     page=None, total=None, section_number=None,
                     category=None, action_title=None):
    """Full-bleed pull quote slide.

    quote:       str — the quote text (no quotation marks needed)
    attribution: str — source / speaker (shown below quote with em-dash)
    context:     str — optional small descriptor tag above the quote
    """
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    dark = _is_dark(style)
    q_color    = MERCK_YELLOW if dark else MERCK_PURPLE
    attr_color = PANEL_LIGHT  if dark else INK_GRAY
    ctx_color  = PANEL_LIGHT  if dark else PURPLE_MUTED

    # Shift quote content below the action-title chrome.  A single-line title
    # clears at ~1.50; a two-line title can reach ~1.95.  Use the subtitle-
    # aware offset so content never overlaps the action title text.
    q_top = Inches(2.10) if not subtitle else Inches(2.50)

    txt(slide, Inches(0.55), q_top, Inches(1.20), Inches(1.40),
        "\u201c", sz=96, color=pal["highlight"], bold=True,
        font=FONT_HEAD, align=PP_ALIGN.LEFT)

    if context:
        txt(slide, Inches(1.40), q_top + Inches(0.15), Inches(10.5), Inches(0.30),
            str(context).upper(), sz=9, color=ctx_color, bold=True,
            font=FONT_BODY, align=PP_ALIGN.LEFT)

    txt(slide, Inches(1.40), q_top + Inches(0.50), Inches(10.5), Inches(2.80),
        str(quote or ''), sz=32, color=q_color, bold=True, italic=True,
        font=FONT_HEAD, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

    attr_y = q_top + Inches(3.55)
    hairline(slide, Inches(1.40), attr_y, Inches(2.80),
             Emu(int(Pt(2))), pal["highlight"])

    if attribution:
        txt(slide, Inches(1.40), attr_y + Inches(0.14), Inches(10.5), Inches(0.42),
            f'\u2014 {attribution}', sz=13, color=attr_color, italic=True,
            font=FONT_BODY, align=PP_ALIGN.LEFT)

    return slide


# ===========================================================================
# Layout: DONUT / RING CHART
# ===========================================================================

def build_donut_chart(prs, meta, action_title=None, segments=None,
                      center_label=None, center_value=None,
                      legend_title=None, takeaway="", source=None,
                      subtitle=None, methodology_note=None,
                      style="merck_executive", page=None, total=None,
                      section_number=None, category=None):
    """Donut ring chart (native pptx DOUGHNUT) with legend.

    segments:     list of dicts — label (str), value (float), color (tuple opt)
    center_value: str — big number shown in the donut hole
    center_label: str — small descriptor below center_value
    legend_title: str — optional legend heading
    """
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    segments = list(segments or [])[:8]
    if not segments:
        return slide

    dark = _is_dark(style)
    bg = pal["bg"]
    DEFAULT_COLORS = [MERCK_PURPLE, MERCK_BLUE, MERCK_GOLD, MERCK_AQUA,
                      GOOD_GREEN, BAD_RED, PURPLE_MUTED, MERCK_YELLOW]
    total_val = sum(float(s.get("value", 0)) for s in segments) or 1.0

    from pptx.chart.data import CategoryChartData as _CCD
    chart_data = _CCD()
    chart_data.categories = [str(s.get("label", "")) for s in segments]
    chart_data.add_series("", [float(s.get("value", 0)) for s in segments])

    outer_r = Inches(1.90)
    cx, cy  = Inches(4.20), Inches(3.85)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        _emu(cx - outer_r), _emu(cy - outer_r),
        _emu(outer_r * 2),  _emu(outer_r * 2),
        chart_data,
    )
    ch = chart_shape.chart
    ch.has_legend = False
    ch.has_title  = False

    series = ch.series[0]
    for i, s in enumerate(segments):
        col = s.get("color") or DEFAULT_COLORS[i % len(DEFAULT_COLORS)]
        pt  = series.points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = _rgb_tuple(col)

    try:
        from lxml import etree as _et
        plot_el = ch.plots[0]._element
        for el in plot_el.iter():
            if el.tag.endswith("}holeSize"):
                el.set("val", "55")
                break
        else:
            ns = "http://schemas.openxmlformats.org/drawingml/2006/chart"
            _et.SubElement(plot_el, f"{{{ns}}}holeSize").set("val", "55")
    except Exception:
        pass

    inner_r = Inches(1.00)
    circle(slide, cx - inner_r, cy - inner_r, inner_r * 2, fill=bg)

    val_color = MERCK_YELLOW if dark else MERCK_PURPLE
    lab_color = PANEL_LIGHT  if dark else INK_GRAY
    if center_value:
        txt(slide, cx - inner_r, cy - Inches(0.46), inner_r * 2, Inches(0.72),
            str(center_value), sz=32, color=val_color, bold=True,
            font=FONT_HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if center_label:
        txt(slide, cx - inner_r, cy + Inches(0.20), inner_r * 2, Inches(0.30),
            str(center_label), sz=10, color=lab_color,
            font=FONT_BODY, align=PP_ALIGN.CENTER)

    leg_x, leg_y = Inches(8.60), Inches(2.10)
    leg_item_h   = Inches(0.44)
    dot_sz       = Inches(0.18)
    leg_color    = PANEL_LIGHT if dark else INK_DARK

    if legend_title:
        txt(slide, leg_x, leg_y - Inches(0.36), Inches(4.5), Inches(0.28),
            str(legend_title).upper(), sz=9, color=MERCK_GOLD, bold=True,
            font=FONT_BODY)

    for i, s in enumerate(segments):
        col = s.get("color") or DEFAULT_COLORS[i % len(DEFAULT_COLORS)]
        iy  = leg_y + i * leg_item_h
        circle(slide, leg_x, iy + (leg_item_h - dot_sz) / 2, dot_sz, fill=col)
        pct = f"{float(s.get('value', 0)) / total_val * 100:.0f}%"
        txt(slide, leg_x + dot_sz + Inches(0.12), iy,
            Inches(4.0), leg_item_h,
            f"{pct}  {s.get('label', '')}",
            sz=11, color=leg_color, font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE)

    return slide


# ===========================================================================
# Layout: KPI DASHBOARD
# ===========================================================================

def build_kpi_dashboard(prs, meta, action_title=None, kpis=None,
                        takeaway="", source=None, subtitle=None,
                        methodology_note=None, style="merck_executive",
                        page=None, total=None, section_number=None,
                        category=None):
    """4-6 KPI metric cards with RAG status and trend direction.

    kpis: list of dicts with keys:
        label     (str) — metric name
        value     (str) — big display number
        unit      (str, opt) — unit suffix shown smaller
        status    (str) — 'green', 'amber', or 'red'
        trend     (str) — 'up', 'down', or 'flat'
        context   (str, opt) — one-line note below the value
        sparkline (list of float, opt) — 5-8 points for mini bar
    """
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    kpis = list(kpis or [])[:6]
    if not kpis:
        return slide

    dark      = _is_dark(style)
    card_fill = PURPLE_DEEP if dark else PANEL_LIGHT
    val_color = MERCK_YELLOW if dark else MERCK_PURPLE
    lab_color = PANEL_LIGHT  if dark else INK_DARK
    ctx_color = PANEL_LIGHT  if dark else INK_GRAY
    RAG       = {"green": GOOD_GREEN, "amber": MERCK_YELLOW, "red": BAD_RED}
    TREND_G   = {"up": "↑", "down": "↓", "flat": "→"}

    n    = len(kpis)
    cols = min(n, 3)
    rows = (n + cols - 1) // cols
    zone_x, zone_y = Inches(0.55), Inches(1.62)
    zone_w, zone_h = Inches(12.2), Inches(4.70)
    gap    = Inches(0.18)
    card_w = (zone_w - gap * (cols - 1)) / cols
    card_h = (zone_h - gap * (rows - 1)) / rows

    for i, kpi in enumerate(kpis):
        row = i // cols
        col = i % cols
        cx  = zone_x + col * (card_w + gap)
        cy  = zone_y + row * (card_h + gap)

        rounded(slide, cx, cy, card_w, card_h, fill=card_fill)

        rag_col = RAG.get(str(kpi.get("status", "")).lower(), PURPLE_MUTED)
        dot_sz  = Inches(0.20)
        circle(slide, cx + card_w - dot_sz - Inches(0.14),
               cy + Inches(0.14), dot_sz, fill=rag_col)

        t_str   = str(kpi.get("trend", "")).lower()
        glyph   = TREND_G.get(t_str, "")
        t_col   = (GOOD_GREEN if t_str == "up" else
                   BAD_RED    if t_str == "down" else PURPLE_MUTED)
        if glyph:
            txt(slide, cx + Inches(0.12), cy + Inches(0.10),
                Inches(0.30), Inches(0.30), glyph,
                sz=14, color=t_col, bold=True, font=FONT_BODY)

        txt(slide, cx + Inches(0.16), cy + Inches(0.12),
            card_w - Inches(0.60), Inches(0.28),
            str(kpi.get("label", "")),
            sz=9, color=lab_color, font=FONT_BODY, align=PP_ALIGN.LEFT)

        val_y = cy + Inches(0.44)
        txt(slide, cx + Inches(0.16), val_y,
            card_w - Inches(0.32), Inches(0.80),
            str(kpi.get("value", "")),
            sz=34, color=val_color, bold=True, font=FONT_HEAD,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

        if kpi.get("unit"):
            txt(slide, cx + Inches(0.16), val_y + Inches(0.62),
                card_w - Inches(0.32), Inches(0.26),
                str(kpi["unit"]), sz=10, color=ctx_color, font=FONT_BODY)

        if kpi.get("context"):
            txt(slide, cx + Inches(0.16), cy + card_h - Inches(0.58),
                card_w - Inches(0.32), Inches(0.44),
                str(kpi["context"]), sz=9, color=ctx_color,
                italic=True, font=FONT_BODY)

        sp = kpi.get("sparkline")
        if sp and len(sp) >= 2:
            sp   = [float(v) for v in sp]
            mn, mx = min(sp), max(sp)
            rng  = (mx - mn) or 1.0
            sp_x = cx + Inches(0.16)
            sp_y = cy + card_h - Inches(0.30)
            sp_w = card_w - Inches(0.32)
            sp_h = Inches(0.16)
            bar_w = sp_w / len(sp) - Inches(0.02)
            for j, v in enumerate(sp):
                bh = max(sp_h * (v - mn) / rng, Emu(28000))
                rect(slide,
                     sp_x + j * (bar_w + Inches(0.02)), sp_y + sp_h - bh,
                     bar_w, bh,
                     fill=pal["highlight"] if j == len(sp) - 1 else PURPLE_MUTED)

    return slide


# ===========================================================================
# Layout: ICON GRID
# ===========================================================================

def build_icon_grid(prs, meta, action_title=None, items=None,
                    columns=3, takeaway="", source=None, subtitle=None,
                    methodology_note=None, style="merck_executive",
                    page=None, total=None, section_number=None,
                    category=None):
    """6 or 9 equal icon-card cells in a 2x3 or 3x3 grid.

    items: list of dicts with keys:
        icon        (str) — single emoji
        title       (str) — bold card title
        body        (str) — 1-2 sentence description
        highlighted (bool, opt) — MERCK_PURPLE card fill
    columns: 2 or 3 (default 3)
    """
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    items = list(items or [])[:9]
    if not items:
        return slide

    cols         = min(int(columns), 3)
    rows_needed  = (len(items) + cols - 1) // cols
    zone_x, zone_y = Inches(0.55), Inches(1.62)
    zone_w, zone_h = Inches(12.2), Inches(4.72)
    gap    = Inches(0.15)
    card_w = (zone_w - gap * (cols - 1)) / cols
    card_h = (zone_h - gap * (rows_needed - 1)) / rows_needed
    icon_sz = Inches(0.55)
    dark    = _is_dark(style)

    for i, item in enumerate(items):
        row = i // cols
        col = i % cols
        cx  = zone_x + col * (card_w + gap)
        cy  = zone_y + row * (card_h + gap)
        highlighted = bool(item.get("highlighted"))
        card_fill   = MERCK_PURPLE if highlighted else PANEL_LIGHT
        icon_bg     = pal["highlight"] if highlighted else MERCK_PURPLE
        title_col   = WHITE        if highlighted else MERCK_PURPLE
        body_col    = PANEL_LIGHT  if highlighted else INK_DARK

        rounded(slide, cx, cy, card_w, card_h, fill=card_fill)
        icon_cx, icon_cy = cx + Inches(0.22), cy + Inches(0.20)
        circle(slide, icon_cx, icon_cy, icon_sz, fill=icon_bg)
        txt(slide, icon_cx, icon_cy, icon_sz, icon_sz,
            str(item.get("icon", "●")),
            sz=18, color=WHITE, font=FONT_BODY,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(slide,
            icon_cx + icon_sz + Inches(0.12), icon_cy,
            card_w - icon_sz - Inches(0.46), icon_sz,
            str(item.get("title", "")),
            sz=12, color=title_col, bold=True, font=FONT_BODY,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        body_y = cy + icon_sz + Inches(0.32)
        txt(slide, cx + Inches(0.18), body_y,
            card_w - Inches(0.36), card_h - icon_sz - Inches(0.46),
            str(item.get("body", "")),
            sz=10, color=body_col, font=FONT_BODY, anchor=MSO_ANCHOR.TOP)

    return slide


# ===========================================================================
# Layout: FUNNEL / CONVERGENCE
# ===========================================================================

def build_funnel(prs, meta, action_title=None, inputs=None,
                 output=None, takeaway="", source=None, subtitle=None,
                 methodology_note=None, style="merck_corporate",
                 page=None, total=None, section_number=None,
                 category=None):
    """2-4 input boxes converging via diagonal lines to a single output.

    inputs: list of dicts with label (str) and body (str)
    output: dict with label (str) and body (str)
    """
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    inputs = list(inputs or [])[:4]
    if not inputs:
        return slide

    zone_x, zone_y = Inches(0.55), Inches(1.62)
    zone_w, zone_h = Inches(12.2), Inches(4.72)
    in_w  = Inches(3.40)
    out_w = Inches(3.80)
    box_h = Inches(0.95)
    gap   = Inches(0.22)
    n     = len(inputs)

    total_in_h = n * box_h + (n - 1) * gap
    in_start_y = zone_y + (zone_h - total_in_h) / 2
    IN_COLORS  = [MERCK_PURPLE, MERCK_BLUE, PURPLE_MUTED, MERCK_AQUA]

    for i, inp in enumerate(inputs):
        iy  = in_start_y + i * (box_h + gap)
        col = IN_COLORS[i % len(IN_COLORS)]
        rounded(slide, zone_x, iy, in_w, box_h, fill=col)
        txt(slide, zone_x + Inches(0.14), iy + Inches(0.08),
            in_w - Inches(0.28), Inches(0.24),
            str(inp.get("label", "")).upper(),
            sz=9, color=MERCK_YELLOW, bold=True, font=FONT_BODY)
        txt(slide, zone_x + Inches(0.14), iy + Inches(0.34),
            in_w - Inches(0.28), box_h - Inches(0.42),
            str(inp.get("body", "")),
            sz=10, color=WHITE, font=FONT_BODY, anchor=MSO_ANCHOR.TOP)

    out   = output or {}
    out_x = zone_x + zone_w - out_w
    out_y = zone_y + (zone_h - Inches(1.80)) / 2
    out_cy = out_y + Inches(0.90)

    for i in range(n):
        iy_c = in_start_y + i * (box_h + gap) + box_h / 2
        line(slide, zone_x + in_w, iy_c, out_x, out_cy,
             color=LIGHT_GRAY, weight=Pt(1.5))

    rounded(slide, out_x, out_y, out_w, Inches(1.80), fill=pal["highlight"])
    txt(slide, out_x + Inches(0.16), out_y + Inches(0.12),
        out_w - Inches(0.32), Inches(0.28),
        str(out.get("label", "")).upper(),
        sz=9, color=PURPLE_DEEP, bold=True, font=FONT_BODY)
    txt(slide, out_x + Inches(0.16), out_y + Inches(0.44),
        out_w - Inches(0.32), Inches(1.20),
        str(out.get("body", "")),
        sz=13, color=PURPLE_DEEP, bold=True, font=FONT_BODY,
        anchor=MSO_ANCHOR.TOP)

    return slide


# ===========================================================================
# Layout: COMPARISON TABLE (feature matrix)
# ===========================================================================

def build_comparison_table(prs, meta, action_title=None, options=None,
                            features=None, takeaway="", source=None,
                            subtitle=None, methodology_note=None,
                            style="merck_executive", page=None, total=None,
                            section_number=None, category=None):
    """Feature matrix with check/cross/partial marks.

    options:  list of str — column headers (2-4 options to compare)
    features: list of dicts with keys:
        label       (str) — row label
        values      (list of str) — one per option: 'yes'/'no'/'partial'/text
        highlighted (bool, opt) — highlight this row
    """
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    options  = list(options  or [])[:4]
    features = list(features or [])[:12]
    n_opts   = max(len(options), 1)

    MARKS = {
        "yes":     ("✓", GOOD_GREEN),
        "no":      ("✗", BAD_RED),
        "partial": ("∼", MERCK_YELLOW),
    }

    zone_x, zone_y = Inches(0.55), Inches(1.62)
    zone_w         = Inches(12.2)
    label_w        = Inches(4.20)
    opt_w          = (zone_w - label_w) / n_opts
    header_h       = Inches(0.52)
    row_h          = max((Inches(4.72) - header_h) / max(len(features), 1),
                        Inches(0.30))

    rect(slide, zone_x, zone_y, label_w, header_h, fill=PURPLE_DEEP)
    for j, opt in enumerate(options):
        px = zone_x + label_w + j * opt_w
        rect(slide, px, zone_y, opt_w, header_h, fill=MERCK_PURPLE)
        txt(slide, px, zone_y, opt_w, header_h,
            str(opt), sz=11, color=WHITE, bold=True,
            font=FONT_BODY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    for i, feat in enumerate(features):
        ry          = zone_y + header_h + i * row_h
        highlighted = bool(feat.get("highlighted"))
        row_bg      = (MERCK_PURPLE if highlighted else
                       PANEL_LIGHT  if i % 2 == 0 else WHITE)
        lbl_color   = WHITE if highlighted else INK_DARK

        rect(slide, zone_x, ry, label_w, row_h, fill=row_bg)
        txt(slide, zone_x + Inches(0.16), ry, label_w - Inches(0.20), row_h,
            str(feat.get("label", "")),
            sz=10, color=lbl_color, bold=highlighted,
            font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE)

        vals = list(feat.get("values", []))
        for j in range(n_opts):
            px = zone_x + label_w + j * opt_w
            rect(slide, px, ry, opt_w, row_h,
                 fill=row_bg, border=LIGHT_GRAY, border_w=Pt(0.25))
            val = vals[j] if j < len(vals) else ""
            mark, m_color = MARKS.get(str(val).lower(), (str(val), INK_DARK))
            if highlighted:
                m_color = WHITE
            txt(slide, px, ry, opt_w, row_h,
                mark, sz=13, color=m_color, bold=True,
                font=FONT_BODY, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    return slide


# ===========================================================================
# Layout: SCORE / RATING TABLE
# ===========================================================================

def build_score_table(prs, meta, action_title=None, rows=None,
                      scale=5, scale_label=None,
                      takeaway="", source=None, subtitle=None,
                      methodology_note=None, style="merck_executive",
                      page=None, total=None, section_number=None,
                      category=None):
    """Audit/maturity table with visual dot ratings per row.

    rows: list of dicts with keys:
        label    (str) — item name
        score    (float 0-scale) — rating
        category (str, opt) — category tag
        note     (str, opt) — brief annotation
    scale: int — max rating (default 5)
    """
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    rows = list(rows or [])[:12]
    if not rows:
        return slide

    zone_x  = Inches(0.55)
    zone_y  = Inches(1.62)
    label_w = Inches(4.80)
    cat_w   = Inches(1.60)
    note_w  = Inches(2.80)
    score_w = Inches(2.90)
    header_h = Inches(0.40)
    row_h    = min((Inches(4.72) - header_h) / len(rows), Inches(0.52))

    col_starts = [zone_x,
                  zone_x + label_w,
                  zone_x + label_w + cat_w,
                  zone_x + label_w + cat_w + note_w]
    col_widths = [label_w, cat_w, note_w, score_w]
    col_labels = ["ITEM", "CATEGORY", "NOTES",
                  scale_label or f"RATING (/{scale})"]

    for j, (cx, cw, cl) in enumerate(zip(col_starts, col_widths, col_labels)):
        rect(slide, cx, zone_y, cw, header_h,
             fill=PURPLE_DEEP if j == 3 else MERCK_PURPLE)
        txt(slide, cx + Inches(0.10), zone_y, cw - Inches(0.10), header_h,
            cl, sz=9, color=WHITE, bold=True, font=FONT_BODY,
            align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.MIDDLE)

    dot_sz, dot_gap = Inches(0.16), Inches(0.06)

    for i, row in enumerate(rows):
        ry = zone_y + header_h + i * row_h
        bg = PANEL_LIGHT if i % 2 == 0 else WHITE

        rect(slide, zone_x, ry, label_w, row_h, fill=bg)
        txt(slide, zone_x + Inches(0.14), ry, label_w - Inches(0.18), row_h,
            str(row.get("label", "")),
            sz=10, color=INK_DARK, font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE)

        rect(slide, col_starts[1], ry, cat_w, row_h,
             fill=bg, border=LIGHT_GRAY, border_w=Pt(0.25))
        if row.get("category"):
            txt(slide, col_starts[1] + Inches(0.06), ry,
                cat_w - Inches(0.12), row_h, str(row["category"]),
                sz=8, color=PURPLE_MUTED, font=FONT_BODY,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        rect(slide, col_starts[2], ry, note_w, row_h,
             fill=bg, border=LIGHT_GRAY, border_w=Pt(0.25))
        if row.get("note"):
            txt(slide, col_starts[2] + Inches(0.10), ry,
                note_w - Inches(0.16), row_h, str(row["note"]),
                sz=9, color=INK_GRAY, italic=True,
                font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE)

        rect(slide, col_starts[3], ry, score_w, row_h,
             fill=bg, border=LIGHT_GRAY, border_w=Pt(0.25))
        score       = float(row.get("score", 0))
        total_dot_w = scale * (dot_sz + dot_gap) - dot_gap
        dot_x0      = col_starts[3] + (score_w - total_dot_w) / 2
        dot_y       = ry + (row_h - dot_sz) / 2
        for d in range(scale):
            filled = (d + 1) <= score
            half   = not filled and d < score < d + 1
            fill_c = (MERCK_PURPLE if filled else
                      MERCK_GOLD   if half   else LIGHT_GRAY)
            circle(slide, dot_x0 + d * (dot_sz + dot_gap), dot_y,
                   dot_sz, fill=fill_c)

    return slide


# ===========================================================================
# Layout: INFLUENCE DIAGRAM
# ===========================================================================

def build_influence_diagram(prs, meta, action_title=None, center=None,
                             forces=None, takeaway="", source=None,
                             subtitle=None, methodology_note=None,
                             style="merck_executive", page=None,
                             total=None, section_number=None,
                             category=None):
    """Central topic box with labeled force arrows pushing in from sides.

    center: dict with label (str) and body (str)
    forces: list of dicts with:
        label (str), body (str),
        side  ('left'/'right'/'top'/'bottom'),
        tone  ('positive'/'negative'/'neutral')
    """
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    center = center or {}
    cx, cy = Inches(6.666), Inches(3.80)
    cw, ch = Inches(3.20),  Inches(1.60)
    rounded(slide, cx - cw / 2, cy - ch / 2, cw, ch, fill=MERCK_PURPLE)
    txt(slide, cx - cw / 2 + Inches(0.12), cy - ch / 2 + Inches(0.10),
        cw - Inches(0.24), Inches(0.30),
        str(center.get("label", "")).upper(),
        sz=9, color=MERCK_YELLOW, bold=True, font=FONT_BODY)
    txt(slide, cx - cw / 2 + Inches(0.12), cy - ch / 2 + Inches(0.38),
        cw - Inches(0.24), ch - Inches(0.46),
        str(center.get("body", "")),
        sz=13, color=WHITE, bold=True, font=FONT_BODY, anchor=MSO_ANCHOR.TOP)

    TONE_COLOR = {"positive": GOOD_GREEN, "negative": BAD_RED,
                  "neutral": MERCK_BLUE}
    SIDE_BASE  = {
        "left":   (Inches(0.55), Inches(2.30)),
        "right":  (Inches(9.80), Inches(2.30)),
        "top":    (Inches(3.80), Inches(1.62)),
        "bottom": (Inches(3.80), Inches(5.60)),
    }
    force_w, force_h, gap = Inches(2.60), Inches(1.00), Inches(0.22)
    processed = {s: 0 for s in SIDE_BASE}

    for f in list(forces or []):
        side = str(f.get("side", "left")).lower()
        if side not in SIDE_BASE:
            side = "left"
        idx  = processed[side]
        col  = TONE_COLOR.get(str(f.get("tone", "neutral")).lower(), MERCK_BLUE)
        bx0, by0 = SIDE_BASE[side]
        if side in ("left", "right"):
            bx, by = bx0, by0 + idx * (force_h + gap)
        else:
            bx, by = bx0 + idx * (force_w + gap), by0

        rounded(slide, bx, by, force_w, force_h, fill=col)
        txt(slide, bx + Inches(0.10), by + Inches(0.08),
            force_w - Inches(0.20), Inches(0.24),
            str(f.get("label", "")).upper(),
            sz=8, color=WHITE, bold=True, font=FONT_BODY)
        txt(slide, bx + Inches(0.10), by + Inches(0.30),
            force_w - Inches(0.20), force_h - Inches(0.36),
            str(f.get("body", "")),
            sz=9, color=WHITE, font=FONT_BODY, anchor=MSO_ANCHOR.TOP)

        if side == "left":
            ax1, ay1 = bx + force_w, by + force_h / 2
            ax2, ay2 = cx - cw / 2,  cy
        elif side == "right":
            ax1, ay1 = bx,           by + force_h / 2
            ax2, ay2 = cx + cw / 2,  cy
        elif side == "top":
            ax1, ay1 = bx + force_w / 2, by + force_h
            ax2, ay2 = cx,               cy - ch / 2
        else:
            ax1, ay1 = bx + force_w / 2, by
            ax2, ay2 = cx,               cy + ch / 2
        line(slide, ax1, ay1, ax2, ay2, color=col, weight=Pt(1.5))
        processed[side] += 1

    return slide


# ===========================================================================
# Layout: WORD CLOUD
# ===========================================================================

def build_word_cloud(prs, meta, action_title=None, words=None,
                     takeaway="", source=None, subtitle=None,
                     methodology_note=None, style="merck_executive",
                     page=None, total=None, section_number=None,
                     category=None):
    """Keyword cluster with size and color mapped to importance.

    words: list of dicts with keys:
        text   (str) — keyword or phrase
        weight (float 1-5) — importance (15pt-35pt)
        color  (tuple, opt) — explicit color override
    """
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    words = list(words or [])[:30]
    if not words:
        return slide

    WORD_COLORS = [MERCK_PURPLE, MERCK_BLUE, MERCK_GOLD,
                   GOOD_GREEN, PURPLE_MUTED, MERCK_AQUA, BAD_RED]
    GRID = [
        (0.75, 1.72), (2.50, 1.72), (4.40, 1.72), (6.30, 1.72), (8.20, 1.72), (10.10, 1.72),
        (0.55, 2.68), (2.20, 2.60), (4.00, 2.55), (5.90, 2.68), (7.80, 2.55), (9.90, 2.60),
        (0.80, 3.50), (2.60, 3.44), (4.50, 3.58), (6.20, 3.44), (8.10, 3.50), (10.00, 3.44),
        (0.60, 4.36), (2.30, 4.28), (4.20, 4.44), (6.00, 4.28), (7.90, 4.36), (9.80, 4.44),
        (1.20, 5.18), (3.10, 5.10), (5.00, 5.24), (6.80, 5.10), (8.60, 5.18), (10.40, 5.10),
    ]

    for i, word in enumerate(words):
        if i >= len(GRID):
            break
        gx, gy = GRID[i]
        weight  = max(1.0, min(5.0, float(word.get("weight") or 2)))
        sz      = int(10 + weight * 5)
        col     = word.get("color") or WORD_COLORS[i % len(WORD_COLORS)]
        txt(slide, Inches(gx), Inches(gy), Inches(2.20), Inches(0.46),
            str(word.get("text", "")),
            sz=sz, color=col, bold=(weight >= 4.0),
            font=FONT_HEAD if weight >= 4 else FONT_BODY,
            align=PP_ALIGN.LEFT)

    return slide


# ===========================================================================
# Layout: PYRAMID
# ===========================================================================

def build_pyramid(prs, meta, action_title=None, tiers=None,
                  orientation="up", takeaway="", source=None, subtitle=None,
                  methodology_note=None, style="merck_executive",
                  page=None, total=None, section_number=None,
                  category=None):
    """3-5 stacked tiers forming a pyramid.

    tiers: list of dicts ordered bottom-to-top with keys:
        label (str) — tier label
        body  (str, opt) — descriptor shown to the right
        color (tuple, opt) — fill override
    orientation: 'up' (default, wide base) or 'down' (inverted)
    """
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    tiers = list(tiers or [])[:5]
    if not tiers:
        return slide

    dark    = _is_dark(style)
    zone_cx = Inches(5.20)
    zone_y  = Inches(1.70)
    zone_h  = Inches(4.60)
    n       = len(tiers)
    tier_h  = zone_h / n
    max_w   = Inches(6.50)
    taper   = max_w / (n + 1)

    TIER_COLORS = [MERCK_PURPLE, MERCK_BLUE, PURPLE_MUTED, MERCK_AQUA, LIGHT_GRAY]

    for i, tier in enumerate(tiers):
        level = i if orientation == "up" else (n - 1 - i)
        width = max_w - level * taper
        ty    = zone_y + i * tier_h
        tx    = zone_cx - width / 2
        col   = tier.get("color") or TIER_COLORS[i % len(TIER_COLORS)]

        rect(slide, tx, ty, width, tier_h - Emu(18000), fill=col)
        txt(slide, tx + Inches(0.12), ty,
            width - Inches(0.24), tier_h - Emu(18000),
            str(tier.get("label", "")),
            sz=12, color=WHITE, bold=True, font=FONT_BODY,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        if tier.get("body"):
            body_x = zone_cx + max_w / 2 + Inches(0.22)
            line(slide, tx + width + Inches(0.10), ty + tier_h / 2,
                 body_x, ty + tier_h / 2, color=LIGHT_GRAY, weight=Pt(0.75))
            txt(slide, body_x, ty + Inches(0.05),
                Inches(4.00), tier_h - Inches(0.10),
                str(tier["body"]),
                sz=10, color=INK_GRAY if not dark else PANEL_LIGHT,
                font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE)

    return slide


# ===========================================================================
# Layout: VENN DIAGRAM
# ===========================================================================

def build_venn(prs, meta, action_title=None, circles=None,
               intersection=None, takeaway="", source=None, subtitle=None,
               methodology_note=None, style="merck_executive",
               page=None, total=None, section_number=None,
               category=None):
    """2 or 3 overlapping translucent circles.

    circles: list of 2-3 dicts with keys:
        label (str) — circle title shown above
        body  (str, opt) — text inside
        color (tuple, opt) — fill color
    intersection: str — text in the overlap center
    """
    import math as _math
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    circles = list(circles or [])[:3]
    if not circles:
        return slide

    dark = _is_dark(style)
    CIRC_COLORS = [MERCK_PURPLE, MERCK_BLUE, MERCK_GOLD]
    r = Inches(2.10)

    n = len(circles)
    if n == 2:
        overlap = Inches(0.80)
        cx_list = [Inches(4.20), Inches(4.20) + r * 2 - overlap]
        cy_list = [Inches(3.80), Inches(3.80)]
    else:
        base_cx, base_cy = Inches(6.10), Inches(3.90)
        tri_r = r * 0.82
        cx_list = [base_cx + tri_r * _math.cos(_math.radians(a))
                   for a in (210, 330, 90)]
        cy_list = [base_cy + tri_r * _math.sin(_math.radians(a))
                   for a in (210, 330, 90)]

    for i, circ in enumerate(circles):
        col = circ.get("color") or CIRC_COLORS[i % len(CIRC_COLORS)]
        ccx, ccy = cx_list[i], cy_list[i]
        shp = oval(slide, ccx - r, ccy - r, r * 2, r * 2, fill=col)

        try:
            sp_pr = shp._element.find(qn("p:spPr"))
            if sp_pr is None:
                sp_pr = shp._element.spPr
            sf = sp_pr.find(f".//{qn('a:solidFill')}")
            if sf is not None:
                sc = sf.find(qn("a:srgbClr"))
                if sc is not None:
                    from lxml import etree as _et
                    al = _et.SubElement(sc, qn("a:alpha"))
                    al.set("val", "65000")
        except Exception:
            pass

        txt(slide, ccx - r, ccy - r - Inches(0.36), r * 2, Inches(0.30),
            str(circ.get("label", "")),
            sz=11, color=WHITE if dark else INK_DARK, bold=True,
            font=FONT_BODY, align=PP_ALIGN.CENTER)

        if circ.get("body"):
            txt(slide, ccx - r * 0.7, ccy - Inches(0.26),
                r * 1.4, Inches(0.60), str(circ["body"]),
                sz=9, color=WHITE, font=FONT_BODY,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    if intersection:
        int_x = sum(cx_list) / n
        int_y = sum(cy_list) / n
        txt(slide, int_x - Inches(0.85), int_y - Inches(0.24),
            Inches(1.70), Inches(0.48), str(intersection),
            sz=10, color=WHITE, bold=True, font=FONT_BODY,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    return slide


# ===========================================================================
# Layout: RISK HEAT MAP
# ===========================================================================

def build_risk_heatmap(prs, meta, action_title=None, risks=None,
                       x_label="LIKELIHOOD", y_label="IMPACT",
                       takeaway="", source=None, subtitle=None,
                       methodology_note=None, style="merck_executive",
                       page=None, total=None, section_number=None,
                       category=None):
    """Probability x Impact heat map with numbered risk dots and legend.

    risks: list of dicts with keys:
        label      (str) — short risk name
        likelihood (int 1-5) — x-axis position
        impact     (int 1-5) — y-axis position
        color      (tuple, opt) — dot color override
    """
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    risks = list(risks or [])[:15]
    dark  = _is_dark(style)

    # grid_y must clear the action title chrome (single-line ~1.50, two-line
    # can reach ~1.90).  Use the same subtitle-aware offset other layouts use.
    GRID_N  = 5
    grid_x  = Inches(1.60)
    grid_y  = Inches(2.55) if not subtitle else Inches(2.95)
    grid_w  = Inches(6.20)
    # Bottom boundary: leave room for takeaway bar (~0.55) + source (~0.35).
    grid_h  = Inches(6.50) - grid_y - Inches(0.90)
    cell_w  = grid_w / GRID_N
    cell_h  = grid_h / GRID_N

    def _heat(row, col):
        s = row + col
        if s <= 3: return (0xC8, 0xE6, 0xC9)
        if s <= 5: return (0xFF, 0xF1, 0x7F)
        if s <= 7: return (0xFF, 0xCC, 0x66)
        return           (0xFF, 0x99, 0x99)

    for row in range(GRID_N):
        for col in range(GRID_N):
            cx = grid_x + col * cell_w
            cy = grid_y + (GRID_N - 1 - row) * cell_h
            rect(slide, cx, cy, cell_w, cell_h,
                 fill=_heat(row, col), border=WHITE, border_w=Pt(0.5))

    ax_col   = PANEL_LIGHT if dark else INK_GRAY
    ax_lbls  = ["Very Low", "Low", "Medium", "High", "Very High"]
    for j in range(GRID_N):
        txt(slide, grid_x + j * cell_w, grid_y + grid_h + Inches(0.04),
            cell_w, Inches(0.24), ax_lbls[j],
            sz=8, color=ax_col, font=FONT_BODY, align=PP_ALIGN.CENTER)
        txt(slide, Inches(0.10), grid_y + (GRID_N - 1 - j) * cell_h,
            Inches(1.40), cell_h, ax_lbls[j],
            sz=8, color=ax_col, font=FONT_BODY,
            align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    txt(slide, grid_x, grid_y + grid_h + Inches(0.30), grid_w, Inches(0.24),
        x_label.upper(), sz=9, color=ax_col, bold=True,
        font=FONT_BODY, align=PP_ALIGN.CENTER)

    RISK_COLORS = [MERCK_PURPLE, MERCK_BLUE, BAD_RED, GOOD_GREEN, MERCK_GOLD,
                   PURPLE_MUTED, MERCK_AQUA, INK_DARK, MERCK_YELLOW,
                   MERCK_BLUE, PURPLE_MUTED, GOOD_GREEN, BAD_RED, MERCK_GOLD, MERCK_PURPLE]
    dot_r = Inches(0.18)
    leg_x = Inches(8.30)
    leg_y = Inches(1.78)
    leg_h = Inches(0.34)
    txt(slide, leg_x, leg_y - Inches(0.28), Inches(4.8), Inches(0.24),
        "RISK REGISTER", sz=9,
        color=MERCK_GOLD if dark else MERCK_PURPLE, bold=True, font=FONT_BODY)

    for i, risk in enumerate(risks):
        col = risk.get("color") or RISK_COLORS[i % len(RISK_COLORS)]
        try:
            lk = max(1, min(5, int(float(risk.get("likelihood", 3)))))
        except (TypeError, ValueError):
            lk = 3
        try:
            im = max(1, min(5, int(float(risk.get("impact", 3)))))
        except (TypeError, ValueError):
            im = 3
        dx  = grid_x + (lk - 1) * cell_w + cell_w / 2
        dy  = grid_y + (GRID_N - im) * cell_h + cell_h / 2
        circle(slide, dx - dot_r, dy - dot_r, dot_r * 2, fill=col)
        txt(slide, dx - dot_r, dy - dot_r, dot_r * 2, dot_r * 2,
            str(i + 1), sz=8, color=WHITE, bold=True, font=FONT_BODY,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        ry = leg_y + i * (leg_h + Inches(0.04))
        circle(slide, leg_x, ry + (leg_h - dot_r * 2) / 2, dot_r * 2, fill=col)
        txt(slide, leg_x + dot_r * 2 + Inches(0.10), ry,
            Inches(4.0), leg_h,
            f"{i + 1}.  {risk.get('label', '')}",
            sz=9, color=PANEL_LIGHT if dark else INK_DARK,
            font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE)

    return slide


# ===========================================================================
# Layout: RADAR / SPIDER CHART (hand-drawn)
# ===========================================================================

def build_radar_chart(prs, meta, action_title=None, axes=None,
                      series=None, takeaway="", source=None, subtitle=None,
                      methodology_note=None, style="merck_executive",
                      page=None, total=None, section_number=None,
                      category=None):
    """Multi-axis spider chart drawn with polygon lines.

    axes:   list of str — axis labels (4-8 axes)
    series: list of dicts with keys:
        label  (str) — series name
        values (list of float 0-100) — one per axis
        color  (tuple, opt)
    """
    import math as _math
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    axes   = list(axes   or [])[:8]
    series = list(series or [])[:4]
    n_axes = len(axes)
    if n_axes < 3:
        return slide

    dark = _is_dark(style)
    cx, cy   = Inches(5.00), Inches(3.85)
    R        = Inches(2.20)
    N_RINGS  = 4
    S_COLORS = [MERCK_PURPLE, MERCK_GOLD, MERCK_BLUE, GOOD_GREEN]

    def _pt(axis_i, frac):
        angle = _math.radians(-90 + 360 / n_axes * axis_i)
        return (int(cx + frac * R * _math.cos(angle)),
                int(cy + frac * R * _math.sin(angle)))

    for ring in range(1, N_RINGS + 1):
        pts = [_pt(j, ring / N_RINGS) for j in range(n_axes)]
        for k in range(n_axes):
            line(slide, pts[k][0], pts[k][1],
                 pts[(k + 1) % n_axes][0], pts[(k + 1) % n_axes][1],
                 color=LIGHT_GRAY, weight=Pt(0.5))

    for j in range(n_axes):
        x2, y2 = _pt(j, 1.0)
        line(slide, int(cx), int(cy), x2, y2, color=LIGHT_GRAY, weight=Pt(0.75))

    for j, lbl in enumerate(axes):
        ax, ay = _pt(j, 1.20)
        txt(slide, ax - Inches(0.70), ay - Inches(0.20),
            Inches(1.40), Inches(0.42), lbl,
            sz=9, color=PANEL_LIGHT if dark else INK_DARK,
            font=FONT_BODY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    for si, ser in enumerate(series):
        vals = [float(v) for v in ser.get("values", [])]
        if len(vals) < n_axes:
            vals += [0.0] * (n_axes - len(vals))
        col = ser.get("color") or S_COLORS[si % len(S_COLORS)]
        pts = [_pt(j, max(0.0, min(1.0, vals[j] / 100.0)))
               for j in range(n_axes)]
        _freeform_poly(slide, pts, fill=col, border=col, border_w=Pt(1.5))

    leg_x, leg_y, leg_h = Inches(10.60), Inches(2.30), Inches(0.36)
    for si, ser in enumerate(series):
        col = ser.get("color") or S_COLORS[si % len(S_COLORS)]
        ry  = leg_y + si * (leg_h + Inches(0.10))
        rect(slide, leg_x, ry + Inches(0.10),
             Inches(0.20), Inches(0.16), fill=col)
        txt(slide, leg_x + Inches(0.28), ry, Inches(2.30), leg_h,
            str(ser.get("label", "")),
            sz=9, color=PANEL_LIGHT if dark else INK_DARK,
            font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE)

    return slide


# ===========================================================================
# Layout: PROS / CONS
# ===========================================================================

def build_pros_cons(prs, meta, action_title=None, pros=None, cons=None,
                    pros_label="ADVANTAGES", cons_label="RISKS",
                    subject=None, takeaway="", source=None, subtitle=None,
                    methodology_note=None, style="merck_executive",
                    page=None, total=None, section_number=None,
                    category=None):
    """Two-panel green pros vs red cons layout.

    pros:    list of str — advantage bullets
    cons:    list of str — risk/disadvantage bullets
    subject: str — optional dark banner above both panels
    """
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    pros = list(pros or [])
    cons = list(cons or [])
    zone_x = Inches(0.55)
    zone_y = Inches(2.10) if not subtitle else Inches(2.55)
    zone_w = Inches(12.2)
    zone_h = Inches(6.50) - zone_y - Inches(1.10)   # respect takeaway/source at bottom
    gap    = Inches(0.14)

    if subject:
        subj_h = Inches(0.46)
        rect(slide, zone_x, zone_y, zone_w, subj_h, fill=PURPLE_DEEP)
        txt(slide, zone_x + Inches(0.20), zone_y, zone_w, subj_h,
            str(subject), sz=13, color=WHITE, bold=True,
            font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE)
        panel_y = zone_y + subj_h + gap
        panel_h = zone_h - subj_h - gap
    else:
        panel_y, panel_h = zone_y, zone_h

    panel_w = (zone_w - gap) / 2
    HDR_H   = Inches(0.46)
    BULLET  = "•  "
    PANEL_BG  = {"pros": (0xE8, 0xF5, 0xE9), "cons": (0xFF, 0xEB, 0xEB)}
    HDR_COLOR = {"pros": GOOD_GREEN, "cons": BAD_RED}
    ICON      = {"pros": "✓", "cons": "✗"}

    for side in ("pros", "cons"):
        items = pros if side == "pros" else cons
        px    = zone_x if side == "pros" else zone_x + panel_w + gap
        lbl   = pros_label if side == "pros" else cons_label
        rect(slide, px, panel_y, panel_w, panel_h, fill=PANEL_BG[side])
        rect(slide, px, panel_y, panel_w, HDR_H,   fill=HDR_COLOR[side])
        txt(slide, px + Inches(0.18), panel_y, panel_w, HDR_H,
            f"{ICON[side]}  {lbl}",
            sz=12, color=WHITE, bold=True, font=FONT_BODY,
            anchor=MSO_ANCHOR.MIDDLE)
        body_y = panel_y + HDR_H + Inches(0.12)
        body_h = panel_h - HDR_H - Inches(0.16)
        item_h = min(body_h / max(len(items), 1), Inches(0.60))
        for i, item in enumerate(items):
            txt(slide, px + Inches(0.18), body_y + i * item_h,
                panel_w - Inches(0.30), item_h,
                BULLET + str(item),
                sz=10, color=INK_DARK, font=FONT_BODY, anchor=MSO_ANCHOR.TOP)

    return slide


# ===========================================================================
# Layout: LAYERED STACK
# ===========================================================================

def build_layered_stack(prs, meta, action_title=None, layers=None,
                        orientation="vertical", takeaway="", source=None,
                        subtitle=None, methodology_note=None,
                        style="merck_executive", page=None, total=None,
                        section_number=None, category=None):
    """Stacked layer architecture diagram (tech stack, org tiers, platform).

    layers: list of dicts (top-to-bottom for vertical) with keys:
        label (str) — layer name
        body  (str, opt) — short description
        icon  (str, opt) — emoji
        color (tuple, opt) — fill override
    orientation: 'vertical' (default) or 'horizontal'
    """
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    layers = list(layers or [])[:7]
    if not layers:
        return slide

    L_COLORS = [MERCK_PURPLE, MERCK_BLUE, PURPLE_MUTED,
                 MERCK_AQUA, GOOD_GREEN, MERCK_GOLD, LIGHT_GRAY]
    zone_x, zone_y = Inches(0.55), Inches(1.65)
    zone_w, zone_h = Inches(12.2), Inches(4.62)
    gap = Inches(0.10)
    n   = len(layers)

    if orientation == "horizontal":
        layer_w = (zone_w - gap * (n - 1)) / n
        for i, layer in enumerate(layers):
            lx  = zone_x + i * (layer_w + gap)
            col = layer.get("color") or L_COLORS[i % len(L_COLORS)]
            rounded(slide, lx, zone_y, layer_w, zone_h, fill=col)
            y_c = zone_y + Inches(0.22)
            if layer.get("icon"):
                txt(slide, lx, y_c, layer_w, Inches(0.55),
                    str(layer["icon"]), sz=20, color=WHITE,
                    font=FONT_BODY, align=PP_ALIGN.CENTER)
                y_c += Inches(0.60)
            txt(slide, lx + Inches(0.10), y_c, layer_w - Inches(0.20), Inches(0.42),
                str(layer.get("label", "")),
                sz=11, color=WHITE, bold=True, font=FONT_BODY,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            if layer.get("body"):
                txt(slide, lx + Inches(0.10), y_c + Inches(0.48),
                    layer_w - Inches(0.20),
                    zone_h - (y_c - zone_y) - Inches(0.56),
                    str(layer["body"]),
                    sz=9, color=WHITE, font=FONT_BODY, anchor=MSO_ANCHOR.TOP)
    else:
        layer_h = (zone_h - gap * (n - 1)) / n
        for i, layer in enumerate(layers):
            ly  = zone_y + i * (layer_h + gap)
            col = layer.get("color") or L_COLORS[i % len(L_COLORS)]
            rounded(slide, zone_x, ly, zone_w, layer_h, fill=col)
            x_c = zone_x + Inches(0.18)
            if layer.get("icon"):
                txt(slide, x_c, ly, Inches(0.55), layer_h,
                    str(layer["icon"]), sz=16, color=WHITE, font=FONT_BODY,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                x_c += Inches(0.60)
            txt(slide, x_c, ly, Inches(3.0), layer_h,
                str(layer.get("label", "")),
                sz=12, color=WHITE, bold=True, font=FONT_BODY,
                anchor=MSO_ANCHOR.MIDDLE)
            if layer.get("body"):
                txt(slide, x_c + Inches(3.10), ly,
                    zone_w - (x_c - zone_x) - Inches(3.20), layer_h,
                    str(layer["body"]),
                    sz=10, color=WHITE, font=FONT_BODY,
                    anchor=MSO_ANCHOR.MIDDLE)

    return slide


# ===========================================================================
# Layout: PHOTO + TEXT (split panel)
# ===========================================================================

def build_photo_text(prs, meta, action_title=None, image_path=None,
                     image_label=None, title=None, bullets=None,
                     image_side="left", takeaway="", source=None,
                     subtitle=None, methodology_note=None,
                     style="merck_executive", page=None, total=None,
                     section_number=None, category=None):
    """Half-slide image panel + half-slide text narrative.

    image_path:  str — path to png/jpg; placeholder drawn if None/missing
    image_label: str — caption inside the placeholder
    title:       str — bold heading in the text panel
    bullets:     list of str — narrative points
    image_side:  'left' (default) or 'right'
    """
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    dark    = _is_dark(style)
    zone_y  = Inches(1.62)
    zone_h  = Inches(4.72)
    panel_w = Inches(5.90)
    text_w  = Inches(12.2) - panel_w - Inches(0.20)

    if image_side == "right":
        img_x  = Inches(0.55) + text_w + Inches(0.20)
        text_x = Inches(0.55)
    else:
        img_x  = Inches(0.55)
        text_x = Inches(0.55) + panel_w + Inches(0.20)

    if image_path and os.path.isfile(str(image_path)):
        try:
            slide.shapes.add_picture(str(image_path),
                                     _emu(img_x), _emu(zone_y),
                                     _emu(panel_w), _emu(zone_h))
        except Exception:
            _draw_img_placeholder(slide, img_x, zone_y, panel_w, zone_h,
                                  image_label, dark)
    else:
        _draw_img_placeholder(slide, img_x, zone_y, panel_w, zone_h,
                              image_label, dark)

    y_cur = zone_y + Inches(0.14)
    if title:
        txt(slide, text_x, y_cur, text_w, Inches(0.58), str(title),
            sz=16, color=MERCK_YELLOW if dark else MERCK_PURPLE,
            bold=True, font=FONT_HEAD, anchor=MSO_ANCHOR.TOP)
        hairline(slide, text_x, y_cur + Inches(0.66),
                 text_w * 0.55, Emu(int(Pt(2))), MERCK_GOLD)
        y_cur += Inches(0.84)

    body_h  = zone_h - (y_cur - zone_y) - Inches(0.10)
    bullets = list(bullets or [])
    if bullets:
        item_h  = min(body_h / len(bullets), Inches(0.70))
        b_color = PANEL_LIGHT if dark else INK_DARK
        for i, b in enumerate(bullets):
            txt(slide, text_x, y_cur + i * item_h, text_w, item_h,
                "•  " + str(b),
                sz=11, color=b_color, font=FONT_BODY, anchor=MSO_ANCHOR.TOP)

    return slide


def _draw_img_placeholder(slide, x, y, w, h, label, dark):
    """Grey placeholder rectangle when no image file is provided."""
    fill = (0x55, 0x55, 0x66) if dark else (0xCC, 0xCC, 0xDD)
    rect(slide, x, y, w, h, fill=fill)
    txt(slide, x, y + h / 2 - Inches(0.28), w, Inches(0.48),
        str(label) if label else "[ Image ]",
        sz=12, color=WHITE, italic=True, font=FONT_BODY,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ===========================================================================
# Layout: FISHBONE (Ishikawa / Cause-and-Effect Diagram)
# ===========================================================================

def build_fishbone(prs, meta, action_title=None, effect=None, bones=None,
                   takeaway=None, source=None, category=None, subtitle=None,
                   style="merck_executive", page=None, total=None,
                   section_number=None, methodology_note=None, content=None):
    """Cause-and-effect (Ishikawa) diagram.

    effect: str - the problem or outcome at the arrow head (right side).
    bones: list of dicts {label, causes} - up to 6 branch categories.
           Alternates above/below the spine.
    """
    if content:
        if "effect"   in content: effect   = content["effect"]
        if "bones"    in content: bones    = content["bones"]
        if "takeaway" in content: takeaway = content["takeaway"]
    style = _style_or_promote(category, style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category, subtitle=subtitle,
                 takeaway=takeaway, source=source, page=page, total=total,
                 palette=style, section_number=section_number,
                 methodology_note=methodology_note)

    dark = _is_dark(style)
    text_color = WHITE if dark else INK_DARK

    zone_top = Inches(2.55) if not subtitle else Inches(2.95)
    zone_h   = Inches(6.40) - zone_top
    mid_y    = zone_top + zone_h / 2

    spine_x1  = Inches(0.80)
    spine_x2  = Inches(10.80)
    effect_w  = Inches(2.00)
    effect_h  = Inches(0.90)

    # Spine
    line(slide, spine_x1, mid_y, spine_x2, mid_y, MERCK_GOLD, Pt(3))

    # Effect box at right end
    rounded(slide, spine_x2 - effect_w / 2, mid_y - effect_h / 2,
            effect_w, effect_h, fill=MERCK_PURPLE)
    txt(slide, spine_x2 - effect_w / 2, mid_y - effect_h / 2,
        effect_w, effect_h, str(effect or "Effect"),
        sz=12, color=WHITE, bold=True, font=FONT_BODY,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    bs = list(bones or [])[:6]
    above = [b for i, b in enumerate(bs) if i % 2 == 0]
    below = [b for i, b in enumerate(bs) if i % 2 == 1]

    n_above = len(above)
    n_below = len(below)
    usable  = spine_x2 - spine_x1 - Inches(1.20)

    for group, sign, n_grp in [(above, -1, n_above), (below, 1, n_below)]:
        if not n_grp:
            continue
        step = usable / n_grp
        for k, bone in enumerate(group):
            bx      = spine_x1 + Inches(0.60) + k * step + step / 2
            bone_y  = mid_y + sign * Inches(1.40)
            meet_x  = bx + Inches(0.40) * sign * -1
            # Diagonal branch to spine
            line(slide, bx, bone_y, meet_x, mid_y, MERCK_PURPLE, Pt(1.5))
            # Category label
            lbl_w = Inches(2.00)
            lbl_h = Inches(0.34)
            lbl_y = bone_y - lbl_h if sign < 0 else bone_y
            rounded(slide, bx - lbl_w / 2, lbl_y, lbl_w, lbl_h, fill=MERCK_PURPLE)
            txt(slide, bx - lbl_w / 2, lbl_y, lbl_w, lbl_h,
                str(bone.get("label", "")),
                sz=9, color=WHITE, bold=True, font=FONT_BODY,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            # Sub-causes
            for ci, cause in enumerate(list(bone.get("causes") or [])[:3]):
                cy2 = bone_y + sign * Inches(0.42 + ci * 0.30)
                txt(slide, bx - Inches(1.10), cy2,
                    Inches(2.20), Inches(0.28),
                    "•  " + str(cause), sz=9, color=text_color,
                    font=FONT_BODY, align=PP_ALIGN.LEFT)

    return slide


# ===========================================================================
# Layout: JOURNEY MAP (multi-row swim-lane)
# ===========================================================================

def build_journey_map(prs, meta, action_title=None, phases=None, rows=None,
                      takeaway=None, source=None, category=None, subtitle=None,
                      style="merck_executive", page=None, total=None,
                      section_number=None, methodology_note=None, content=None):
    """Customer journey map: phases as columns, swim-lane rows as labeled tracks.

    phases: list of str - column headers (e.g. ["Aware","Consider","Buy","Use"]).
    rows: list of dicts {label, cells} - cells is a list of str, one per phase.
    """
    if content:
        if "phases"   in content: phases   = content["phases"]
        if "rows"     in content: rows     = content["rows"]
        if "takeaway" in content: takeaway = content["takeaway"]
    style = _style_or_promote(category, style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category, subtitle=subtitle,
                 takeaway=takeaway, source=source, page=page, total=total,
                 palette=style, section_number=section_number,
                 methodology_note=methodology_note)

    dark = _is_dark(style)
    text_color = WHITE if dark else INK_DARK

    phases    = list(phases or [])
    rows      = list(rows   or [])
    n_phases  = max(len(phases), 1)
    n_rows    = max(len(rows),   1)

    zone_top = Inches(2.55) if not subtitle else Inches(2.95)
    zone_h   = Inches(6.40) - zone_top
    label_w  = Inches(1.60)
    grid_x   = Inches(0.65) + label_w
    grid_w   = Inches(12.33) - grid_x
    col_w    = grid_w / n_phases
    hdr_h    = Inches(0.44)
    row_h    = (zone_h - hdr_h) / n_rows

    # Phase header row
    for j, ph in enumerate(phases):
        cx         = grid_x + j * col_w
        hdr_fill   = MERCK_PURPLE if j % 2 == 0 else PURPLE_DEEP
        rect(slide, cx, zone_top, col_w, hdr_h, fill=hdr_fill)
        txt(slide, cx, zone_top, col_w, hdr_h, str(ph),
            sz=10, color=WHITE, bold=True, font=FONT_BODY,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    row_label_colors = [MERCK_PURPLE, LY_CYAN, MERCK_GOLD, OP_LIME, PURPLE_MUTED, GOOD_GREEN]

    for i, row in enumerate(rows):
        ry = zone_top + hdr_h + i * row_h
        lc = row_label_colors[i % len(row_label_colors)]
        # Row label
        rounded(slide, Inches(0.65), ry, label_w, row_h, fill=lc)
        txt(slide, Inches(0.65), ry, label_w, row_h,
            str(row.get("label", "")),
            sz=9, color=WHITE, bold=True, font=FONT_BODY,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Cells
        cells = list(row.get("cells") or [])
        for j in range(n_phases):
            cx        = grid_x + j * col_w
            cell_fill = PANEL_LIGHT if (i + j) % 2 == 0 else WHITE
            cell      = rounded(slide, cx, ry, col_w, row_h, fill=cell_fill)
            _apply_border(cell, LIGHT_GRAY, Pt(0.5))
            cell_text = cells[j] if j < len(cells) else ""
            if cell_text:
                txt(slide, cx + Inches(0.10), ry + Inches(0.08),
                    col_w - Inches(0.20), row_h - Inches(0.16),
                    str(cell_text), sz=9, color=text_color,
                    font=FONT_BODY, anchor=MSO_ANCHOR.TOP)

    return slide
