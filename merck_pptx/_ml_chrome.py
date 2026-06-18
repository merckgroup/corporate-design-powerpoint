from __future__ import annotations
import sys
from typing import Optional
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from ._ml_constants import (
    FONT_HEAD, FONT_BODY, SLIDE_W, SLIDE_H,
    MERCK_PURPLE, MERCK_GOLD, MERCK_YELLOW, MERCK_BLUE,
    WHITE, INK_DARK, INK_GRAY, LIGHT_GRAY, PANEL_LIGHT,
    PURPLE_DEEP, PURPLE_MUTED, BAD_RED, GOOD_GREEN,
    ACT_PURPLE, LY_CYAN, OP_LIME, PHASE_1_COLOR, PHASE_2_COLOR, PHASE_3_COLOR,
    GOLD_RULE_X, GOLD_RULE_Y, GOLD_RULE_W, GOLD_RULE_H,
    CLASS_BADGE_X, CLASS_BADGE_Y, CLASS_BADGE_W,
    BREADCRUMB_X, BREADCRUMB_Y, BREADCRUMB_W, BREADCRUMB_H,
    SECTION_CIRCLE_X, SECTION_CIRCLE_Y, SECTION_CIRCLE_D,
    SECTION_TAG_X, SECTION_TAG_Y, SECTION_TAG_W, SECTION_TAG_H,
    TITLE_X, TITLE_Y_NUMBERED, TITLE_Y_UNNUMBERED, TITLE_W, TITLE_H,
    SUB_X, SUB_W, SUB_H, SUB_GAP,
    CONTENT_X, CONTENT_Y, CONTENT_Y_SUBTITLE, CONTENT_W, CONTENT_H,
    SOURCE_Y, SOURCE_H, TAKEAWAY_Y, TAKEAWAY_H, PHASE_Y, PHASE_H,
    FOOTER_Y, FOOTER_H, FOOTER_TEXT_Y,
    PALETTES, _palette_for, _rgb_tuple, _is_dark,
)
from ._ml_primitives import (
    rect, rounded, oval, circle, line, hairline, txt, _add_run,
    _freeform_poly, _emu, _apply_fill, _apply_border,
)
from ._ml_icons import draw_icon, ICON_REGISTRY

# ===========================================================================
# Tracked-letter helper (simulates letter-spacing)
# ===========================================================================

def _tracked(text: str) -> str:
    """Insert spaces between letters to simulate uppercase tracking."""
    if not text:
        return ""
    words = str(text).upper().split()
    spaced = [" ".join(list(w)) for w in words]
    return "   ".join(spaced)


def _track_letters(text: str) -> str:
    """Alias for _tracked, exposed for layouts that prefer this name."""
    return _tracked(text)


def _estimate_card_content_h(col: dict, default_h):
    """Estimate vertical height needed for a card based on label + title + items.

    col supports keys: label, title (or body for two/three col), body (for stat
    strip), items (list), milestone, timestamp. Returns an Emu value.
    """
    h = Inches(0.40)  # padding top + label
    title_text = col.get("title") or col.get("body_title") or col.get("body") or ""
    if title_text:
        # ~1 line per 32 chars at 16pt.
        title_lines = max(1, (len(str(title_text)) // 32) + (1 if len(str(title_text)) % 32 else 0))
        h += Inches(0.42) * title_lines
    body_text = col.get("desc") or ""
    if body_text:
        body_lines = max(1, (len(str(body_text)) // 50) + (1 if len(str(body_text)) % 50 else 0))
        h += Inches(0.24) * body_lines
    items = col.get("items") or []
    if items:
        h += Inches(0.05)
        h += Inches(0.32) * len(items)
    if col.get("milestone") or col.get("timestamp"):
        h += Inches(0.50)
    h += Inches(0.30)  # bottom padding
    if h < default_h:
        return default_h
    return h


def _compute_row_card_height(cols, default_h):
    """Equal-height enforcement: return the max needed across all cards.

    Cards in a row must visually look like a row of equals; cards with less
    content show whitespace at the bottom rather than auto-shrinking.
    """
    if not cols:
        return default_h
    heights = [_estimate_card_content_h(c or {}, default_h) for c in cols]
    return max(heights)


def _pad_int(n) -> str:
    try:
        return f"{int(n):02d}"
    except Exception:
        return str(n)


def _format_section_number(num):
    """Render the section-circle text. Pure ints zero-pad to 2 digits to match
    the agenda chips ('01', '02', ..., '10'). Strings with letters pass through
    verbatim to support sub-labels ('3A', '3B', '12B')."""
    if num is None:
        return ""
    s = str(num).strip()
    if not s:
        return ""
    try:
        return f"{int(s):02d}"
    except (ValueError, TypeError):
        return s


# ===========================================================================
# Universal chrome
# ===========================================================================

def _top_chrome(slide, meta, category, palette, top_bar=False,
                page=None, total=None):
    """Top chrome: classification badge + thin progress bar.

    The previous gold rule + breadcrumb ('DECK_LABEL • CATEGORY • MONTH YEAR')
    duplicated information that's already in the footer band (deck_label +
    category) and in the section circle (category tag) — three copies of
    the same words at the top of every content slide. Replaced with:

      - Classification badge (top-right, compliance requirement)
      - Thin progress bar across the top (page / total), filled in
        MERCK_GOLD on light palettes / MERCK_YELLOW on dark. Subtle
        signal of "how far through the deck am I" — the only navigation
        question audiences actually have. Skipped on the cover (when the
        full-width gold top_bar is used instead).

    When meta['cover_top_bar'] is True OR top_bar=True, a full-width gold
    bar renders at y=0. Used for cover slides only; progress bar is
    suppressed in that case.
    """
    pal = _palette_for(palette)
    meta = meta or {}
    dark = _is_dark(palette)

    # NOTE: meta['cover_top_bar'] is intentionally NOT read here — only the
    # explicit top_bar=True argument enables the cover gold-bar treatment.
    # build_cover passes top_bar=True when meta['cover_top_bar'] is set, so
    # content slides never accidentally inherit the cover bar even when the
    # meta key is present in the deck-wide meta dict.
    if bool(top_bar):
        # Full-width gold bar at the very top of the slide (cover only).
        hairline(slide, Inches(0.0), Inches(0.0), SLIDE_W, Inches(0.20),
                 MERCK_GOLD)
    elif page is not None and total:
        # Thin progress bar at the very top: light track + gold fill.
        try:
            p = int(page)
            t = int(total)
        except (TypeError, ValueError):
            p, t = 0, 0
        if t > 0 and p > 0:
            bar_x = Inches(0.65)
            bar_y = Inches(0.10)
            bar_w = Inches(12.0)
            bar_h = Inches(0.05)
            track_color = PURPLE_MUTED if dark else LIGHT_GRAY
            fill_color = pal["hot"] if dark else pal["highlight"]
            hairline(slide, bar_x, bar_y, bar_w, bar_h, track_color)
            # Filled portion proportional to page / total. Min 1px visible.
            frac = max(0.0, min(1.0, p / t))
            fill_w_emu = int(bar_w * frac)
            if fill_w_emu > 0:
                hairline(slide, bar_x, bar_y, fill_w_emu, bar_h, fill_color)

    # Classification badge (top-right, two stacked lines). Kept for
    # compliance — internal vs confidential should be visible at a glance.
    classification = meta.get("classification") or "Internal"
    badge = slide.shapes.add_textbox(CLASS_BADGE_X, CLASS_BADGE_Y,
                                     CLASS_BADGE_W, Inches(0.50))
    tf = badge.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    _add_run(p1, "Classification:", sz=9, color=MERCK_GOLD, font=FONT_BODY)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    _add_run(p2, str(classification).upper(), sz=9, color=MERCK_GOLD,
             bold=True, font=FONT_BODY)


def _bottom_chrome(slide, meta, category, page, total, palette):
    """Footer text boxes: deck label + category (left) and page number (right).

    No band is drawn. Both elements use Rich Purple at 8 pt per CD guidelines.
    Custom text boxes are used (not native placeholder injection) because the
    content is per-slide dynamic.
    """
    meta = meta or {}

    deck_label = meta.get("deck_label", "")
    cat = category or ""
    left_parts = [p for p in [deck_label, cat] if p]
    # On dark-background slides (merck_storytelling) Rich Purple is invisible —
    # use a light off-white so the footer remains readable.
    _ftr_dark = _is_dark(palette)
    footer_color = MERCK_PURPLE if not _ftr_dark else PANEL_LIGHT

    if left_parts:
        txt(slide, Inches(0.65), FOOTER_TEXT_Y, Inches(10.0), Inches(0.20),
            "   •   ".join(left_parts), sz=8,
            color=footer_color, bold=False,
            font=FONT_BODY, anchor=MSO_ANCHOR.TOP)

    if page is not None:
        page_text = _pad_int(page)
        if total:
            page_text = f"{_pad_int(page)} / {_pad_int(total)}"
        txt(slide, Inches(12.0), FOOTER_TEXT_Y, Inches(1.0), Inches(0.20),
            page_text, sz=8, color=footer_color,
            bold=False, font=FONT_BODY, align=PP_ALIGN.RIGHT)


_ICON_DISPATCH = {
    "lightbulb": "icon_lightbulb", "alert": "icon_alert",
    "info": "icon_info", "target": "icon_target",
    "check": "icon_check", "flag": "icon_flag",
    "shield": "icon_shield", "lock": "icon_lock",
    "gear": "icon_gear", "users": "icon_users",
    "chart_bar": "icon_chart_bar", "trending_up": "icon_trending_up",
}

def _section_marker(slide, number, category, palette, icon=None):
    """Numbered circle + uppercase category tag.

    The circle is only drawn when it has visible content (a number or icon) —
    an empty filled circle with nothing inside is omitted. The category tag
    always renders independently of the circle.
    """
    dark = _is_dark(palette)
    circle_fill = MERCK_PURPLE if not dark else MERCK_YELLOW
    number_color = WHITE if not dark else INK_DARK
    tag_color = PURPLE_MUTED if not dark else MERCK_GOLD

    has_number = number is not None and str(number) != ""
    if has_number or icon:
        circle(slide, SECTION_CIRCLE_X, SECTION_CIRCLE_Y,
               SECTION_CIRCLE_D, fill=circle_fill)
        if has_number:
            num_sz = 22 if len(str(number)) <= 2 else 14
            txt(slide, SECTION_CIRCLE_X, SECTION_CIRCLE_Y,
                SECTION_CIRCLE_D, SECTION_CIRCLE_D,
                str(number), sz=num_sz, color=number_color, bold=True,
                font=FONT_BODY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if category:
        txt(slide, SECTION_TAG_X, SECTION_TAG_Y,
            SECTION_TAG_W, SECTION_TAG_H,
            _tracked(category), sz=11, color=tag_color, bold=True,
            font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE)


def _render_action_title(slide, x, y, w, h, content, palette,
                         size=22, italic_color=None, base_color=None):
    """Render action title that may be a string OR list of (text, italic_bool) tuples."""
    pal = _palette_for(palette)
    dark = _is_dark(palette)
    if base_color is None:
        base_color = WHITE if dark else MERCK_PURPLE
    if italic_color is None:
        italic_color = pal["hot"]

    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    # Auto-shrink font if the text is too long to fit — prevents visible clipping.
    try:
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    if isinstance(content, (list, tuple)) and content and \
       all(isinstance(seg, (list, tuple)) and len(seg) == 2 for seg in content):
        for seg_text, seg_italic in content:
            color = italic_color if seg_italic else base_color
            _add_run(p, seg_text, sz=size, color=color, bold=False,
                     italic=bool(seg_italic), font=(FONT_HEAD if seg_italic else FONT_BODY))
    else:
        text = "" if content is None else str(content)
        _add_run(p, text, sz=size, color=base_color, bold=False,
                 italic=False, font=FONT_BODY)
    return box


def _subtitle(slide, x, y, w, h, text, palette):
    if not text:
        return None
    pal = _palette_for(palette)
    color = pal["ink_2"] if not _is_dark(palette) else PANEL_LIGHT
    return txt(slide, x, y, w, h, text, sz=13, color=color, italic=True,
               font=FONT_BODY, anchor=MSO_ANCHOR.TOP)


def _source_line(slide, x, y, w, h, text, palette):
    if not text:
        return None
    pal = _palette_for(palette)
    color = pal["ink_2"] if not _is_dark(palette) else PANEL_LIGHT
    body = str(text)
    if not body.lower().startswith("source"):
        body = f"Source: {body}"
    return txt(slide, x, y, w, h, body, sz=8, color=color, italic=True,
               font=FONT_BODY)


def _methodology_note(slide, x, y, w, h, text, palette):
    """Render a small italic methodology note in INK_GRAY Verdana 9pt."""
    if not text:
        return None
    pal = _palette_for(palette)
    color = pal["ink_2"] if not _is_dark(palette) else PANEL_LIGHT
    body = str(text)
    return txt(slide, x, y, w, h, body, sz=9, color=color, italic=True,
               font=FONT_BODY)


def _takeaway_band(slide, text, palette):
    """Bottom takeaway band placed just above the footer band.

    Single consistent render across all slides: 9pt Verdana bold, mixed case,
    no letter-tracking, no uppercase. The earlier letter-tracked-uppercase
    treatment (the "McKinsey caps look") inflated visual width ~2.2x and
    silently clipped takeaways over ~50 chars at the right edge. Bands are
    decorated by their colored fill (yellow / purple / gold) — typography
    doesn't need to also shout.

    Band geometry: starts at x=0.65, width=11.1 → ends at x=11.75. This is
    0.07" short of the Merck "M" logo placed at x=11.82 in the master slide
    chrome; the previous 12.0" width covered the logo on every content
    slide.

    Author contract: takeaway should be 40-90 characters. Up to ~120 chars
    fit at sz=9 mixed case across the 10.7" usable band width (11.1 minus
    2x 0.20" margins); beyond that, word-wrap pushes to a second line that
    the 0.22" band height clips. The runner validates this.
    """
    if not text:
        return None
    pal  = _palette_for(palette)
    dark = _is_dark(palette)
    # Band fill: use the theme's primary accent colour so each theme's
    # takeaway band is visually distinct.  Dark themes use "hot" (their
    # bright accent on dark bg); light themes use "accent".
    fill = pal["hot"] if dark else pal["accent"]
    # Text colour: pick white or near-black based on perceived luminance of
    # the fill so readability is maintained for all 6 themes.
    r, g, b = fill
    luminance = 0.299 * r + 0.587 * g + 0.114 * b
    text_color = PURPLE_DEEP if luminance > 140 else WHITE
    band = rounded(slide, Inches(0.65), TAKEAWAY_Y, Inches(11.1), TAKEAWAY_H,
                   fill=fill, adj=50000)
    tf = band.text_frame
    tf.margin_left = Inches(0.20)
    tf.margin_right = Inches(0.20)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _add_run(p, str(text).strip(), sz=9, color=text_color, bold=True,
             font=FONT_BODY)
    return band


# ===========================================================================
# Backward-compatible chrome surface
# ===========================================================================

def merck_stub(slide, palette):
    """Legacy: no-op (kept for API parity)."""
    return None


def category_tag(slide, text, palette):
    """Legacy: small uppercase tag at the breadcrumb position."""
    if not text:
        return None
    color = MERCK_GOLD
    return txt(slide, BREADCRUMB_X, BREADCRUMB_Y, BREADCRUMB_W, BREADCRUMB_H,
               _tracked(text), sz=10, color=color, bold=True, font=FONT_BODY,
               anchor=MSO_ANCHOR.MIDDLE)


def action_title(slide, text, palette):
    """Legacy entry point: render action title at TITLE_Y_NUMBERED."""
    return _render_action_title(slide, TITLE_X, TITLE_Y_NUMBERED,
                                TITLE_W, TITLE_H, text, palette)


def subtitle_line(slide, text, palette):
    """Legacy entry point: render subtitle below the action title."""
    return _subtitle(slide, SUB_X, Inches(3.05), SUB_W, SUB_H, text, palette)


def takeaway_band(slide, text, palette):
    return _takeaway_band(slide, text, palette)


def source_line(slide, text, palette):
    return _source_line(slide, Inches(0.65), Inches(7.22),
                        Inches(11.0), Inches(0.20), text, palette)


def footnotes_block(slide, items, palette, y=None):
    """Render numbered footnotes just above the source line.

    items: list of (number, text) tuples. Renders as a single textbox with
    each footnote on its own paragraph. Numbers are rendered as small
    superscript-style runs. Verdana italic 9pt INK_GRAY.
    """
    if not items:
        return None
    pal = _palette_for(palette)
    color = pal["ink_2"] if not _is_dark(palette) else PANEL_LIGHT
    if y is None:
        # Place above SOURCE_Y, but stack upward if many footnotes.
        n = len(items)
        block_h = Inches(0.18) * n
        y = SOURCE_Y - block_h - Inches(0.02)
    box = slide.shapes.add_textbox(Inches(0.65), y, Inches(11.5),
                                   Inches(0.20) * max(len(items), 1) + Inches(0.05))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for i, item in enumerate(items):
        try:
            number, text = item[0], item[1]
        except Exception:
            continue
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(1)
        # Small superscript-like number then italic body text.
        _add_run(p, f"{number}  ", sz=8, color=MERCK_GOLD, bold=True,
                 font=FONT_BODY)
        _add_run(p, str(text), sz=9, color=color, italic=True, font=FONT_BODY)
    return box


def _superscript(slide, x, y, number, color=None):
    """Tiny superscript number used next to chart data labels."""
    txt(slide, x, y, Inches(0.30), Inches(0.18),
        str(number), sz=8, color=color or MERCK_GOLD, bold=True,
        font=FONT_BODY)


def page_number(slide, n, total, palette):
    if n is None:
        return None
    page_text = _pad_int(n)
    if total:
        page_text = f"{_pad_int(n)} / {_pad_int(total)}"
    return txt(slide, Inches(12.0), FOOTER_TEXT_Y, Inches(1.0), Inches(0.20),
               page_text, sz=8, color=MERCK_PURPLE, bold=False, font=FONT_BODY,
               align=PP_ALIGN.RIGHT)


def apply_chrome(slide, meta, action_title, category=None, subtitle=None,
                 takeaway=None, source=None, page=None, total=None,
                 palette="merck_executive", section_number=None,
                 show_section_marker=True, methodology_note=None,
                 section_icon=None):
    """Lay down the universal chrome.

    action_title may be either a string or a list of (text, italic_bool)
    tuples. When show_section_marker is True (default) and category is set,
    a section number marker is rendered.

    section_number: number to show in the circle. Each content slide should
        have a UNIQUE sequential number (1, 2, 3, ...) that matches its entry
        in the agenda. Not the page number.

    section_icon: if provided (e.g. "lightbulb"), the circle renders an icon
        instead of a number. Used for exec_summary and other slides where the
        circle anchors on a concept rather than a sequence position.

    methodology_note: optional italic note rendered just above the source
    line (or in its place if source is empty). 9pt INK_GRAY italic.
    """
    _top_chrome(slide, meta, category, palette, page=page, total=total)
    _bottom_chrome(slide, meta, category, page, total, palette)

    title_y = TITLE_Y_NUMBERED
    if show_section_marker and category:
        if section_icon:
            _section_marker(slide, None, category, palette,
                            icon=section_icon)
            title_y = TITLE_Y_NUMBERED
        else:
            num = section_number
            if num is not None:
                _section_marker(slide,
                                _format_section_number(num),
                                category, palette)
                title_y = TITLE_Y_NUMBERED
            else:
                title_y = TITLE_Y_UNNUMBERED
    else:
        title_y = TITLE_Y_UNNUMBERED

    if action_title:
        _render_action_title(slide, TITLE_X, title_y, TITLE_W, TITLE_H,
                             action_title, palette)
    if subtitle:
        sub_y = title_y + SUB_GAP
        _subtitle(slide, SUB_X, sub_y, SUB_W, SUB_H, subtitle, palette)

    if takeaway:
        _takeaway_band(slide, takeaway, palette)
    # Source + methodology note placement: when both are supplied, the
    # methodology note sits ABOVE the source line (so it doesn't collide with
    # the takeaway band at TAKEAWAY_Y). When only methodology_note is
    # supplied, it occupies the source line slot.
    if source and methodology_note:
        _methodology_note(slide, Inches(0.65),
                          SOURCE_Y - Inches(0.22),
                          Inches(12.0), Inches(0.20),
                          methodology_note, palette)
        _source_line(slide, Inches(0.65), SOURCE_Y,
                     Inches(12.0), SOURCE_H, source, palette)
    elif source:
        _source_line(slide, Inches(0.65), SOURCE_Y,
                     Inches(12.0), SOURCE_H, source, palette)
    elif methodology_note:
        _methodology_note(slide, Inches(0.65), SOURCE_Y,
                          Inches(12.0), SOURCE_H, methodology_note, palette)
    return slide


# ===========================================================================
# Craft helpers
# ===========================================================================

def stub_and_flag(slide, anchor_x, anchor_y, label, palette,
                  direction="up_right",
                  label_offset_x=Inches(0.4),
                  label_offset_y=Inches(-0.4)):
    """Draw a dot at the anchor, an elbow connector, and a small label."""
    pal = _palette_for(palette)
    color = pal["highlight"] if not _is_dark(palette) else pal["hot"]
    dot_size = Inches(0.08)
    circle(slide, anchor_x - dot_size / 2, anchor_y - dot_size / 2,
           dot_size, fill=color)

    if direction == "up_left":
        label_x = anchor_x - Inches(2.0)
        label_y = anchor_y - Inches(0.40)
    elif direction == "down_right":
        label_x = anchor_x + Inches(0.40)
        label_y = anchor_y + Inches(0.30)
    elif direction == "down_left":
        label_x = anchor_x - Inches(2.0)
        label_y = anchor_y + Inches(0.30)
    else:
        label_x = anchor_x + Inches(0.40)
        label_y = anchor_y - Inches(0.40)

    conn = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW,
                                      _emu(anchor_x), _emu(anchor_y),
                                      _emu(label_x), _emu(label_y))
    conn.shadow.inherit = False
    conn.line.color.rgb = _rgb_tuple(color)
    conn.line.width = Pt(0.75)
    box_w = Inches(2.0)
    box_h = Inches(0.30)
    txt(slide, label_x, label_y - box_h / 2,
        box_w, box_h, label, sz=9, color=color, bold=True, font=FONT_BODY)
    return conn


def decimal_align(text_value: str):
    if text_value is None:
        return ("", "")
    s = str(text_value).strip()
    if not s:
        return ("", "")
    if "." in s:
        whole, dec = s.split(".", 1)
        return (whole, "." + dec)
    return (s, "")


# ===========================================================================
# Card pattern (cream PANEL_LIGHT + gold top stripe)
# ===========================================================================



def _draw_card(slide, x, y, w, h, palette, highlighted=False,
               stripe_color=None):
    """Draw a Recall-style card. Returns (top_inside_y, body_color, label_color)."""
    pal = _palette_for(palette)
    if highlighted:
        card_fill = PURPLE_DEEP if _is_dark(palette) else MERCK_PURPLE
        body = rounded(slide, x, y, w, h, fill=card_fill)
        stripe = stripe_color or pal["hot"]
        rounded(slide, x, y, w, Inches(0.06), fill=stripe, adj=50000)
        return WHITE, PANEL_LIGHT, stripe
    # Default card: panel fill + theme-accent stripe + hairline border.
    body = rounded(slide, x, y, w, h, fill=PANEL_LIGHT)
    _apply_border(body, LIGHT_GRAY, Pt(0.5))
    stripe = stripe_color or pal["highlight"]
    rounded(slide, x, y, w, Inches(0.06), fill=stripe, adj=50000)
    return MERCK_PURPLE, INK_GRAY, stripe


def _gold_square_bullet(slide, x, y, size=Inches(0.10), color=None):
    rect(slide, x, y, size, size, fill=color or MERCK_GOLD)


def _bulleted_list(slide, x, y, w, h, items, palette,
                   text_color=None, sz=11, bullet_color=None):
    """Render a list with bullet hierarchy and CD-compliant paragraph spacing.

    Indent detection: items starting with 2+ spaces are rendered as level-3
    sub-bullets (en-dash, slightly smaller, indented) per CD bullet hierarchy.
    All other items are level-2 bullets (▪).  Level-1 (plain text, no bullet)
    is not used here — these are always body list items.

    Paragraph spacing follows CD spec: 3 pt before each paragraph (except the
    first), 3 pt after every paragraph, line spacing 1.05×.
    """
    if not items:
        return
    pal = _palette_for(palette)
    if text_color is None:
        text_color = pal["ink"]
    bcol = bullet_color or pal["highlight"]

    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)

    for i, it in enumerate(items):
        raw = str(it)
        is_sub = raw.startswith("  ")          # 2+ leading spaces → level-3
        text   = raw.lstrip() if is_sub else raw

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(3)             # CD: 3 pt before non-first paragraphs
        p.space_after = Pt(3)                  # CD: 3 pt after every paragraph

        # CD line spacing: 1.05× multiple via XML (no python-pptx shortcut for multiples)
        from lxml import etree as _et
        pPr = p._p.get_or_add_pPr()
        for existing in pPr.findall(qn("a:lnSpc")):
            pPr.remove(existing)
        lnSpc = _et.SubElement(pPr, qn("a:lnSpc"))
        spcPct = _et.SubElement(lnSpc, qn("a:spcPct"))
        spcPct.set("val", "105000")            # 105 % = 1.05×

        if is_sub:
            # Level 3: indented en-dash bullet, slightly smaller
            sub_sz = max(sz - 1, 9)
            _add_run(p, "    –  ", sz=sub_sz, color=bcol, bold=False, font=FONT_BODY)
            _add_run(p, text,     sz=sub_sz, color=text_color, font=FONT_BODY)
        else:
            # Level 2: filled square bullet
            _add_run(p, "▪  ", sz=sz, color=bcol, bold=True, font=FONT_BODY)
            _add_run(p, text, sz=sz, color=text_color, font=FONT_BODY)


# ===========================================================================
# Statement card primitive (left or top stripe)
# ===========================================================================

def statement_card(slide, x, y, w, h, label, body, palette,
                   accent="left", stripe_color=None):
    """Pale-lavender rounded statement card with a theme-accent stripe.

    accent="left" draws a vertical stripe down the left edge; accent="top"
    draws a horizontal stripe across the top. Label sits at the top-left
    inside the card in theme accent colour uppercase tracked Verdana 10pt; body in
    INK_DARK Verdana 14pt bold (auto-shrinks to 13pt for longer bodies).
    """
    pal = _palette_for(palette)
    fill = PANEL_LIGHT
    stripe = stripe_color or pal["highlight"]
    card = rounded(slide, x, y, w, h, fill=fill)
    _apply_border(card, LIGHT_GRAY, Pt(0.5))

    if accent == "top":
        rounded(slide, x, y, w, Inches(0.06), fill=stripe, adj=50000)
        text_inset_x = x + Inches(0.30)
        text_w = w - Inches(0.60)
        label_y = y + Inches(0.18)
        body_y = y + Inches(0.50)
    else:
        # Default: left stripe
        rect(slide, x, y, Inches(0.06), h, fill=stripe)
        text_inset_x = x + Inches(0.24)
        text_w = w - Inches(0.44)
        label_y = y + Inches(0.16)
        body_y = y + Inches(0.48)

    if label:
        txt(slide, text_inset_x, label_y, text_w, Inches(0.24),
            _track_letters(label), sz=10, color=pal["highlight"], bold=True,
            font=FONT_BODY)
    if body:
        body_sz = 14 if len(str(body)) <= 180 else 13
        body_h = (y + h) - body_y - Inches(0.16)
        if body_h < Inches(0.30):
            body_h = Inches(0.30)
        txt(slide, text_inset_x, body_y, text_w, body_h,
            str(body), sz=body_sz, color=INK_DARK, bold=True,
            font=FONT_BODY)
    return card


def in_slide_section(slide, y, label, palette, x=None, w=None, align="center"):
    """Render an in-slide sub-section heading.

    Small theme-accent uppercase tracked Verdana 11pt label used to break a
    content zone into sub-sections (e.g. four cards, then THREE PILLARS
    heading, then three more cards).
    """
    if not label:
        return None
    pal = _palette_for(palette)
    x_use = x if x is not None else Inches(0.65)
    w_use = w if w is not None else Inches(12.0)
    align_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }
    a = align_map.get(align, PP_ALIGN.CENTER)
    return txt(slide, x_use, y, w_use, Inches(0.28),
               _track_letters(label), sz=11, color=pal["highlight"], bold=True,
               font=FONT_BODY, align=a, anchor=MSO_ANCHOR.MIDDLE)


# ===========================================================================
# Phase progress strip (used on cover)
# ===========================================================================

def _phase_progress(slide, phases, palette):
    """Phase progress indicator above the footer band.

    phases: list of dicts {label, desc, current(bool)}.
    """
    if not phases:
        return
    pal = _palette_for(palette)
    dark = _is_dark(palette)
    n = len(phases)
    margin = Inches(0.65)
    avail = SLIDE_W - 2 * margin
    seg = avail / n

    # Draw connecting hairline first.
    cy = PHASE_Y + Inches(0.05)
    line(slide, margin + seg / 2, cy, margin + seg * (n - 0.5), cy,
         PURPLE_MUTED, Pt(0.5))

    for i, ph in enumerate(phases):
        current = bool(ph.get("current"))
        cx = margin + seg * i + seg / 2
        dot_d = Inches(0.20)
        dot_color = pal["highlight"] if current else PURPLE_MUTED
        circle(slide, cx - dot_d / 2, cy - dot_d / 2, dot_d, fill=dot_color)

        label = ph.get("label", "")
        desc = ph.get("desc", "")
        label_color = (WHITE if dark else INK_DARK) if current else PURPLE_MUTED
        desc_color = pal["highlight"] if current else PURPLE_MUTED

        # Label below the dot.
        label_y = cy + Inches(0.18)
        txt(slide, cx - seg / 2 + Inches(0.10), label_y,
            seg - Inches(0.20), Inches(0.22),
            _tracked(label), sz=9, color=label_color, bold=True,
            font=FONT_BODY, align=PP_ALIGN.CENTER)
        # Description.
        if desc:
            txt(slide, cx - seg / 2 + Inches(0.10), label_y + Inches(0.22),
                seg - Inches(0.20), Inches(0.20),
                desc, sz=8, color=desc_color, font=FONT_BODY,
                align=PP_ALIGN.CENTER)


