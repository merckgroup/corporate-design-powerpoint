from __future__ import annotations
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt
from ._ml_constants import FONT_BODY, WHITE, _rgb_tuple

# ===========================================================================
# Primitives
# ===========================================================================

def _apply_fill(shape, fill):
    if fill is None:
        shape.fill.background()
        return
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb_tuple(fill)


def _apply_border(shape, color, weight=None):
    line_obj = shape.line
    if color is None:
        line_obj.fill.background()
        return
    line_obj.color.rgb = _rgb_tuple(color)
    line_obj.width = weight if weight is not None else Pt(0.5)


def _emu(v):
    """Coerce a coordinate or extent value to an integer EMU.

    OOXML's xsd:long requires plain integers; lxml will serialize a Python
    float as '2138340.0', which PowerPoint rejects and shows as a 'needs
    repair' warning on open. python-pptx's add_shape / add_textbox happen
    to auto-cast, but add_connector does NOT — and any arithmetic involving
    `/` (Python 3 true division) silently produces a float that propagates
    all the way to the XML. This helper is the single coercion point used
    by every primitive below.
    """
    if v is None:
        return 0
    try:
        return int(v)
    except (TypeError, ValueError):
        return 0


def rect(slide, x, y, w, h, fill=None, border=None, border_w=None):
    """Sharp rectangle. Default has no border, no fill."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 _emu(x), _emu(y), _emu(w), _emu(h))
    shp.shadow.inherit = False
    _apply_fill(shp, fill)
    _apply_border(shp, border, border_w if border_w is not None else Pt(0.5))
    if shp.has_text_frame:
        shp.text_frame.text = ""
        shp.text_frame.margin_left = Inches(0.05)
        shp.text_frame.margin_right = Inches(0.05)
        shp.text_frame.margin_top = Inches(0.02)
        shp.text_frame.margin_bottom = Inches(0.02)
    return shp


def rounded(slide, x, y, w, h, fill=None, adj=6000):
    """Rounded rectangle. adj controls corner radius (0 to 50000)."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 _emu(x), _emu(y), _emu(w), _emu(h))
    shp.shadow.inherit = False
    _apply_fill(shp, fill)
    _apply_border(shp, None)
    try:
        shp.adjustments[0] = max(0, min(adj, 50000)) / 100000.0
    except Exception:
        pass
    if shp.has_text_frame:
        shp.text_frame.margin_left = Inches(0.1)
        shp.text_frame.margin_right = Inches(0.1)
        shp.text_frame.margin_top = Inches(0.05)
        shp.text_frame.margin_bottom = Inches(0.05)
    return shp


def oval(slide, x, y, w, h, fill=None):
    """Oval; no border by default."""
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 _emu(x), _emu(y), _emu(w), _emu(h))
    shp.shadow.inherit = False
    _apply_fill(shp, fill)
    _apply_border(shp, None)
    return shp


def circle(slide, x, y, size, fill=None):
    """Circle of given side length."""
    return oval(slide, x, y, size, size, fill=fill)


def line(slide, x1, y1, x2, y2, color, weight=Pt(0.5)):
    """Straight line connector. Coords are forced to int EMU — add_connector
    does not auto-cast, and float coords produce 'needs repair' warnings."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      _emu(x1), _emu(y1),
                                      _emu(x2), _emu(y2))
    conn.shadow.inherit = False
    conn.line.color.rgb = _rgb_tuple(color)
    conn.line.width = weight
    return conn


def hairline(slide, x, y, w, h, color):
    """Thin separator drawn as a no-border filled rectangle."""
    shp = rect(slide, x, y, w, h, fill=color)
    _apply_border(shp, None)
    return shp


def txt(slide, x, y, w, h, text, sz=14, color=None, bold=False, italic=False,
        align=PP_ALIGN.LEFT, font=FONT_BODY, anchor=MSO_ANCHOR.TOP):
    """Transparent textbox with one formatted run."""
    box = slide.shapes.add_textbox(_emu(x), _emu(y), _emu(w), _emu(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = "" if text is None else str(text)
    run.font.name = font
    run.font.size = Pt(sz)
    run.font.bold = bool(bold)
    run.font.italic = bool(italic)
    if color is not None:
        run.font.color.rgb = _rgb_tuple(color)
    return box


def _add_run(paragraph, text, sz=14, color=None, bold=False, italic=False,
             font=FONT_BODY):
    run = paragraph.add_run()
    run.text = "" if text is None else str(text)
    run.font.name = font
    run.font.size = Pt(sz)
    run.font.bold = bool(bold)
    run.font.italic = bool(italic)
    if color is not None:
        run.font.color.rgb = _rgb_tuple(color)
    return run


# ===========================================================================
# Freeform helper
# ===========================================================================

def _freeform_poly(slide, points, fill=None, border=None, border_w=None):
    """Closed freeform polygon from a list of (x, y) EMU points."""
    if not points:
        return None
    builder = slide.shapes.build_freeform(points[0][0], points[0][1])
    for px, py in points[1:]:
        builder.add_line_segments([(px, py)], close=False)
    builder.add_line_segments([(points[0][0], points[0][1])], close=True)
    shp = builder.convert_to_shape()
    shp.shadow.inherit = False
    _apply_fill(shp, fill)
    if border is not None:
        _apply_border(shp, border, border_w if border_w is not None else Pt(0.5))
    else:
        _apply_border(shp, None)
    return shp


