from __future__ import annotations
import math
import sys
from typing import Optional
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from ._ml_constants import (
    FONT_HEAD, FONT_BODY,
    MERCK_PURPLE, MERCK_GOLD, MERCK_YELLOW, MERCK_BLUE, MERCK_AQUA,
    WHITE, INK_DARK, INK_GRAY, LIGHT_GRAY, PANEL_LIGHT,
    PURPLE_DEEP, PURPLE_MUTED, BAD_RED, GOOD_GREEN,
    ACT_PURPLE, LY_CYAN, OP_LIME, FC_PINK, DEV_POS_BLUE, DEV_NEG_RED,
    CHART_PALETTE, PHASE_1_COLOR, PHASE_2_COLOR, PHASE_3_COLOR,
    SLIDE_W,
    CONTENT_X, CONTENT_Y, CONTENT_Y_SUBTITLE, CONTENT_W, CONTENT_H,
    SOURCE_Y, SOURCE_H, TAKEAWAY_Y, TAKEAWAY_H, PHASE_Y, PHASE_H,
    FOOTER_Y, FOOTER_H, FOOTER_TEXT_Y,
    TITLE_X, TITLE_Y_NUMBERED, TITLE_Y_UNNUMBERED, TITLE_W, TITLE_H,
    SUB_X, SUB_W, SUB_H, SUB_GAP,
    AUTO_PROMOTE_EXECUTIVE,
    _palette_for, _rgb_tuple, _is_dark,
)
from ._ml_primitives import (
    rect, rounded, oval, circle, line, hairline, txt, _add_run,
    _freeform_poly, _emu, _apply_fill, _apply_border,
)
from ._ml_icons import draw_icon, ICON_REGISTRY
from ._ml_deck import (
    _new_slide, _intro_layout, _divider_layout,
    _cover_picture_layout, _populate_placeholder,
)
from ._ml_chrome import (
    apply_chrome, _phase_progress, _draw_card, _bulleted_list,
    stub_and_flag, footnotes_block, _section_marker,
    _top_chrome, _bottom_chrome, _gold_square_bullet,
    _tracked, _track_letters, _format_section_number, _pad_int,
    _render_action_title, _source_line,
    statement_card, in_slide_section,
    _takeaway_band, _superscript, _callout_block, _chrome, _content_y,
)
from ._ml_charts import (
    add_slope_chart, add_dot_plot, add_marimekko, add_waterfall,
    add_small_multiples, add_simple_bar, _render_chart,
)
from ._ml_helpers import _style_or_promote, _tone_color, _rag_color, _norm_key, _named_color

_CIRCULAR_PRESETS = {
    "pdca": {
        "phases": [
            {"label": "PLAN",  "body": "Identify the problem and plan the change"},
            {"label": "DO",    "body": "Implement the change on a small scale"},
            {"label": "CHECK", "body": "Analyse results and identify learnings"},
            {"label": "ACT",   "body": "Scale up or adjust based on findings"},
        ],
    },
    "dmaic": {
        "phases": [
            {"label": "DEFINE",   "body": "Define the problem and goals"},
            {"label": "MEASURE",  "body": "Measure current process performance"},
            {"label": "ANALYSE",  "body": "Identify root causes"},
            {"label": "IMPROVE",  "body": "Implement and verify solutions"},
            {"label": "CONTROL",  "body": "Sustain the improvements"},
        ],
    },
}


# ===========================================================================
# Layout: PHASE PROCESS
# ===========================================================================

def build_phase_process(prs, meta, action_title=None, phases=None, takeaway=None,
                        source=None, category=None, style="merck_executive",
                        page=None, total=None, section_number=None,
                        highlight_index=None, show_arrows=True,
                        methodology_note=None, subtitle=None, content=None):
    """Horizontal phased flow as a series of cards with arrow connectors.

    Each phase dict: {label, title, body, milestone, highlighted}.
    highlight_index (0-based int): when set, only this index is highlighted
        and overrides per-phase 'highlighted' flags.
    show_arrows (bool): draw a MERCK_GOLD right-arrow in the gutter between
        adjacent cards (default True).

    All cards are rendered at the same height (equal-height enforcement) so
    they line up as a row of equals.
    """
    if content:
        if "phases"          in content: phases          = content["phases"]
        if "highlight_index" in content: highlight_index = content["highlight_index"]
        if "show_arrows"     in content: show_arrows     = content["show_arrows"]
        if "takeaway"        in content: takeaway        = content["takeaway"]
    style = _style_or_promote(category, style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 subtitle=subtitle,
                 takeaway=takeaway, source=source, page=page, total=total,
                 palette=style, section_number=section_number,
                 methodology_note=methodology_note)

    phs = (phases or [])[:5]
    if not phs:
        return slide
    n = len(phs)
    zone_top = _content_y(meta, subtitle=bool(subtitle))
    zone_h = Inches(6.50) - zone_top
    # Reserve a top strip for the "STEP a -> STEP b -> ..." breadcrumb.
    breadcrumb_h = Inches(0.32)
    cards_top = zone_top + breadcrumb_h + Inches(0.10)
    base_cards_h = zone_h - breadcrumb_h - Inches(0.10)
    arrow_w = Inches(0.36) if show_arrows else Inches(0.18)
    total_w = Inches(12.0)
    card_w = (total_w - arrow_w * (n - 1)) / n

    # Equal-height enforcement based on content estimates.
    estimated = []
    for ph in phs:
        # Per-card content height estimate based on title length, body, milestone.
        est_h = Inches(0.50)  # top label
        title_len = len(str(ph.get("title", "")))
        if title_len:
            title_lines = max(1, (title_len // 22) + (1 if title_len % 22 else 0))
            est_h += Inches(0.32) * title_lines
        body_len = len(str(ph.get("body", "")))
        if body_len:
            body_lines = max(1, (body_len // 32) + (1 if body_len % 32 else 0))
            est_h += Inches(0.22) * body_lines
        if ph.get("milestone"):
            est_h += Inches(0.50)
        est_h += Inches(0.40)  # bottom padding
        estimated.append(est_h)
    cards_h = max(base_cards_h, max(estimated))
    # Clamp so cards don't run past the source line.
    max_h_avail = Inches(6.40) - cards_top
    if cards_h > max_h_avail:
        cards_h = max_h_avail

    # Breadcrumb across the top: "STEP 1 -> STEP 2 -> ...".
    crumb_parts = []
    for i, ph in enumerate(phs):
        lbl = ph.get("label", f"STEP {i + 1}")
        crumb_parts.append(_tracked(lbl))
    crumb = "   →   ".join(crumb_parts)
    txt(slide, Inches(0.65), zone_top, total_w, breadcrumb_h,
        crumb, sz=10, color=pal["highlight"], bold=True, font=FONT_BODY,
        anchor=MSO_ANCHOR.MIDDLE)

    for i, ph in enumerate(phs):
        cx = Inches(0.65) + i * (card_w + arrow_w)
        cy = cards_top
        if "highlighted" in ph:
            highlighted = bool(ph["highlighted"])
        elif ph.get("status") == "current":
            highlighted = True
        elif highlight_index is not None:
            highlighted = (i == int(highlight_index))
        else:
            highlighted = False
        _draw_card(slide, cx, cy, card_w, cards_h, style,
                   highlighted=highlighted)

        pad = Inches(0.20)
        # Tiny step label inside card top.
        label_color = pal["hot"] if highlighted else pal["highlight"]
        sub_color = WHITE if highlighted else INK_GRAY
        title_color = WHITE if highlighted else MERCK_PURPLE
        txt(slide, cx + pad, cy + Inches(0.20), card_w - pad * 2,
            Inches(0.22),
            _tracked(ph.get("label", f"STEP {i + 1}")), sz=9,
            color=label_color, bold=True, font=FONT_BODY)

        title = ph.get("title", "")
        if title:
            txt(slide, cx + pad, cy + Inches(0.50), card_w - pad * 2,
                Inches(0.95),
                title, sz=15, color=title_color, bold=True, font=FONT_BODY)

        # Decide milestone treatment before drawing body so we can shrink the
        # body region when a long milestone caption needs vertical room.
        ms = ph.get("milestone")
        ms_text = str(ms) if ms else ""
        # Milestones ≤32 chars: single-line pill (height 0.32").
        # Milestones 33–60 chars: taller pill (height 0.48") with 9pt font to
        #   keep all milestone cards visually consistent (solid colored block).
        # Milestones >60 chars: italic caption fallback (can't fit in a pill).
        ms_len = len(ms_text)
        ms_is_long = bool(ms_text) and ms_len > 60
        ms_is_medium = bool(ms_text) and 32 < ms_len <= 60
        body_bottom_margin = Inches(1.00) if ms_is_long else Inches(0.65)

        body = ph.get("body", "")
        if body:
            body_h = cards_h - Inches(1.50) - body_bottom_margin
            if body_h < Inches(0.30):
                body_h = Inches(0.30)
            txt(slide, cx + pad, cy + Inches(1.50),
                card_w - pad * 2, body_h,
                body, sz=10, color=sub_color, font=FONT_BODY)

        # Milestone at the bottom. Short text (<=32 chars) renders as a
        # colored pill. Long text renders as an italic gold caption (no fill)
        # so it can wrap to 2-3 lines without overlapping body text — the
        # previous fixed 0.32" pill clipped long milestones into the body,
        # producing strikethrough-style render artifacts.
        if ms:
            block_w = card_w - pad * 2
            block_x = cx + pad
            if ms_is_long:
                cap_h = Inches(0.80)
                cap_y = cy + cards_h - cap_h - Inches(0.10)
                cap_color = pal["hot"] if highlighted else pal["highlight"]
                txt(slide, block_x, cap_y, block_w, cap_h,
                    ms_text, sz=9, color=cap_color, italic=True, bold=True,
                    font=FONT_BODY, align=PP_ALIGN.LEFT,
                    anchor=MSO_ANCHOR.TOP)
            else:
                # Medium milestones (33-60 chars) use a taller pill so all
                # milestone treatments remain a consistent solid colored block.
                pill_h  = Inches(0.48) if ms_is_medium else Inches(0.32)
                pill_sz = 9            if ms_is_medium else 10
                pill_y  = cy + cards_h - pill_h - Inches(0.14)
                if highlighted:
                    if _is_dark(style):
                        pill_fill = PURPLE_DEEP
                        pill_text = WHITE
                    else:
                        pill_fill = pal["hot"]
                        pill_text = PURPLE_DEEP
                else:
                    if _is_dark(style):
                        pill_fill = PURPLE_DEEP
                    else:
                        pill_fill = MERCK_PURPLE
                    pill_text = WHITE
                rect(slide, block_x, pill_y, block_w, pill_h, fill=pill_fill)
                txt(slide, block_x, pill_y, block_w, pill_h,
                    ms_text, sz=pill_sz, color=pill_text, bold=True,
                    font=FONT_BODY, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)

        # Arrow connector in the gutter to the next card.
        if i < n - 1 and show_arrows:
            ax = cx + card_w
            ay = cy + cards_h / 2
            ar_w = Inches(0.30)
            ar_h = Inches(0.40)
            ar_x = ax + (arrow_w - ar_w) / 2
            ar_y = ay - ar_h / 2
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                           ar_x, ar_y, ar_w, ar_h)
            arrow.shadow.inherit = False
            _apply_fill(arrow, pal["highlight"])
            _apply_border(arrow, None)
    return slide



# ===========================================================================
# Layout: GANTT
# ===========================================================================

def build_gantt(prs, meta, action_title=None, rows=None, quarters=None, takeaway=None,
                source=None, category=None, style="merck_executive",
                page=None, total=None, section_number=None,
                methodology_note=None, subtitle=None, content=None):
    if content:
        if "rows"     in content: rows     = content["rows"]
        if "quarters" in content: quarters = content["quarters"]
        if "takeaway" in content: takeaway = content["takeaway"]
    style = _style_or_promote(category, style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 subtitle=subtitle,
                 takeaway=takeaway, source=source, page=page, total=total,
                 palette=style, section_number=section_number, methodology_note=methodology_note)

    zone_top = _content_y(meta, subtitle=bool(subtitle))
    zone_h = Inches(6.50) - zone_top
    qs = list(quarters or ["Q1", "Q2", "Q3", "Q4"])
    qn_count = max(len(qs), 1)
    label_w = Inches(2.6)
    grid_x = Inches(0.65) + label_w
    grid_w = Inches(12.0) - label_w
    header_h = Inches(0.40)
    header_y = zone_top

    col_w = grid_w / qn_count
    for i, q in enumerate(qs):
        txt(slide, grid_x + i * col_w, header_y, col_w, header_h,
            _tracked(str(q)), sz=10, color=pal["highlight"], bold=True,
            font=FONT_BODY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        line(slide, grid_x + i * col_w, header_y + header_h,
             grid_x + i * col_w, zone_top + zone_h,
             LIGHT_GRAY, Pt(0.5))
    line(slide, grid_x + qn_count * col_w, header_y + header_h,
         grid_x + qn_count * col_w, zone_top + zone_h,
         LIGHT_GRAY, Pt(0.5))

    body_y = header_y + header_h
    body_h = zone_h - header_h
    rs = rows or []
    n = max(len(rs), 1)
    row_h = body_h / n
    for i, r in enumerate(rs):
        ry = body_y + i * row_h
        txt(slide, Inches(0.65), ry, label_w - Inches(0.10), row_h,
            r.get("label", ""), sz=11, color=INK_DARK, font=FONT_BODY,
            anchor=MSO_ANCHOR.MIDDLE)
        hairline(slide, grid_x, ry + row_h - Emu(int(Pt(0.5))),
                 grid_w, Emu(int(Pt(0.5))), LIGHT_GRAY)
        try:
            sq = float(r.get("start_q", 1))
        except (TypeError, ValueError):
            sq = 1.0
        try:
            du = float(r.get("duration_q", 1))
        except (TypeError, ValueError):
            du = 1.0
        bar_x = grid_x + (sq - 1) * col_w
        bar_w = du * col_w
        if bar_x < grid_x:
            bar_w = bar_w - (grid_x - bar_x)
            bar_x = grid_x
        if bar_x + bar_w > grid_x + grid_w:
            bar_w = grid_x + grid_w - bar_x
        bar_w = max(bar_w, Inches(0.10))
        raw_tone = r.get("tone", "")
        if raw_tone and raw_tone.lower() not in ("neutral", ""):
            tone_color = _tone_color(raw_tone, style)
        else:
            tone_color = pal["accent"]
        pad = Inches(0.10)
        rect(slide, bar_x, ry + pad, bar_w, row_h - pad * 2, fill=tone_color)
    return slide



# ===========================================================================
# Tier A: New layout functions
# ===========================================================================

def _draw_check_mark(slide, cx, cy, size, color):
    """Draw a check-mark inside a circle by composing two line segments."""
    half = size / 2
    arm = size * 0.30
    # Down-stroke of check
    s1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    _emu(cx + half - arm),
                                    _emu(cy + half),
                                    _emu(cx + half - arm * 0.15),
                                    _emu(cy + half + arm * 0.6))
    s1.line.color.rgb = _rgb_tuple(color)
    s1.line.width = Pt(2.25)
    # Up-stroke of check
    s2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    _emu(cx + half - arm * 0.15),
                                    _emu(cy + half + arm * 0.6),
                                    _emu(cx + half + arm * 0.95),
                                    _emu(cy + half - arm * 0.7))
    s2.line.color.rgb = _rgb_tuple(color)
    s2.line.width = Pt(2.25)


def build_milestone_timeline(prs, meta, action_title=None, milestones=None, takeaway=None,
                             source=None, methodology_note=None,
                             category=None, subtitle=None,
                             style="merck_executive", page=None, total=None,
                             section_number=None, content=None):
    """Horizontal milestone timeline.

    milestones = list of {"date","title","body","status"} where status is one
    of "done", "current", "future".
    """
    if content:
        if "milestones" in content: milestones = content["milestones"]
        if "takeaway"   in content: takeaway   = content["takeaway"]
    style = _style_or_promote(category, style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 subtitle=subtitle, takeaway=takeaway, source=source,
                 page=page, total=total, palette=style,
                 section_number=section_number,
                 methodology_note=methodology_note)

    items = list(milestones or [])[:6]
    n = max(len(items), 1)
    zone_x = CONTENT_X
    zone_w = CONTENT_W
    col_w = zone_w / n
    line_y = Inches(4.10)
    circle_d = Inches(0.42)
    date_y = Inches(3.25)
    title_y = Inches(3.55)
    body_y = Inches(4.75)

    # Pre-pass: figure out last done index for line coloring.
    statuses = [str(it.get("status", "future")).lower() for it in items]
    last_done = -1
    for i, s in enumerate(statuses):
        if s == "done":
            last_done = i

    # Draw connecting lines between adjacent circles.
    for i in range(n - 1):
        cx1 = zone_x + col_w * i + col_w / 2
        cx2 = zone_x + col_w * (i + 1) + col_w / 2
        s_left = statuses[i]
        s_right = statuses[i + 1]
        if s_left == "done" and (s_right in ("done", "current")):
            line_color = pal["accent"]
        else:
            line_color = LIGHT_GRAY
        rect(slide, cx1 + circle_d / 2, line_y + circle_d / 2 - Emu(int(Pt(1.0))),
             cx2 - cx1 - circle_d, Emu(int(Pt(2.0))),
             fill=line_color)

    # Draw each milestone column.
    for i, m in enumerate(items):
        col_cx = zone_x + col_w * i + col_w / 2
        st = statuses[i]

        # Date label.
        txt(slide, zone_x + col_w * i, date_y, col_w, Inches(0.26),
            _track_letters(str(m.get("date", "")).upper()),
            sz=10, color=pal["highlight"], bold=True, font=FONT_BODY,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Title — use white on dark-background styles so text is readable.
        dark = _is_dark(style)
        title_color = pal["highlight"] if st == "current" else (WHITE if dark else MERCK_PURPLE)
        txt(slide, zone_x + col_w * i, title_y, col_w, Inches(0.40),
            str(m.get("title", "")), sz=13, color=title_color, bold=True,
            font=FONT_BODY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Circle.
        circle_x = col_cx - circle_d / 2
        circle_y = line_y
        if st == "done":
            circle(slide, circle_x, circle_y, circle_d, fill=pal["accent"])
            _draw_check_mark(slide, circle_x, circle_y, circle_d, WHITE)
        elif st == "current":
            circle(slide, circle_x, circle_y, circle_d, fill=pal["highlight"])
        else:
            outer = circle(slide, circle_x, circle_y, circle_d, fill=pal["bg"])
            _apply_border(outer, LIGHT_GRAY, Pt(1.25))

        # Body text.
        txt(slide, zone_x + col_w * i + Inches(0.10), body_y,
            col_w - Inches(0.20), Inches(1.50),
            str(m.get("body", "")), sz=10, color=INK_GRAY, font=FONT_BODY,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    return slide



def build_circular_flow(prs, meta, action_title, phases, takeaway="", source=None,
                        subtitle=None, methodology_note=None,
                        style="merck_corporate", page=None, total=None,
                        section_number=None, category=None):
    """Circular process flow with 2-8 phases arranged around a centre point.

    phases: list of dicts with keys: label (str), body (str), icon (optional str icon name).
    Returns the slide.
    """
    import math

    style = _style_or_promote(category, style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 subtitle=subtitle, takeaway=takeaway, source=source,
                 page=page, total=total, palette=style,
                 section_number=section_number,
                 methodology_note=methodology_note)

    n = max(2, min(8, len(phases)))
    cx = SLIDE_W / 2
    cy = Inches(4.0)
    radius = Inches(2.0)
    circle_size = Inches(0.9)
    half = circle_size / 2

    # Compute centre positions for each phase
    positions = []
    for i in range(n):
        angle = (2 * math.pi * i / n) - math.pi / 2
        px = cx + radius * math.cos(angle)
        py = cy + radius * math.sin(angle)
        positions.append((px, py))

    # Draw connector lines between adjacent phase circles
    for i in range(n):
        x1, y1 = positions[i]
        x2, y2 = positions[(i + 1) % n]
        # Shorten line endpoints to circle edge
        dx = x2 - x1
        dy = y2 - y1
        dist = math.hypot(dx, dy)
        if dist > 0:
            ux, uy = dx / dist, dy / dist
            edge = half + Inches(0.05)
            lx1 = x1 + ux * edge
            ly1 = y1 + uy * edge
            lx2 = x2 - ux * edge
            ly2 = y2 - uy * edge
            line(slide, lx1, ly1, lx2, ly2,
                 color=pal.get("accent", MERCK_PURPLE), weight=Pt(1.5))

    # Draw phase circles and labels
    body_box_h = Inches(0.55)
    for i, phase in enumerate(phases[:n]):
        px, py = positions[i]
        cx_shape = px - half
        cy_shape = py - half

        # Filled circle
        circle(slide, cx_shape, cy_shape, circle_size,
               fill=pal.get("accent", MERCK_PURPLE))

        # Label inside circle (number + short label)
        label_text = str(phase.get("label", str(i + 1)))
        txt(slide, cx_shape, cy_shape, circle_size, circle_size,
            label_text, sz=10, color=WHITE, bold=True,
            align=PP_ALIGN.CENTER, font=FONT_BODY,
            anchor=MSO_ANCHOR.MIDDLE)

        # Body text below circle
        body_text = str(phase.get("body", ""))
        if body_text:
            txt(slide, cx_shape - Inches(0.3), py + half + Inches(0.06),
                circle_size + Inches(0.6), body_box_h,
                body_text, sz=11, color=pal.get("ink", INK_DARK),
                align=PP_ALIGN.CENTER, font=FONT_BODY,
                anchor=MSO_ANCHOR.TOP)

    return slide



def build_arrow_chain(prs, meta, action_title, steps, consequence=None,
                      takeaway="", source=None, subtitle=None,
                      methodology_note=None, style="merck_corporate",
                      page=None, total=None, section_number=None, category=None,
                      content=None):
    """Horizontal arrow chain: 3-5 steps + optional consequence box.

    steps: list of dicts with keys:
        label (str) — uppercase short label
        body (str) — 1-2 sentence description
        highlighted (bool, optional) — if True, uses MERCK_PURPLE fill + white text
    consequence: dict with label and body (optional right-end conclusion box,
                 MERCK_YELLOW fill)
    Returns the slide.
    """
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    steps = list(steps)[:5]
    n_steps = max(len(steps), 1)
    n_arrows = n_steps - 1
    has_consequence = consequence is not None

    # Layout geometry
    zone_x = Inches(0.55)
    zone_y = Inches(2.3)
    zone_w = Inches(12.2)
    box_h = Inches(2.4)

    arrow_w = Inches(0.35)
    # consequence box is 1.0 * box_w wide; preceded by a slightly larger arrow
    consequence_w_ratio = 1.1
    consequence_arrow_w = Inches(0.42)

    # Total "units" of width consumed:
    #   n_steps boxes  +  n_arrows regular arrows
    #   + (consequence_arrow_w + consequence_w_ratio * box_w) if consequence
    # Solve for box_w:
    #   zone_w = n_steps*box_w + n_arrows*arrow_w
    #            [+ consequence_arrow_w + consequence_w_ratio*box_w]
    if has_consequence:
        # zone_w = (n_steps + consequence_w_ratio)*box_w
        #        + n_arrows*arrow_w + consequence_arrow_w
        box_w = (zone_w - n_arrows * arrow_w - consequence_arrow_w) / (
            n_steps + consequence_w_ratio)
        consequence_box_w = box_w * consequence_w_ratio
    else:
        box_w = (zone_w - n_arrows * arrow_w) / n_steps
        consequence_box_w = 0.0

    # Helper: draw a right-pointing filled triangle (arrow connector)
    def _draw_triangle(ax, ay, aw, ah, fill_color):
        """Draw filled right-pointing triangle at (ax, ay) with size (aw x ah)."""
        tip_x = ax + aw
        mid_y = ay + ah / 2
        pts = [
            (int(ax), int(ay)),
            (int(ax), int(ay + ah)),
            (int(tip_x), int(mid_y)),
        ]
        _freeform_poly(slide, pts, fill=fill_color)

    cursor_x = zone_x
    arrow_h = Inches(0.55)  # triangle height (vertical span)
    arrow_y = zone_y + (box_h - arrow_h) / 2  # vertically centred

    for i, step in enumerate(steps):
        highlighted = bool(step.get("highlighted", False))

        # Box fill and text colors
        box_fill = MERCK_PURPLE if highlighted else PANEL_LIGHT
        label_color = WHITE if highlighted else pal["highlight"]
        body_color = WHITE if highlighted else INK_DARK

        # Draw box
        rounded(slide, cursor_x, zone_y, box_w, box_h, fill=box_fill)

        # Label — uppercase, tracked at sz=9
        label_text = str(step.get("label", "")).upper()
        label_pad_x = Inches(0.12)
        label_pad_y = Inches(0.14)
        label_h = Inches(0.25)
        if label_text:
            txt(slide,
                cursor_x + label_pad_x,
                zone_y + label_pad_y,
                box_w - label_pad_x * 2,
                label_h,
                label_text,
                sz=9, color=label_color, bold=False,
                align=PP_ALIGN.LEFT, font=FONT_BODY,
                anchor=MSO_ANCHOR.TOP)

        # Body text
        body_top = zone_y + label_pad_y + label_h + Inches(0.08)
        body_h = box_h - label_pad_y - label_h - Inches(0.08) - Inches(0.12)
        txt(slide,
            cursor_x + label_pad_x,
            body_top,
            box_w - label_pad_x * 2,
            body_h,
            str(step.get("body", "")),
            sz=11, color=body_color, bold=False,
            align=PP_ALIGN.LEFT, font=FONT_BODY,
            anchor=MSO_ANCHOR.TOP)

        cursor_x += box_w

        # Arrow between boxes (not after last step unless consequence follows)
        if i < n_steps - 1:
            _draw_triangle(cursor_x, arrow_y, arrow_w, arrow_h, MERCK_PURPLE)
            cursor_x += arrow_w
        elif has_consequence:
            # Larger arrow before consequence box
            _draw_triangle(cursor_x, arrow_y - Inches(0.04),
                           consequence_arrow_w, arrow_h + Inches(0.08),
                           MERCK_PURPLE)
            cursor_x += consequence_arrow_w

    # Consequence box
    if has_consequence:
        con_label = str(consequence.get("label", "")).upper()
        con_body = str(consequence.get("body", ""))

        rounded(slide, cursor_x, zone_y, consequence_box_w, box_h,
             fill=pal["hot"])

        label_pad_x = Inches(0.12)
        label_pad_y = Inches(0.14)
        label_h = Inches(0.25)
        if con_label:
            txt(slide,
                cursor_x + label_pad_x,
                zone_y + label_pad_y,
                consequence_box_w - label_pad_x * 2,
                label_h,
                con_label,
                sz=9, color=PURPLE_DEEP, bold=True,
                align=PP_ALIGN.LEFT, font=FONT_BODY,
                anchor=MSO_ANCHOR.TOP)

        body_top = zone_y + label_pad_y + label_h + Inches(0.08)
        body_h = box_h - label_pad_y - label_h - Inches(0.08) - Inches(0.12)
        txt(slide,
            cursor_x + label_pad_x,
            body_top,
            consequence_box_w - label_pad_x * 2,
            body_h,
            con_body,
            sz=11, color=PURPLE_DEEP, bold=True,
            align=PP_ALIGN.LEFT, font=FONT_BODY,
            anchor=MSO_ANCHOR.TOP)

    callout = (content or {}).get("callout")
    if isinstance(callout, dict):
        _callout_block(slide, callout.get("type", "conclusion"),
                       callout.get("text", ""), pal,
                       has_takeaway=bool(takeaway))
    return slide



# ===========================================================================
# Layout: FUNNEL / CONVERGENCE
# ===========================================================================

def build_funnel(prs, meta, action_title=None, inputs=None,
                 output=None, takeaway="", source=None, subtitle=None,
                 methodology_note=None, style="merck_corporate",
                 page=None, total=None, section_number=None,
                 category=None):
    """2-4 input boxes converging via diagonal lines to a single output.

    inputs: list of dicts with label (str) and body (str)
    output: dict with label (str) and body (str)
    """
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    # Cap at 4 inputs — the funnel geometry is designed for 2-4 boxes.
    inputs = list(inputs or [])[:4]
    if not inputs:
        return slide

    zone_x, zone_y = Inches(0.55), Inches(1.62)
    zone_w, zone_h = Inches(12.2), Inches(4.72)
    in_w  = Inches(3.40)
    out_w = Inches(3.80)
    box_h = Inches(0.95)
    gap   = Inches(0.22)
    n     = len(inputs)

    total_in_h = n * box_h + (n - 1) * gap
    in_start_y = zone_y + (zone_h - total_in_h) / 2
    # Default progression: muted purple → teal → green → yellow, visually distinct per step.
    IN_COLORS  = [PURPLE_MUTED, LY_CYAN, GOOD_GREEN, MERCK_YELLOW]

    for i, inp in enumerate(inputs):
        iy  = in_start_y + i * (box_h + gap)
        # Per-input color: inp["color"] name overrides the default progression.
        col = _named_color(inp.get("color"), IN_COLORS[i % len(IN_COLORS)])
        rounded(slide, zone_x, iy, in_w, box_h, fill=col)
        # Ensure label text is readable: switch to dark text on yellow-fill boxes.
        label_text_color = MERCK_PURPLE if col == MERCK_YELLOW else MERCK_YELLOW
        txt(slide, zone_x + Inches(0.14), iy + Inches(0.08),
            in_w - Inches(0.28), Inches(0.24),
            str(inp.get("label", "")).upper(),
            sz=9, color=label_text_color, bold=True, font=FONT_BODY)
        txt(slide, zone_x + Inches(0.14), iy + Inches(0.34),
            in_w - Inches(0.28), box_h - Inches(0.42),
            str(inp.get("body", "")),
            sz=10, color=WHITE, font=FONT_BODY, anchor=MSO_ANCHOR.TOP)

    out   = output or {}
    out_x = zone_x + zone_w - out_w
    out_y = zone_y + (zone_h - Inches(1.80)) / 2
    out_cy = out_y + Inches(0.90)

    for i in range(n):
        iy_c = in_start_y + i * (box_h + gap) + box_h / 2
        line(slide, zone_x + in_w, iy_c, out_x, out_cy,
             color=LIGHT_GRAY, weight=Pt(1.5))

    rounded(slide, out_x, out_y, out_w, Inches(1.80), fill=pal["highlight"])
    txt(slide, out_x + Inches(0.16), out_y + Inches(0.12),
        out_w - Inches(0.32), Inches(0.28),
        str(out.get("label", "")).upper(),
        sz=9, color=PURPLE_DEEP, bold=True, font=FONT_BODY)
    txt(slide, out_x + Inches(0.16), out_y + Inches(0.44),
        out_w - Inches(0.32), Inches(1.20),
        str(out.get("body", "")),
        sz=13, color=PURPLE_DEEP, bold=True, font=FONT_BODY,
        anchor=MSO_ANCHOR.TOP)

    return slide



# ===========================================================================
# Layout: JOURNEY MAP (multi-row swim-lane)
# ===========================================================================

def build_journey_map(prs, meta, action_title=None, phases=None, rows=None,
                      takeaway=None, source=None, category=None, subtitle=None,
                      style="merck_executive", page=None, total=None,
                      section_number=None, methodology_note=None, content=None):
    """Customer journey map: phases as columns, swim-lane rows as labeled tracks.

    phases: list of str - column headers (e.g. ["Aware","Consider","Buy","Use"]).
    rows: list of dicts {label, cells} - cells is a list of str, one per phase.
    """
    if content:
        if "phases"   in content: phases   = content["phases"]
        if "rows"     in content: rows     = content["rows"]
        if "takeaway" in content: takeaway = content["takeaway"]
    style = _style_or_promote(category, style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category, subtitle=subtitle,
                 takeaway=takeaway, source=source, page=page, total=total,
                 palette=style, section_number=section_number,
                 methodology_note=methodology_note)

    dark = _is_dark(style)
    text_color = WHITE if dark else INK_DARK

    phases    = list(phases or [])
    rows      = list(rows   or [])
    n_phases  = max(len(phases), 1)
    n_rows    = max(len(rows),   1)

    zone_top = _content_y(meta, subtitle=bool(subtitle))
    zone_h   = Inches(6.40) - zone_top
    label_w  = Inches(1.60)
    grid_x   = Inches(0.65) + label_w
    grid_w   = Inches(12.33) - grid_x
    col_w    = grid_w / n_phases
    hdr_h    = Inches(0.44)
    row_h    = (zone_h - hdr_h) / n_rows

    # Phase header row
    for j, ph in enumerate(phases):
        cx         = grid_x + j * col_w
        hdr_fill   = MERCK_PURPLE if j % 2 == 0 else PURPLE_DEEP
        rect(slide, cx, zone_top, col_w, hdr_h, fill=hdr_fill)
        txt(slide, cx, zone_top, col_w, hdr_h, str(ph),
            sz=10, color=WHITE, bold=True, font=FONT_BODY,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    row_label_colors = [MERCK_PURPLE, LY_CYAN, MERCK_GOLD, OP_LIME, PURPLE_MUTED, GOOD_GREEN]

    for i, row in enumerate(rows):
        ry = zone_top + hdr_h + i * row_h
        lc = row_label_colors[i % len(row_label_colors)]
        # Row label
        rounded(slide, Inches(0.65), ry, label_w, row_h, fill=lc)
        txt(slide, Inches(0.65), ry, label_w, row_h,
            str(row.get("label", "")),
            sz=9, color=WHITE, bold=True, font=FONT_BODY,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Cells
        cells = list(row.get("cells") or [])
        for j in range(n_phases):
            cx        = grid_x + j * col_w
            cell_fill = PANEL_LIGHT if (i + j) % 2 == 0 else WHITE
            cell      = rounded(slide, cx, ry, col_w, row_h, fill=cell_fill)
            _apply_border(cell, LIGHT_GRAY, Pt(0.5))
            cell_text = cells[j] if j < len(cells) else ""
            if cell_text:
                txt(slide, cx + Inches(0.10), ry + Inches(0.08),
                    col_w - Inches(0.20), row_h - Inches(0.16),
                    str(cell_text), sz=9, color=text_color,
                    font=FONT_BODY, anchor=MSO_ANCHOR.TOP)

    return slide


# ===========================================================================
# Layout: ROAD TO SUCCESS
# ===========================================================================

def build_road_to_success(prs, meta, action_title=None, stages=None,
                          milestones=None, takeaway="", source=None,
                          subtitle=None, methodology_note=None,
                          style="merck_executive", page=None, total=None,
                          section_number=None, category=None, content=None):
    """Horizontal timeline path with milestone dots and stage columns below.

    content keys:
        stages     (list of {title, body}) — 2-4 stage columns below the path
        milestones (list of str)           — 2-6 labels along the path (alternating above/below)
    """
    if content:
        stages     = stages     or content.get("stages")
        milestones = milestones or content.get("milestones")
    style = _style_or_promote(category, style)
    pal   = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    PATH_Y  = Inches(3.60)
    PATH_X0 = CONTENT_X
    PATH_W  = CONTENT_W
    ACC     = pal["accent"]
    INK     = pal["ink"]
    INK2    = pal["ink_2"]

    # Path: thick accent-coloured hairline spanning slide width
    hairline(slide, PATH_X0, PATH_Y, PATH_W, Emu(int(Pt(3))), ACC)

    # Milestone dots along the path
    mss = list(milestones or [])[:6]
    if mss:
        dot_d = Inches(0.15)
        n_ms  = len(mss)
        for i, ms in enumerate(mss):
            frac  = i / max(n_ms - 1, 1)
            dot_x = PATH_X0 + (PATH_W - dot_d) * frac
            circle(slide, dot_x, PATH_Y - dot_d / 2, dot_d, fill=ACC)
            if ms:
                lbl_y = (PATH_Y - Inches(0.32) if i % 2 == 0
                         else PATH_Y + dot_d + Inches(0.04))
                txt(slide, dot_x - Inches(0.38), lbl_y,
                    Inches(0.90), Inches(0.22),
                    str(ms), sz=9, color=INK2, font=FONT_BODY,
                    align=PP_ALIGN.CENTER)

    # Stage columns below path
    stgs = list(stages or [])[:4]
    n    = len(stgs)
    if n:
        col_w     = PATH_W / n
        col_y_top = PATH_Y + Inches(0.48)
        col_h     = CONTENT_Y + CONTENT_H - col_y_top
        for i, stage in enumerate(stgs):
            cx = PATH_X0 + i * col_w
            # Colour accent bar at top of column
            hairline(slide, cx, col_y_top,
                     col_w - Inches(0.15), Emu(int(Pt(2))), ACC)
            # Title
            txt(slide, cx, col_y_top + Inches(0.10),
                col_w - Inches(0.15), Inches(0.38),
                str(stage.get("title", "")), sz=13, color=ACC,
                bold=True, font=FONT_BODY)
            # Body
            txt(slide, cx, col_y_top + Inches(0.52),
                col_w - Inches(0.15), col_h - Inches(0.52),
                str(stage.get("body", "")), sz=11, color=INK,
                font=FONT_BODY, anchor=MSO_ANCHOR.TOP)
    return slide
