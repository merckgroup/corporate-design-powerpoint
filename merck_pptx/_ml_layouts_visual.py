from __future__ import annotations
import math
import os
import sys
from typing import Optional
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from ._ml_constants import (
    FONT_HEAD, FONT_BODY, SLIDE_W,
    MERCK_PURPLE, MERCK_GOLD, MERCK_YELLOW, MERCK_BLUE, MERCK_AQUA,
    WHITE, INK_DARK, INK_GRAY, LIGHT_GRAY, PANEL_LIGHT,
    PURPLE_DEEP, PURPLE_MUTED, BAD_RED, GOOD_GREEN,
    ACT_PURPLE, LY_CYAN, OP_LIME, FC_PINK, DEV_POS_BLUE, DEV_NEG_RED,
    CHART_PALETTE, PHASE_1_COLOR, PHASE_2_COLOR, PHASE_3_COLOR,
    CONTENT_X, CONTENT_Y, CONTENT_Y_SUBTITLE, CONTENT_W, CONTENT_H,
    SOURCE_Y, SOURCE_H, TAKEAWAY_Y, TAKEAWAY_H, PHASE_Y, PHASE_H,
    FOOTER_Y, FOOTER_H, FOOTER_TEXT_Y,
    TITLE_X, TITLE_Y_NUMBERED, TITLE_Y_UNNUMBERED, TITLE_W, TITLE_H,
    SUB_X, SUB_W, SUB_H, SUB_GAP,
    AUTO_PROMOTE_EXECUTIVE,
    _palette_for, _rgb_tuple, _is_dark,
)
from ._ml_primitives import (
    rect, rounded, oval, circle, line, hairline, txt, _add_run,
    _freeform_poly, _emu, _apply_fill, _apply_border,
)
from ._ml_icons import draw_icon, ICON_REGISTRY
from ._ml_deck import (
    _new_slide, _intro_layout, _divider_layout,
    _cover_picture_layout, _populate_placeholder,
)
from ._ml_chrome import (
    apply_chrome, _phase_progress, _draw_card, _bulleted_list,
    stub_and_flag, footnotes_block, _section_marker,
    _top_chrome, _bottom_chrome, _gold_square_bullet,
    _tracked, _track_letters, _format_section_number, _pad_int,
    _render_action_title, _source_line,
    statement_card, in_slide_section,
    _takeaway_band, _superscript,
)
from ._ml_charts import (
    add_slope_chart, add_dot_plot, add_marimekko, add_waterfall,
    add_small_multiples, add_simple_bar, _render_chart,
)
from ._ml_helpers import _style_or_promote, _tone_color, _rag_color, _norm_key

# ===========================================================================
# Layout: PULL QUOTE / STATEMENT
# ===========================================================================

def build_pull_quote(prs, meta, quote=None, attribution=None,
                     context=None, takeaway="", source=None, subtitle=None,
                     methodology_note=None, style="merck_storytelling",
                     page=None, total=None, section_number=None,
                     category=None, action_title=None):
    """Full-bleed pull quote slide.

    quote:       str — the quote text (no quotation marks needed)
    attribution: str — source / speaker (shown below quote with em-dash)
    context:     str — optional small descriptor tag above the quote
    """
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    dark = _is_dark(style)
    q_color    = MERCK_YELLOW if dark else MERCK_PURPLE
    attr_color = PANEL_LIGHT  if dark else INK_GRAY
    ctx_color  = PANEL_LIGHT  if dark else PURPLE_MUTED

    # Shift quote content below the action-title chrome.  A single-line title
    # clears at ~1.50; a two-line title can reach ~1.95.  Use the subtitle-
    # aware offset so content never overlaps the action title text.
    q_top = Inches(2.10) if not subtitle else Inches(2.50)

    txt(slide, Inches(0.55), q_top, Inches(1.20), Inches(1.40),
        "\u201c", sz=96, color=pal["highlight"], bold=True,
        font=FONT_HEAD, align=PP_ALIGN.LEFT)

    if context:
        txt(slide, Inches(1.40), q_top + Inches(0.15), Inches(10.5), Inches(0.30),
            str(context).upper(), sz=9, color=ctx_color, bold=True,
            font=FONT_BODY, align=PP_ALIGN.LEFT)

    txt(slide, Inches(1.40), q_top + Inches(0.50), Inches(10.5), Inches(2.80),
        str(quote or ''), sz=32, color=q_color, bold=True, italic=True,
        font=FONT_HEAD, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

    attr_y = q_top + Inches(3.55)
    hairline(slide, Inches(1.40), attr_y, Inches(2.80),
             Emu(int(Pt(2))), pal["highlight"])

    if attribution:
        txt(slide, Inches(1.40), attr_y + Inches(0.14), Inches(10.5), Inches(0.42),
            f'\u2014 {attribution}', sz=13, color=attr_color, italic=True,
            font=FONT_BODY, align=PP_ALIGN.LEFT)

    return slide



# ===========================================================================
# Layout: WORD CLOUD
# ===========================================================================

def build_word_cloud(prs, meta, action_title=None, words=None,
                     takeaway="", source=None, subtitle=None,
                     methodology_note=None, style="merck_executive",
                     page=None, total=None, section_number=None,
                     category=None):
    """Keyword cluster with size and color mapped to importance.

    words: list of dicts with keys:
        text   (str) — keyword or phrase
        weight (float 1-5) — importance (15pt-35pt)
        color  (tuple, opt) — explicit color override
    """
    pal = _palette_for(style)
    ACC  = pal["accent"]
    HOT  = pal["hot"]
    HLT  = pal["highlight"]
    INK  = pal["ink"]
    INK2 = pal["ink_2"]
    dark = _is_dark(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    words = list(words or [])[:30]
    if not words:
        return slide

    WORD_COLORS = [ACC, HLT, HOT, pal["accent_2"], pal["accent_3"], pal["muted"]]

    # Dynamic grid: fill the full content zone (not hard-coded y=1.72)
    _zone_top = Inches(2.95) if subtitle else Inches(2.55)
    _zone_bot = SOURCE_Y - Inches(0.10)
    _zone_h   = _zone_bot - _zone_top
    _GRID_ROWS, _GRID_COLS = 5, 6
    _row_h   = _zone_h / _GRID_ROWS
    _col_w   = CONTENT_W / _GRID_COLS
    GRID = [
        (CONTENT_X + c * _col_w + _col_w * 0.08 + (r % 2) * _col_w * 0.06,
         _zone_top  + r * _row_h + _row_h  * 0.12)
        for r in range(_GRID_ROWS)
        for c in range(_GRID_COLS)
    ]

    for i, word in enumerate(words):
        if i >= len(GRID):
            break
        gx, gy = GRID[i]
        # Accept plain strings or dicts
        if isinstance(word, str):
            word_text = word
            weight = 2.0
            col = WORD_COLORS[i % len(WORD_COLORS)]
        else:
            weight  = max(1.0, min(5.0, float(word.get("weight") or 2)))
            col     = word.get("color") or WORD_COLORS[i % len(WORD_COLORS)]
            word_text = str(word.get("text", ""))
        sz = int(10 + weight * 5)
        txt(slide, gx, gy, Inches(2.20), Inches(0.46),
            word_text,
            sz=sz, color=col, bold=(weight >= 4.0),
            font=FONT_HEAD if weight >= 4 else FONT_BODY,
            align=PP_ALIGN.LEFT)

    return slide


# ===========================================================================
# Layout: PYRAMID
# ===========================================================================

def build_pyramid(prs, meta, action_title=None, tiers=None,
                  orientation="up", takeaway="", source=None, subtitle=None,
                  methodology_note=None, style="merck_executive",
                  page=None, total=None, section_number=None,
                  category=None):
    """3-5 stacked tiers forming a pyramid.

    tiers: list of dicts ordered bottom-to-top with keys:
        label (str) — tier label
        body  (str, opt) — descriptor shown to the right
        color (tuple, opt) — fill override
    orientation: 'up' (default, wide base) or 'down' (inverted)
    """
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    tiers = list(tiers or [])[:5]
    if not tiers:
        return slide

    dark    = _is_dark(style)
    zone_cx = Inches(5.20)
    zone_y  = Inches(1.70)
    zone_h  = Inches(4.60)
    n       = len(tiers)
    tier_h  = zone_h / n
    max_w   = Inches(6.50)
    taper   = max_w / (n + 1)

    TIER_COLORS = [pal["accent"], pal["accent_2"], pal["accent_3"], pal["muted"], pal["rule"]]

    for i, tier in enumerate(tiers):
        level = i if orientation == "up" else (n - 1 - i)
        width = max_w - level * taper
        ty    = zone_y + i * tier_h
        tx    = zone_cx - width / 2
        col   = tier.get("color") or TIER_COLORS[i % len(TIER_COLORS)]

        rect(slide, tx, ty, width, tier_h - Emu(18000), fill=col)
        txt(slide, tx + Inches(0.12), ty,
            width - Inches(0.24), tier_h - Emu(18000),
            str(tier.get("label", "")),
            sz=12, color=WHITE, bold=True, font=FONT_BODY,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        if tier.get("body"):
            body_x = zone_cx + max_w / 2 + Inches(0.22)
            line(slide, tx + width + Inches(0.10), ty + tier_h / 2,
                 body_x, ty + tier_h / 2, color=LIGHT_GRAY, weight=Pt(0.75))
            txt(slide, body_x, ty + Inches(0.05),
                Inches(4.00), tier_h - Inches(0.10),
                str(tier["body"]),
                sz=10, color=INK_GRAY if not dark else PANEL_LIGHT,
                font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE)

    return slide


# ===========================================================================
# Layout: VENN DIAGRAM
# ===========================================================================

def build_venn(prs, meta, action_title=None, circles=None,
               intersection=None, takeaway="", source=None, subtitle=None,
               methodology_note=None, style="merck_executive",
               page=None, total=None, section_number=None,
               category=None):
    """2 or 3 overlapping translucent circles.

    circles: list of 2-3 dicts with keys:
        label (str) — circle title shown above
        body  (str, opt) — text inside
        color (tuple, opt) — fill color
    intersection: str — text in the overlap center
    """
    import math as _math
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    circles = list(circles or [])[:3]
    if not circles:
        return slide

    dark = _is_dark(style)
    CIRC_COLORS = [pal["accent"], pal["accent_2"], pal["accent_3"], pal["highlight"]]
    r = Inches(2.10)

    n = len(circles)
    if n == 2:
        overlap = Inches(0.80)
        cx_list = [Inches(4.20), Inches(4.20) + r * 2 - overlap]
        cy_list = [Inches(3.80), Inches(3.80)]
    else:
        base_cx, base_cy = Inches(6.10), Inches(3.90)
        tri_r = r * 0.82
        cx_list = [base_cx + tri_r * _math.cos(_math.radians(a))
                   for a in (210, 330, 90)]
        cy_list = [base_cy + tri_r * _math.sin(_math.radians(a))
                   for a in (210, 330, 90)]

    label_color = WHITE if dark else INK_DARK
    group_cx = sum(cx_list[:n]) / n
    group_cy = sum(cy_list[:n]) / n

    for i, circ in enumerate(circles):
        col = circ.get("color") or CIRC_COLORS[i % len(CIRC_COLORS)]
        ccx, ccy = cx_list[i], cy_list[i]
        shp = oval(slide, ccx - r, ccy - r, r * 2, r * 2, fill=col)

        try:
            sp_pr = shp._element.find(qn("p:spPr"))
            if sp_pr is None:
                sp_pr = shp._element.spPr
            sf = sp_pr.find(f".//{qn('a:solidFill')}")
            if sf is not None:
                sc = sf.find(qn("a:srgbClr"))
                if sc is not None:
                    from lxml import etree as _et
                    al = _et.SubElement(sc, qn("a:alpha"))
                    al.set("val", "65000")
        except Exception:
            pass

        # Radial label: push outward from the group centroid so the top-circle
        # label never flies off the slide (RCA-10).
        import math as _m
        dx = ccx - group_cx
        dy = ccy - group_cy
        dist = _m.hypot(dx, dy)
        if dist > 0:
            ux, uy = dx / dist, dy / dist
        else:
            ux, uy = 0.0, -1.0
        lbl_cx = ccx + ux * (r + Inches(0.12))
        lbl_cy = ccy + uy * (r + Inches(0.04)) - Inches(0.15)
        # clamp to slide safe zone
        lbl_cx = max(CONTENT_X, min(lbl_cx, CONTENT_X + CONTENT_W - r * 2))
        lbl_cy = max(Inches(1.20), min(lbl_cy, SOURCE_Y - Inches(0.32)))
        txt(slide, lbl_cx - r, lbl_cy, r * 2, Inches(0.30),
            str(circ.get("label", "")),
            sz=10, color=label_color, bold=True, font=FONT_BODY,
            align=PP_ALIGN.CENTER)

        if circ.get("body"):
            txt(slide, ccx - r * 0.7, ccy - Inches(0.26),
                r * 1.4, Inches(0.60), str(circ["body"]),
                sz=9, color=WHITE, font=FONT_BODY,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    if intersection:
        int_x = sum(cx_list) / n
        int_y = sum(cy_list) / n
        txt(slide, int_x - Inches(0.85), int_y - Inches(0.24),
            Inches(1.70), Inches(0.48), str(intersection),
            sz=10, color=WHITE, bold=True, font=FONT_BODY,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    return slide



# ===========================================================================
# Layout: LAYERED STACK
# ===========================================================================

def build_layered_stack(prs, meta, action_title=None, layers=None,
                        orientation="vertical", takeaway="", source=None,
                        subtitle=None, methodology_note=None,
                        style="merck_executive", page=None, total=None,
                        section_number=None, category=None):
    """Stacked layer architecture diagram (tech stack, org tiers, platform).

    layers: list of dicts (top-to-bottom for vertical) with keys:
        label (str) — layer name
        body  (str, opt) — short description
        icon  (str, opt) — emoji
        color (tuple, opt) — fill override
    orientation: 'vertical' (default) or 'horizontal'
    """
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    layers = list(layers or [])[:7]
    if not layers:
        return slide

    L_COLORS = [pal["accent"], pal["accent_2"], pal["accent_3"],
                 pal["highlight"], pal["muted"], pal["hot"], pal["rule"]]
    zone_x, zone_y = Inches(0.55), Inches(1.65)
    zone_w, zone_h = Inches(12.2), Inches(4.62)
    gap = Inches(0.10)
    n   = len(layers)

    if orientation == "horizontal":
        layer_w = (zone_w - gap * (n - 1)) / n
        for i, layer in enumerate(layers):
            lx  = zone_x + i * (layer_w + gap)
            col = layer.get("color") or L_COLORS[i % len(L_COLORS)]
            rounded(slide, lx, zone_y, layer_w, zone_h, fill=col)
            y_c = zone_y + Inches(0.22)
            if layer.get("icon"):
                txt(slide, lx, y_c, layer_w, Inches(0.55),
                    str(layer["icon"]), sz=20, color=WHITE,
                    font=FONT_BODY, align=PP_ALIGN.CENTER)
                y_c += Inches(0.60)
            txt(slide, lx + Inches(0.10), y_c, layer_w - Inches(0.20), Inches(0.42),
                str(layer.get("label", "")),
                sz=11, color=WHITE, bold=True, font=FONT_BODY,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            if layer.get("body"):
                txt(slide, lx + Inches(0.10), y_c + Inches(0.48),
                    layer_w - Inches(0.20),
                    zone_h - (y_c - zone_y) - Inches(0.56),
                    str(layer["body"]),
                    sz=9, color=WHITE, font=FONT_BODY, anchor=MSO_ANCHOR.TOP)
    else:
        layer_h = (zone_h - gap * (n - 1)) / n
        for i, layer in enumerate(layers):
            ly  = zone_y + i * (layer_h + gap)
            col = layer.get("color") or L_COLORS[i % len(L_COLORS)]
            rounded(slide, zone_x, ly, zone_w, layer_h, fill=col)
            x_c = zone_x + Inches(0.18)
            if layer.get("icon"):
                txt(slide, x_c, ly, Inches(0.55), layer_h,
                    str(layer["icon"]), sz=16, color=WHITE, font=FONT_BODY,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                x_c += Inches(0.60)
            txt(slide, x_c, ly, Inches(3.0), layer_h,
                str(layer.get("label", "")),
                sz=12, color=WHITE, bold=True, font=FONT_BODY,
                anchor=MSO_ANCHOR.MIDDLE)
            if layer.get("body"):
                txt(slide, x_c + Inches(3.10), ly,
                    zone_w - (x_c - zone_x) - Inches(3.20), layer_h,
                    str(layer["body"]),
                    sz=10, color=WHITE, font=FONT_BODY,
                    anchor=MSO_ANCHOR.MIDDLE)

    return slide


# ===========================================================================
# Layout: PHOTO + TEXT (split panel)
# ===========================================================================

def build_photo_text(prs, meta, action_title=None, image_path=None,
                     image_label=None, title=None, bullets=None,
                     image_side="left", takeaway="", source=None,
                     subtitle=None, methodology_note=None,
                     style="merck_executive", page=None, total=None,
                     section_number=None, category=None):
    """Half-slide image panel + half-slide text narrative.

    image_path:  str — path to png/jpg; placeholder drawn if None/missing
    image_label: str — caption inside the placeholder
    title:       str — bold heading in the text panel
    bullets:     list of str — narrative points
    image_side:  'left' (default) or 'right'
    """
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    dark    = _is_dark(style)
    zone_y  = Inches(1.62)
    zone_h  = Inches(4.72)
    panel_w = Inches(5.90)
    text_w  = Inches(12.2) - panel_w - Inches(0.20)

    if image_side == "right":
        img_x  = Inches(0.55) + text_w + Inches(0.20)
        text_x = Inches(0.55)
    else:
        img_x  = Inches(0.55)
        text_x = Inches(0.55) + panel_w + Inches(0.20)

    if image_path and os.path.isfile(str(image_path)):
        try:
            slide.shapes.add_picture(str(image_path),
                                     _emu(img_x), _emu(zone_y),
                                     _emu(panel_w), _emu(zone_h))
        except Exception:
            _draw_img_placeholder(slide, img_x, zone_y, panel_w, zone_h,
                                  image_label, dark, pal=pal)
    else:
        _draw_img_placeholder(slide, img_x, zone_y, panel_w, zone_h,
                              image_label, dark, pal=pal)

    y_cur = zone_y + Inches(0.14)
    if title:
        txt(slide, text_x, y_cur, text_w, Inches(0.58), str(title),
            sz=16, color=MERCK_YELLOW if dark else MERCK_PURPLE,
            bold=True, font=FONT_HEAD, anchor=MSO_ANCHOR.TOP)
        hairline(slide, text_x, y_cur + Inches(0.66),
                 text_w * 0.55, Emu(int(Pt(2))), MERCK_GOLD)
        y_cur += Inches(0.84)

    body_h  = zone_h - (y_cur - zone_y) - Inches(0.10)
    bullets = list(bullets or [])
    if bullets:
        item_h  = min(body_h / len(bullets), Inches(0.70))
        b_color = PANEL_LIGHT if dark else INK_DARK
        for i, b in enumerate(bullets):
            txt(slide, text_x, y_cur + i * item_h, text_w, item_h,
                "•  " + str(b),
                sz=11, color=b_color, font=FONT_BODY, anchor=MSO_ANCHOR.TOP)

    return slide


def _draw_img_placeholder(slide, x, y, w, h, label, dark, pal=None):
    """Branded placeholder rectangle when no image file is provided."""
    fill = pal["panel"] if pal else ((0x55, 0x55, 0x66) if dark else (0xCC, 0xCC, 0xDD))
    text_color = pal["muted"] if pal else WHITE
    rect(slide, x, y, w, h, fill=fill)
    display_label = ("[Image]  " + str(label)) if label else "[Image]"
    txt(slide, x, y + h / 2 - Inches(0.28), w, Inches(0.48),
        display_label,
        sz=13, color=text_color, italic=True, font=FONT_BODY,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ===========================================================================
# Layout: FISHBONE (Ishikawa / Cause-and-Effect Diagram)
# ===========================================================================

def build_fishbone(prs, meta, action_title=None, effect=None, bones=None,
                   takeaway=None, source=None, category=None, subtitle=None,
                   style="merck_executive", page=None, total=None,
                   section_number=None, methodology_note=None, content=None):
    """Cause-and-effect (Ishikawa) diagram.

    effect: str - the problem or outcome at the arrow head (right side).
    bones: list of dicts {label, causes} - up to 6 branch categories.
           Alternates above/below the spine.
    """
    if content:
        if "effect"   in content: effect   = content["effect"]
        if "bones"    in content: bones    = content["bones"]
        if "takeaway" in content: takeaway = content["takeaway"]
    style = _style_or_promote(category, style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category, subtitle=subtitle,
                 takeaway=takeaway, source=source, page=page, total=total,
                 palette=style, section_number=section_number,
                 methodology_note=methodology_note)

    dark = _is_dark(style)
    text_color = WHITE if dark else INK_DARK

    zone_top = Inches(2.55) if not subtitle else Inches(2.95)
    zone_h   = Inches(6.40) - zone_top
    mid_y    = zone_top + zone_h / 2

    spine_x1  = Inches(0.80)
    spine_x2  = Inches(10.80)
    effect_w  = Inches(2.00)
    effect_h  = Inches(0.90)

    # Spine
    line(slide, spine_x1, mid_y, spine_x2, mid_y, pal["highlight"], Pt(3))

    # Effect box at right end
    rounded(slide, spine_x2 - effect_w / 2, mid_y - effect_h / 2,
            effect_w, effect_h, fill=pal["accent"])
    txt(slide, spine_x2 - effect_w / 2, mid_y - effect_h / 2,
        effect_w, effect_h, str(effect or "Effect"),
        sz=12, color=WHITE, bold=True, font=FONT_BODY,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    bs = list(bones or [])[:6]
    above = [b for i, b in enumerate(bs) if i % 2 == 0]
    below = [b for i, b in enumerate(bs) if i % 2 == 1]

    n_above = len(above)
    n_below = len(below)
    usable  = spine_x2 - spine_x1 - Inches(1.20)

    for group, sign, n_grp in [(above, -1, n_above), (below, 1, n_below)]:
        if not n_grp:
            continue
        step = usable / n_grp
        # Dynamic label width: never wider than step minus a 0.20" gutter (RCA-09)
        lbl_w = min(Inches(2.00), step - Inches(0.20))
        lbl_h = Inches(0.36)
        for k, bone in enumerate(group):
            bx      = spine_x1 + Inches(0.60) + k * step + step / 2
            bone_y  = mid_y + sign * Inches(1.40)
            meet_x  = bx + Inches(0.40) * sign * -1
            # Diagonal branch to spine
            line(slide, bx, bone_y, meet_x, mid_y, pal["highlight"], Pt(1.5))
            # Category label
            lbl_y = bone_y - lbl_h if sign < 0 else bone_y
            rounded(slide, bx - lbl_w / 2, lbl_y, lbl_w, lbl_h, fill=pal["accent"])
            txt(slide, bx - lbl_w / 2, lbl_y, lbl_w, lbl_h,
                str(bone.get("label", "")),
                sz=9, color=WHITE, bold=True, font=FONT_BODY,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            # Sub-causes
            for ci, cause in enumerate(list(bone.get("causes") or [])[:3]):
                cy2 = bone_y + sign * Inches(0.42 + ci * 0.30)
                txt(slide, bx - Inches(1.10), cy2,
                    Inches(2.20), Inches(0.28),
                    "•  " + str(cause), sz=9, color=text_color,
                    font=FONT_BODY, align=PP_ALIGN.LEFT)

    return slide


# ===========================================================================
# Layout: KEY QUESTION
# ===========================================================================

def build_key_question(prs, meta, action_title=None, question=None, context=None,
                       takeaway="", source=None, subtitle=None,
                       methodology_note=None, style="merck_executive",
                       page=None, total=None, section_number=None,
                       category=None, content=None):
    """Centred question-mark icon with four directional arrows and question text.

    content keys:
        question (str) — the key question text
        context  (str) — optional supporting context line
    """
    if content:
        question = question or content.get("question")
        context  = context  or content.get("context")
    style = _style_or_promote(category, style)
    pal   = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title, category=category,
                 takeaway=takeaway, source=source, subtitle=subtitle,
                 methodology_note=methodology_note,
                 page=page, total=total, section_number=section_number,
                 palette=style)

    # Icon circle centred on slide
    CX = SLIDE_W / 2
    CY = Inches(3.30)
    R  = Inches(0.60)
    ACC = pal["accent"]
    HOT = pal["hot"]

    circle(slide, CX - R, CY - R, R * 2, fill=ACC)
    txt(slide, CX - R, CY - R, R * 2, R * 2, "?",
        sz=36, color=WHITE, bold=True, font=FONT_HEAD,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Four directional arrows pointing toward the icon
    AW, AH, GAP = Inches(0.50), Inches(0.28), Inches(0.18)
    arrow_specs = [
        (MSO_SHAPE.LEFT_ARROW,  CX - R - GAP - AW,  CY - AH / 2,  AW, AH),
        (MSO_SHAPE.RIGHT_ARROW, CX + R + GAP,        CY - AH / 2,  AW, AH),
        (MSO_SHAPE.UP_ARROW,    CX - AH / 2, CY - R - GAP - AW,   AH, AW),
        (MSO_SHAPE.DOWN_ARROW,  CX - AH / 2, CY + R + GAP,         AH, AW),
    ]
    for shape_type, ax, ay, aw, ah in arrow_specs:
        shp = slide.shapes.add_shape(shape_type,
                                     int(ax), int(ay), int(aw), int(ah))
        shp.shadow.inherit = False
        _apply_fill(shp, HOT)
        _apply_border(shp, None)

    # Question text
    if question:
        txt(slide, Inches(1.5), CY + R + Inches(0.50), Inches(10.33), Inches(0.70),
            str(question), sz=22, color=pal["ink"], bold=True, font=FONT_HEAD,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    if context:
        txt(slide, Inches(2.0), CY + R + Inches(1.30), Inches(9.33), Inches(0.50),
            str(context), sz=13, color=pal["ink_2"], font=FONT_BODY,
            align=PP_ALIGN.CENTER)
    return slide


