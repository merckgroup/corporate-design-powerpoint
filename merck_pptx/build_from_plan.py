"""Canonical Merck deck builder: JSON plan → PowerPoint.

Reads a slide plan (dict or .json file) and calls the correct build_*
function in merck_layouts.py for every slide. Handles style resolution,
auto-promote, agenda auto-fill, appendix slides, hyperlink wiring, and
post-build speaker notes / image placement.

Usage:
    from merck_pptx import build_from_plan
    build_from_plan("plan.json", "output/deck.pptx")

Or from a dict:
    from merck_pptx import build_from_plan
    build_from_plan(plan_dict, "output/deck.pptx")
"""

from __future__ import annotations

import json
import pathlib
import sys
from typing import Optional

try:
    from lxml import etree as _ET
    _HAS_LXML = True
except ImportError:
    import xml.etree.ElementTree as _ET
    _HAS_LXML = False
    import warnings as _lxml_warnings
    _lxml_warnings.warn(
        "lxml is not installed. The 'synthetic' and 'electronics' color themes will "
        "render with an incorrect background color. Install lxml to fix: pip install lxml",
        ImportWarning,
        stacklevel=2,
    )

from merck_pptx.merck_layouts import (
    open_deck, save_deck, AUTO_PROMOTE_EXECUTIVE,
    add_speaker_notes, add_slide_jump_hyperlink, add_image,
    build_cover, build_exec_summary, build_agenda, build_section_divider,
    build_chart_slide, build_two_column, build_three_column, build_2x2_matrix,
    build_phase_process, build_vertical_numbered, build_waterfall_slide,
    build_decision_rows, build_gantt, build_hero_stat, build_close,
    build_stat_strip, build_before_after, build_milestone_timeline,
    build_status_table, build_hub_spoke, build_pillar_detail,
    build_four_column, build_label_rows,
    build_circular_flow, build_org_chart, build_topic_set, build_arrow_chain,
    build_pull_quote, build_donut_chart, build_kpi_dashboard,
    build_fishbone, build_journey_map,
    build_icon_grid, build_funnel,
    build_comparison_table, build_score_table, build_influence_diagram,
    build_word_cloud, build_pyramid, build_venn, build_risk_heatmap,
    build_radar_chart, build_pros_cons, build_layered_stack,
    build_photo_text, build_key_question, build_road_to_success,
)
from merck_pptx.validate_plan import validate_plan, ValidationError

_PKG_DIR = pathlib.Path(__file__).parent
_TEMPLATE_DIR = _PKG_DIR / "templates"

# ---------------------------------------------------------------------------
# Template file selection — one base file per division (not per theme).
# Color themes are applied programmatically via _apply_color_theme() after
# the template is loaded, so only a small set of division template files is
# needed regardless of how many color_theme variants exist.
#
# Supported division values (meta.division):
#   merck (default) | emd_serono | emd_electronics | millipore_sigma | merck_asia
# Supported region values (meta.region):
#   eu (default) | usa / us / canada
# ---------------------------------------------------------------------------
_DIVISION_TEMPLATES = {
    # (region, division) → filename in _TEMPLATE_DIR
    # EU / global
    ("eu",  "merck"):           "EU_Merck_Themed.pptx",
    ("eu",  "emd_serono"):      "EU_EMDSerono_Themed.pptx",
    ("eu",  "emd_electronics"): "EU_EMDElectronics_Themed.pptx",
    ("eu",  "millipore_sigma"): "EU_MilliporeSigma_Themed.pptx",
    ("eu",  "merck_asia"):      "EU_MerckAsia_Themed.pptx",
    # USA / Canada
    ("usa", "merck"):           "USA_Merck_Themed_Base_v1.pptx",
    ("usa", "emd_serono"):      "USA_EMDSerono_Themed.pptx",
    ("usa", "emd_electronics"): "USA_EMDElectronics_Themed.pptx",
    ("usa", "millipore_sigma"): "USA_MilliporeSigma_Themed.pptx",
    ("usa", "merck_asia"):      "USA_MerckAsia_Themed.pptx",
}

_TEMPLATE_DEFAULT = _TEMPLATE_DIR / "EU_Merck_Themed.pptx"

# Backward-compat alias used by legacy callers.
_TEMPLATES = {
    "eu":     _TEMPLATE_DIR / "EU_Merck_Themed.pptx",
    "usa":    _TEMPLATE_DIR / "USA_Merck_Themed_Base_v1.pptx",
    "us":     _TEMPLATE_DIR / "USA_Merck_Themed_Base_v1.pptx",
    "canada": _TEMPLATE_DIR / "USA_Merck_Themed_Base_v1.pptx",
}


def _resolve_template(region: str, color_theme: str,
                      division: str = "merck") -> pathlib.Path:
    """Return the best-matching base template for (region, division).

    color_theme is intentionally ignored here — themes are applied
    programmatically by _apply_color_theme() after the template is loaded,
    so only one file per (region, division) pair is needed.

    Lookup order:
    1. (region, division) from _DIVISION_TEMPLATES (if file exists on disk)
    2. (region, "merck")  — region default
    3. global default     — EU_Merck_Themed.pptx
    """
    r = str(region   or "eu").lower().strip()
    d = str(division or "merck").lower().strip().replace("-", "_").replace(" ", "_")
    if r in ("us", "canada"):
        r = "usa"

    # 1. Division-specific template.
    fname = _DIVISION_TEMPLATES.get((r, d))
    if fname:
        candidate = _TEMPLATE_DIR / fname
        if candidate.exists():
            return candidate

    # 2. Region Merck default.
    merck_default = _DIVISION_TEMPLATES.get((r, "merck"))
    if merck_default:
        candidate = _TEMPLATE_DIR / merck_default
        if candidate.exists():
            return candidate

    # 3. Global fallback.
    return _TEMPLATE_DEFAULT


# ---------------------------------------------------------------------------
# Programmatic color-theme application
# ---------------------------------------------------------------------------

# Merck Corporate Design exact hex values.
_MC = {
    "violet":     "503291",
    "pink":       "EB3C96",
    "red":        "E61E50",
    "teal":       "2DBECD",
    "lightgreen": "A5CD50",
    "yellow":     "FFC832",
    "cream":      "FFDCB9",
    "palegreen":  "B4DC96",   # hardcoded in the current EU template shapes
}

# Per-theme overrides for the PowerPoint theme XML color slots.
# Slots used by the organic blob shapes in the EU template:
#   accent1 → main large blob (cover + divider left panel bg)
#   dk2     → background-covering rectangle (slide bg visual layer)
#   accent5 → accent/secondary blob (divider right panel + cover accent)
#   lt2/bg2 → used in some divider sub-shapes
# _BG_HEX: the hardcoded #B4DC96 pale-green that forms the "light" panel
#           in the cover and divider — replaced with theme background color.
_THEME_SCHEME_OVERRIDES = {
    # Per-theme overrides for (accent1, dk2, accent5, lt2).
    #
    # accent1 → main organic blob group shapes (cover + divider large panel)
    # dk2     → Rechteck 51: the full-slide background rectangle below the blobs
    # accent5 → secondary accent blobs in the divider layout
    # lt2     → secondary uses in divider sub-shapes
    #
    # For light themes (bg = lime green or cream):
    #   dk2 = accent colour so the background rect peeks through the freeform gap
    # For dark themes (bg = violet):
    #   dk2 = "503291" so the background rect is violet (dark) → dark slide feel
    #
    #            accent1      dk2         accent5     lt2
    "plastic":   ("503291", "EB3C96",  "EB3C96",  "B4DC96"),
    "functional":("503291", "2DBECD",  "2DBECD",  "A5CD50"),
    "organic":   ("503291", "E61E50",  "E61E50",  "FFDCB9"),
    "synthetic": ("FFC832", "503291",  "FFC832",  "B4DC96"),
    # technical: light-theme (cream bg) — dk2/accent5 must be teal so the
    # background rectangle shows through the freeform gap in teal, not violet.
    "technical": ("2DBECD", "2DBECD",  "2DBECD",  "FFDCB9"),
    "electronics":("FFC832","503291",  "FFC832",  "B4DC96"),
}

# Hardcoded background hex values found in the EU and USA template layout shapes.
# The EU template uses pale-green (#B4DC96) for the "light panel" freeform; the
# USA template uses lime-green (#A5CD50) for the equivalent shape.  Both must be
# patched so organic/technical (cream) and synthetic/electronics (transparent)
# themes render correctly across both regional templates.
_TEMPLATE_HARDCODED_BG   = "B4DC96"   # EU template light-panel colour
_TEMPLATE_HARDCODED_BG_2 = "A5CD50"   # USA template light-panel colour

# Background replacement color per theme (replaces the hardcoded pale-green
# freeform that forms the "light panel" on covers and dividers).
# None = make the freeform transparent so dark-theme blob shapes show through.
_THEME_BG_HEX = {
    "plastic":    "A5CD50",   # lime green
    "functional": "A5CD50",   # lime green
    "organic":    "FFDCB9",   # cream / paleyellow
    "synthetic":  None,        # transparent → yellow blobs visible on violet bg
    "technical":  "FFDCB9",   # cream
    "electronics":None,        # transparent → yellow blobs visible on violet bg
}


def _apply_color_theme(prs, color_theme: str) -> None:
    """Apply a Merck Corporate Design color theme to the loaded presentation.

    Modifies two things IN MEMORY (original file untouched):
    1. The slide master's theme color scheme XML — so shapes using
       scheme:accent1 / scheme:dk2 / scheme:accent5 etc. automatically
       render in the correct theme colors.
    2. All layout shape XML — replaces the hardcoded #B4DC96 pale-green
       that forms the "light panel" background on covers and dividers with
       the appropriate theme background color.

    Must be called immediately after open_deck() and before any slides are
    added.  Themes map exactly to the 6 official Merck Corporate Design
    variants in the empower library.
    """
    theme_lower = str(color_theme or "plastic").lower().strip()
    overrides = _THEME_SCHEME_OVERRIDES.get(theme_lower)
    bg_hex    = _THEME_BG_HEX.get(theme_lower, _TEMPLATE_HARDCODED_BG)

    if not overrides:
        return   # unknown theme — leave template colors unchanged

    a1, dk2, a5, lt2 = overrides

    # ------------------------------------------------------------------
    # 1. Modify the theme color scheme in the slide master's theme part.
    # ------------------------------------------------------------------
    _NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

    def _set_scheme_color(clr_scheme, slot_tag: str, hex_val: str):
        """Find <a:{slot_tag}> inside <a:clrScheme> and set its srgbClr."""
        for child in clr_scheme:
            tag = child.tag.split("}")[1] if "}" in child.tag else child.tag
            if tag != slot_tag:
                continue
            srgb = child.find(f"{{{_NS}}}srgbClr")
            if srgb is not None:
                srgb.set("val", hex_val)
            else:
                # Replace any sysClr with srgbClr
                for old in list(child):
                    child.remove(old)
                srgb = _ET.SubElement(child, f"{{{_NS}}}srgbClr")
                srgb.set("val", hex_val)
            return

    master = prs.slide_masters[0]
    try:
        # Locate the theme part via relationship.
        # The theme part is a raw opc.Part (bytes), not a parsed XML part,
        # so we access it via .blob, modify, and write back.
        THEME_RELTYPE = (
            "http://schemas.openxmlformats.org/officeDocument/2006/"
            "relationships/theme"
        )
        theme_part = None
        for rel in master.part.rels.values():
            if not rel.is_external and THEME_RELTYPE in str(rel.reltype):
                theme_part = rel.target_part
                break

        if theme_part is not None:
            # Parse the theme XML from its raw bytes.
            theme_xml  = theme_part.blob
            theme_elem = _ET.fromstring(theme_xml)
            clr_scheme = theme_elem.find(f".//{{{_NS}}}clrScheme")
            if clr_scheme is not None:
                _set_scheme_color(clr_scheme, "accent1", a1)
                _set_scheme_color(clr_scheme, "dk2",     dk2)
                _set_scheme_color(clr_scheme, "accent5", a5)
                _set_scheme_color(clr_scheme, "lt2",     lt2)
                # Write the modified XML back as bytes.
                # standalone=True is lxml-only; stdlib ET ignores the kwarg.
                _ts_kw = {"standalone": True} if _HAS_LXML else {}
                theme_part.blob = _ET.tostring(
                    theme_elem, xml_declaration=True,
                    encoding="UTF-8", **_ts_kw,
                )
    except Exception as exc:
        print(f"WARNING: could not apply theme color scheme: {exc}",
              file=sys.stderr)

    # ------------------------------------------------------------------
    # 2. Replace hardcoded #B4DC96 hex in all layout shapes.
    #    bg_hex = new color hex  → replace the fill value
    #    bg_hex = None           → make the shape transparent (noFill)
    # ------------------------------------------------------------------
    if bg_hex is not None and bg_hex.upper() == _TEMPLATE_HARDCODED_BG.upper():
        return   # no change needed for plastic (default)

    # Hex values that are "baseline background" colours baked into template shapes.
    # B4DC96 = EU template pale-green panel; A5CD50 = USA template lime-green panel.
    _BG_TARGETS = {_TEMPLATE_HARDCODED_BG.upper(), "A5CD50"}
    # For plastic/functional the lime-green IS the intended background — skip it.
    if bg_hex is not None and bg_hex.upper() == "A5CD50":
        _BG_TARGETS.discard("A5CD50")

    def _replace_bg_srgb(srgb):
        """Replace or noFill one srgbClr element; returns True if acted on."""
        solidFill = srgb.getparent()
        if solidFill is None:
            return False
        parent = solidFill.getparent()
        if parent is None:
            return False
        # Context guard: only patch shape/group fills, not text-run colours.
        _ptag = parent.tag.split("}")[1] if "}" in parent.tag else parent.tag
        if _ptag in ("rPr", "pPr", "defRPr", "endParaRPr", "lstStyle"):
            return False
        if bg_hex is not None:
            srgb.set("val", bg_hex.upper())
        elif _HAS_LXML:
            # Replace <a:solidFill><a:srgbClr.../></a:solidFill> with <a:noFill/>
            # so the freeform becomes transparent and dark-theme blobs show through.
            idx = list(parent).index(solidFill)
            parent.remove(solidFill)
            noFill = _ET.Element(f"{{{_NS}}}noFill")
            parent.insert(idx, noFill)
        else:
            # Without lxml, paint with dark-theme violet as closest approximation.
            srgb.set("val", "503291")
        return True

    try:
        elements_to_patch = list(master.element.iter(f"{{{_NS}}}srgbClr"))
        for layout in master.slide_layouts:
            elements_to_patch.extend(layout.element.iter(f"{{{_NS}}}srgbClr"))

        for srgb in elements_to_patch:
            if srgb.get("val", "").upper() in _BG_TARGETS:
                _replace_bg_srgb(srgb)
    except Exception as exc:
        print(f"WARNING: could not replace hardcoded bg color: {exc}",
              file=sys.stderr)

# Permitted image file extensions — prevents a crafted plan from embedding
# arbitrary files (e.g. databases, documents) by disguising them as images.
_SAFE_IMAGE_EXTENSIONS = frozenset({
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif", ".webp", ".emf", ".wmf",
})


class BuildError(RuntimeError):
    """Wraps a layout-builder exception with page + layout context."""


# Layouts that carry no section-number circle.
# Must stay in sync with _STRUCTURAL_LAYOUTS in validate_plan.py.
_NO_CIRCLE_LAYOUTS = {
    "cover", "agenda", "section_divider", "close", "exec_summary", "hero_stat", "pull_quote"
}


def _is_appendix(slide: dict) -> bool:
    return bool(slide.get("appendix"))


def _short_title(text, max_chars: int = 40) -> str:
    if text is None:
        return ""
    if isinstance(text, (list, tuple)):
        text = "".join(seg[0] if isinstance(seg, (list, tuple)) else str(seg) for seg in text)
    s = str(text).strip()
    if len(s) <= max_chars:
        return s
    return s[:max_chars].rsplit(" ", 1)[0] + "…"


def _place_image(slide, img_spec: dict) -> None:
    from pptx.util import Inches as _In
    raw_path = img_spec.get("path")
    if not raw_path:
        return

    # Security: validate the path before embedding.
    # (1) Extension check — prevents arbitrary file types from being embedded.
    # (2) Directory jail — restricts images to the current working directory tree
    #     to prevent a crafted plan from exfiltrating arbitrary local files
    #     (e.g. via prompt injection setting image.path to a sensitive path).
    img_path = pathlib.Path(raw_path).resolve()
    if img_path.suffix.lower() not in _SAFE_IMAGE_EXTENSIONS:
        print(
            f"WARNING: image not placed — '{img_path.name}' is not a recognised image type. "
            f"Allowed: {sorted(_SAFE_IMAGE_EXTENSIONS)}",
            file=sys.stderr,
        )
        return
    _allowed_image_root = pathlib.Path.cwd().resolve()
    try:
        img_path.relative_to(_allowed_image_root)
    except ValueError:
        print(
            f"WARNING: image not placed — '{img_path}' is outside the allowed "
            f"directory ({_allowed_image_root}). Move the image file inside the "
            f"project directory.",
            file=sys.stderr,
        )
        return

    placement = (img_spec.get("placement") or "right_panel").lower()
    presets = {
        "hero":        (6.83, 1.30, 6.00, 5.00),
        "right_panel": (8.50, 3.10, 4.50, 3.50),
        "background":  (0.00, 0.00, 13.33, 7.50),
    }
    x, y, w, h = presets.get(placement, presets["right_panel"])
    w = img_spec.get("w") or w
    h = img_spec.get("h") or h
    try:
        add_image(slide, str(img_path), _In(x), _In(y),
                  w=_In(w) if w else None, h=_In(h) if h else None)
    except Exception as exc:
        print(f"WARNING: image not placed: {exc}", file=sys.stderr)


def _wire_agenda_hyperlinks(prs, ordered_slides: list) -> None:
    sn_to_idx: dict[str, int] = {}
    for i, s in enumerate(ordered_slides):
        if s.get("layout") in _NO_CIRCLE_LAYOUTS or _is_appendix(s):
            continue
        sn = s.get("section_number")
        if sn is None:
            continue
        sn_str = str(sn).strip()
        sn_to_idx[sn_str] = i
        try:
            sn_to_idx[f"{int(sn_str):02d}"] = i
        except (ValueError, TypeError):
            pass

    for i, s in enumerate(ordered_slides):
        if s.get("layout") != "agenda":
            continue
        agenda_slide = prs.slides[i]
        for shp in agenda_slide.shapes:
            name = getattr(shp, "name", "") or ""
            if not name.startswith("AgendaChapter_"):
                continue
            key = name.split("_", 1)[1]
            target_idx = sn_to_idx.get(key)
            if target_idx is None:
                stripped = key.lstrip("0")
                if stripped:
                    target_idx = sn_to_idx.get(stripped)
            if target_idx is None:
                print(f"WARNING: agenda hyperlink '{name}' — no slide with section_number '{key}'",
                      file=sys.stderr)
                continue
            try:
                add_slide_jump_hyperlink(shp, prs.slides[target_idx])
            except Exception as exc:
                print(f"WARNING: agenda hyperlink '{name}' not set: {exc}", file=sys.stderr)


def _autofill_agenda(plan: dict) -> None:
    slides = plan.get("slides") or []
    content_slides = [
        s for s in slides
        if s.get("layout") not in _NO_CIRCLE_LAYOUTS and not _is_appendix(s)
    ]
    for s in slides:
        if s.get("layout") != "agenda":
            continue
        content = s.get("content") or {}
        if content.get("chapters"):
            continue
        derived = []
        for cs in content_slides:
            sn = cs.get("section_number")
            num = ""
            if sn is not None:
                try:
                    num = f"{int(sn):02d}"
                except (ValueError, TypeError):
                    num = str(sn)
            derived.append({
                "number": num,
                "title": _short_title(cs.get("action_title"), 40),
                "subtitle": _short_title(cs.get("category") or "", 30),
            })
        if derived:
            s["content"] = {**content, "chapters": derived}


# ---------------------------------------------------------------------------
# Style helpers
# ---------------------------------------------------------------------------

_VALID_COLOR_THEMES = frozenset({
    "functional", "organic", "plastic", "synthetic", "technical", "electronics",
})


def _resolve_style(slide: dict, meta: dict) -> str:
    deck_style  = meta.get("deck_style")
    color_theme = str(meta.get("color_theme") or "").lower().strip()

    # merck_storytelling decks are visually coherent ONLY when every slide
    # uses the same dark palette.  Ignore per-slide style overrides and
    # auto-promote entirely so that Risk/Recommendation/Decision slides don't
    # revert to white inside an otherwise dark purple deck.
    if deck_style == "merck_storytelling":
        return "merck_storytelling"

    style = slide.get("style", "inherit")
    if style == "inherit" or not style:
        # Use deck_style if set; otherwise fall back to color_theme so the
        # 6 theme palettes are automatically applied without setting deck_style.
        if deck_style:
            style = deck_style
        elif color_theme in _VALID_COLOR_THEMES:
            style = color_theme
        else:
            style = "merck_executive"

    # Auto-promote: category is the canonical trigger field; page_function is
    # the fallback so plans that set only page_function still work.
    _cat = slide.get("category") or slide.get("page_function")
    if _cat in AUTO_PROMOTE_EXECUTIVE:
        style = "merck_executive"
    return style


def _resolve_category(slide: dict) -> Optional[str]:
    """Auto-set category to match page_function for auto-promote slides.
    The library triggers on category, not page_function."""
    pf = slide.get("page_function")
    if pf in AUTO_PROMOTE_EXECUTIVE:
        return pf
    return slide.get("category")


def _common_kwargs(slide: dict, meta: dict, total) -> dict:
    return {
        "style":          _resolve_style(slide, meta),
        "page":           slide.get("page"),
        "total":          total,
        "section_number": slide.get("section_number"),
        "category":       _resolve_category(slide),
    }


def _content(slide: dict) -> dict:
    c = slide.get("content")
    return c if isinstance(c, dict) else {}


# ---------------------------------------------------------------------------
# Per-layout builder wrappers
# ---------------------------------------------------------------------------

def _build_cover(prs, meta, slide, total):
    c = _content(slide)
    return build_cover(
        prs, meta,
        action_title=slide.get("action_title"),
        subtitle=slide.get("subtitle", ""),
        style=_resolve_style(slide, meta),
        key_messages=c.get("key_messages"),
        phases=c.get("phases"),
        authors=c.get("authors"),
        top_bar=bool(meta.get("cover_top_bar", False)),
        page=slide.get("page"),
        total=total,
        color_theme=str(meta.get("color_theme") or "").lower().strip() or None,
    )


def _build_exec_summary(prs, meta, slide, total):
    c = _content(slide)
    return build_exec_summary(
        prs, meta,
        action_title=slide.get("action_title", ""),
        key_messages=c.get("key_messages", []),
        takeaway=slide.get("takeaway", ""),
        source=slide.get("source"),
        subtitle=slide.get("subtitle"),
        methodology_note=c.get("methodology_note"),
        **_common_kwargs(slide, meta, total),
    )


def _build_agenda(prs, meta, slide, total):
    c = _content(slide)
    return build_agenda(
        prs, meta,
        chapters=c.get("chapters", []),
        style=_resolve_style(slide, meta),
        action_title=slide.get("action_title"),
        page=slide.get("page"),
        total=total,
    )


def _build_section_divider(prs, meta, slide, total):
    number = slide.get("section_number") or _content(slide).get("number") or ""
    return build_section_divider(
        prs, meta,
        number=number,
        title=slide.get("action_title"),
        style=_resolve_style(slide, meta),
        page=slide.get("page"),
        total=total,
        takeaway=slide.get("takeaway"),
        source=slide.get("source"),
    )


def _build_chart_slide(prs, meta, slide, total):
    c = _content(slide)
    return build_chart_slide(
        prs, meta,
        action_title=slide.get("action_title", ""),
        chart=c.get("chart", {}),
        takeaway=slide.get("takeaway", ""),
        source=slide.get("source"),
        callouts=c.get("callouts"),
        subtitle=slide.get("subtitle"),
        footnotes=c.get("footnotes"),
        methodology_note=c.get("methodology_note"),
        **_common_kwargs(slide, meta, total),
    )


def _build_two_column(prs, meta, slide, total):
    c = _content(slide)
    return build_two_column(
        prs, meta,
        action_title=slide.get("action_title", ""),
        left=c.get("left", {}),
        right=c.get("right", {}),
        takeaway=slide.get("takeaway", ""),
        source=slide.get("source"),
        subtitle=slide.get("subtitle"),
        methodology_note=c.get("methodology_note"),
        **_common_kwargs(slide, meta, total),
    )


def _build_three_column(prs, meta, slide, total):
    c = _content(slide)
    return build_three_column(
        prs, meta,
        action_title=slide.get("action_title", ""),
        columns=c.get("columns", []),
        takeaway=slide.get("takeaway", ""),
        source=slide.get("source"),
        subtitle=slide.get("subtitle"),
        methodology_note=c.get("methodology_note"),
        **_common_kwargs(slide, meta, total),
    )


def _build_2x2_matrix(prs, meta, slide, total):
    c = _content(slide)
    return build_2x2_matrix(
        prs, meta,
        action_title=slide.get("action_title", ""),
        x_axis=c.get("x_axis", {}),
        y_axis=c.get("y_axis", {}),
        quadrants=c.get("quadrants", {}),
        takeaway=slide.get("takeaway", ""),
        source=slide.get("source"),
        subtitle=slide.get("subtitle"),
        methodology_note=c.get("methodology_note"),
        **_common_kwargs(slide, meta, total),
    )


def _build_phase_process(prs, meta, slide, total):
    c = _content(slide)
    return build_phase_process(
        prs, meta,
        action_title=slide.get("action_title", ""),
        phases=c.get("phases", []),
        takeaway=slide.get("takeaway", ""),
        source=slide.get("source"),
        subtitle=slide.get("subtitle"),
        highlight_index=c.get("highlight_index"),
        show_arrows=c.get("show_arrows", True),
        methodology_note=c.get("methodology_note"),
        **_common_kwargs(slide, meta, total),
    )


def _build_vertical_numbered(prs, meta, slide, total):
    c = _content(slide)
    return build_vertical_numbered(
        prs, meta,
        action_title=slide.get("action_title", ""),
        items=c.get("items", []),
        takeaway=slide.get("takeaway", ""),
        source=slide.get("source"),
        methodology_note=c.get("methodology_note"),
        **_common_kwargs(slide, meta, total),
    )


def _normalise_waterfall_bars(bars: list) -> list:
    """Translate plan-schema bar types to the layout function's internal types.

    Plan schema uses:  "total" | "positive" | "negative"
    Layout function expects: "start" | "end" | "up" | "down"

    Rules:
    - "total"    → "start" for the first bar in a run of totals,
                   "end" for the last; all others in the middle become "end"
                   (the function treats "end" and "start" identically in terms
                   of running-total logic, both anchor to zero).
    - "positive" → "up"   (value kept as-is, must be > 0)
    - "negative" → "down" (value stored as positive; add_waterfall uses abs())

    Unknown types are left unchanged so hand-crafted plans using the internal
    type names ("up", "down", "start", "end") continue to work.
    """
    _ALIAS = {"positive": "up", "negative": "down", "total": "end"}
    normalised = []
    for i, b in enumerate(bars):
        b = dict(b)
        t = str(b.get("type", "up")).lower()
        if t == "total":
            # First bar anchors the baseline; subsequent totals close the chart.
            b["type"] = "start" if i == 0 else "end"
        elif t == "negative":
            b["type"] = "down"
            # Value must be positive; layout function uses abs() for "down" bars.
            try:
                b["value"] = abs(float(b.get("value", 0)))
            except (TypeError, ValueError):
                pass
        elif t in _ALIAS:
            b["type"] = _ALIAS[t]
        normalised.append(b)
    return normalised


def _build_waterfall_slide(prs, meta, slide, total):
    c = _content(slide)
    chart = c.get("chart") or {}
    # Schema: content.chart.data.bars — extract bars for the layout function.
    bars = chart.get("data", {}).get("bars") or chart.get("bars") or []
    bars = _normalise_waterfall_bars(bars)
    return build_waterfall_slide(
        prs, meta,
        action_title=slide.get("action_title", ""),
        bars=bars,
        takeaway=slide.get("takeaway", ""),
        source=slide.get("source"),
        footnotes=c.get("footnotes"),
        methodology_note=c.get("methodology_note"),
        subtitle=slide.get("subtitle"),
        **_common_kwargs(slide, meta, total),
    )


def _normalise_decision(d: dict) -> dict:
    """Translate plan-schema decision dict to the layout function's expected keys.

    Plan schema:  {tone, number, owner, text}
    Layout reads: {tone, number, owner, title (opt), desc/body/description}

    Maps "text" → "body" whenever body/desc/description are absent, regardless
    of whether "title" is present.  Explicit body/desc/description keys are left
    untouched so plans already using internal naming continue to work.
    """
    d = dict(d)
    has_body = bool(d.get("desc") or d.get("body") or d.get("description"))
    if not has_body and d.get("text") is not None:
        d["body"] = d.pop("text")
    return d


def _build_decision_rows(prs, meta, slide, total):
    c = _content(slide)
    decisions = [_normalise_decision(d) for d in c.get("decisions", [])]
    return build_decision_rows(
        prs, meta,
        action_title=slide.get("action_title", ""),
        decisions=decisions,
        takeaway=slide.get("takeaway", ""),
        source=slide.get("source"),
        subtitle=slide.get("subtitle"),
        methodology_note=c.get("methodology_note"),
        **_common_kwargs(slide, meta, total),
    )


def _build_gantt(prs, meta, slide, total):
    c = _content(slide)
    return build_gantt(
        prs, meta,
        action_title=slide.get("action_title", ""),
        quarters=c.get("quarters", []),
        rows=c.get("rows", []),
        takeaway=slide.get("takeaway", ""),
        source=slide.get("source"),
        methodology_note=c.get("methodology_note"),
        **_common_kwargs(slide, meta, total),
    )


def _build_hero_stat(prs, meta, slide, total):
    c = _content(slide)
    return build_hero_stat(
        prs, meta,
        stat=c.get("stat", {}),
        context=c.get("context", ""),
        source=c.get("source") or slide.get("source"),
        style=_resolve_style(slide, meta),
        page=slide.get("page"),
        total=total,
        category=_resolve_category(slide),
    )


def _build_close(prs, meta, slide, total):
    c = _content(slide)
    return build_close(
        prs, meta,
        action_statement=c.get("action_statement") or slide.get("action_title"),
        style=_resolve_style(slide, meta),
        page=slide.get("page"),
        total=total,
        takeaway=slide.get("takeaway"),
        source=slide.get("source"),
    )


def _build_stat_strip(prs, meta, slide, total):
    c = _content(slide)
    return build_stat_strip(
        prs, meta,
        action_title=slide.get("action_title", ""),
        stats=c.get("stats", []),
        takeaway=slide.get("takeaway", ""),
        source=slide.get("source"),
        subtitle=slide.get("subtitle"),
        methodology_note=c.get("methodology_note"),
        **_common_kwargs(slide, meta, total),
    )


def _build_before_after(prs, meta, slide, total):
    c = _content(slide)
    before = c.get("before") or {}
    after  = c.get("after")  or {}
    # Labels can live at content level (before_label / after_label) OR nested
    # inside the before/after dicts as "label".  Content-level takes precedence
    # so that plans which explicitly set before_label continue to work.
    before_label = (c.get("before_label")
                    or (before.get("label") if isinstance(before, dict) else None)
                    or "TODAY")
    after_label  = (c.get("after_label")
                    or (after.get("label") if isinstance(after, dict) else None)
                    or "TOMORROW")
    return build_before_after(
        prs, meta,
        action_title=slide.get("action_title", ""),
        before=before,
        after=after,
        takeaway=slide.get("takeaway", ""),
        source=slide.get("source"),
        subtitle=slide.get("subtitle"),
        before_label=before_label,
        after_label=after_label,
        methodology_note=c.get("methodology_note"),
        **_common_kwargs(slide, meta, total),
    )


_MILESTONE_STATUS_ALIASES = {
    "upcoming": "future",
    "future":   "future",
    "pending":  "future",
    "planned":  "future",
    "completed": "done",
    "complete":  "done",
    "done":      "done",
    "finished":  "done",
    "active":    "current",
    "current":   "current",
    "in_progress": "current",
    "inprogress":  "current",
}


def _normalise_milestone(m: dict) -> dict:
    """Translate plan-schema milestone dict to the layout function's expected keys.

    Plan schema:  {date, label, description, status}
    Layout reads: {date, title,  body,        status ("done"|"current"|"future")}
    """
    m = dict(m)
    # Key aliases: label → title, description → body
    if "title" not in m and "label" in m:
        m["title"] = m.pop("label")
    if "body" not in m and "description" in m:
        m["body"] = m.pop("description")
    # Status normalisation
    raw_status = str(m.get("status", "future")).lower().replace(" ", "_")
    m["status"] = _MILESTONE_STATUS_ALIASES.get(raw_status, "future")
    return m


def _build_milestone_timeline(prs, meta, slide, total):
    c = _content(slide)
    milestones = [_normalise_milestone(m) for m in c.get("milestones", [])]
    return build_milestone_timeline(
        prs, meta,
        action_title=slide.get("action_title", ""),
        milestones=milestones,
        takeaway=slide.get("takeaway", ""),
        source=slide.get("source"),
        subtitle=slide.get("subtitle"),
        methodology_note=c.get("methodology_note"),
        **_common_kwargs(slide, meta, total),
    )


_STATUS_TABLE_COL_ORDER = [
    "program", "phase", "status", "rag", "health",
    "milestone", "owner", "due", "date", "comment", "notes",
]


def _derive_status_columns(rows: list) -> list:
    """Derive column headers from the first row's keys when columns are absent.

    Key ordering follows _STATUS_TABLE_COL_ORDER so common fields (program,
    phase, status …) always appear in a sensible sequence.  Unknown keys are
    appended alphabetically after the known ones.
    """
    if not rows:
        return []
    first = rows[0] if isinstance(rows[0], dict) else {}
    keys = list(first.keys())
    ordered = [k for k in _STATUS_TABLE_COL_ORDER if k in keys]
    remainder = sorted(k for k in keys if k not in ordered)
    return [k.replace("_", " ").title() for k in (ordered + remainder)]


def _build_status_table(prs, meta, slide, total):
    c = _content(slide)
    rows = c.get("rows", [])
    columns = c.get("columns") or _derive_status_columns(rows)
    return build_status_table(
        prs, meta,
        action_title=slide.get("action_title", ""),
        columns=columns,
        rows=rows,
        takeaway=slide.get("takeaway", ""),
        source=slide.get("source"),
        subtitle=slide.get("subtitle"),
        methodology_note=c.get("methodology_note"),
        **_common_kwargs(slide, meta, total),
    )


def _build_hub_spoke(prs, meta, slide, total):
    c = _content(slide)
    return build_hub_spoke(
        prs, meta,
        action_title=slide.get("action_title", ""),
        hub=c.get("hub", {}),
        spokes=c.get("spokes", []),
        takeaway=slide.get("takeaway", ""),
        source=slide.get("source"),
        subtitle=slide.get("subtitle"),
        methodology_note=c.get("methodology_note"),
        **_common_kwargs(slide, meta, total),
    )


def _build_pillar_detail(prs, meta, slide, total):
    c = _content(slide)
    return build_pillar_detail(
        prs, meta,
        action_title=slide.get("action_title", ""),
        pillar_number=c.get("pillar_number", ""),
        pillar_label=c.get("pillar_label", "PILLAR"),
        owner=c.get("owner"),
        sections=c.get("sections", []),
        takeaway=slide.get("takeaway", ""),
        source=slide.get("source"),
        subtitle=slide.get("subtitle"),
        methodology_note=c.get("methodology_note"),
        **_common_kwargs(slide, meta, total),
    )


def _build_four_column(prs, meta, slide, total):
    c = _content(slide)
    return build_four_column(
        prs, meta,
        action_title=slide.get("action_title", ""),
        columns=c.get("columns", []),
        takeaway=slide.get("takeaway", ""),
        source=slide.get("source"),
        subtitle=slide.get("subtitle"),
        methodology_note=c.get("methodology_note"),
        **_common_kwargs(slide, meta, total),
    )


def _build_label_rows(prs, meta, slide, total):
    c = _content(slide)
    return build_label_rows(
        prs, meta,
        action_title=slide.get("action_title", ""),
        rows=c.get("rows", []),
        takeaway=slide.get("takeaway", ""),
        source=slide.get("source"),
        subtitle=slide.get("subtitle"),
        label_color=c.get("label_color"),
        methodology_note=c.get("methodology_note"),
        **_common_kwargs(slide, meta, total),
    )


def _build_circular_flow(prs, meta, slide, total):
    c = _content(slide)
    return build_circular_flow(
        prs, meta,
        action_title=slide.get("action_title", ""),
        phases=c.get("phases", []),
        takeaway=slide.get("takeaway", ""),
        source=slide.get("source"),
        subtitle=slide.get("subtitle"),
        methodology_note=c.get("methodology_note"),
        **_common_kwargs(slide, meta, total),
    )


def _build_org_chart(prs, meta, slide, total):
    c = _content(slide)
    return build_org_chart(
        prs, meta,
        action_title=slide.get("action_title", ""),
        root=c.get("root", {}),
        children=c.get("children", []),
        takeaway=slide.get("takeaway", ""),
        source=slide.get("source"),
        subtitle=slide.get("subtitle"),
        methodology_note=c.get("methodology_note"),
        **_common_kwargs(slide, meta, total),
    )


def _build_topic_set(prs, meta, slide, total):
    c = _content(slide)
    return build_topic_set(
        prs, meta,
        action_title=slide.get("action_title", ""),
        topics=c.get("topics", []),
        takeaway=slide.get("takeaway", ""),
        source=slide.get("source"),
        subtitle=slide.get("subtitle"),
        methodology_note=c.get("methodology_note"),
        **_common_kwargs(slide, meta, total),
    )


def _build_arrow_chain(prs, meta, slide, total):
    c = _content(slide)
    return build_arrow_chain(
        prs, meta,
        action_title=slide.get("action_title", ""),
        steps=c.get("steps", []),
        consequence=c.get("consequence"),
        takeaway=slide.get("takeaway", ""),
        source=slide.get("source"),
        subtitle=slide.get("subtitle"),
        methodology_note=c.get("methodology_note"),
        content=c,
        **_common_kwargs(slide, meta, total),
    )


def _build_pull_quote(prs, meta, slide, total):
    c = _content(slide)
    return build_pull_quote(
        prs, meta,
        quote=c.get("quote") or slide.get("quote"),
        attribution=c.get("attribution") or slide.get("attribution"),
        context=c.get("context"),
        action_title=slide.get("action_title"),
        **_common_kwargs(slide, meta, total),
    )


def _build_donut_chart(prs, meta, slide, total):
    c = _content(slide)
    return build_donut_chart(
        prs, meta,
        action_title=slide.get("action_title", ""),
        segments=c.get("segments", []),
        center_value=c.get("center_value"),
        center_label=c.get("center_label"),
        legend_title=c.get("legend_title"),
        **_common_kwargs(slide, meta, total),
    )


def _build_kpi_dashboard(prs, meta, slide, total):
    c = _content(slide)
    return build_kpi_dashboard(
        prs, meta,
        action_title=slide.get("action_title", ""),
        kpis=c.get("kpis", []),
        **_common_kwargs(slide, meta, total),
    )


def _build_icon_grid(prs, meta, slide, total):
    c = _content(slide)
    return build_icon_grid(
        prs, meta,
        action_title=slide.get("action_title", ""),
        items=c.get("items", []),
        columns=c.get("columns", 3),
        **_common_kwargs(slide, meta, total),
    )


def _normalise_journey_rows(raw_rows: list) -> tuple:
    """Normalise journey map rows to the layout function's {label, cells} format.

    Accepted input formats:
      A) Native layout format: [{label, cells: [str, ...]}]
      B) Rich step format:     [{actor, steps: [{stage, action, emotion}, ...]}]
         or                    [{name,  steps: [{stage, action, emotion}, ...]}]

    For Format B, cells are derived from step.action (with stage as fallback).
    Phase headers are derived from step.stage values of the first row when
    content.phases is absent.

    Returns (phases, rows) where phases is a list of str column headers and
    rows is a list of {label, cells} dicts.
    """
    phases_derived: list = []
    normalised: list = []
    for i, row in enumerate(raw_rows):
        if not isinstance(row, dict):
            continue
        # Determine row label — accept "label", "name", or "actor".
        label = (row.get("label") or row.get("name") or row.get("actor") or "")

        steps = row.get("steps")
        cells = row.get("cells")

        if steps is not None:
            # Rich step format: extract action text from each step.
            cells = [str(s.get("action") or s.get("stage") or "") for s in steps]
            if i == 0 and not phases_derived:
                phases_derived = [str(s.get("stage", "")) for s in steps]
        elif cells is None:
            cells = []

        normalised.append({"label": str(label), "cells": [str(c) for c in cells]})
    return phases_derived, normalised


def _build_journey_map(prs, meta, slide, total):
    c = _content(slide)
    raw_rows = c.get("rows") or c.get("actors") or []
    phases_from_content = c.get("phases") or []
    phases_derived, rows = _normalise_journey_rows(raw_rows)
    # Prefer explicitly declared phases; fall back to those derived from steps.
    phases = phases_from_content or phases_derived
    return build_journey_map(
        prs, meta,
        action_title=slide.get("action_title", ""),
        phases=phases,
        rows=rows,
        takeaway=slide.get("takeaway", ""),
        source=slide.get("source"),
        subtitle=slide.get("subtitle"),
        methodology_note=c.get("methodology_note"),
        **_common_kwargs(slide, meta, total),
    )


def _build_funnel(prs, meta, slide, total):
    c = _content(slide)
    return build_funnel(
        prs, meta,
        action_title=slide.get("action_title", ""),
        inputs=c.get("inputs", []),      # schema key: inputs (not stages)
        output=c.get("output"),
        **_common_kwargs(slide, meta, total),
    )


def _normalise_comparison_table(c: dict) -> tuple:
    """Return (options, features) normalised from whatever format the plan uses.

    Two accepted formats:
      A) Native layout format (used by LLM when it knows the schema):
           options:  ["Option A", "Option B", ...]
           features: [{"label": "Criterion", "values": ["yes", "no", ...], "highlighted": false}, ...]

      B) Simple matrix format (natural LLM/human output):
           headers: ["Criterion", "Option A", "Option B", ...]   ← first cell is the row-label column
           rows:    [["Timeline", "Full recovery", "Partial", ...], ...]

    Format B is normalised to Format A.  If both "options" and "headers" are
    present, "options" wins (backwards compat for plans already using Format A).
    """
    options  = c.get("options")
    features = c.get("features")
    if options or features:
        return list(options or []), list(features or [])

    headers = c.get("headers") or []
    rows    = c.get("rows") or []

    if not headers and not rows:
        return [], []

    # First header is the row-label column; remaining are option columns.
    derived_options  = [str(h) for h in headers[1:]] if len(headers) > 1 else []
    derived_features = []
    for row in rows:
        if isinstance(row, (list, tuple)):
            label  = str(row[0]) if row else ""
            values = [str(v) for v in row[1:]]
        elif isinstance(row, dict):
            label  = str(row.get("label", ""))
            values = [str(v) for v in row.get("values", [])]
        else:
            continue
        derived_features.append({"label": label, "values": values})
    return derived_options, derived_features


def _build_comparison_table(prs, meta, slide, total):
    c = _content(slide)
    options, features = _normalise_comparison_table(c)
    return build_comparison_table(
        prs, meta,
        action_title=slide.get("action_title", ""),
        options=options,
        features=features,
        **_common_kwargs(slide, meta, total),
    )


def _build_score_table(prs, meta, slide, total):
    c = _content(slide)
    return build_score_table(
        prs, meta,
        action_title=slide.get("action_title", ""),
        rows=c.get("rows", []),          # schema keys: rows + scale
        scale=c.get("scale", 5),
        scale_label=c.get("scale_label"),
        content=c,
        **_common_kwargs(slide, meta, total),
    )


def _build_influence_diagram(prs, meta, slide, total):
    c = _content(slide)
    return build_influence_diagram(
        prs, meta,
        action_title=slide.get("action_title", ""),
        center=c.get("center", {}),
        forces=c.get("forces", []),      # schema key: forces (not nodes)
        **_common_kwargs(slide, meta, total),
    )


def _build_word_cloud(prs, meta, slide, total):
    c = _content(slide)
    return build_word_cloud(
        prs, meta,
        action_title=slide.get("action_title", ""),
        words=c.get("words", []),
        **_common_kwargs(slide, meta, total),
    )


def _build_pyramid(prs, meta, slide, total):
    c = _content(slide)
    return build_pyramid(
        prs, meta,
        action_title=slide.get("action_title", ""),
        tiers=c.get("tiers", []),
        orientation=c.get("orientation", "up"),
        **_common_kwargs(slide, meta, total),
    )


def _build_venn(prs, meta, slide, total):
    c = _content(slide)
    return build_venn(
        prs, meta,
        action_title=slide.get("action_title", ""),
        circles=c.get("circles", []),
        intersection=c.get("intersection"),  # schema key: intersection
        **_common_kwargs(slide, meta, total),
    )


def _build_risk_heatmap(prs, meta, slide, total):
    c = _content(slide)
    return build_risk_heatmap(
        prs, meta,
        action_title=slide.get("action_title", ""),
        risks=c.get("risks", []),
        x_label=c.get("x_label", "LIKELIHOOD"),
        y_label=c.get("y_label", "IMPACT"),
        **_common_kwargs(slide, meta, total),
    )


def _build_radar_chart(prs, meta, slide, total):
    c = _content(slide)
    return build_radar_chart(
        prs, meta,
        action_title=slide.get("action_title", ""),
        axes=c.get("axes", []),
        series=c.get("series", []),
        **_common_kwargs(slide, meta, total),
    )


def _build_pros_cons(prs, meta, slide, total):
    c = _content(slide)
    return build_pros_cons(
        prs, meta,
        action_title=slide.get("action_title", ""),
        pros=c.get("pros", []),
        cons=c.get("cons", []),
        pros_label=c.get("pros_label", "ADVANTAGES"),
        cons_label=c.get("cons_label", "RISKS"),
        subject=c.get("subject"),
        **_common_kwargs(slide, meta, total),
    )


def _build_layered_stack(prs, meta, slide, total):
    c = _content(slide)
    return build_layered_stack(
        prs, meta,
        action_title=slide.get("action_title", ""),
        layers=c.get("layers", []),
        orientation=c.get("orientation", "vertical"),
        **_common_kwargs(slide, meta, total),
    )


def _build_photo_text(prs, meta, slide, total):
    c = _content(slide)
    return build_photo_text(
        prs, meta,
        action_title=slide.get("action_title", ""),
        image_path=c.get("image_path"),      # schema keys: image_path, image_side
        image_label=c.get("image_label"),
        title=c.get("title"),
        bullets=c.get("bullets", []),
        image_side=c.get("image_side", "left"),
        **_common_kwargs(slide, meta, total),
    )


def _build_fishbone(prs, meta, slide, total):
    c = _content(slide)
    return build_fishbone(
        prs, meta,
        action_title=slide.get("action_title", ""),
        effect=c.get("effect", ""),
        bones=c.get("bones", []),            # schema key: bones (not causes)
        takeaway=slide.get("takeaway", ""),
        source=slide.get("source"),
        subtitle=slide.get("subtitle"),
        methodology_note=c.get("methodology_note"),
        **_common_kwargs(slide, meta, total),
    )


def _build_key_question(prs, meta, slide, total):
    c = _content(slide)
    return build_key_question(
        prs, meta,
        action_title=slide.get("action_title", ""),
        question=c.get("question"),
        context=c.get("context"),
        content=c,
        **_common_kwargs(slide, meta, total),
    )


def _build_road_to_success(prs, meta, slide, total):
    c = _content(slide)
    return build_road_to_success(
        prs, meta,
        action_title=slide.get("action_title", ""),
        stages=c.get("stages"),
        milestones=c.get("milestones"),
        content=c,
        **_common_kwargs(slide, meta, total),
    )


def _build_columns_auto(prs, meta, slide, total):
    """Auto-dispatch to two/three/four_column based on column count."""
    c = _content(slide)
    cols = c.get("columns", [])
    n = len(cols)
    if n == 0:
        print(
            f"WARNING: 'columns' layout on page {slide.get('page', '?')} has no "
            f"column data — slide will render blank.",
            file=sys.stderr,
        )
    layout = "three_column" if n == 3 else ("two_column" if n <= 2 else "four_column")
    dispatch = {
        "two_column":   _build_two_column,
        "three_column": _build_three_column,
        "four_column":  _build_four_column,
    }
    patched = {**slide, "layout": layout}
    return dispatch[layout](prs, meta, patched, total)


_DISPATCH = {
    "cover":              _build_cover,
    "exec_summary":       _build_exec_summary,
    "agenda":             _build_agenda,
    "section_divider":    _build_section_divider,
    "chart_slide":        _build_chart_slide,
    "two_column":         _build_two_column,
    "three_column":       _build_three_column,
    "2x2_matrix":         _build_2x2_matrix,
    "phase_process":      _build_phase_process,
    "vertical_numbered":  _build_vertical_numbered,
    "waterfall_slide":    _build_waterfall_slide,
    "decision_rows":      _build_decision_rows,
    "gantt":              _build_gantt,
    "hero_stat":          _build_hero_stat,
    "close":              _build_close,
    "stat_strip":         _build_stat_strip,
    "before_after":       _build_before_after,
    "milestone_timeline": _build_milestone_timeline,
    "status_table":       _build_status_table,
    "hub_spoke":          _build_hub_spoke,
    "pillar_detail":      _build_pillar_detail,
    "four_column":        _build_four_column,
    "label_rows":         _build_label_rows,
    "circular_flow":      _build_circular_flow,
    "org_chart":          _build_org_chart,
    "topic_set":          _build_topic_set,
    "arrow_chain":        _build_arrow_chain,
    "columns":            _build_columns_auto,
    "pull_quote":         _build_pull_quote,
    "donut_chart":        _build_donut_chart,
    "kpi_dashboard":      _build_kpi_dashboard,
    "icon_grid":          _build_icon_grid,
    "journey_map":        _build_journey_map,
    "funnel":             _build_funnel,
    "comparison_table":   _build_comparison_table,
    "score_table":        _build_score_table,
    "influence_diagram":  _build_influence_diagram,
    "word_cloud":         _build_word_cloud,
    "pyramid":            _build_pyramid,
    "venn":               _build_venn,
    "risk_heatmap":       _build_risk_heatmap,
    "radar_chart":        _build_radar_chart,
    "pros_cons":          _build_pros_cons,
    "layered_stack":      _build_layered_stack,
    "photo_text":         _build_photo_text,
    "fishbone":           _build_fishbone,
    "key_question":       _build_key_question,
    "road_to_success":    _build_road_to_success,
}


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------
# Plan sanitization helpers
# ---------------------------------------------------------------------------

def _trunc(s: object, limit: int) -> object:
    """Truncate a string to `limit` chars at a word boundary, appending '…'."""
    if not isinstance(s, str) or len(s) <= limit:
        return s
    if limit <= 3:
        return s[:limit]
    # Cut at last space to avoid mid-word truncation (looks unprofessional on slides).
    candidate = s[:limit - 1]  # leave room for ellipsis char
    last_space = candidate.rfind(" ")
    if last_space > limit // 3:
        candidate = candidate[:last_space]
    return candidate.rstrip(",;: ") + "…"


def _sanitize_obj(obj: object, rules: dict) -> None:
    """Apply truncation rules to a dict, mutating in place."""
    if not isinstance(obj, dict):
        return
    for field, limit in rules.items():
        if isinstance(obj.get(field), str):
            obj[field] = _trunc(obj[field], limit)


# Per-field character limits (must align with rendering engine constraints).
# Top-level slide fields:
_SLIDE_LIMITS = {"action_title": 80}
# Content fields — applied at every dict level (top-level content and list items).
# Using one dict for both avoids divergence and closes the gap where list-item
# fields like "context" (used in kpi_dashboard) were not being truncated.
_CONTENT_FIELD_LIMITS = {
    "takeaway": 120,
    "body":     130,
    "desc":     160,
    "context":  130,
    "note":     80,
}


def _sanitize_content(content: object) -> None:
    """Recursively sanitize content dict: truncate known long-text fields."""
    if not isinstance(content, dict):
        return
    _sanitize_obj(content, _CONTENT_FIELD_LIMITS)
    for val in content.values():
        if isinstance(val, list):
            for item in val:
                if isinstance(item, dict):
                    _sanitize_obj(item, _CONTENT_FIELD_LIMITS)
                    _sanitize_content(item)  # recurse into nested dicts/lists
        elif isinstance(val, dict):
            _sanitize_content(val)


def _sanitize_plan(plan: dict) -> None:
    """Truncate overflow-prone fields across all slides before validation."""
    meta = plan.get("meta") or {}
    for slide in plan.get("slides") or []:
        # Takeaway lives both at slide level and inside content — handle first
        # so _sanitize_content below doesn't double-truncate it.
        for loc in (slide, slide.get("content") or {}):
            if isinstance(loc.get("takeaway"), str) and len(loc["takeaway"]) > 120:
                loc["takeaway"] = _trunc(loc["takeaway"], 120)

        # Cover slide: ensure a subtitle is present (either "Title; Subtitle" in
        # action_title, or a standalone subtitle field).  The LLM occasionally
        # omits it; auto-patch using the deck_label rather than hard-failing.
        # Do this BEFORE _sanitize_obj so the combined string is truncated together.
        if slide.get("layout") == "cover":
            raw = (slide.get("action_title") or "").rstrip().rstrip(";").rstrip()
            has_sub = (";" in (slide.get("action_title") or "") and
                       (slide.get("action_title") or "").split(";", 1)[1].strip()) or \
                      (slide.get("subtitle") or "").strip() or \
                      ((slide.get("content") or {}).get("subtitle") or "").strip()
            if not has_sub:
                deck_label = meta.get("deck_label") or "Merck Presentation"
                combined = (raw + "; " + deck_label) if raw else deck_label
                # Keep within 60-char cover-title limit (merck_layouts enforces this).
                slide["action_title"] = _trunc(combined, 60)

        _sanitize_obj(slide, _SLIDE_LIMITS)
        _sanitize_content(slide.get("content") or {})


# ---------------------------------------------------------------------------

def build_from_plan(plan, output_path, base_pptx: Optional[str] = None,
                    strict: bool = True) -> str:
    """
    Build a Merck PowerPoint deck from a slide plan.

    Parameters
    ----------
    plan : dict | str | pathlib.Path
        A plan dict, or a path to a JSON file.
    output_path : str | pathlib.Path
        Destination .pptx path.
    base_pptx : str | None
        Explicit path to a Merck theme .pptx. If None, selected from
        meta.region using the package-bundled templates.
    strict : bool
        If True (default), raise ValidationError on hard errors.

    Returns
    -------
    str
        Absolute path of the saved file.
    """
    if isinstance(plan, (str, pathlib.Path)):
        with open(plan, encoding="utf-8") as f:
            plan = json.load(f)

    # Security: validate output path ends with .pptx — sanity guard against
    # accidental overwrites of non-pptx files (e.g. output_path typo).
    output_path = pathlib.Path(output_path)
    if output_path.suffix.lower() != ".pptx":
        raise ValueError(
            f"output_path must end with .pptx, got: '{output_path.suffix}'"
        )

    # Auto-sanitize: truncate fields that overflow fixed-dimension text boxes in the
    # rendering engine. Truncation is a silent safety net — the LLM prompt instructs
    # the model to stay within these limits, but the build should not fail on overshoot.
    _sanitize_plan(plan)

    # Auto-fill agenda chapters before validation sees the plan.
    _autofill_agenda(plan)

    warnings = validate_plan(plan)
    for w in warnings:
        print(f"WARNING: {w}")

    meta = plan.get("meta") or {}
    slides = plan.get("slides") or []
    main_slides = [s for s in slides if not _is_appendix(s)]
    appendix_slides = [s for s in slides if _is_appendix(s)]
    total = len(main_slides)

    # Renumber pages so progress bar is consistent
    for i, s in enumerate(main_slides, start=1):
        s["page"] = i
    for i, s in enumerate(appendix_slides, start=1):
        s["page"] = f"A{i}"

    if base_pptx is None:
        region      = str(meta.get("region",   "eu")).lower().strip()
        color_theme = str(meta.get("color_theme") or "plastic").lower().strip()
        division    = str(meta.get("division")    or "merck").lower().strip()
        template_path = _resolve_template(region, color_theme, division)
    else:
        # Security: restrict base_pptx to the package templates directory so
        # a caller cannot force-load an arbitrary or maliciously crafted .pptx.
        bp = pathlib.Path(base_pptx).resolve()
        if not bp.is_file() or bp.suffix.lower() != ".pptx":
            raise ValueError(f"base_pptx must be an existing .pptx file: {bp}")
        try:
            bp.relative_to(_TEMPLATE_DIR.resolve())
        except ValueError:
            raise ValueError(
                f"base_pptx must be inside the package templates directory "
                f"({_TEMPLATE_DIR}). Got: {bp}"
            )
        template_path = bp
        color_theme = str(meta.get("color_theme") or "plastic").lower().strip()

    prs = open_deck(str(template_path))

    # Apply color theme programmatically — modifies the master's theme XML and
    # replaces hardcoded shape fill colors so ALL cover/divider blob shapes
    # render in the selected theme's accent colors.  Must happen before slides.
    _apply_color_theme(prs, color_theme)
    ordered = main_slides + appendix_slides

    for s in ordered:
        layout = s.get("layout")
        builder = _DISPATCH.get(layout)
        if builder is None:
            raise BuildError(
                f"page {s.get('page', '?')}: unknown layout {layout!r}"
            )

        is_apx = _is_appendix(s)
        if is_apx:
            if not s.get("category"):
                s["category"] = "APPENDIX"
            slide_total = None
        else:
            slide_total = total

        try:
            built_slide = builder(prs, meta, s, slide_total)
        except Exception as exc:
            raise BuildError(
                f"page {s.get('page', '?')} ({layout}): {exc}"
            ) from exc

        # Speaker notes
        notes = s.get("notes")
        if notes and built_slide is not None:
            try:
                add_speaker_notes(built_slide, notes)
            except Exception as exc:
                print(
                    f"WARNING: speaker notes not written for page {s.get('page')}: {exc}",
                    file=sys.stderr,
                )

        # Image placement
        img = _content(s).get("image")
        if img and built_slide is not None and img.get("path"):
            _place_image(built_slide, img)

    _wire_agenda_hyperlinks(prs, ordered)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    save_deck(prs, str(output_path))
    return str(output_path.resolve())
