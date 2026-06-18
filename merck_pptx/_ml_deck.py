from __future__ import annotations
import os
import sys
from typing import Optional
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from ._ml_constants import SLIDE_W, SLIDE_H, _rgb_tuple
from ._ml_primitives import _emu

# ===========================================================================
# Deck lifecycle
# ===========================================================================

def open_deck(base_path: Optional[str] = None) -> Presentation:
    """Open a Presentation. Reuses the Merck themed base pptx if it can be
    found at the given path or any common location; else falls back to a
    default python-pptx deck and warns to stderr.

    Common search paths tried when base_path is missing or not found:
      - the path passed in
      - Merck_Themed_Base_v1.pptx in the current working directory
      - /mnt/data/Merck_Themed_Base_v1.pptx (MyGPT sandbox upload location)
      - the directory containing merck_layouts.py
    """
    candidates = []
    if base_path:
        candidates.append(base_path)
    # Common fallback locations so a wrong-path call still finds the base.
    fname = "Merck_Themed_Base_v1.pptx"
    candidates.extend([
        fname,
        os.path.join("/mnt/data", fname),
        os.path.join(os.path.dirname(os.path.abspath(__file__)), fname),
    ])

    prs = None
    used = None
    for cand in candidates:
        if not cand:
            continue
        if os.path.isfile(cand):
            try:
                prs = Presentation(cand)
                used = cand
                # Strip pre-existing slides so the base provides only theme
                # and layouts; the caller adds slides from scratch.
                sldIdLst = prs.slides._sldIdLst  # noqa
                for sldId in list(sldIdLst):
                    rId = sldId.get(qn("r:id"))
                    try:
                        prs.part.drop_rel(rId)
                    except Exception:
                        pass
                    sldIdLst.remove(sldId)
                break
            except Exception:
                prs = None
                continue

    if prs is None:
        import sys
        print("WARNING: Merck_Themed_Base_v1.pptx not found in any search "
              "path. Falling back to default python-pptx layouts. Slides may "
              "look generic and PowerPoint may show a 'repairs' warning on "
              "open. Place the base pptx next to merck_layouts.py or pass "
              "the correct path to open_deck().",
              file=sys.stderr)
        prs = Presentation()

    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def save_deck(prs: Presentation, output_path: str) -> str:
    """Persist the deck and return the path."""
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    return output_path


def _blank_layout(prs: Presentation):
    """Pick the most placeholder-free layout we can find.

    Order of preference:
      1. ANY layout whose name is exactly "Blank" (case-insensitive).
         Merck_Themed_Base_v1.pptx has this at index 9; default python-pptx
         has it at index 6. Searching by name handles both cleanly.
      2. The layout with the FEWEST placeholders (avoids "Vertical Title and
         Text" and other heavy default layouts that trigger PowerPoint
         repair warnings when chrome is drawn on top).
      3. Last resort: the first layout. Never the last layout (which is the
         default python-pptx "Vertical Title and Text" and looked wrong on
         every slide of the prior run).
    """
    layouts = list(prs.slide_layouts)
    if not layouts:
        return None
    # 1. Search by name.
    for layout in layouts:
        if str(layout.name).strip().lower() == "blank":
            return layout
    # 2. Fewest placeholders.
    best = None
    best_count = None
    for layout in layouts:
        try:
            count = len(list(layout.placeholders))
        except Exception:
            continue
        if best_count is None or count < best_count:
            best = layout
            best_count = count
            if count == 0:
                return layout
    if best is not None:
        return best
    # 3. First layout (NOT last).
    return layouts[0]


def _new_slide(prs: Presentation, bg_color=None):
    slide = prs.slides.add_slide(_blank_layout(prs))
    if bg_color is not None:
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _rgb_tuple(bg_color)
    return slide


def _intro_layout(prs: Presentation):
    """Pick the themed cover layout from Merck_Themed_Base_v1.pptx.

    Returns the 'Title' layout (index 1 in the Merck base template) which
    provides the diagonal green/yellow cover design, the auto-disclaimer
    text, EMD logos, and three placeholders: title, subtitle, name/date.

    Falls back to None when running against a generic deck (no themed base)
    so build_cover can route to its chrome-from-scratch path.
    """
    if prs is None:
        return None
    # 1. Match by exact name 'Title' (Merck template).
    for layout in prs.slide_layouts:
        if str(layout.name).strip().lower() == "title":
            # Sanity: must have all 3 cover placeholders.
            try:
                placeholders = list(layout.placeholders)
            except Exception:
                continue
            has_title = any(
                ph.placeholder_format.idx == 0 for ph in placeholders
            )
            has_subtitle = any(
                ph.placeholder_format.idx == 1 for ph in placeholders
            )
            if has_title and has_subtitle:
                return layout
    # 2. No themed cover available — caller will fall back.
    return None


def _divider_layout(prs: Presentation):
    """Return the template's native 'Divider' layout, or None.

    The Merck base template includes a branded section-divider layout (index 10,
    named 'Divider') that carries the organic blob shapes and colour scheme.
    Using it makes dividers visually consistent with covers from the same template.

    Layout must have placeholder idx=0 (TITLE — section number) and idx=13
    (BODY — chapter title).  'Divider plain' (no number) is intentionally skipped.
    """
    if prs is None:
        return None
    for layout in prs.slide_layouts:
        name = str(layout.name).strip().lower()
        if name != "divider":          # skip "divider plain" and others
            continue
        try:
            placeholders = list(layout.placeholders)
        except Exception:
            continue
        has_title = any(ph.placeholder_format.idx == 0  for ph in placeholders)
        has_body  = any(ph.placeholder_format.idx == 13 for ph in placeholders)
        if has_title and has_body:
            return layout
    return None


def _cover_picture_layout(prs: Presentation):
    """Return the 'Title with picture' cover layout for EMD Electronics, or None.

    This layout (index 2 in the Merck base template, named 'Title with picture')
    includes a PICTURE placeholder (idx=20) that the user fills in PowerPoint
    after generation.  It is only used when color_theme == 'electronics'.
    """
    if prs is None:
        return None
    for layout in prs.slide_layouts:
        name = str(layout.name).strip().lower()
        if "picture" not in name:
            continue
        try:
            placeholders = list(layout.placeholders)
        except Exception:
            continue
        has_title    = any(ph.placeholder_format.idx == 0  for ph in placeholders)
        has_subtitle = any(ph.placeholder_format.idx == 1  for ph in placeholders)
        has_picture  = any(ph.placeholder_format.idx == 20 for ph in placeholders)
        if has_title and has_subtitle and has_picture:
            return layout
    return None


def add_image(slide, image_path, x, y, w=None, h=None):
    """Insert a user-provided image. Returns the picture shape, or None
    if the file is missing or unreadable. Aspect ratio preserved when only
    one of w/h is supplied.

    Path is validated to be a local file (not a URL) for safety. The
    canonical runner passes user-uploaded paths from `content.image.path`.
    """
    if not image_path:
        return None
    if not os.path.isfile(image_path):
        return None
    try:
        if w is not None and h is not None:
            return slide.shapes.add_picture(image_path,
                                            _emu(x), _emu(y),
                                            width=_emu(w), height=_emu(h))
        if w is not None:
            return slide.shapes.add_picture(image_path,
                                            _emu(x), _emu(y),
                                            width=_emu(w))
        if h is not None:
            return slide.shapes.add_picture(image_path,
                                            _emu(x), _emu(y),
                                            height=_emu(h))
        return slide.shapes.add_picture(image_path, _emu(x), _emu(y))
    except Exception:
        return None


def add_slide_jump_hyperlink(shape, target_slide):
    """Make a shape's text a clickable jump to target_slide. Works by
    creating a slide-jump relationship on the source slide part and
    wrapping each <a:r> text run with an <a:hlinkClick> element pointing
    at that relationship via r:id, plus action='ppaction://hlinksldjump'.

    Used by the canonical runner to make agenda chapter rows clickable.
    """
    from lxml import etree
    if not shape.has_text_frame:
        return
    src_part = shape.part
    try:
        rId = src_part.relate_to(
            target_slide.part,
            "http://schemas.openxmlformats.org/officeDocument/2006/"
            "relationships/slide"
        )
    except Exception:
        return
    tx_body = shape.text_frame._txBody
    for r in tx_body.iter(qn("a:r")):
        rPr = r.find(qn("a:rPr"))
        if rPr is None:
            rPr = etree.SubElement(r, qn("a:rPr"))
            # rPr must be the FIRST child of <a:r> per OOXML schema.
            r.insert(0, rPr)
        # Remove any existing hlinkClick before adding ours.
        for existing in rPr.findall(qn("a:hlinkClick")):
            rPr.remove(existing)
        hlink = etree.SubElement(rPr, qn("a:hlinkClick"))
        hlink.set(qn("r:id"), rId)
        hlink.set("action", "ppaction://hlinksldjump")


def add_speaker_notes(slide, notes):
    """Add speaker notes to a slide. Multi-line text becomes multiple
    paragraphs. Safe to call with None / empty string (no-op).

    Used by the canonical runner — agents emit `notes` per slide and the
    runner lifts them through to PowerPoint's Notes pane for presenters.
    """
    if not notes:
        return
    text = str(notes).strip()
    if not text:
        return
    try:
        notes_tf = slide.notes_slide.notes_text_frame
    except Exception:
        return
    notes_tf.text = ""
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = notes_tf.paragraphs[0] if i == 0 else notes_tf.add_paragraph()
        p.text = line


def _populate_placeholder(layout_ph_idx, slide, text, *, font=None, sz=None,
                          color=None, bold=None, italic=None):
    """Find a placeholder by idx on the slide and set its text + formatting.

    Returns the placeholder shape (or None if not found). When text contains
    newlines, each line becomes a separate paragraph.
    """
    if text is None:
        return None
    target = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == layout_ph_idx:
            target = ph
            break
    if target is None or not target.has_text_frame:
        return None
    tf = target.text_frame
    tf.clear()
    lines = str(text).split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        if font is not None:
            run.font.name = font
        if sz is not None:
            run.font.size = Pt(sz)
        if bold is not None:
            run.font.bold = bool(bold)
        if italic is not None:
            run.font.italic = bool(italic)
        if color is not None:
            run.font.color.rgb = _rgb_tuple(color)
    return target


