from __future__ import annotations
import math
import sys
from typing import Optional
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt
from ._ml_constants import (
    FONT_HEAD, FONT_BODY,
    MERCK_PURPLE, MERCK_GOLD, MERCK_YELLOW, MERCK_BLUE, MERCK_AQUA,
    WHITE, INK_DARK, INK_GRAY, LIGHT_GRAY, PANEL_LIGHT,
    PURPLE_DEEP, PURPLE_MUTED, BAD_RED, GOOD_GREEN,
    ACT_PURPLE, LY_CYAN, OP_LIME, FC_PINK, DEV_POS_BLUE, DEV_NEG_RED,
    CHART_PALETTE, PHASE_1_COLOR, PHASE_2_COLOR, PHASE_3_COLOR,
    CONTENT_X, CONTENT_Y, CONTENT_Y_SUBTITLE, CONTENT_W, CONTENT_H,
    SOURCE_Y, SOURCE_H, TAKEAWAY_Y, TAKEAWAY_H, PHASE_Y, PHASE_H,
    FOOTER_Y, FOOTER_H, FOOTER_TEXT_Y,
    TITLE_X, TITLE_Y_NUMBERED, TITLE_Y_UNNUMBERED, TITLE_W, TITLE_H,
    SUB_X, SUB_W, SUB_H, SUB_GAP,
    AUTO_PROMOTE_EXECUTIVE,
    _palette_for, _rgb_tuple, _is_dark,
)
from ._ml_primitives import (
    rect, rounded, oval, circle, line, hairline, txt, _add_run,
    _freeform_poly, _emu, _apply_fill, _apply_border,
)
from ._ml_icons import draw_icon, ICON_REGISTRY
from ._ml_deck import (
    _new_slide, _intro_layout, _divider_layout,
    _cover_picture_layout, _populate_placeholder,
)
from ._ml_chrome import (
    apply_chrome, _phase_progress, _draw_card, _bulleted_list,
    stub_and_flag, footnotes_block, _section_marker,
    _top_chrome, _bottom_chrome, _gold_square_bullet,
    _tracked, _track_letters, _format_section_number, _pad_int,
    _render_action_title, _source_line,
    statement_card, in_slide_section,
    _takeaway_band, _superscript,
)
from ._ml_charts import (
    add_slope_chart, add_dot_plot, add_marimekko, add_waterfall,
    add_small_multiples, add_simple_bar, _render_chart,
)
from ._ml_helpers import _style_or_promote, _tone_color, _rag_color, _norm_key

# ===========================================================================
# Layout: COVER
# ===========================================================================

def _draw_cover_keymessages_grid(slide, cards, grid_top, grid_h, card_style):
    """N-aware key_messages grid for cover slides.

    - 1 card: full-width single row
    - 2 cards: two columns, one row
    - 3 cards: three columns, one row (centered, not 2x2 with orphan)
    - 4 cards: 2x2 grid

    Avoids the asymmetric look of 3 cards in a 2x2 grid (top-left, top-right,
    bottom-left, with the bottom-right empty).
    """
    if not cards:
        return
    cards = list(cards)[:4]
    n = len(cards)
    gap = Inches(0.30)
    pad = Inches(0.20)
    left = Inches(0.65)

    if n == 4:
        cols, rows = 2, 2
    elif n == 3:
        cols, rows = 3, 1
    elif n == 2:
        cols, rows = 2, 1
    else:
        cols, rows = 1, 1

    col_w = (Inches(12.0) - gap * (cols - 1)) / cols
    row_h = (grid_h - gap * (rows - 1)) / rows

    for i, m in enumerate(cards):
        r = i // cols
        c = i % cols
        cx = left + c * (col_w + gap)
        cy = grid_top + r * (row_h + gap)
        _draw_card(slide, cx, cy, col_w, row_h, card_style, highlighted=False)
        _gold_square_bullet(slide, cx + pad, cy + Inches(0.22))
        title_box = slide.shapes.add_textbox(
            cx + pad + Inches(0.20), cy + Inches(0.15),
            col_w - pad * 2 - Inches(0.25), Inches(0.36))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.0)
        tf.margin_bottom = Inches(0.0)
        p = tf.paragraphs[0]
        _add_run(p, m.get("label", ""), sz=14, color=MERCK_PURPLE,
                 bold=True, font=FONT_BODY)
        txt(slide, cx + pad, cy + Inches(0.55), col_w - pad * 2,
            row_h - Inches(0.65),
            m.get("body", ""), sz=11, color=INK_GRAY, font=FONT_BODY)


def _draw_cover_authors(slide, authors, dark=False):
    """Stack a list of {name, title} byline entries below the subtitle area.

    Authors render one per line from y ~5.20 to y ~6.10. Name in bold (white
    on dark covers, INK_DARK on light), title in lighter color on the same
    line, separated by a small gap.
    """
    if not authors:
        return
    rows = list(authors)[:5]
    line_h = Inches(0.32)
    start_y = Inches(5.20)
    if dark:
        name_color = WHITE
        title_color = PANEL_LIGHT
    else:
        name_color = INK_DARK
        title_color = INK_GRAY
    for i, a in enumerate(rows):
        ry = start_y + line_h * i
        name = a.get("name", "") if isinstance(a, dict) else str(a)
        atitle = a.get("title", "") if isinstance(a, dict) else ""
        box = slide.shapes.add_textbox(Inches(0.65), ry, Inches(12.0), line_h)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.02)
        tf.margin_right = Inches(0.02)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        if name:
            _add_run(p, name, sz=12, color=name_color, bold=True,
                     font=FONT_BODY)
        if atitle:
            _add_run(p, "   ", sz=12, color=name_color, font=FONT_BODY)
            _add_run(p, atitle, sz=11, color=title_color, font=FONT_BODY)


def build_cover(prs, meta, title=None, subtitle="", style="merck_executive",
                key_messages=None, phases=None, action_title=None,
                page=None, total=None, section_number=None, category=None,
                authors=None, top_bar=False, methodology_note=None,
                content=None, color_theme=None):
    """Cover slide.

    Preferred path: when the Merck themed base is loaded, this uses the
    'Title' layout (index 1) directly. The layout provides the diagonal
    green/yellow design, EMD/MilliporeSigma/EMD-Electronics logos, the auto
    disclaimer ('The businesses of Merck KGaA...'), and the classification
    badge. The function populates three placeholders:
      - Title 1 (idx 0): action_title (keep short, ≤60 chars)
      - Subtitle 2 (idx 1): subtitle (main subtitle, multi-line OK)
      - Name/Date (idx 10): name + month_year on two lines

    Electronics theme path: uses the 'Title with picture' layout (index 2).
    The PICTURE placeholder (idx 20) is left empty so the user can add their
    own image in PowerPoint after generation.

    Fallback path: when no themed layout is found (generic python-pptx
    deck), draws the cover chrome from scratch — the legacy behavior.
    key_messages / phases / authors-as-grid render only in fallback mode.

    title (also accepted as action_title) may be a string OR a list of
    (text, italic_bool) tuples — italic runs render in MERCK_YELLOW (only
    honored on the fallback path; the themed layout uses its own type).
    """
    if content:
        if "authors"      in content: authors      = content["authors"]
        if "key_messages" in content: key_messages = content["key_messages"]
        if "phases"       in content: phases       = content["phases"]
        if "title"        in content: title        = content["title"]
        if "subtitle"     in content: subtitle     = content["subtitle"]
    if title is None:
        title = action_title if action_title is not None else ""

    # For Electronics theme, try the 'Title with picture' layout first.
    theme_lower = str(color_theme or "").lower()
    if theme_lower == "electronics":
        pic_layout = _cover_picture_layout(prs)
        if pic_layout is not None:
            import sys
            print(
                "INFO: Electronics cover — image placeholder (idx 20) left "
                "empty for manual editing in PowerPoint.",
                file=sys.stderr,
            )
            intro_layout = pic_layout
        else:
            intro_layout = _intro_layout(prs)
    else:
        intro_layout = _intro_layout(prs)

    # Themed-layout path: drop the visual chrome onto the template.
    if intro_layout is not None:
        slide = prs.slides.add_slide(intro_layout)

        # Dark color themes (synthetic, electronics) need a dark slide background.
        # The template's "light panel" freeform is made transparent by
        # _apply_color_theme(), so the slide's own background fills the gap.
        # This must be set before placeholders are populated.
        if theme_lower in ("synthetic", "electronics"):
            try:
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = _rgb_tuple(MERCK_PURPLE)
            except Exception:
                pass

        # Flatten tuple titles to a single string; the template's Title 1
        # uses its own typography so we don't try to honor italic runs.
        if isinstance(title, (list, tuple)) and title and \
           all(isinstance(seg, (list, tuple)) and len(seg) == 2
               for seg in title):
            title_text = "".join(seg[0] for seg in title)
        else:
            title_text = str(title) if title else ""

        # Split title at the schema-convention ';' separator: everything before
        # the first ';' is the deck title; everything after is the subtitle.
        # This prevents long combined strings from overflowing the title box.
        title_part, _, subtitle_part = title_text.partition(";")
        title_part    = title_part.strip()
        subtitle_part = subtitle_part.strip()

        ph_title = _populate_placeholder(0, slide, title_part)
        if ph_title is not None and ph_title.has_text_frame:
            try:
                from pptx.enum.text import MSO_AUTO_SIZE
                ph_title.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except Exception:
                pass

        # Theme-dependent body text color for subtitle / name-date placeholders.
        # The USA template can inherit teal from its theme master; explicitly
        # override so these placeholders always use a CD-compliant neutral.
        _cover_dark = theme_lower in ("synthetic", "electronics")
        _body_color = PANEL_LIGHT if _cover_dark else INK_GRAY

        # Populate subtitle placeholder (idx 1) with subtitle_part if the
        # slide had no explicit subtitle and the title contained a ';' split.
        if subtitle_part and not subtitle:
            subtitle = subtitle_part
        if subtitle:
            ph_sub = _populate_placeholder(1, slide, subtitle,
                                           font=FONT_BODY, color=_body_color)
            if ph_sub is not None and ph_sub.has_text_frame:
                ph_sub.text_frame.word_wrap = True
                try:
                    from pptx.enum.text import MSO_AUTO_SIZE
                    ph_sub.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                except Exception:
                    pass
        # Name/Date placeholder (idx 10): two lines (author name; month/year).
        name_line = ""
        if authors:
            first = authors[0] if isinstance(authors, list) else authors
            if isinstance(first, dict):
                name_line = first.get("name", "")
        if not name_line:
            name_line = (meta or {}).get("author") or ""
        date_line = (meta or {}).get("month_year", "")
        nd_text = "\n".join(s for s in [name_line, date_line] if s)
        if nd_text:
            _populate_placeholder(10, slide, nd_text,
                                  font=FONT_BODY, color=_body_color)
        return slide

    # Fallback: no themed base. Build the cover from scratch (legacy path).
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])

    if style == "merck_storytelling":
        # Dark cover: full-bleed purple + top chrome + huge italic hero.
        _top_chrome(slide, meta, category, style,
                    top_bar=top_bar or bool((meta or {}).get("cover_top_bar")))
        # Hero title.
        _render_action_title(slide, Inches(0.65), Inches(1.50),
                             Inches(12.0), Inches(1.60), title, style,
                             size=44, italic_color=MERCK_YELLOW,
                             base_color=WHITE)
        # Subtitle.
        if subtitle:
            txt(slide, Inches(0.65), Inches(3.10), Inches(12.0), Inches(0.50),
                subtitle, sz=15, color=PANEL_LIGHT, italic=True,
                font=FONT_BODY)

        # N-aware key_messages grid (1/2/3/4 cards).
        if key_messages:
            _draw_cover_keymessages_grid(slide, key_messages,
                                         Inches(4.10), Inches(2.10),
                                         "merck_executive")

        # Authors byline (stacked, name bold + title in lighter color).
        if authors:
            _draw_cover_authors(slide, authors, dark=True)

        # Phase progress strip.
        if phases:
            _phase_progress(slide, phases, style)

        # Bottom footer band.
        _bottom_chrome(slide, meta, None, None, None, style)

    else:
        # Light cover: white bg.
        _top_chrome(slide, meta, category, style,
                    top_bar=top_bar or bool((meta or {}).get("cover_top_bar")))
        # Hero title in purple bold (with optional italic emphasis in yellow).
        _render_action_title(slide, Inches(0.65), Inches(1.80),
                             Inches(12.0), Inches(1.70), title, style,
                             size=44, italic_color=MERCK_YELLOW,
                             base_color=MERCK_PURPLE)
        # Theme-accent rule under title.
        hairline(slide, Inches(0.65), Inches(3.55), Inches(2.6),
                 Emu(int(Pt(2.5))), pal["highlight"])
        # Subtitle gray italic.
        if subtitle:
            txt(slide, Inches(0.65), Inches(3.75), Inches(12.0), Inches(0.50),
                subtitle, sz=16, color=INK_GRAY, italic=True, font=FONT_BODY)

        # N-aware key_messages grid (1/2/3/4 cards).
        if key_messages:
            _draw_cover_keymessages_grid(slide, key_messages,
                                         Inches(4.20), Inches(1.80), style)

        # Authors byline (stacked, name bold + title in lighter color).
        if authors:
            _draw_cover_authors(slide, authors, dark=False)

        # Deck label and month/year strip near bottom.
        deck_label = (meta or {}).get("deck_label", "")
        month_year = (meta or {}).get("month_year", "")
        if deck_label or month_year:
            parts = []
            if deck_label:
                parts.append(_tracked(deck_label))
            if month_year:
                parts.append(_tracked(month_year))
            sep = "   •   "
            txt(slide, Inches(0.65), Inches(5.95), Inches(12.0), Inches(0.30),
                sep.join(parts), sz=10, color=PURPLE_MUTED, bold=True,
                font=FONT_BODY)

        # Phase strip.
        if phases:
            _phase_progress(slide, phases, style)

        # Footer band.
        _bottom_chrome(slide, meta, None, None, None, style)

    return slide


# ===========================================================================
# Layout: EXEC SUMMARY
# ===========================================================================

def build_exec_summary(prs, meta, action_title=None, key_messages=None, takeaway=None,
                       source=None, page=None, total=None, style=None,
                       category=None, section_number=None,
                       methodology_note=None, subtitle=None, content=None):
    """Executive Summary: numbered rows with label + body.
    Auto-promoted to executive style (unless overridden by an exec-eligible
    category). Category defaults to 'Executive Summary'.
    """
    if content:
        if "key_messages" in content: key_messages = content["key_messages"]
        if "takeaway"     in content: takeaway     = content["takeaway"]
    style = _style_or_promote(category or "Executive Summary", style)
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    apply_chrome(slide, meta, action_title,
                 category="Executive Summary",
                 subtitle=subtitle,
                 takeaway=takeaway,
                 source=source, page=page, total=total, palette=style,
                 section_number=section_number,
                 methodology_note=methodology_note)

    items = (key_messages or [])[:5]
    n = max(len(items), 1)
    # Content area below the title block.
    zone_top = Inches(2.95) if subtitle else Inches(2.55)
    zone_bot = Inches(6.50)
    zone_h = zone_bot - zone_top
    row_gap = Inches(0.10)
    row_h = (zone_h - row_gap * (n - 1)) / n
    for i, msg in enumerate(items):
        ry = zone_top + i * (row_h + row_gap)
        num_w = Inches(0.65)
        txt(slide, Inches(0.65), ry + Inches(0.02), num_w, row_h,
            _pad_int(i + 1), sz=24, color=MERCK_PURPLE, bold=True,
            font=FONT_BODY, anchor=MSO_ANCHOR.TOP)
        hairline(slide, Inches(0.65) + num_w + Inches(0.10),
                 ry + row_h - Emu(int(Pt(0.5))),
                 Inches(12.0) - num_w - Inches(0.10),
                 Emu(int(Pt(0.5))), LIGHT_GRAY)
        label = msg.get("label", "")
        body = msg.get("body", "")
        body_x = Inches(0.65) + num_w + Inches(0.20)
        body_w = Inches(12.0) - num_w - Inches(0.30)
        if label:
            txt(slide, body_x, ry + Inches(0.02), body_w, Inches(0.30),
                label, sz=13, color=MERCK_PURPLE, bold=True, font=FONT_BODY)
        txt(slide, body_x, ry + Inches(0.34), body_w, row_h - Inches(0.40),
            body, sz=10, color=INK_GRAY, font=FONT_BODY)
    return slide


# ===========================================================================
# Layout: AGENDA
# ===========================================================================

def build_agenda(prs, meta, chapters=None, style="merck_executive",
                 page=None, total=None, action_title=None,
                 section_number=None, category=None,
                 methodology_note=None, content=None):
    # action_title is honored when provided; otherwise the title defaults to
    # "INDEX". methodology_note is accepted for caller consistency and ignored.
    """Agenda / chapters slide (no section marker; INDEX-style title)."""
    if content:
        if "chapters" in content: chapters = content["chapters"]
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    _top_chrome(slide, meta, None, style, page=page, total=total)
    _bottom_chrome(slide, meta, "Index", page, total, style)

    # Big "INDEX" title (or supplied action_title).
    dark = _is_dark(style)
    title_color = WHITE if dark else MERCK_PURPLE
    title_text = action_title if action_title else "INDEX"
    heading_box = txt(slide, Inches(0.65), Inches(1.30), Inches(12.0), Inches(1.20),
                      str(title_text).upper(), sz=44, color=title_color, bold=True,
                      font=FONT_HEAD)
    try:
        from pptx.enum.text import MSO_AUTO_SIZE
        heading_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass

    items = (chapters or [])[:12]
    if not items:
        return slide
    # Two-column layout (left half + right half).
    half = len(items) // 2 + len(items) % 2
    col_top = Inches(2.80)  # heading box ends at 1.30+1.20=2.50; 0.30 gap
    col_h = Inches(3.95)
    col_w = (Inches(12.0) - Inches(0.50)) / 2
    row_gap = Inches(0.12)
    rows_left = items[:half]
    rows_right = items[half:]

    def _draw_col(rows, x0):
        if not rows:
            return
        n = len(rows)
        compact = (n >= 5)
        title_sz = 12 if compact else 14
        sub_sz = 9 if compact else 10
        row_h = (col_h - row_gap * (n - 1)) / n
        for i, ch in enumerate(rows):
            ry = col_top + i * (row_h + row_gap)
            number = ch.get("number") or _pad_int(i + 1)
            title = ch.get("title", "")
            sub = ch.get("subtitle", "")
            # Number in purple bold. Named with the chapter number so the
            # canonical runner can attach a slide-jump hyperlink in a
            # post-pass after all target slides exist.
            num_color   = MERCK_PURPLE if not dark else pal["hot"]
            num_shape = txt(slide, x0, ry, Inches(0.60), Inches(0.40),
                str(number), sz=title_sz, color=num_color, bold=True,
                font=FONT_BODY)
            try:
                num_shape.name = f"AgendaChapter_{number}"
            except Exception:
                pass
            # Title in INK_DARK bold (also named for hyperlink targeting).
            # Height is dynamic so long titles can wrap to 2 lines without clipping.
            title_h = min(row_h * 0.60, Inches(0.52))
            title_shape = txt(slide, x0 + Inches(0.55), ry,
                col_w - Inches(0.55),
                title_h,
                title, sz=title_sz,
                color=INK_DARK if not dark else WHITE,
                bold=True, font=FONT_BODY)
            try:
                title_shape.name = f"AgendaChapter_{number}"
            except Exception:
                pass
            if sub:
                sub_y = ry + title_h + Inches(0.02)
                sub_h = max(row_h - title_h - Inches(0.06), Inches(0.18))
                txt(slide, x0 + Inches(0.55), sub_y,
                    col_w - Inches(0.55),
                    sub_h,
                    sub, sz=sub_sz,
                    color=INK_GRAY if not dark else PANEL_LIGHT,
                    italic=True, font=FONT_BODY)
    _draw_col(rows_left, Inches(0.65))
    _draw_col(rows_right, Inches(0.65) + col_w + Inches(0.50))
    return slide


# ===========================================================================
# Layout: SECTION DIVIDER
# ===========================================================================

def build_section_divider(prs, meta, number=None, title=None, style="merck_executive",
                          page=None, total=None, action_title=None,
                          section_number=None, category=None,
                          takeaway=None, source=None,
                          methodology_note=None, content=None):
    """Section divider with a section number and chapter title.

    Preferred path: uses the template's native 'Divider' layout (branded
    organic blob shapes, colour scheme from the theme template file).  The
    section number goes into placeholder idx=0 and the chapter title into
    idx=13.

    Fallback path (no 'Divider' layout found): draws the divider from scratch
    using the programmatic chrome — a large serif number at left and an italic
    title at right with a gold rule below.

    title (also accepted as action_title) is the section heading.
    methodology_note is accepted for caller consistency and ignored.
    """
    if content:
        if "number" in content: number = content["number"]
        if "title"  in content: title  = content["title"]
    if title is None:
        title = action_title if action_title is not None else ""

    # Format the section number as a zero-padded two-digit string.
    num_str = ""
    if number is not None:
        try:
            num_str = f"{int(number):02d}"
        except (TypeError, ValueError):
            num_str = str(number)

    pal = _palette_for(style)

    # ------------------------------------------------------------------
    # Preferred path: native template Divider layout.
    # ------------------------------------------------------------------
    # For dark themes (synthetic, electronics) the native Divider layout's
    # full-slide background rectangle uses scheme:accent1, which is yellow
    # for those themes — this would make the divider entirely yellow.
    # Use the programmatic fallback instead so the palette bg (violet) is
    # applied correctly.
    divider_layout = _divider_layout(prs) if not _is_dark(style) else None
    if divider_layout is not None:
        slide = prs.slides.add_slide(divider_layout)
        # Populate number placeholder (idx 0).
        if num_str:
            ph0 = _populate_placeholder(0, slide, num_str)
            if ph0 is not None and ph0.has_text_frame:
                ph0.text_frame.word_wrap = True
                try:
                    from pptx.enum.text import MSO_AUTO_SIZE
                    ph0.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                except Exception:
                    pass
        # Populate chapter title placeholder (idx 13).
        ph13 = _populate_placeholder(13, slide, str(title))
        if ph13 is not None and ph13.has_text_frame:
            ph13.text_frame.word_wrap = True
            try:
                from pptx.enum.text import MSO_AUTO_SIZE
                ph13.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except Exception:
                pass
        # Add bottom chrome (classification + page number).
        _bottom_chrome(slide, meta, "Section", page, total, style)
        return slide

    # ------------------------------------------------------------------
    # Fallback: programmatic renderer (no themed template available).
    # ------------------------------------------------------------------
    slide = _new_slide(prs, bg_color=pal["bg"])
    _top_chrome(slide, meta, None, style, page=page, total=total)
    _bottom_chrome(slide, meta, "Section", page, total, style)

    dark = _is_dark(style)
    number_color  = pal["hot"]       if dark else pal["accent"]
    title_color   = WHITE            if dark else INK_DARK
    accent_color  = pal["highlight"]            # was hardcoded MERCK_GOLD

    # Huge serif number at left.
    _num_box = txt(slide, Inches(0.65), Inches(1.80), Inches(5.0), Inches(3.20),
                   str(num_str or number or ""), sz=96, color=number_color, bold=True,
                   font=FONT_HEAD, anchor=MSO_ANCHOR.MIDDLE)
    try:
        from pptx.enum.text import MSO_AUTO_SIZE as _AS
        _num_box.text_frame.auto_size = _AS.TEXT_TO_FIT_SHAPE
    except Exception:
        pass
    # Title to the right (uppercase tracking-wide).
    txt(slide, Inches(5.50), Inches(3.10), Inches(7.5), Inches(0.40),
        _tracked(title), sz=13, color=accent_color, bold=True, font=FONT_BODY,
        anchor=MSO_ANCHOR.TOP)
    # Italic phrase below.
    txt(slide, Inches(5.50), Inches(3.55), Inches(7.5), Inches(1.20),
        str(title), sz=32, color=title_color, italic=True, bold=True,
        font=FONT_HEAD)
    # Accent rule — colour from palette (was hardcoded MERCK_GOLD).
    hairline(slide, Inches(5.50), Inches(5.05), Inches(2.4),
             Emu(int(Pt(2.5))), accent_color)

    # Emit takeaway band and/or source line if supplied.
    if takeaway:
        _takeaway_band(slide, takeaway, style)
    if source:
        _source_line(slide, Inches(0.65), SOURCE_Y, Inches(12.0), SOURCE_H,
                     source, style)
    return slide



# ===========================================================================
# Layout: CLOSE
# ===========================================================================

def build_close(prs, meta, action_statement=None, style="merck_executive",
                page=None, total=None, section_number=None,
                category=None, takeaway=None, source=None,
                action_title=None, methodology_note=None, content=None):
    """Close slide. Action statement in big italic Merck Web, gold rule.
    page and total are accepted for caller consistency; close slides typically
    omit the page number visually but the kwargs prevent a TypeError on callers
    that pass page=N, total=N as part of their consistent slide-call pattern.
    action_title is accepted as an alias for action_statement.
    methodology_note is accepted for caller consistency and ignored.
    """
    if content:
        if "action_statement" in content: action_statement = content["action_statement"]
        if "takeaway"         in content: takeaway         = content["takeaway"]
    if action_statement is None:
        action_statement = action_title
    pal = _palette_for(style)
    slide = _new_slide(prs, bg_color=pal["bg"])
    # Close slides: progress bar shows full (deck ends here) when total known.
    close_p = total if total is not None else page
    _top_chrome(slide, meta, None, style, page=close_p, total=total)
    _bottom_chrome(slide, meta, "Close", None, None, style)

    dark = _is_dark(style)
    rule_color = pal["highlight"]
    text_color = WHITE if dark else MERCK_PURPLE

    # Gold accent rule above.
    hairline(slide, Inches(0.65), Inches(2.80), Inches(2.6),
             Emu(int(Pt(2.5))), rule_color)

    # Action statement text box — use the full safe zone from rule to source line.
    stmt_top = Inches(3.10)
    stmt_h   = SOURCE_Y - Inches(0.15) - stmt_top  # ≈ 3.30 in (3.10 → 6.40)

    # Adaptive font: scale down for long statements so text fits without clipping.
    _stmt = action_statement or ""
    _n = (len(_stmt) if isinstance(_stmt, str)
          else sum(len(s) for s, _ in _stmt if isinstance(s, str)))
    sz_stmt = 38 if _n <= 100 else (28 if _n <= 200 else 20)

    # Action statement in italic display font. Supports either a string OR a list
    # of (text, italic_bool) tuples for mixed-emphasis runs.
    if isinstance(action_statement, (list, tuple)) and action_statement and \
       all(isinstance(seg, (list, tuple)) and len(seg) == 2
           for seg in action_statement):
        box = slide.shapes.add_textbox(Inches(0.65), stmt_top,
                                       Inches(12.0), stmt_h)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.02)
        tf.margin_right = Inches(0.02)
        p = tf.paragraphs[0]
        for seg_text, seg_italic in action_statement:
            seg_color = pal["hot"] if seg_italic else text_color
            _add_run(p, seg_text, sz=sz_stmt, color=seg_color, bold=True,
                     italic=True, font=FONT_HEAD)
    else:
        txt(slide, Inches(0.65), stmt_top, Inches(12.0), stmt_h,
            action_statement or "", sz=sz_stmt, color=text_color, italic=True,
            bold=True, font=FONT_HEAD)
    # Small tagline beneath: deck label + month/year.
    deck_label = (meta or {}).get("deck_label", "")
    month_year = (meta or {}).get("month_year", "")
    parts = []
    if deck_label:
        parts.append(_tracked(deck_label))
    if month_year:
        parts.append(_tracked(month_year))
    if parts:
        sep = "   •   "
        tag_color = PANEL_LIGHT if dark else PURPLE_MUTED
        # When a takeaway band is present, the takeaway sits at 6.83; move the
        # tagline above the gold rule (which is at 2.80) so it doesn't collide.
        tag_y = Inches(2.50) if takeaway else Inches(6.40)
        txt(slide, Inches(0.65), tag_y, Inches(12.0), Inches(0.30),
            sep.join(parts), sz=10, color=tag_color, bold=True,
            font=FONT_BODY)

    # Emit takeaway band and/or source line if supplied.
    if takeaway:
        _takeaway_band(slide, takeaway, style)
    if source:
        _source_line(slide, Inches(0.65), SOURCE_Y, Inches(12.0), SOURCE_H,
                     source, style)
    return slide


